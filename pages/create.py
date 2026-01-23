import streamlit as st
import streamlit.components.v1 as components
import io
import json
from pptx import Presentation
import time
from datetime import datetime
from io import BytesIO
from docx import Document
import base64
import hashlib
from docx.enum.text import WD_ALIGN_PARAGRAPH
from templates.templateconfig import SYSTEM_TEMPLATES, ATS_COLORS, load_css_template
from utils import (
    generate_generic_html, get_user_resume, get_score_color, render_skills_section,
    get_score_label, ai_ats_score, extract_temp_from_docx,render_generic_section,
    should_regenerate_resume, generate_enhanced_resume, extract_template_from_html,save_user_templates,
    save_and_improve, add_new_item,delete_user_ppt_template,render_basic_details,load_user_doc_templates,
    docx_to_html_preview,save_user_doc_templates,analyze_slide_structure,generate_ppt_sections,save_user_doc_templates,
    match_generated_to_original,clear_and_replace_text,save_user_ppt_templates,ask_ai_for_mapping,auto_process_docx,
    load_user_templates,load_user_ppt_templates,delete_user_doc_template,convert_html_to_docx_spire,clear_template_state
)


st.set_page_config(
    page_title="CVmate Resume Creator",
    layout="wide",
    initial_sidebar_state="collapsed"
)



if st.query_params.get("logout") == "true":
    for key in list(st.session_state.keys()):
        del st.session_state[key]
    st.query_params.clear()
    st.switch_page("app.py")
    st.stop()


if 'logged_in_user' not in st.session_state:
    st.session_state.logged_in_user = None


if st.session_state.logged_in_user is None:
    logged_user = st.query_params.get("user")
    if logged_user:
        st.session_state.logged_in_user = logged_user
    else:
        st.warning("Please login first!")
        st.switch_page("app.py")
        st.stop()

if st.session_state.logged_in_user:
    current_param = st.query_params.get("user")
    if current_param != st.session_state.logged_in_user:
        st.query_params["user"] = st.session_state.logged_in_user


email = st.session_state.logged_in_user

# ============= SESSIONS =============


if 'selected_template' not in st.session_state:
    st.session_state.selected_template = None
if 'template_preview_html' not in st.session_state:
    st.session_state.template_preview_html = None
if 'template_preview_css' not in st.session_state:
    st.session_state.template_preview_css = None
if 'show_template_selector' not in st.session_state:
    st.session_state.show_template_selector = True
if 'show_visual_editor' not in st.session_state:
    st.session_state.show_visual_editor = False
if 'enhanced_resume' not in st.session_state:
    st.session_state.enhanced_resume = None
if 'final_resume_data' not in st.session_state:
    st.session_state.final_resume_data = None
if 'editor_content' not in st.session_state:
    st.session_state.editor_content = ""
if 'ats_result' not in st.session_state:
    st.session_state.ats_result = {}
if 'edit_mode' not in st.session_state:
    st.session_state.edit_mode = False
if 'current_edit_section' not in st.session_state:
    st.session_state.current_edit_section = None
if 'show_upload_modal' not in st.session_state:
    st.session_state.show_upload_modal = False
if 'uploaded_templates' not in st.session_state:
    st.session_state.uploaded_templates = {}
if 'generated_docx' not in st.session_state:
    st.session_state.generated_docx = None
if 'generated_docx_temp' not in st.session_state:
    st.session_state.generated_docx_temp = None
if 'generated_ppt' not in st.session_state:
    st.session_state.generated_ppt = None
if 'job_description' not in st.session_state:
    st.session_state.job_description = None


if st.session_state.job_description is None:
    jd_param = st.query_params.get("jd")
    if jd_param:
        try:
            decoded_bytes = base64.b64decode(jd_param)
            decoded_json = decoded_bytes.decode('utf-8')
            st.session_state.job_description = json.loads(decoded_json)
        except Exception as e:
            st.session_state.job_description = {}
    else:
        st.session_state.job_description = {}


if st.session_state.job_description and isinstance(st.session_state.job_description, dict) and len(st.session_state.job_description) > 0:
    try:
        jd_json = json.dumps(st.session_state.job_description)
        encoded_bytes = base64.b64encode(jd_json.encode('utf-8'))
        encoded_str = encoded_bytes.decode('utf-8')
        st.query_params["jd"] = encoded_str
    except:
        pass

if st.session_state.final_resume_data is None:
    user_resume = get_user_resume(email)
    if user_resume:
        st.session_state.final_resume_data = user_resume
else:
    user_resume = st.session_state.final_resume_data


if st.session_state.enhanced_resume is None and user_resume:
    st.session_state.enhanced_resume = user_resume.copy()


resume_data = st.session_state.get('enhanced_resume')
# st.write(resume_data)
jd_data = st.session_state.get('job_description')


if 'enhanced_resume' not in st.session_state or st.session_state.enhanced_resume is None:
    user_resume = get_user_resume(email)
    if user_resume:
        st.session_state.enhanced_resume = user_resume.copy()


jd_data = st.session_state.get('job_description')


if jd_data and isinstance(jd_data, dict) and len(jd_data) > 0:
    if 'last_enhanced_jd' not in st.session_state:
        st.session_state.last_enhanced_jd = None
  
    jd_hash = hashlib.md5(str(jd_data).encode()).hexdigest()
    
  
    if st.session_state.last_enhanced_jd != jd_hash or st.session_state.get('force_enhance', False):
        resume_data = st.session_state.enhanced_resume
        generate_enhanced_resume(resume_data, jd_data)
        st.session_state.last_enhanced_jd = jd_hash
        st.session_state.force_enhance = False

resume_data = st.session_state.get('enhanced_resume')


if resume_data and jd_data and isinstance(jd_data, dict) and len(jd_data) > 0:
    try:
        if not st.session_state.ats_result or st.session_state.ats_result.get('overall_score', 0) == 0:
            st.session_state.ats_result = ai_ats_score(st.session_state.enhanced_resume, jd_data)
    except Exception as e:
        st.warning(f"Could not calculate ATS score: {str(e)}")
        st.session_state.ats_result = {}
else:
    st.session_state.ats_result = {}






RESUME_ORDER = ["education", "experience", "skills", "projects", "certifications", "achievements"]
st.markdown("""
<style>
    [data-testid="stSidebar"], [data-testid="collapsedControl"], 
    [data-testid="stSidebarNav"] {display: none;}
    #MainMenu, footer, header, button[kind="header"] {visibility: hidden;}
    
    /* Orange and White Theme */
    :root {
        --primary-orange: #FF6B35;
        --primary-orange-light: #FF8B5C;
        --primary-orange-dark: #D45A2B;
        --white: #FFFFFF;
        --light-gray: #F8F9FA;
        --medium-gray: #E9ECEF;
        --dark-gray: #495057;
        --text-dark: #212529;
        --success-green: #28A745;
    }
    
    .stApp {
        background: var(--white) !important;
    }
    
    /* Main container */
    .main .block-container {
        padding-top: 0.5rem !important;
        padding-left: 1rem !important;
        padding-right: 1rem !important;
        max-width: 100% !important;
    }
    
    /* Template Gallery - Orange Theme */
    .template-header {
        margin-bottom: 1.5rem;
    }
    
    .template-title {
        font-size: 2.2rem;
        font-weight: 700;
        color: var(--text-dark);
        margin-bottom: 0.3rem;
    }
    
    .template-subtitle {
        font-size: 1rem;
        color: var(--dark-gray);
        margin-bottom: 1.5rem;
    }
    
    .template-card {
        background: var(--white);
        border-radius: 12px;
        overflow: hidden;
        box-shadow: 0 4px 12px rgba(255, 107, 53, 0.1);
        transition: all 0.3s ease;
        border: 1px solid var(--medium-gray);
        height: 100%;
        cursor: pointer;
        margin-bottom: 1rem;
    }
    
    .template-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 12px 24px rgba(255, 107, 53, 0.2);
        border-color: var(--primary-orange);
    }
    
    .create-blank-card {
        background: var(--white);
        border-radius: 12px;
        padding: 1.5rem;
        box-shadow: 0 4px 12px rgba(255, 107, 53, 0.1);
        transition: all 0.3s ease;
        border: 2px dashed var(--primary-orange);
        height: 380px;
        display: flex;
        flex-direction: column;
        align-items: center;
        justify-content: center;
        cursor: pointer;
        margin-bottom: 1rem;
    }
    
    .create-blank-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 12px 24px rgba(255, 107, 53, 0.2);
        border-color: var(--primary-orange-dark);
        background: linear-gradient(135deg, rgba(255, 107, 53, 0.05) 0%, rgba(255, 107, 53, 0.02) 100%);
    }
    
    .plus-icon {
        width: 60px;
        height: 60px;
        border-radius: 50%;
        background: linear-gradient(135deg, var(--primary-orange) 0%, var(--primary-orange-light) 100%);
        display: flex;
        align-items: center;
        justify-content: center;
        font-size: 2rem;
        color: var(--white);
        margin-bottom: 1rem;
    }
    
    .template-preview {
        width: 100%;
        height: 200px;
        background: linear-gradient(135deg, var(--light-gray) 0%, var(--medium-gray) 100%);
        display: flex;
        align-items: center;
        justify-content: center;
        position: relative;
        overflow: hidden;
    }
    
    .template-badge {
        position: absolute;
        top: 0.5rem;
        right: 0.5rem;
        background: var(--primary-orange);
        color: var(--white);
        padding: 0.3rem 0.8rem;
        border-radius: 12px;
        font-size: 0.7rem;
        font-weight: 600;
    }
    
    /* Visual Editor Container */
 .visual-editor-container {
        width: 100%;
        height: calc(100vh - 120px);
        border: none;
        overflow: auto;
        margin-top: 0.5rem;
        background: var(--light-gray);
        border-radius: 12px;
        border: 1px solid var(--medium-gray);
    }
    
  
    iframe {
        overflow: auto !important;
    }
    /* Resume Tools Panel - Orange Theme */
    .resume-tools-panel {
        background: var(--white);
        border-radius: 12px;
        padding: 1.5rem;
        box-shadow: 0 4px 12px rgba(255, 107, 53, 0.1);
        height: fit-content;
        margin-top: 0;
        border: 1px solid var(--medium-gray);
    }
    
    .ats-score-display {
        background: linear-gradient(135deg, rgba(255, 107, 53, 0.1) 0%, rgba(255, 107, 53, 0.05) 100%);
        border-radius: 12px;
        padding: 1.5rem;
        text-align: center;
        margin-bottom: 1.5rem;
        border: 1px solid rgba(255, 107, 53, 0.2);
    }
    
    .ats-score-number {
        font-size: 3rem;
        font-weight: 800;
        color: var(--primary-orange);
        line-height: 1;
        margin-bottom: 0.5rem;
    }
    
    .ats-score-label {
        font-size: 1.1rem;
        color: var(--text-dark);
        font-weight: 700;
        margin-bottom: 0.3rem;
    }
    
    .ats-score-subtext {
        font-size: 0.9rem;
        color: var(--dark-gray);
    }
    
   
.stButton > button[kind="primary"] {
    background: linear-gradient(135deg, #FF6B35 0%, #FF8B5C 100%) !important;
    color: #FFFFFF !important;
    border: none !important;
}

/* Hover */
.stButton > button[kind="primary"]:hover {
    background: linear-gradient(135deg, #D45A2B 0%, #FF6B35 100%) !important;
    box-shadow: 0 6px 20px rgba(255, 107, 53, 0.35) !important;
    transform: translateY(-2px);
}


.stButton > button[kind="primary"]:active {
    background: #D45A2B !important;
}


    .editor-header {
        background: var(--white);
        padding: 1rem;
        border-radius: 12px;
        margin-bottom: 1rem;
        box-shadow: 0 2px 8px rgba(255, 107, 53, 0.1);
        border: 1px solid var(--medium-gray);
    }
    
 
    .inline-editor {
        position: absolute;
        background: var(--white);
        border: 2px solid var(--primary-orange);
        border-radius: 8px;
        padding: 1rem;
        box-shadow: 0 8px 24px rgba(0, 0, 0, 0.15);
        z-index: 10000;
        min-width: 300px;
        max-width: 500px;
    }
    
    .inline-editor-header {
        display: flex;
        justify-content: space-between;
        align-items: center;
        margin-bottom: 1rem;
        padding-bottom: 0.5rem;
        border-bottom: 1px solid var(--medium-gray);
    }
    
    .inline-editor-title {
        font-weight: 600;
        color: var(--primary-orange);
        font-size: 1.1rem;
    }
    
    .inline-editor-close {
        background: none;
        border: none;
        color: var(--dark-gray);
        cursor: pointer;
        font-size: 1.2rem;
        padding: 0;
        width: 24px;
        height: 24px;
        display: flex;
        align-items: center;
        justify-content: center;
        border-radius: 4px;
    }
    
    .inline-editor-close:hover {
        background: var(--light-gray);
        color: var(--primary-orange);
    }
    
    /* Template grid */
    [data-testid="column"] {
        padding: 0.5rem;
    }
    
    /* Expander styling */
    .streamlit-expanderHeader {
        background-color: rgba(255, 107, 53, 0.05) !important;
        border-radius: 8px !important;
        border: 1px solid rgba(255, 107, 53, 0.1) !important;
    }
    
    .streamlit-expanderHeader:hover {
        background-color: rgba(255, 107, 53, 0.1) !important;
    }
    

    .empty-section {
        display: none !important;
    }
    

    .stTextArea textarea {
        border-radius: 8px !important;
        border: 2px solid var(--medium-gray) !important;
        transition: all 0.3s ease !important;
    }
    
    .stTextArea textarea:focus {
        border-color: var(--primary-orange) !important;
        box-shadow: 0 0 0 3px rgba(255, 107, 53, 0.1) !important;
    }
    
    .stTextInput input {
        border-radius: 8px !important;
        border: 2px solid var(--medium-gray) !important;
        transition: all 0.3s ease !important;
    }
    
    .stTextInput input:focus {
        border-color: var(--primary-orange) !important;
        box-shadow: 0 0 0 3px rgba(255, 107, 53, 0.1) !important;
    }
    :root {
        --primary: #FF6B35;
        --primary-dark: #E85A28;
        --primary-light: #FF8C5A;
        --accent: #FFA500;
        --bg-primary: #FAFAFA;
        --bg-secondary: #FFFFFF;
        --text-primary: #1A1A1A;
        --text-secondary: #666666;
        --text-light: #999999;
        --border: #E5E5E5;
        --shadow: rgba(255, 107, 53, 0.12);
    }     

     html, body, .stApp {
        font-family: 'Inter', sans-serif;
        background: var(--bg-primary);
        color: var(--text-primary);
        scroll-behavior: smooth;
    }

    /* ==================== NAVIGATION ==================== */
    .nav-wrapper {
        position: fixed;
        top: 0;
        left: 0;
        right: 0;
        z-index: 1000;
        background: rgba(255, 255, 255, 0.95);
        backdrop-filter: blur(20px);
        border-bottom: 1px solid var(--border);
        animation: slideDown 0.6s ease-out;
    }

    @keyframes slideDown {
        from {
            transform: translateY(-100%);
            opacity: 0;
        }
        to {
            transform: translateY(0);
            opacity: 1;
        }
    }

    .nav-container {
        max-width: 1400px;
        margin: 0 auto;
        padding: 0 3rem;
        display: flex;
        align-items: center;
        justify-content: space-between;
        height: 80px;
    }

    .logo {
        font-family: 'Archivo', sans-serif;
        font-size: 28px;
        font-weight: 900;
        background: linear-gradient(135deg, var(--primary) 0%, var(--accent) 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
        letter-spacing: -1px;
    }

    .nav-menu {
        display: flex;
        gap: 2rem;
        align-items: center;
    }

    .nav-link {
        color: var(--text-secondary) !important;
        text-decoration: none !important;
        font-size: 15px;
        font-weight: 500;
        padding: 10px 20px;
        border-radius: 8px;
        transition: all 0.3s ease;
        position: relative;
    }

    .nav-link:hover {
        color: var(--primary) !important;
        background: rgba(255, 107, 53, 0.08);
    }

    .nav-link.btn-primary {
        background: linear-gradient(135deg, var(--primary) 0%, var(--primary-dark) 100%);
        color: white !important;
        font-weight: 600;
        padding: 12px 28px;
        box-shadow: 0 4px 12px var(--shadow);
    }

    .nav-link.btn-primary:hover {
        transform: translateY(-2px);
        box-shadow: 0 6px 20px var(--shadow);
    }
</style>
""", unsafe_allow_html=True)



if "logged_in_user" not in st.session_state or st.session_state.logged_in_user is None:
    logged_user = st.query_params.get("user")
    if logged_user:
        st.session_state.logged_in_user = logged_user

current_user = st.session_state.get('logged_in_user', '')
is_logged_in = bool(current_user)


if is_logged_in and current_user:
    home_url = f"/?user={current_user}"
    ats_url = f"ats?user={current_user}"
    qu_url = f"qu?user={current_user}"
    change_url = f"change?user={current_user}"
else:
    home_url = "/"
    ats_url = "#ats"
    qu_url = "#qu"
    change_url = "#change"

if is_logged_in:
    auth_button = '<a class="nav-link" href="?logout=true" target="_self">⏻</a>'
else:
    auth_button = '<a class="nav-link" href="#Login" target="_self">Login</a>'


st.markdown(f"""
<div class="nav-wrapper">
    <div class="nav-container">
        <div class="logo">CVmate</div>
        <div class="nav-menu">
            <a class="nav-link" href="{home_url}" target="_self">Home</a>
            <a class="nav-link" href="main?&user={current_user}" target="_self">Create Resume</a>
            <a class="nav-link" href="job?&user={current_user}" target="_self">Add New JD</a>
            <a class="nav-link" href="{ats_url}" target="_self">ATS Checker</a>
            <a class="nav-link" href="{qu_url}" target="_self">AI Assistant</a>
            <a class="nav-link" href="{change_url}" target="_self">Change Template</a>
            {auth_button}
        </div>
    </div>
</div>
""", unsafe_allow_html=True)



def process_html_upload(uploaded_file):
    """Process uploaded HTML file."""
    import chardet
    resume_data = st.session_state.get('final_resume_data') or {}
    raw_data = uploaded_file.read()
    detected = chardet.detect(raw_data)
    encoding = detected["encoding"] or "utf-8"
    content = raw_data.decode(encoding, errors="ignore")

    parsed_template = extract_template_from_html(content)

    st.session_state.temp_upload_config = {
        'name': f"Uploaded_{uploaded_file.name.split('.')[0]}",
        'css': parsed_template.get('css', ''),
        'html': parsed_template.get('html', ''),
        'original_filename': uploaded_file.name
    }
    

    preview_html = f"""
        <style>{parsed_template.get('css', '')}</style>
        <div class="ats-page">{generate_generic_html(resume_data)}</div>
    """
    st.session_state.temp_upload_preview = preview_html
    st.session_state.selected_template_config = st.session_state.temp_upload_config
    st.session_state.template_source = 'temp_upload'

def show_inline_doc_mapping_editor():
    """Show inline editor for Word document mapping in visual editor."""
    st.markdown("---")
    st.markdown("### ✏️ Edit Document Mapping")
    
  
    if 'mapping' not in st.session_state or 'doc_original_bytes' not in st.session_state:
        st.warning("⚠️ No mapping data available. Generate document first.")
        return
    
    mapping = st.session_state['mapping'].copy()
    
    st.info("💡 Edit the values below to change how your resume data maps to the document")
    
   
    st.markdown("#### Current Mappings")
    
    if mapping:
       
        for idx, (template_key, resume_value) in enumerate(mapping.items()):
            with st.container():
                col1, col2, col3 = st.columns([2, 3, 1])
                
                with col1:
                    st.text_input(
                        "Template",
                        value=template_key,
                        disabled=True,
                        key=f"inline_template_key_{idx}",
                        label_visibility="collapsed"
                    )
                
                with col2:
                    new_value = st.text_area(
                        "Value",
                        value=resume_value,
                        height=80,
                        key=f"inline_resume_value_{idx}",
                        label_visibility="collapsed"
                    )
                    mapping[template_key] = new_value
                
                with col3:
                    st.write("")
                    st.write("")
                    if st.button("🗑️", key=f"inline_delete_{idx}", use_container_width=True):
                        del mapping[template_key]
                        st.session_state['mapping'] = mapping
                        st.rerun()
        
        st.session_state['mapping'] = mapping
    else:
        st.warning("No mappings found.")
    
    
    st.markdown("---")
    st.markdown("#### Add New Mapping")
    
    col1, col2, col3 = st.columns([2, 3, 1])
    
    with col1:
        new_key = st.text_input("Template Text", key="inline_new_template_key")
    
    with col2:
        new_value = st.text_area("Resume Value", height=80, key="inline_new_resume_value")
    
    with col3:
        st.write("")
        st.write("")
        if st.button("➕ Add", use_container_width=True):
            if new_key and new_value:
                st.session_state['mapping'][new_key] = new_value
                st.success("✅ Added!")
                st.rerun()
    
    
    st.markdown("---")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("✨ Apply Changes & Regenerate", type="primary", use_container_width=True):
            try:
             
                
                import io
                original_bytes = st.session_state['doc_original_bytes']
                
                with st.spinner("Regenerating document..."):
                    output_doc = auto_process_docx(
                        io.BytesIO(original_bytes),
                        st.session_state['mapping']
                    )
           
                st.session_state['generated_docx'] = output_doc.getvalue()
                
                st.success("✅ Document regenerated with new mapping!")
                time.sleep(1)
                st.rerun()
                
            except Exception as e:
                st.error(f"❌ Error regenerating: {str(e)}")
                import traceback
                st.error(traceback.format_exc())
    
    with col2:
        if st.button("🔄 Reset to Original", type="secondary", use_container_width=True):
            try:
                import io

                original_bytes = st.session_state['doc_original_bytes']
                uploadtext = extract_temp_from_docx(io.BytesIO(original_bytes))
                
                resume_data = st.session_state.get('enhanced_resume') or st.session_state.get('final_resume_data') or {}
                
                with st.spinner("Resetting to AI mapping..."):
                    mapped_data = ask_ai_for_mapping(uploadtext, resume_data)
                    
                    if isinstance(mapped_data, list):
                        mapped_data = {
                            item["template"]: item["new"]
                            for item in mapped_data
                            if "template" in item and "new" in item
                        }
                    
                    st.session_state['mapping'] = mapped_data
                
                st.success("✅ Reset to original mapping!")
                st.rerun()
                
            except Exception as e:
                st.error(f"❌ Error resetting: {str(e)}")


def show_inline_ppt_content_editor():
    """Show inline editor for PowerPoint content in visual editor - FIXED VERSION."""
    st.markdown("---")
    st.markdown("### ✏️ Edit Presentation Content")
    

    if 'temp_ppt_text_elements' not in st.session_state or not st.session_state.get('temp_ppt_text_elements'):

        if st.session_state.get('generated_ppt'):
            try:
                import io
                from pptx import Presentation
                
                prs = Presentation(io.BytesIO(st.session_state['generated_ppt']))
                
                text_elements = []
                for slide_idx, slide in enumerate(prs.slides):
                    for shape_idx, shape in enumerate(slide.shapes):
                        if shape.has_text_frame and shape.text.strip():
                            text_elements.append({
                                'slide': slide_idx + 1,
                                'shape': shape_idx,
                                'original_text': shape.text.strip(),
                                'shape_type': type(shape).__name__
                            })
                
                st.session_state['temp_ppt_text_elements'] = text_elements
                
            except Exception as e:
                st.error(f"❌ Could not extract presentation content: {str(e)}")
                st.info("💡 Try regenerating the presentation using 'Save & Auto-Improve'")
                return
        

        elif st.session_state.get('selected_ppt_template'):
            template_data = st.session_state['selected_ppt_template']
            if 'text_elements' in template_data:
                st.session_state['temp_ppt_text_elements'] = template_data['text_elements']
            else:
                st.warning("⚠️ No content data available. Please regenerate the presentation.")
                return
        else:
            st.warning("⚠️ No presentation data found. Generate presentation first.")
            return
    

    if 'temp_ppt_edits' not in st.session_state:
        st.session_state['temp_ppt_edits'] = {}
    
    edits = st.session_state['temp_ppt_edits'].copy()
    text_elements = st.session_state['temp_ppt_text_elements']
    
    if not text_elements:
        st.warning("⚠️ No text elements found in presentation.")
        return
    
    st.info("💡 Edit the content for each slide below")

    slides_dict = {}
    for element in text_elements:
        slide_num = element['slide']
        if slide_num not in slides_dict:
            slides_dict[slide_num] = []
        slides_dict[slide_num].append(element)
    

    for slide_num in sorted(slides_dict.keys()):
        with st.expander(f"📊 Slide {slide_num}", expanded=(slide_num==1)):
            for element in slides_dict[slide_num]:
                key = f"{element['slide']}_{element['shape']}"
                current_text = edits.get(key, element['original_text'])
                
                col1, col2 = st.columns([4, 1])
                
                with col1:
                    new_text = st.text_area(
                        f"Shape {element['shape']}",
                        value=current_text,
                        height=100,
                        key=f"inline_ppt_{key}"
                    )
                    edits[key] = new_text
                
                with col2:
                    st.write("")
                    st.write("")
                    if st.button("🔄", key=f"inline_reset_{key}", use_container_width=True, help="Reset to original"):
                        edits[key] = element['original_text']
                        st.rerun()
    
    st.session_state['temp_ppt_edits'] = edits
    

    st.markdown("---")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("✨ Apply Changes & Regenerate", type="primary", use_container_width=True):
            try:
                import io
                from pptx import Presentation
                
      
                ppt_data = None
                
        
                if st.session_state.get('generated_ppt'):
                    ppt_data = st.session_state['generated_ppt']
   
                elif st.session_state.get('selected_ppt_template'):
                    ppt_data = st.session_state['selected_ppt_template'].get('ppt_data')
                
                if not ppt_data:
                    st.error("❌ No presentation data found to regenerate.")
                    return
                
                working_prs = Presentation(io.BytesIO(ppt_data))
                
        
                with st.spinner("Regenerating presentation..."):
                    for element in text_elements:
                        key = f"{element['slide']}_{element['shape']}"
                        if key in edits:
                            slide_idx = element['slide'] - 1
                            shape_idx = element['shape']
                            
                            if slide_idx < len(working_prs.slides):
                                slide = working_prs.slides[slide_idx]
                                if shape_idx < len(slide.shapes):
                                    shape = slide.shapes[shape_idx]
                                    if shape.has_text_frame:
                                        clear_and_replace_text(shape, edits[key])
                
                # Save
                output = io.BytesIO()
                working_prs.save(output)
                output.seek(0)
                
                st.session_state['generated_ppt'] = output.getvalue()
                
                st.success("✅ Presentation regenerated!")
                
                time.sleep(1)
                st.rerun()
                
            except Exception as e:
                st.error(f"❌ Error regenerating: {str(e)}")
                import traceback
                st.error(traceback.format_exc())
    
    with col2:
        if st.button("🔄 Reset All", type="secondary", use_container_width=True):

            for element in text_elements:
                key = f"{element['slide']}_{element['shape']}"
                edits[key] = element['original_text']
            
            st.session_state['temp_ppt_edits'] = edits
            st.success("✅ Reset to original content!")
            st.rerun()


def show_inline_doc_mapping_editor():
    """Show inline editor for Word document mapping in visual editor - FIXED VERSION."""
    st.markdown("---")
    st.markdown("### ✏️ Edit Document Mapping")
    
  
    if 'mapping' not in st.session_state or not st.session_state.get('mapping'):
        if st.session_state.get('doc_original_bytes'):
            try:
                import io
                
     
                uploadtext = extract_temp_from_docx(io.BytesIO(st.session_state['doc_original_bytes']))
                
         
                resume_data = st.session_state.get('enhanced_resume') or st.session_state.get('final_resume_data') or {}
                
                with st.spinner("Generating mapping..."):
                    mapped_data = ask_ai_for_mapping(uploadtext, resume_data)
                    
                    if isinstance(mapped_data, list):
                        mapped_data = {
                            item["template"]: item["new"]
                            for item in mapped_data
                            if "template" in item and "new" in item
                        }
                    
                    st.session_state['mapping'] = mapped_data
                    st.session_state['template_text'] = uploadtext
                    
            except Exception as e:
                st.error(f"❌ Could not generate mapping: {str(e)}")
                st.info("💡 Try regenerating the document using 'Save & Auto-Improve'")
                return
        else:
            st.warning("⚠️ No document data found. Generate document first.")
            return
    

    if 'mapping' not in st.session_state or 'doc_original_bytes' not in st.session_state:
        st.warning("⚠️ No mapping data available. Generate document first.")
        return
    
    mapping = st.session_state['mapping'].copy()
    
    st.info("💡 Edit the values below to change how your resume data maps to the document")
    

    st.markdown("#### Current Mappings")
    
    if mapping:

        for idx, (template_key, resume_value) in enumerate(mapping.items()):
            with st.container():
                col1, col2, col3 = st.columns([2, 3, 1])
                
                with col1:
                    st.text_input(
                        "Template",
                        value=template_key,
                        disabled=True,
                        key=f"inline_template_key_{idx}",
                        label_visibility="collapsed"
                    )
                
                with col2:
                    new_value = st.text_area(
                        "Value",
                        value=resume_value,
                        height=80,
                        key=f"inline_resume_value_{idx}",
                        label_visibility="collapsed"
                    )
                    mapping[template_key] = new_value
                
                with col3:
                    st.write("")
                    st.write("")
                    if st.button("🗑️", key=f"inline_delete_{idx}", use_container_width=True):
                        del mapping[template_key]
                        st.session_state['mapping'] = mapping
                        st.rerun()
        
        st.session_state['mapping'] = mapping
    else:
        st.warning("No mappings found.")
    

    st.markdown("---")
    st.markdown("#### Add New Mapping")
    
    col1, col2, col3 = st.columns([2, 3, 1])
    
    with col1:
        new_key = st.text_input("Template Text", key="inline_new_template_key")
    
    with col2:
        new_value = st.text_area("Resume Value", height=80, key="inline_new_resume_value")
    
    with col3:
        st.write("")
        st.write("")
        if st.button("➕ Add", use_container_width=True):
            if new_key and new_value:
                st.session_state['mapping'][new_key] = new_value
                st.success("✅ Added!")
                st.rerun()
    

    st.markdown("---")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("✨ Apply Changes & Regenerate", type="primary", use_container_width=True):
            try:
                import io
      
                original_bytes = st.session_state['doc_original_bytes']
                
  
                with st.spinner("Regenerating document..."):
                    output_doc = auto_process_docx(
                        io.BytesIO(original_bytes),
                        st.session_state['mapping']
                    )
                
  
                st.session_state['generated_docx'] = output_doc.getvalue()
                
                st.success("✅ Document regenerated with new mapping!")
                
                time.sleep(1)
                st.rerun()
                
            except Exception as e:
                st.error(f"❌ Error regenerating: {str(e)}")
                import traceback
                st.error(traceback.format_exc())
    
    with col2:
        if st.button("🔄 Reset to Original", type="secondary", use_container_width=True):
            try:
                import io
                
      
                original_bytes = st.session_state['doc_original_bytes']
                uploadtext = extract_temp_from_docx(io.BytesIO(original_bytes))
                
                resume_data = st.session_state.get('enhanced_resume') or st.session_state.get('final_resume_data') or {}
                
                with st.spinner("Resetting to AI mapping..."):
                    mapped_data = ask_ai_for_mapping(uploadtext, resume_data)
                    
                    if isinstance(mapped_data, list):
                        mapped_data = {
                            item["template"]: item["new"]
                            for item in mapped_data
                            if "template" in item and "new" in item
                        }
                    
                    st.session_state['mapping'] = mapped_data
                
                st.success("✅ Reset to original mapping!")
                st.rerun()
                
            except Exception as e:
                st.error(f"❌ Error resetting: {str(e)}")


def save_uploaded_template(template_name, file_type):
    """Save the uploaded template to user's templates."""
    current_user = st.session_state.logged_in_user
    
    try:
        if file_type == 'html':
            if 'uploaded_templates' not in st.session_state:
                st.session_state.uploaded_templates = {}
            
            template_id = f"html_{datetime.now().strftime('%Y%m%d%H%M%S')}"
            temp_config = st.session_state.get('temp_upload_config', {})
            
            st.session_state.uploaded_templates[template_id] = {
                'name': template_name,
                'css': temp_config.get('css', ''),
                'html': temp_config.get('html', ''),
                'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'original_filename': st.session_state.get('uploaded_file_name', 'unknown.html'),
                'type': 'html'
            }
  
            save_user_templates(current_user, st.session_state.uploaded_templates)
        
        elif file_type in ['docx', 'doc']:
            if 'doc_templates' not in st.session_state:
                st.session_state.doc_templates = {}
            
            template_id = f"doc_{datetime.now().strftime('%Y%m%d%H%M%S')}"
            temp_config = st.session_state.get('temp_upload_config', {})
            
            st.session_state.doc_templates[template_id] = {
                'name': template_name,
                'doc_data': temp_config.get('doc_data', st.session_state.get('temp_doc_data')),
                'structure': temp_config.get('structure', st.session_state.get('temp_doc_structure')),
                'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'original_filename': st.session_state.uploaded_file_name,
                'type': 'word',
                'sections_detected': temp_config.get('sections_detected', [])
            }
            
 

            save_user_doc_templates(current_user, st.session_state.doc_templates)
        
        elif file_type in ['pptx', 'ppt']:
            if 'ppt_templates' not in st.session_state:
                st.session_state.ppt_templates = {}
            
            template_id = f"ppt_{datetime.now().strftime('%Y%m%d%H%M%S')}"
            temp_config = st.session_state.get('temp_upload_config', {})
            
            st.session_state.ppt_templates[template_id] = {
                'name': template_name,
                'ppt_data': temp_config.get('ppt_data', st.session_state.get('temp_ppt_data')),
                'structure': temp_config.get('structure', st.session_state.get('temp_ppt_structure')),
                'edits': temp_config.get('edits', st.session_state.get('temp_ppt_edits', {})),
                'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'original_filename': st.session_state.uploaded_file_name,
                'type': 'powerpoint',
                'text_elements': temp_config.get('text_elements', [])
            }
            
            save_user_ppt_templates(current_user, st.session_state.ppt_templates)
        
    except Exception as e:
        st.error(f"Error saving template: {str(e)}")
        raise e


def show_word_doc_preview_and_save(uploaded_file, final_data):
    """Display Word document preview with save template section."""
    try:
 
        st.markdown("---")
        st.markdown("### 💾 Save Word Template")
        
        col1, col2 = st.columns([2, 1])
        
        with col1:
            doc_template_name = st.text_input(
                "Template Name:",
                value=f"DocTemplate_{uploaded_file.name.split('.')[0]}",
                key="doc_template_name"
            )
        
        with col2:
            st.write("")
            st.write("")
            if st.button("💾 Save Template", width=True, type="primary", key="save_doc_btn"):
                current_user = st.session_state.logged_in_user
                
                if 'doc_templates' not in st.session_state:
                    st.session_state.doc_templates = {}
                
                template_id = f"doc_{datetime.now().strftime('%Y%m%d%H%M%S')}"
                structure = st.session_state.temp_doc_structure
                
                st.session_state.doc_templates[template_id] = {
                    'name': doc_template_name,
                    'doc_data': st.session_state.temp_doc_data,
                    'structure': structure,
                    'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                    'original_filename': uploaded_file.name,
                    'sections_detected': [s['section'] for s in structure],
                    'type': 'word'
                }
                
                if save_user_doc_templates(current_user, st.session_state.doc_templates):
                    st.success(f"✅ Template '{doc_template_name}' saved!")
                    
                    time.sleep(1)
                    st.rerun()
                else:
                    st.error("Failed to save template. Please try again.")
        
        # ========== PREVIEW SECTION ==========
        st.markdown("---")
        st.markdown("### 🔍 Document Preview (Not Saved Yet)")
        
        if st.session_state.get('generated_doc'):
            from docx import Document
            import io
            
            doc_stream = io.BytesIO(st.session_state.generated_doc)
            processed_doc = Document(doc_stream)
            
            st.markdown("""
            <style>
            .doc-preview {
                border: 2px solid #e0e0e0;
                border-radius: 10px;
                padding: 40px;
                background: white;
                min-height: 600px;
                max-height: 800px;
                overflow-y: auto;
                font-family: 'Calibri', 'Arial', sans-serif;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            }
            .doc-name { 
                font-size: 24px; 
                font-weight: bold; 
                margin-bottom: 5px;
                color: #1a1a1a;
            }
            .doc-title { 
                font-size: 14px; 
                margin-bottom: 5px;
                color: #4a4a4a;
            }
            .doc-contact { 
                font-size: 12px; 
                margin-bottom: 20px;
                color: #666;
            }
            .doc-heading { 
                font-size: 16px; 
                font-weight: bold; 
                margin: 20px 0 10px 0;
                border-bottom: 2px solid #333;
                padding-bottom: 5px;
                color: #1a1a1a;
            }
            .doc-text { 
                font-size: 11pt; 
                line-height: 1.6;
                margin: 8px 0;
                color: #333;
                white-space: pre-wrap;
            }
            </style>
            """, unsafe_allow_html=True)
            
            html_content = '<div class="doc-preview">'
            para_count = 0
            
            for para in processed_doc.paragraphs:
                if not para.text.strip():
                    continue
                
                text = para.text.strip()
                
                if para_count == 0:
                    html_content += f'<div class="doc-name">{text}</div>'
                elif para_count == 1:
                    html_content += f'<div class="doc-title">{text}</div>'
                elif para_count == 2:
                    html_content += f'<div class="doc-contact">{text}</div>'
                elif para.style.name.startswith('Heading') or (para.runs and para.runs[0].bold and len(text.split()) <= 10):
                    html_content += f'<div class="doc-heading">{text}</div>'
                else:
                    formatted_text = text.replace('\n', '<br>')
                    html_content += f'<div class="doc-text">{formatted_text}</div>'
                
                para_count += 1
            
            html_content += '</div>'
            st.markdown(html_content, unsafe_allow_html=True)
        
        # ========== DOWNLOAD SECTION ==========
        st.markdown("---")
        filename = f"{final_data.get('name', 'Resume').replace(' ', '_')}_Final.docx"
        
        col1, col2 = st.columns([3, 1])
        
        with col1:
            st.download_button(
                label="📥 Download Document (Not Saved Yet)",
                data=st.session_state.generated_doc,
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True,
                type="primary"
            )
        
        with col2:
            if st.button("🔄 Reset", use_container_width=True):
                for key in ['generated_doc', 'temp_doc_data', 'temp_doc_filename', 'temp_doc_structure']:
                    if key in st.session_state:
                        del st.session_state[key]
                st.rerun()
                
    except Exception as e:
        st.error(f"Preview error: {str(e)}")

def process_word_upload(uploaded_file, final_data):
    """Process uploaded Word document and generate preview - FIXED VERSION."""
    try:
        # Check if already processed to avoid reprocessing
        if not st.session_state.get('doc_already_processed', False):
            # ✅ Read the file ONCE and store it
            uploaded_file.seek(0)
            original_file_bytes = uploaded_file.read()
            
            # ✅ Extract template text
            uploaded_file.seek(0)
            uploadtext = extract_temp_from_docx(uploaded_file)
            
            # ✅ Generate mapping
            mapped_data = ask_ai_for_mapping(uploadtext, final_data)

            # ✅ FORCE mapping to be a dictionary
            if isinstance(mapped_data, list):
                mapped_data = {
                    item["template"]: item["new"]
                    for item in mapped_data
                    if "template" in item and "new" in item
                }

            st.session_state['mapping'] = mapped_data
            st.session_state['template_text'] = uploadtext

            # ✅ Generate updated DOCX using the original bytes
            output_doc = auto_process_docx(
                io.BytesIO(original_file_bytes),
                st.session_state['mapping']
            )

            # ✅ Store everything properly
            st.session_state['generated_docx'] = output_doc.getvalue()
            st.session_state['doc_original_filename'] = uploaded_file.name
            st.session_state['doc_original_bytes'] = original_file_bytes
            
            # ✅ CRITICAL: Set template source to trigger col3 preview
            st.session_state['template_source'] = 'doc_saved'
            
            # Mark as processed
            st.session_state['doc_already_processed'] = True

            st.success("✅ DOCX processed successfully! Check preview on the right →")

        # ==============================
        # ✅ SAVE TEMPLATE SECTION
        # ==============================
        st.divider()
        st.subheader("💾 Save Template")

        # Get the template name (preserve it if already set)
        default_name = f"Doc_{st.session_state.get('doc_original_filename', 'template').split('.')[0]}"
        
        col1, col2 = st.columns([2, 1])

        with col1:
            doc_template_name = st.text_input(
                "Template Name:",
                value=st.session_state.get('doc_template_name_input', default_name),
                key="doc_template_name_input_field"
            )
            # Store in session to persist
            st.session_state['doc_template_name_input'] = doc_template_name

        with col2:
            st.write("")
            st.write("")
            if st.button("💾 Save Template", width='stretch', type="primary", key="save_doc_template_btn"):
                current_user = st.session_state.logged_in_user
                
                if 'doc_templates' not in st.session_state:
                    st.session_state.doc_templates = load_user_doc_templates(current_user)

                template_id = f"doc_{datetime.now().strftime('%Y%m%d%H%M%S')}"

                # ✅ Save with ORIGINAL file bytes
                st.session_state.doc_templates[template_id] = {
                    'name': doc_template_name,
                    'doc_data': st.session_state['doc_original_bytes'],
                    'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                    'original_filename': st.session_state['doc_original_filename']
                }

                try:
                    result = save_user_doc_templates(
                        current_user,
                        st.session_state.doc_templates
                    )

                    if result:
                        st.success(f"✅ Template '{doc_template_name}' saved successfully!")
                        
                        
                        # Set as active template
                        st.session_state.selected_doc_template_id = template_id
                        st.session_state.selected_doc_template = st.session_state.doc_templates[template_id]
                        
                        # Clear the upload states
                        st.session_state['doc_already_processed'] = False
                        st.session_state.pop('doc_template_name_input', None)
                        
                        # Wait a moment to show success
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error("❌ Failed to save template (function returned False)")
                
                except Exception as e:
                    st.error("❌ Save function crashed!")
                    st.exception(e)

        # ==============================
        # ✅ DOWNLOAD BUTTON
        # ==============================
        st.divider()
        st.download_button(
            label="⬇️ Download Updated Resume",
            data=st.session_state['generated_docx'],
            file_name=f"Resume_{final_data.get('name', 'User').replace(' ', '_')}.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            width='stretch',
            key="download_doc_from_upload"
        )

    except Exception as e:
        st.error(f"❌ Error processing document: {str(e)}")
        st.exception(e)



def process_ppt_upload(uploaded_file, final_data):
    """Process uploaded PowerPoint file and generate preview - FIXED WITH DOWNLOAD."""
    try:
        # Check if already processed to avoid reprocessing
        if not st.session_state.get('ppt_already_processed', False):
            # Store original file
            ppt_data = uploaded_file.getvalue()
            st.session_state.temp_ppt_data = ppt_data
            st.session_state.temp_ppt_filename = uploaded_file.name
            
            # Load presentation
            prs = Presentation(io.BytesIO(ppt_data))
            
            # Extract slide structure
            slide_texts = []
            for slide_idx, slide in enumerate(prs.slides):
                text_blocks = []
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.has_text_frame and shape.text.strip():
                        text_blocks.append({
                            "index": shape_idx,
                            "text": shape.text.strip(),
                            "position": {"x": shape.left, "y": shape.top}
                        })
                
                if text_blocks:
                    text_blocks.sort(key=lambda x: (x["position"]["y"], x["position"]["x"]))
                    slide_texts.append({
                        "slide_number": slide_idx + 1,
                        "text_blocks": text_blocks
                    })
            
            st.session_state.temp_ppt_structure = slide_texts
            
            # Generate content with user's resume data
            print("inside the ppt funtion")
            structured_slides = analyze_slide_structure(slide_texts)
            generated_sections = generate_ppt_sections(final_data, structured_slides)
            print("every ai model analysis is completed")
            # Create text elements list
            text_elements = []
            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.has_text_frame and shape.text.strip():
                        text_elements.append({
                            'slide': slide_idx + 1,
                            'shape': shape_idx,
                            'original_text': shape.text.strip(),
                            'shape_type': type(shape).__name__
                        })
            
            st.session_state.temp_ppt_text_elements = text_elements
            
            # Match content
            content_mapping, heading_shapes, basic_info_shapes = match_generated_to_original(
                text_elements, generated_sections, prs)
            print("compled content mapping")
            # Generate preview
            working_prs = Presentation(io.BytesIO(ppt_data))
            edits = {}
            
            for element in text_elements:
                key = f"{element['slide']}_{element['shape']}"
                if key not in heading_shapes:
                    edits[key] = content_mapping.get(key, element['original_text'])
            
            # Apply edits to working presentation
            for element in text_elements:
                key = f"{element['slide']}_{element['shape']}"
                if key not in heading_shapes and key in edits:
                    slide_idx = element['slide'] - 1
                    shape_idx = element['shape']
                    
                    if slide_idx < len(working_prs.slides):
                        slide = working_prs.slides[slide_idx]
                        if shape_idx < len(slide.shapes):
                            shape = slide.shapes[shape_idx]
                            if shape.has_text_frame:
                                clear_and_replace_text(shape, edits[key])
                                print(f"Updated Slide {element['slide']} Shape {shape_idx} with new text.")
            # Save output
            output = io.BytesIO()
            working_prs.save(output)
            output.seek(0)
            
            # ✅ CRITICAL: Store generated PPT and set template source
            st.session_state['generated_ppt'] = output.getvalue()
            st.session_state['template_source'] = 'ppt_saved'
            st.session_state.temp_ppt_edits = edits
            
            # Mark as processed
            st.session_state['ppt_already_processed'] = True
            
            st.success("✅ PowerPoint processed successfully! Check preview on the right →")
        
        # ==============================
        # ✅ SAVE TEMPLATE SECTION
        # ==============================
        st.divider()
        st.subheader("💾 Save Template")
        
        # Get the template name (preserve it if already set)
        default_name = f"PPT_{st.session_state.get('temp_ppt_filename', 'template').split('.')[0]}"
        
        col1, col2 = st.columns([2, 1])
        
        with col1:
            ppt_template_name = st.text_input(
                "Template Name:",
                value=st.session_state.get('ppt_template_name_input', default_name),
                key="ppt_template_name_input_field1"
            )
            # Store in session to persist
            st.session_state['ppt_template_name_input'] = ppt_template_name
        
        with col2:
            st.write("")
            st.write("")
            if st.button("💾 Save Template", width='stretch', type="primary", key="save_ppt_template_btn1"):
                current_user = st.session_state.logged_in_user
                
                if 'ppt_templates' not in st.session_state:
                    st.session_state.ppt_templates = {}
                
                template_id = f"ppt_{datetime.now().strftime('%Y%m%d%H%M%S')}"
                
                st.session_state.ppt_templates[template_id] = {
                    'name': ppt_template_name,
                    'ppt_data': st.session_state.temp_ppt_data,
                    'structure': st.session_state.temp_ppt_structure,
                    'edits': st.session_state.temp_ppt_edits,
                    'text_elements': st.session_state.temp_ppt_text_elements,
                    'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                    'original_filename': st.session_state.temp_ppt_filename,
                    'type': 'powerpoint'
                }
                
                if save_user_ppt_templates(current_user, st.session_state.ppt_templates):
                    st.success(f"✅ Template '{ppt_template_name}' saved successfully!")
                    
                    
                    # Set as active template
                    st.session_state.selected_ppt_template_id = template_id
                    st.session_state.selected_ppt_template = st.session_state.ppt_templates[template_id]
                    
                    # Clear the upload states
                    st.session_state['ppt_already_processed'] = False
                    st.session_state.pop('ppt_template_name_input', None)
                    
                    # Wait a moment to show success
                    time.sleep(1)
                    st.rerun()
                else:
                    st.error("Failed to save template. Please try again.")
        
        # ==============================
        # ✅ DOWNLOAD BUTTON
        # ==============================
        st.divider()
        st.download_button(
            label="⬇️ Download Updated Presentation",
            data=st.session_state['generated_ppt'],
            file_name=f"Resume_{final_data.get('name', 'User').replace(' ', '_')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            width='stretch',
            key="download_ppt_from_upload1"
        )
        
        return True
        
    except Exception as e:
        st.error(f"Error processing PowerPoint file: {str(e)}")
        st.exception(e)
        return False

def show_word_preview_inline(doc_bytes):
    """Display Word document preview inline with proper formatting."""
    from docx import Document
    import io
    
    try:
        doc = Document(io.BytesIO(doc_bytes))
        
        st.markdown("""
        <style>
        .doc-preview-inline {
            border: 2px solid #e0e0e0;
            border-radius: 12px;
            padding: 40px;
            background: white;
            max-height: 800px;
            overflow-y: auto;
            font-family: 'Calibri', 'Arial', sans-serif;
            box-shadow: 0 4px 12px rgba(0,0,0,0.1);
        }
        .doc-name { 
            font-size: 24px; 
            font-weight: bold; 
            margin-bottom: 5px;
            color: #1a1a1a;
        }
        .doc-title { 
            font-size: 14px; 
            margin-bottom: 5px;
            color: #4a4a4a;
        }
        .doc-contact { 
            font-size: 12px; 
            margin-bottom: 20px;
            color: #666;
        }
        .doc-heading { 
            font-size: 16px; 
            font-weight: bold; 
            margin: 20px 0 10px 0;
            border-bottom: 2px solid #333;
            padding-bottom: 5px;
            color: #1a1a1a;
        }
        .doc-text { 
            font-size: 11pt; 
            line-height: 1.6;
            margin: 8px 0;
            color: #333;
            white-space: pre-wrap;
        }
        </style>
        """, unsafe_allow_html=True)
        
        html_content = '<div class="doc-preview-inline">'
        para_count = 0
        
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            
            text = para.text.strip()
            
            if para_count == 0:
                html_content += f'<div class="doc-name">{text}</div>'
            elif para_count == 1:
                html_content += f'<div class="doc-title">{text}</div>'
            elif para_count == 2:
                html_content += f'<div class="doc-contact">{text}</div>'
            elif para.style.name.startswith('Heading') or (para.runs and para.runs[0].bold and len(text.split()) <= 10):
                html_content += f'<div class="doc-heading">{text}</div>'
            else:
                formatted_text = text.replace('\n', '<br>')
                html_content += f'<div class="doc-text">{formatted_text}</div>'
            
            para_count += 1
        
        html_content += '</div>'
        st.markdown(html_content, unsafe_allow_html=True)
        
    except Exception as e:
        st.error(f"Preview error: {str(e)}")


def show_ppt_preview_inline(ppt_bytes):
    """Display PowerPoint preview inline with proper formatting."""
    from pptx import Presentation
    import io
    
    try:
        prs = Presentation(io.BytesIO(ppt_bytes))
        
        st.markdown("""
        <style>
        .ppt-slide-container {
            border: 2px solid #e0e0e0;
            border-radius: 12px;
            padding: 20px;
            margin-bottom: 20px;
            background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%);
            box-shadow: 0 4px 12px rgba(0,0,0,0.1);
        }
        .ppt-slide-header {
            background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 10px 15px;
            border-radius: 8px;
            font-weight: bold;
            margin-bottom: 15px;
            font-size: 16px;
        }
        .ppt-content-box {
            background: white;
            padding: 15px;
            border-radius: 8px;
            margin: 10px 0;
            border-left: 4px solid #667eea;
            box-shadow: 0 2px 4px rgba(0,0,0,0.05);
        }
        .ppt-title-text {
            font-size: 18px;
            font-weight: bold;
            color: #2c3e50;
            margin-bottom: 10px;
        }
        .ppt-body-text {
            font-size: 14px;
            color: #34495e;
            line-height: 1.6;
            white-space: pre-wrap;
        }
        .ppt-bullet {
            color: #667eea;
            margin-right: 8px;
        }
        </style>
        """, unsafe_allow_html=True)
        
        for slide_idx, slide in enumerate(prs.slides):
            st.markdown(f"""
            <div class="ppt-slide-container">
                <div class="ppt-slide-header">📊 Slide {slide_idx + 1}</div>
            """, unsafe_allow_html=True)
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Check if it's a title
                    is_title = False
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.runs:
                                first_run = paragraph.runs[0]
                                if first_run.font.size and first_run.font.size.pt > 18:
                                    is_title = True
                                    break
                                if first_run.font.bold:
                                    is_title = True
                                    break
                    
                    # Format text
                    if '\n' in text:
                        lines = text.split('\n')
                        formatted_lines = [f'<span class="ppt-bullet">●</span> {line.strip()}' 
                                         for line in lines if line.strip()]
                        formatted_text = '<br>'.join(formatted_lines)
                    else:
                        formatted_text = text
                    
                    css_class = "ppt-title-text" if is_title else "ppt-body-text"
                    st.markdown(f"""
                    <div class="ppt-content-box">
                        <div class="{css_class}">{formatted_text}</div>
                    </div>
                    """, unsafe_allow_html=True)
            
            st.markdown("</div>", unsafe_allow_html=True)
            
    except Exception as e:
        st.error(f"Preview error: {str(e)}")


def save_uploaded_template(template_name, file_type):
    """Save the uploaded template to user's templates - CORRECTED VERSION."""
    current_user = st.session_state.logged_in_user
    
    try:
        if file_type == 'html':
            if 'uploaded_templates' not in st.session_state:
                st.session_state.uploaded_templates = {}
            
            template_id = f"html_{datetime.now().strftime('%Y%m%d%H%M%S')}"
            temp_config = st.session_state.get('temp_upload_config', {})
            
            st.session_state.uploaded_templates[template_id] = {
                'name': template_name,
                'css': temp_config.get('css', ''),
                'html': temp_config.get('html', ''),
                'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'original_filename': st.session_state.get('uploaded_file_name', 'unknown.html'),
                'type': 'html'
            }
            
            save_user_templates(current_user, st.session_state.uploaded_templates)
            return True
        
        elif file_type in ['docx', 'doc']:
            if 'doc_templates' not in st.session_state:
                st.session_state.doc_templates = {}
            
            template_id = f"doc_{datetime.now().strftime('%Y%m%d%H%M%S')}"
            structure = st.session_state.get('temp_doc_structure', [])
            
            st.session_state.doc_templates[template_id] = {
                'name': template_name,
                'doc_data': st.session_state.temp_doc_data,
                'structure': structure,
                'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'original_filename': st.session_state.uploaded_file_name,
                'type': 'word',
                'sections_detected': [s['section'] for s in structure]
            }
            
            save_user_doc_templates(current_user, st.session_state.doc_templates)
            return True
        
        elif file_type in ['pptx', 'ppt']:
            if 'ppt_templates' not in st.session_state:
                st.session_state.ppt_templates = {}
            
            template_id = f"ppt_{datetime.now().strftime('%Y%m%d%H%M%S')}"
            
            st.session_state.ppt_templates[template_id] = {
                'name': template_name,
                'ppt_data': st.session_state.temp_ppt_data,
                'structure': st.session_state.get('temp_ppt_structure', []),
                'edits': st.session_state.get('temp_ppt_edits', {}),
                'text_elements': st.session_state.get('temp_ppt_text_elements', []),
                'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'original_filename': st.session_state.uploaded_file_name,
                'type': 'powerpoint'
            }
            
            save_user_ppt_templates(current_user, st.session_state.ppt_templates)
            return True
        
        return False
        
    except Exception as e:
        st.error(f"Error saving template: {str(e)}")
        import traceback
        st.error(f"Details: {traceback.format_exc()}")
        return False


def use_uploaded_template_now(template_name, file_type):
    """Use the uploaded template immediately without saving."""
    try:
        if file_type == 'html':
            resume_data = generate_resume_for_template()
            
            css = st.session_state.temp_upload_css
            html_content = generate_generic_html(resume_data)
            
            full_html = f"""
            <div class="ats-page">
                {html_content}
            </div>
            """
            
            st.session_state.selected_template = template_name
            st.session_state.template_preview_html = full_html
            st.session_state.template_preview_css = css
            st.session_state.final_resume_data = resume_data
            st.session_state.show_upload_interface = False
            st.session_state.show_template_selector = False
            st.session_state.show_visual_editor = True
            st.rerun()
        
        elif file_type in ['docx', 'doc']:
            st.info("Word template functionality will redirect to download after use.")
            # Implement Word template usage logic here
        
        elif file_type in ['pptx', 'ppt']:
            st.info("PowerPoint template functionality will redirect to download after use.")
            # Implement PowerPoint template usage logic here
            
    except Exception as e:
        st.error(f"Error using template: {str(e)}")

def show_upload_interface():
    """Show the full-page upload interface with inline editing - ENHANCED VERSION."""
    final_data = st.session_state.get('final_resume_data') or st.session_state.get('enhanced_resume') or {}
    
    # Header with back button
    st.markdown('<div class="editor-header">', unsafe_allow_html=True)
    col1, col2 = st.columns([1, 5])
    with col1:
        if st.button("← Back to Templates", type="primary", use_container_width=True):
            st.session_state.show_upload_interface = False
            # Clear ALL temporary upload states INCLUDING processing flags
            for key in ['temp_upload_html', 'temp_upload_css', 'temp_upload_preview', 
                       'uploaded_file_name', 'uploaded_file_type', 'temp_upload_config',
                       'temp_doc_data', 'temp_doc_structure', 'temp_doc_filename',
                       'generated_docx', 'doc_structure', 'doc_replaced', 'doc_removed',
                       'temp_ppt_data', 'temp_ppt_structure', 'temp_ppt_edits', 
                       'temp_ppt_text_elements', 'generated_ppt', 'doc_original_bytes',
                       'doc_original_filename', 'mapping', 'template_text',
                       'doc_already_processed', 'ppt_already_processed',
                       'doc_template_name_input', 'ppt_template_name_input',
                       'show_upload_doc_editor', 'show_upload_ppt_editor']:
                if key in st.session_state:
                    del st.session_state[key]
            st.rerun()
    
    with col2:
        st.markdown("<h3 style='margin: 0; color: #212529;'>📤 Upload Custom Resume Template</h3>", 
                   unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Main upload section
    st.markdown("### 📁 Select Template File")
    
    uploaded_file = st.file_uploader(
        "Upload your template file",
        type=['html', 'pptx', 'docx', 'doc'],
        key="template_file_upload",
        help="Supported formats: HTML, Word (.docx, .doc), PowerPoint (.pptx)"
    )
    
    if uploaded_file is not None:
        st.success(f"✅ File uploaded: {uploaded_file.name}")
        
        file_type = uploaded_file.name.split('.')[-1].lower()
        st.session_state.uploaded_file_name = uploaded_file.name
        st.session_state.uploaded_file_type = file_type
        
        # Create two columns: controls + editor (col1) and preview (col2)
        col1, col2 = st.columns([1, 1.5])
        
        with col1:
            # ========== PROCESS FILES ==========
            if file_type in ['docx', 'doc'] and not st.session_state.get('doc_already_processed'):
                with st.spinner("Processing template..."):
                    process_word_upload(uploaded_file, final_data)
            elif file_type in ['pptx', 'ppt'] and not st.session_state.get('ppt_already_processed'):
                with st.spinner("Processing template..."):
                    process_ppt_upload(uploaded_file, final_data)
            elif file_type == 'html':
                with st.spinner("Processing template..."):
                    process_html_upload(uploaded_file)
            
            # ========== WORD DOCUMENT SECTION ==========
            if file_type in ['docx', 'doc'] and st.session_state.get('doc_already_processed'):
                st.markdown("---")
                st.markdown("### 📝 Word Document Controls")
                
                # Template name input
                default_name = f"Doc_{st.session_state.get('doc_original_filename', 'template').split('.')[0]}"
                doc_template_name = st.text_input(
                    "Template Name:",
                    value=st.session_state.get('doc_template_name_input', default_name),
                    key="doc_template_name_input_id"
                )
                st.session_state['doc_template_name_input'] = doc_template_name
                
                # Toggle editor button
                st.markdown("---")
                current_state = st.session_state.get('show_upload_doc_editor', False)
                button_label = "✏️ Hide Editor" if current_state else "✏️ Edit Document Mapping"
                
                if st.button(button_label, use_container_width=True, type="primary" if not current_state else "secondary", key="toggle_upload_doc_editor"):
                    st.session_state.show_upload_doc_editor = not current_state
                    st.rerun()
                
                # Show inline editor if enabled
                if st.session_state.get('show_upload_doc_editor', False):
                    st.markdown("---")
                    show_upload_doc_editor()
                
                # Save and Download buttons
                st.markdown("---")
                col_save, col_download = st.columns(2)
                
                with col_save:
                    if st.button("💾 Save Template", type="primary", use_container_width=True, key="save_doc_template_btn1"):
                        current_user = st.session_state.logged_in_user
                        
                        if 'doc_templates' not in st.session_state:
                            st.session_state.doc_templates = load_user_doc_templates(current_user)

                        template_id = f"doc_{datetime.now().strftime('%Y%m%d%H%M%S')}"

                        st.session_state.doc_templates[template_id] = {
                            'name': doc_template_name,
                            'doc_data': st.session_state['doc_original_bytes'],
                            'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                            'original_filename': st.session_state['doc_original_filename']
                        }

                        try:
                            result = save_user_doc_templates(current_user, st.session_state.doc_templates)

                            if result:
                                st.success(f"✅ Template '{doc_template_name}' saved successfully!")
                                
                                
                                st.session_state.selected_doc_template_id = template_id
                                st.session_state.selected_doc_template = st.session_state.doc_templates[template_id]
                                
                                st.session_state['doc_already_processed'] = False
                                st.session_state.pop('doc_template_name_input', None)
                                st.session_state.pop('show_upload_doc_editor', None)
                                
                                time.sleep(1.5)
                                st.session_state.show_upload_interface = False
                                st.rerun()
                            else:
                                st.error("❌ Failed to save template")
                        except Exception as e:
                            st.error("❌ Save function crashed!")
                            st.exception(e)
                
                with col_download:
                    st.download_button(
                        label="⬇️ Download",
                        data=st.session_state['generated_docx'],
                        file_name=f"{doc_template_name}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True,
                        key="download_doc_from_upload1"
                    )
            
            # ========== POWERPOINT SECTION ==========
            elif file_type in ['pptx', 'ppt'] and st.session_state.get('ppt_already_processed'):
                st.markdown("---")
                st.markdown("### 📊 PowerPoint Controls")
                
                # Template name input
                default_name = f"PPT_{st.session_state.get('temp_ppt_filename', 'template').split('.')[0]}"
                ppt_template_name = st.text_input(
                    "Template Name:",
                    value=st.session_state.get('ppt_template_name_input', default_name),
                    key="ppt_template_name_input_field"
                )
                st.session_state['ppt_template_name_input'] = ppt_template_name
                
                # Toggle editor button
                st.markdown("---")
                current_state = st.session_state.get('show_upload_ppt_editor', False)
                button_label = "✏️ Hide Editor" if current_state else "✏️ Edit Presentation Content"
                
                if st.button(button_label, use_container_width=True, type="primary" if not current_state else "secondary", key="toggle_upload_ppt_editor"):
                    st.session_state.show_upload_ppt_editor = not current_state
                    st.rerun()
                
                # Show inline editor if enabled
                if st.session_state.get('show_upload_ppt_editor', False):
                    st.markdown("---")
                    show_upload_ppt_editor()
                
                # Save and Download buttons
                st.markdown("---")
                col_save, col_download = st.columns(2)
                
                with col_save:
                    if st.button("💾 Save Template", type="primary", use_container_width=True, key="save_ppt_template_btn2"):
                        current_user = st.session_state.logged_in_user
                        
                        if 'ppt_templates' not in st.session_state:
                            st.session_state.ppt_templates = {}
                        
                        template_id = f"ppt_{datetime.now().strftime('%Y%m%d%H%M%S')}"
                        
                        st.session_state.ppt_templates[template_id] = {
                            'name': ppt_template_name,
                            'ppt_data': st.session_state.temp_ppt_data,
                            'structure': st.session_state.temp_ppt_structure,
                            'edits': st.session_state.temp_ppt_edits,
                            'text_elements': st.session_state.temp_ppt_text_elements,
                            'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                            'original_filename': st.session_state.temp_ppt_filename,
                            'type': 'powerpoint'
                        }
                        
                        if save_user_ppt_templates(current_user, st.session_state.ppt_templates):
                            st.success(f"✅ Template '{ppt_template_name}' saved successfully!")
                            
                            
                            st.session_state.selected_ppt_template_id = template_id
                            st.session_state.selected_ppt_template = st.session_state.ppt_templates[template_id]
                            
                            st.session_state['ppt_already_processed'] = False
                            st.session_state.pop('ppt_template_name_input', None)
                            st.session_state.pop('show_upload_ppt_editor', None)
                            
                            time.sleep(1.5)
                            st.session_state.show_upload_interface = False
                            st.rerun()
                        else:
                            st.error("Failed to save template. Please try again.")
                
                with col_download:
                    st.download_button(
                        label="⬇️ Download",
                        data=st.session_state['generated_ppt'],
                        file_name=f"{ppt_template_name}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                        key="download_ppt_from_upload"
                    )
            
            # ========== HTML SECTION ==========
            elif file_type == 'html':
                st.markdown("---")
                st.markdown("### 📄 HTML Template Controls")
                
                template_name = st.text_input(
                    "Template Name:",
                    value=f"Custom_{uploaded_file.name.split('.')[0]}",
                    key="html_template_name"
                )
                
                col_save, col_download = st.columns(2)
                
                with col_save:
                    if st.button("💾 Save", type="primary", use_container_width=True, key="save_html"):
                        if save_uploaded_template(template_name, file_type):
                            st.success("✅ Saved!")
                            
                            time.sleep(1)
                            st.session_state.show_upload_interface = False
                            st.rerun()
                
                with col_download:
                    if st.session_state.get('temp_upload_preview'):
                        st.download_button(
                            label="📥 Download",
                            data=st.session_state.temp_upload_preview,
                            file_name=f"{template_name}.html",
                            mime="text/html",
                            use_container_width=True,
                            type="secondary"
                        )
        
        # ========== PREVIEW COLUMN ==========
        with col2:
            st.markdown("---")
            st.markdown("### 👁️ Template Preview")
            
            if file_type == 'html' and st.session_state.get('temp_upload_preview'):
                st.components.v1.html(st.session_state.temp_upload_preview, height=800, scrolling=True)
            
            elif file_type in ['docx', 'doc'] and st.session_state.get('generated_docx'):
                try:
                    preview_html = docx_to_html_preview(io.BytesIO(st.session_state['generated_docx']))
                    st.components.v1.html(preview_html, height=800, scrolling=True)
                except Exception as e:
                    st.error(f"Preview error: {str(e)}")
            
            elif file_type in ['pptx', 'ppt'] and st.session_state.get('generated_ppt'):
                show_ppt_preview_inline(st.session_state['generated_ppt'])
    else:
        # Clear processing flags when no file is uploaded
        st.session_state.pop('doc_already_processed', None)
        st.session_state.pop('ppt_already_processed', None)


def show_upload_doc_editor():
    """Inline Word document editor for upload interface."""
    st.markdown("#### ✏️ Document Mapping Editor")
    
    if 'mapping' not in st.session_state or not st.session_state.get('mapping'):
        st.warning("⚠️ No mapping data available.")
        return
    
    mapping = st.session_state['mapping'].copy()
    
    st.info("💡 Edit how your resume data maps to the document template")
    
    # Show mappings
    for idx, (template_key, resume_value) in enumerate(mapping.items()):
        with st.container():
            col1, col2, col3 = st.columns([2, 3, 1])
            
            with col1:
                st.text_input(
                    "Template",
                    value=template_key,
                    disabled=True,
                    key=f"upload_doc_key_{idx}",
                    label_visibility="collapsed"
                )
            
            with col2:
                new_value = st.text_area(
                    "Value",
                    value=resume_value,
                    height=60,
                    key=f"upload_doc_value_{idx}",
                    label_visibility="collapsed"
                )
                mapping[template_key] = new_value
            
            with col3:
                if st.button("🗑️", key=f"upload_doc_delete_{idx}", use_container_width=True):
                    del mapping[template_key]
                    st.session_state['mapping'] = mapping
                    st.rerun()
    
    st.session_state['mapping'] = mapping
    
    # Apply button
    if st.button("✨ Apply Changes", type="primary", use_container_width=True, key="apply_upload_doc_changes"):
        try:
            import io
            
            original_bytes = st.session_state['doc_original_bytes']
            
            with st.spinner("Applying changes..."):
                output_doc = auto_process_docx(
                    io.BytesIO(original_bytes),
                    st.session_state['mapping']
                )
            
            st.session_state['generated_docx'] = output_doc.getvalue()
            st.success("✅ Changes applied!")
            st.rerun()
            
        except Exception as e:
            st.error(f"❌ Error: {str(e)}")


def show_upload_ppt_editor():
    """Inline PowerPoint editor for upload interface."""
    st.markdown("#### ✏️ Presentation Content Editor")
    
    if 'temp_ppt_edits' not in st.session_state or 'temp_ppt_text_elements' not in st.session_state:
        st.warning("⚠️ No content data available.")
        return
    
    edits = st.session_state['temp_ppt_edits'].copy()
    text_elements = st.session_state['temp_ppt_text_elements']
    
    st.info("💡 Edit the content for each slide")
    
    # Group by slide
    slides_dict = {}
    for element in text_elements:
        slide_num = element['slide']
        if slide_num not in slides_dict:
            slides_dict[slide_num] = []
        slides_dict[slide_num].append(element)
    
    # Display by slide
    for slide_num in sorted(slides_dict.keys()):
        with st.expander(f"📊 Slide {slide_num}", expanded=(slide_num==1)):
            for element in slides_dict[slide_num]:
                key = f"{element['slide']}_{element['shape']}"
                current_text = edits.get(key, element['original_text'])
                
                col1, col2 = st.columns([4, 1])
                
                with col1:
                    new_text = st.text_area(
                        f"Shape {element['shape']}",
                        value=current_text,
                        height=80,
                        key=f"upload_ppt_{key}"
                    )
                    edits[key] = new_text
                
                with col2:
                    if st.button("🔄", key=f"upload_ppt_reset_{key}", use_container_width=True):
                        edits[key] = element['original_text']
                        st.rerun()
    
    st.session_state['temp_ppt_edits'] = edits
    
    # Apply button
    if st.button("✨ Apply Changes", type="primary", use_container_width=True, key="apply_upload_ppt_changes"):
        try:
            import io
            from pptx import Presentation
            
            ppt_data = st.session_state['temp_ppt_data']
            working_prs = Presentation(io.BytesIO(ppt_data))
            
            with st.spinner("Applying changes..."):
                for element in text_elements:
                    key = f"{element['slide']}_{element['shape']}"
                    if key in edits:
                        slide_idx = element['slide'] - 1
                        shape_idx = element['shape']
                        
                        if slide_idx < len(working_prs.slides):
                            slide = working_prs.slides[slide_idx]
                            if shape_idx < len(slide.shapes):
                                shape = slide.shapes[shape_idx]
                                if shape.has_text_frame:
                                    clear_and_replace_text(shape, edits[key])
            
            output = io.BytesIO()
            working_prs.save(output)
            output.seek(0)
            
            st.session_state['generated_ppt'] = output.getvalue()
            st.success("✅ Changes applied!")
            st.rerun()
            
        except Exception as e:
            st.error(f"❌ Error: {str(e)}")





def regenerate_live_preview():
    """Regenerate the live preview with updated resume data."""
    if not st.session_state.get('selected_template'):
        return None, None
    
    # Get current resume data from enhanced_resume
    resume_data = st.session_state.get('enhanced_resume', {})
    
    # Get template configuration
    template_config = SYSTEM_TEMPLATES.get(st.session_state.selected_template)
    if not template_config:
        return None, None
    
    # Generate fresh HTML and CSS
    selected_color = ATS_COLORS["Professional Blue (Default)"]
    css = load_css_template(
        template_config['css_template'],
        selected_color
    )
    html_content = template_config['html_generator'](resume_data)
    
    # Create full HTML
    full_html = f"""
    <div class="ats-page">
        {html_content}
    </div>
    """
    
    return full_html, css



def get_standard_keys():
    """Return set of standard resume keys."""
    return {
        "name", "email", "phone", "location", "url", "summary", "job_title",
        "education", "experience", "skills", "projects", "certifications", 
        "total_experience_count", "input_method", "achievements", "languages"
    }

def clean_resume_data(data):
    """Clean and prepare resume data for templates."""
    if not data:
        return {}
    
    cleaned = data.copy()
    
    # Remove empty values
    for key in list(cleaned.keys()):
        if cleaned[key] is None or cleaned[key] == "":
            del cleaned[key]
        elif isinstance(cleaned[key], list) and len(cleaned[key]) == 0:
            del cleaned[key]
        elif isinstance(cleaned[key], dict) and len(cleaned[key]) == 0:
            del cleaned[key]
    
    # Ensure required sections exist
    if 'experience' not in cleaned:
        cleaned['experience'] = []
    if 'education' not in cleaned:
        cleaned['education'] = []
    if 'skills' not in cleaned:
        cleaned['skills'] = []
    
    return cleaned

def generate_resume_for_template():
    """Generate enhanced resume data for template."""
    # Return enhanced resume if it exists
    if st.session_state.enhanced_resume:
        return clean_resume_data(st.session_state.enhanced_resume)
    
    # Check if we should regenerate (only if job description exists)
    jd_data = st.session_state.get('job_description')
    has_jd = jd_data is not None and isinstance(jd_data, dict) and len(jd_data) > 0
    
    # Only regenerate if JD exists and edit mode is off
    if has_jd and should_regenerate_resume() and not st.session_state.get('edit_toggle', False):
        try:
            generate_enhanced_resume(resume_data,jd_data)
        except Exception as e:
            st.warning(f"Could not enhance resume: {str(e)}")
            # Fall back to original resume
            pass
    
    # Get resume data - enhanced if available, otherwise original
    resume_data = st.session_state.get('enhanced_resume') or st.session_state.get('final_resume_data') or user_resume
    # st.write(resume_data)
    if not resume_data:
        st.error("No resume data found. Please create a resume first.")
        return {}
    
    resume_data = clean_resume_data(resume_data)
    
    
    return resume_data

def save_custom_sections():
    """Save custom section edits."""
    if 'enhanced_resume' in st.session_state:
        data = st.session_state.enhanced_resume
        standard_keys = get_standard_keys()
        
        for key in list(data.keys()):
            if key not in standard_keys and isinstance(data.get(key), str):
                edit_key = f"edit_custom_{key}"
                if edit_key in st.session_state:
                    data[key] = st.session_state[edit_key].strip()
        
        st.session_state.enhanced_resume = data

VISUAL_EDITOR_HTML = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }

        html, body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            background: #F8F9FA;
            overflow: visible !important;
            height: auto !important;
            min-height: 100vh;
        }

        body {
            padding: 1rem;
        }

        .canvas-container {
            width: 100%;
            background: #F8F9FA;
            display: flex;
            align-items: flex-start;
            justify-content: center;
            padding: 1rem;
            min-height: 100vh;
        }

        .canvas {
            width: 8.5in;
            min-height: 11in;
            background: white;
            box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
            position: relative;
            overflow: visible !important;
            margin: 0 auto;
            padding: 0.5in;
        }

        /* Make elements editable */
        .editable {
            cursor: pointer;
            transition: all 0.2s ease;
            position: relative;
        }

        .editable:hover {
            background: rgba(255, 107, 53, 0.05) !important;
            outline: 2px dashed rgba(255, 107, 53, 0.3) !important;
            outline-offset: 2px;
        }

        .editable:active {
            background: rgba(255, 107, 53, 0.1) !important;
        }

        .editable-section {
            border-radius: 4px;
            padding: 4px;
        }

        /* Inline Editor Modal */
        .inline-editor-modal {
            display: none;
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            background: rgba(0, 0, 0, 0.5);
            z-index: 10000;
            justify-content: center;
            align-items: center;
        }

        .inline-editor-content {
            background: white;
            border-radius: 12px;
            padding: 1.5rem;
            width: 90%;
            max-width: 600px;
            max-height: 80vh;
            overflow-y: auto;
            box-shadow: 0 20px 40px rgba(0, 0, 0, 0.2);
            border: 2px solid #FF6B35;
        }

        .inline-editor-header {
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-bottom: 1.5rem;
            padding-bottom: 1rem;
            border-bottom: 2px solid #F8F9FA;
        }

        .inline-editor-title {
            font-size: 1.3rem;
            font-weight: 700;
            color: #FF6B35;
        }

        .inline-editor-close {
            background: none;
            border: none;
            font-size: 1.5rem;
            color: #495057;
            cursor: pointer;
            width: 32px;
            height: 32px;
            display: flex;
            align-items: center;
            justify-content: center;
            border-radius: 6px;
            transition: all 0.2s ease;
        }

        .inline-editor-close:hover {
            background: #F8F9FA;
            color: #FF6B35;
        }

        .inline-editor-body {
            margin-bottom: 1.5rem;
        }

        .inline-editor-textarea {
            width: 100%;
            min-height: 200px;
            padding: 1rem;
            border: 2px solid #E9ECEF;
            border-radius: 8px;
            font-family: inherit;
            font-size: 1rem;
            line-height: 1.5;
            resize: vertical;
            transition: all 0.3s ease;
        }

        .inline-editor-textarea:focus {
            outline: none;
            border-color: #FF6B35;
            box-shadow: 0 0 0 3px rgba(255, 107, 53, 0.1);
        }

        .inline-editor-actions {
            display: flex;
            gap: 1rem;
            justify-content: flex-end;
        }

        .inline-editor-btn {
            padding: 0.75rem 1.5rem;
            border: none;
            border-radius: 8px;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.3s ease;
        }

        .inline-editor-save {
            background: linear-gradient(135deg, #FF6B35 0%, #FF8B5C 100%);
            color: white;
        }

        .inline-editor-save:hover {
            transform: translateY(-2px);
            box-shadow: 0 6px 20px rgba(255, 107, 53, 0.3);
        }

        .inline-editor-cancel {
            background: #F8F9FA;
            color: #495057;
            border: 2px solid #E9ECEF;
        }

        .inline-editor-cancel:hover {
            background: #E9ECEF;
        }

        /* Remove any hidden elements */
        [style*="display: none"],
        [style*="visibility: hidden"],
        .empty-section {
            display: none !important;
        }
        
        /* Ensure content is visible */
        .ats-page {
            width: 100% !important;
            height: auto !important;
            min-height: 11in !important;
            padding: 0.5in !important;
            overflow: visible !important;
        }
        
        /* Make sure all content shows */
        div, p, h1, h2, h3, h4, li, span {
            overflow: visible !important;
            white-space: normal !important;
        }
        
    </style>
    <style id="template-css"></style>
</head>
<body>
    <!-- Inline Editor Modal -->
    <div class="inline-editor-modal" id="inlineEditor">
        <div class="inline-editor-content">
            <div class="inline-editor-header">
                <div class="inline-editor-title" id="editorTitle">Edit Section</div>
                <button class="inline-editor-close" onclick="closeEditor()">×</button>
            </div>
            <div class="inline-editor-body">
                <textarea class="inline-editor-textarea" id="editorTextarea" placeholder="Enter your content here..."></textarea>
            </div>
            <div class="inline-editor-actions">
                <button class="inline-editor-btn inline-editor-cancel" onclick="closeEditor()">Cancel</button>
                <button class="inline-editor-btn inline-editor-save" onclick="saveContent()">Save Changes</button>
            </div>
        </div>
    </div>

    <div class="canvas-container">
        <div class="canvas" id="canvas">
            <!-- Template content will be injected here -->
        </div>
    </div>

    <script>
        let currentEditElement = null;
        let currentEditPath = '';
        
        // Load template content
        function loadTemplateContent() {
            const templateData = window.templateContent || '';
            const templateCss = window.templateCss || '';
            
            if (templateData) {
                document.getElementById('canvas').innerHTML = templateData;
            }
            
            if (templateCss) {
                document.getElementById('template-css').innerHTML = templateCss;
            }
            
            // Make all text elements editable
            setTimeout(() => {
                makeElementsEditable();
                removeEmptySections();
            }, 100);
        }
        
        function makeElementsEditable() {
            // Select elements that should be editable
            const editableSelectors = [
                'h1', 'h2', 'h3', 'h4', 'h5', 'h6',
                'p', 'span', 'li', 'div[class*="name"]',
                'div[class*="title"]', 'div[class*="summary"]',
                'div[class*="description"]', 'div[class*="skill"]'
            ];
            
            editableSelectors.forEach(selector => {
                document.querySelectorAll(`#canvas ${selector}`).forEach(el => {
                    if (el.textContent.trim() && !el.hasAttribute('data-non-editable')) {
                        el.classList.add('editable', 'editable-section');
                        el.setAttribute('title', 'Click to edit');
                        el.addEventListener('click', (e) => {
                            e.stopPropagation();
                            openEditor(el);
                        });
                    }
                });
            });
        }
        
        function removeEmptySections() {
            // Remove completely empty sections
            const sections = document.querySelectorAll('#canvas > div, #canvas section, #canvas .section');
            sections.forEach(section => {
                if (section.textContent.trim() === '' || section.innerHTML.trim() === '<br>') {
                    section.style.display = 'none';
                }
            });
        }
        
        
        
        // Initialize
        document.addEventListener('DOMContentLoaded', loadTemplateContent);
        window.addEventListener('load', loadTemplateContent);
        
        // Listen for messages from parent
        window.addEventListener('message', (event) => {
            if (event.data.type === 'update_content') {
                document.getElementById('canvas').innerHTML = event.data.content;
                const styleEl = document.getElementById('template-css');
                if (styleEl && event.data.css) {
                    styleEl.innerHTML = event.data.css;
                }
                setTimeout(() => {
                    makeElementsEditable();
                    removeEmptySections();
                }, 100);
            }
            
            if (event.data.type === 'refresh_editor') {
                loadTemplateContent();
            }
        });
        
        // Close editor on ESC key
        document.addEventListener('keydown', (e) => {
            if (e.key === 'Escape') {
                closeEditor();
            }
        });
        
        // Close editor when clicking outside
        document.getElementById('inlineEditor').addEventListener('click', (e) => {
            if (e.target.classList.contains('inline-editor-modal')) {
                closeEditor();
            }
        });
    </script>
</body>
</html>
"""




def render_template_card(template_name, template_config, template_categories):
    """Render a single template card."""
    template_type = template_categories.get(template_name, "Template")
    
    st.markdown(f"""
    <div class="template-card">
        <div class="template-preview">
            <div style="width: 85%; height: 90%; background: white; 
                 border-radius: 4px; padding: 1rem; box-shadow: 0 2px 8px rgba(0,0,0,0.1);">
                <div style="width: 60%; height: 10px; background: #FF6B35; 
                     border-radius: 2px; margin-bottom: 0.5rem;"></div>
                <div style="width: 40%; height: 6px; background: #E9ECEF; 
                     border-radius: 2px; margin-bottom: 1rem;"></div>
                <div style="width: 80%; height: 6px; background: #E9ECEF; 
                     border-radius: 2px; margin-bottom: 0.5rem;"></div>
            </div>
            <div class="template-badge">{template_type}</div>
        </div>
        <div style="padding: 1rem;">
            <div style="font-size: 1rem; font-weight: 600; color: #212529; 
                 margin-bottom: 0.5rem;">📄 {template_name}</div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    if st.button("Use Template", key=f"use_{template_name}", type="primary", use_container_width=True):
        # Generate resume data
        clear_template_state()
        resume_data = generate_resume_for_template()
        
        # Generate template HTML
        selected_color = ATS_COLORS["Professional Blue (Default)"]
        css = load_css_template(template_config['css_template'], selected_color)
        html_content = template_config['html_generator'](resume_data)
        
        # Create full HTML
        full_html = f"""
        <div class="ats-page">
            {html_content}
        </div>
        """
        
        # Store in session state
        st.session_state.selected_template = template_name
        st.session_state.template_preview_html = full_html
        st.session_state.template_preview_css = css
        st.session_state.final_resume_data = resume_data
        st.session_state['template_source'] = 'system'
        st.session_state.show_template_selector = False
        st.session_state.show_visual_editor = True
        st.rerun()

def show_upload_modal():
    """Show the upload modal for custom templates."""
    st.markdown("---")
    st.markdown("### 📤 Upload Custom Template")
    
    # Create tabs for different file types
    tab1, tab2, tab3 = st.tabs(["📄 HTML Template", "📝 Word Document", "📊 PowerPoint"])
    
    with tab1:
        show_html_upload()
    
    with tab2:
        show_word_upload()
    
    with tab3:
        show_powerpoint_upload()
    
    # Close button
    col1, col2, col3 = st.columns([1, 1, 1])
    with col2:
        if st.button("❌ Close Upload", type="secondary", use_container_width=True):
            st.session_state.show_upload_modal = False
            st.rerun()


def show_html_upload():
    """Handle HTML template upload."""
    st.markdown("Upload an HTML file with embedded CSS styles")
    
    uploaded_file = st.file_uploader(
        "Choose HTML file",
        type=['html'],
        key="html_upload",
        help="Upload an HTML file with your custom resume template"
    )
    
    if uploaded_file is not None:
        try:
            import chardet
            raw_data = uploaded_file.read()
            detected = chardet.detect(raw_data)
            encoding = detected["encoding"] or "utf-8"
            content = raw_data.decode(encoding, errors="ignore")
            
            # Extract template structure (you'll need to implement this based on your needs)
            from templates.templateconfig import extract_template_from_html
            parsed_template = extract_template_from_html(content)
            
            # FIXED: Changed key from 'html_template_name' to be more unique
            template_name = st.text_input(
                "Template Name:",
                value=f"Custom_{uploaded_file.name.split('.')[0]}",
                key="modal_html_template_name"  # Changed key here
            )
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("💾 Save Template", type="primary", use_container_width=True, key="save_html_template_btn"):
                    save_html_template(template_name, parsed_template, uploaded_file.name)
                    st.success(f"✅ Template '{template_name}' saved!")
                    
                    time.sleep(1)
                    st.session_state.show_upload_modal = False
                    st.rerun()
            
            with col2:
                if st.button("👁️ Preview & Use", type="secondary", use_container_width=True, key="preview_html_template_btn"):
                    use_uploaded_html_template(template_name, parsed_template)
            
            # Show preview
            st.markdown("#### Preview")
            resume_data = st.session_state.get('final_resume_data') or {}
            preview_html = f"""
                <style>{parsed_template.get('css', '')}</style>
                <div class="ats-page">{generate_generic_html(resume_data)}</div>
            """
            st.components.v1.html(preview_html, height=600, scrolling=True)
            
        except Exception as e:
            st.error(f"Error processing HTML file: {str(e)}")


def show_word_upload():
    """Handle Word document upload."""
    st.markdown("Upload a Word document (.docx) template")
    
    uploaded_file = st.file_uploader(
        "Choose Word file",
        type=['docx', 'doc'],
        key="word_upload",
        help="Upload a Word document with your custom resume template"
    )
    
    if uploaded_file is not None:
        try:
            from templates.templateconfig import extract_document_structure, replace_content
            
            # Extract structure
            uploaded_file.seek(0)
            doc, structure = extract_document_structure(uploaded_file)
            
            # Store template data
            uploaded_file.seek(0)
            doc_data = uploaded_file.read()
            
            # FIXED: Changed key to be unique
            template_name = st.text_input(
                "Template Name:",
                value=f"Word_{uploaded_file.name.split('.')[0]}",
                key="modal_word_template_name"  # Changed key here
            )
            
            # Get resume data
            resume_data = st.session_state.get('final_resume_data') or {}
            
            # Generate preview
            uploaded_file.seek(0)
            doc, _ = extract_document_structure(uploaded_file)
            output, replaced, removed = replace_content(doc, structure, resume_data)
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("💾 Save Template", type="primary", use_container_width=True, key="save_word_template_btn"):
                    save_word_template(template_name, doc_data, structure, uploaded_file.name)
                    st.success(f"✅ Template '{template_name}' saved!")
                    
                    time.sleep(1)
                    st.session_state.show_upload_modal = False
                    st.rerun()
            
            with col2:
                if st.button("📥 Download Preview", type="secondary", use_container_width=True, key="download_word_preview_btn"):
                    st.download_button(
                        label="Download Word Document",
                        data=output.getvalue(),
                        file_name=f"{template_name}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )
            
            # Show preview
            st.markdown("#### Preview")
            show_word_preview(output.getvalue())
            
        except Exception as e:
            st.error(f"Error processing Word file: {str(e)}")


def show_powerpoint_upload():
    """Handle PowerPoint upload."""
    st.markdown("Upload a PowerPoint presentation (.pptx) template")
    
    uploaded_file = st.file_uploader(
        "Choose PowerPoint file",
        type=['pptx', 'ppt'],
        key="ppt_upload",
        help="Upload a PowerPoint presentation with your custom resume template"
    )
    
    if uploaded_file is not None:
        try:
            import io
            from pptx import Presentation
            from templates.templateconfig import analyze_slide_structure, generate_ppt_sections
            
            # Load presentation
            ppt_data = uploaded_file.getvalue()
            prs = Presentation(io.BytesIO(ppt_data))
            
            # Extract structure
            slide_texts = []
            for slide_idx, slide in enumerate(prs.slides):
                text_blocks = []
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.has_text_frame and shape.text.strip():
                        text_blocks.append({
                            "index": shape_idx,
                            "text": shape.text.strip(),
                            "position": {"x": shape.left, "y": shape.top}
                        })
                
                if text_blocks:
                    text_blocks.sort(key=lambda x: (x["position"]["y"], x["position"]["x"]))
                    slide_texts.append({
                        "slide_number": slide_idx + 1,
                        "text_blocks": text_blocks
                    })
            
            # FIXED: Changed key to be unique
            template_name = st.text_input(
                "Template Name:",
                value=f"PPT_{uploaded_file.name.split('.')[0]}",
                key="modal_ppt_template_name"  # Changed key here
            )
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("💾 Save Template", type="primary", use_container_width=True, key="save_ppt_template_btn3"):
                    save_ppt_template(template_name, ppt_data, slide_texts, uploaded_file.name)
                    st.success(f"✅ Template '{template_name}' saved!")
                    
                    time.sleep(1)
                    st.session_state.show_upload_modal = False
                    st.rerun()
            
            with col2:
                if st.button("👁️ Preview", type="secondary", use_container_width=True, key="preview_ppt_btn"):
                    st.session_state.show_ppt_preview = True
            
            # Show preview if requested
            if st.session_state.get('show_ppt_preview', False):
                st.markdown("#### Preview")
                show_ppt_preview(prs)
            
        except Exception as e:
            st.error(f"Error processing PowerPoint file: {str(e)}")


def save_html_template(name, parsed_template, filename):
    """Save HTML template to user's templates."""
    if 'uploaded_templates' not in st.session_state:
        st.session_state.uploaded_templates = {}
    
    template_id = f"html_{datetime.now().strftime('%Y%m%d%H%M%S')}"
    st.session_state.uploaded_templates[template_id] = {
        'name': name,
        'css': parsed_template.get('css', ''),
        'html': parsed_template.get('html', ''),
        'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
        'original_filename': filename,
        'type': 'html'
    }
    
    # Save to database (implement this based on your storage)
    # save_user_templates(st.session_state.logged_in_user, st.session_state.uploaded_templates)


def save_word_template(name, doc_data, structure, filename):
    """Save Word template to user's templates."""
    if 'doc_templates' not in st.session_state:
        st.session_state.doc_templates = {}
    
    template_id = f"doc_{datetime.now().strftime('%Y%m%d%H%M%S')}"
    st.session_state.doc_templates[template_id] = {
        'name': name,
        'doc_data': doc_data,
        'structure': structure,
        'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
        'original_filename': filename,
        'type': 'word'
    }
    
    # Save to database
    # save_user_doc_templates(st.session_state.logged_in_user, st.session_state.doc_templates)


def save_ppt_template(name, ppt_data, slide_texts, filename):
    """Save PowerPoint template to user's templates."""
    if 'ppt_templates' not in st.session_state:
        st.session_state.ppt_templates = {}
    
    template_id = f"ppt_{datetime.now().strftime('%Y%m%d%H%M%S')}"
    st.session_state.ppt_templates[template_id] = {
        'name': name,
        'ppt_data': ppt_data,
        'slide_texts': slide_texts,
        'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
        'original_filename': filename,
        'type': 'powerpoint'
    }
    
    # Save to database
    # save_user_ppt_templates(st.session_state.logged_in_user, st.session_state.ppt_templates)


def show_word_preview(doc_bytes):
    """Display Word document preview."""
    from docx import Document
    import io
    
    doc = Document(io.BytesIO(doc_bytes))
    
    st.markdown("""
    <style>
    .doc-preview {
        border: 2px solid #e0e0e0;
        border-radius: 10px;
        padding: 40px;
        background: white;
        max-height: 600px;
        overflow-y: auto;
    }
    </style>
    """, unsafe_allow_html=True)
    
    html_content = '<div class="doc-preview">'
    for para in doc.paragraphs:
        if para.text.strip():
            html_content += f'<p>{para.text}</p>'
    html_content += '</div>'
    
    st.markdown(html_content, unsafe_allow_html=True)


def show_ppt_preview(prs):
    """Display PowerPoint preview."""
    st.markdown("""
    <style>
    .ppt-slide {
        border: 2px solid #e0e0e0;
        border-radius: 10px;
        padding: 20px;
        margin-bottom: 20px;
        background: white;
    }
    </style>
    """, unsafe_allow_html=True)
    
    for slide_idx, slide in enumerate(prs.slides):
        st.markdown(f'<div class="ppt-slide"><h4>Slide {slide_idx + 1}</h4>', unsafe_allow_html=True)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                st.markdown(f'<p>{shape.text}</p>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)


def use_uploaded_html_template(name, parsed_template):
    """Use the uploaded HTML template immediately."""
    resume_data = generate_resume_for_template()
    
    css = parsed_template.get('css', '')
    html_content = generate_generic_html(resume_data)
    
    full_html = f"""
    <div class="ats-page">
        {html_content}
    </div>
    """
    
    st.session_state.selected_template = name
    st.session_state.template_preview_html = full_html
    st.session_state.template_preview_css = css
    st.session_state.final_resume_data = resume_data
    st.session_state.show_template_selector = False
    st.session_state.show_visual_editor = True
    st.session_state.show_upload_modal = False
    st.rerun()

def show_template_selector():
    """Show template selection gallery with custom upload option and saved templates."""
    # Check if we should show upload interface
   
    if st.session_state.get('show_upload_interface', False):
        show_upload_interface()
        return
    st.markdown("""
    <style>
    /* ================= HERO HEADER ================= */
    .ats-main-wrapper {
        background: linear-gradient(135deg, #fff9f5 0%, #ffffff 50%, #fff5f0 100%);
        padding: 60px 0 30px;
        position: center;
    }

    .ats-hero {
        text-align: center;
        max-width: 900px;
        margin: 0 auto;
    }

    .ats-hero-badge {
        display: inline-block;
        background: linear-gradient(135deg, #fff5f0 0%, #ffe8d6 100%);
        color: #e87532;
        padding: 6px 10px;
        border-radius: 999px;
        font-size: 12px;
        font-weight: 600;
        margin-bottom: 16px;
        border: 1px solid rgba(232, 117, 50, 0.25);
    }

    .ats-main-title {
        font-size: 2.6rem;
        font-weight: 800;
        color: #0f172a;
        margin-bottom: 12px;
        letter-spacing: -1px;
    }

    .ats-main-title .highlight {
        background: linear-gradient(135deg, #e87532 0%, #ff8c42 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }

    .ats-hero-description {
        font-size: 1rem;
        color: #64748b;
        line-height: 1.6;
        max-width: 620px;
        margin-left: 120px !important;
    }
    </style>
    """, unsafe_allow_html=True)

    st.markdown("""
    <div class="ats-main-wrapper">
        <div class="ats-hero">
            <div class="ats-hero-badge">Step 3 of 3</div>
            <h1 class="ats-main-title">Download your <span class="highlight">Resume</span></h1>
            <p class="ats-hero-description">
               Choose from our professionally designed templates or upload your own
            </p>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    # Search bar
    search_term = st.text_input("🔍 Search templates...", 
                               placeholder="Search templates...",
                               label_visibility="collapsed",
                               key="template_search")
    
    # Template categories
    template_categories = {
        "Minimalist (ATS Best)": "ATS",
        "Horizontal Line": "Modern",
        "Bold Title Accent": "Bold",
        "Date Below": "Clean",
        "Section Box Header": "Boxed",
        "Times New Roman Classic": "Classic",
        "Sophisticated Minimal": "Minimal",
        "Clean Look": "Modern",
        "Elegant": "Elegant",
        "Modern Minimal": "Minimal",
        "Two Coloumn": "Layout",
    }
    
    # Load all saved templates
    current_user = st.session_state.logged_in_user
    if 'uploaded_templates' not in st.session_state:
        st.session_state.uploaded_templates = load_user_templates(current_user)
    
    if 'doc_templates' not in st.session_state:
        st.session_state.doc_templates = load_user_doc_templates(current_user)
    
    if 'ppt_templates' not in st.session_state:
        st.session_state.ppt_templates = load_user_ppt_templates(current_user)
    
    # Filter system templates
    filtered_templates = {
        name: config for name, config in SYSTEM_TEMPLATES.items()
        if search_term.lower() in name.lower() or 
        search_term.lower() in template_categories.get(name, "").lower()
    }
    
    # Filter saved templates
    filtered_html = {
        tid: data for tid, data in st.session_state.uploaded_templates.items()
        if search_term.lower() in data['name'].lower()
    }
    
    filtered_doc = {
        tid: data for tid, data in st.session_state.doc_templates.items()
        if search_term.lower() in data['name'].lower()
    }
    
    filtered_ppt = {
        tid: data for tid, data in st.session_state.ppt_templates.items()
        if search_term.lower() in data['name'].lower()
    }
    
    # Create tabs for different template types
    tab1, tab2, tab3, tab4 = st.tabs(["📋 System Templates", "📄 HTML Templates", "📝 Word Templates", "📊 PowerPoint Templates"])
    
    # ========== SYSTEM TEMPLATES TAB ==========
    with tab1:
        templates_list = list(filtered_templates.items())
        
        # First row with Upload Custom Template card
        cols = st.columns(3)
        
        # Upload Custom Template Card
        with cols[0]:
            st.markdown("""
            <div class="create-blank-card" id="upload-template-card">
                <div class="plus-icon">📤</div>
                <div style="font-size: 1.2rem; font-weight: 700; color: #212529; margin-bottom: 0.5rem;">
                    Upload Custom Template
                </div>
                <div style="font-size: 0.9rem; color: #6c757d; text-align: center;">
                    Upload your own HTML, Word, or PowerPoint template
                </div>
            </div>
            """, unsafe_allow_html=True)
            
            if st.button("📤 Upload Template", key="upload_custom_btn", type="primary", use_container_width=True):
                st.session_state.show_upload_interface = True
                st.rerun()
        
        # Fill remaining slots in first row
        for col_idx in range(1, min(3, len(templates_list) + 1)):
            if col_idx - 1 < len(templates_list):
                template_name, template_config = templates_list[col_idx - 1]
                with cols[col_idx]:
                    render_template_card(template_name, template_config, template_categories)
        
        # Render remaining templates in rows of 3
        if len(templates_list) > 2:
            remaining_templates = templates_list[2:]
            for i in range(0, len(remaining_templates), 3):
                cols = st.columns(3)
                row_templates = remaining_templates[i:i + 3]
                
                for col_idx, (template_name, template_config) in enumerate(row_templates):
                    with cols[col_idx]:
                        render_template_card(template_name, template_config, template_categories)
    
    # ========== HTML TEMPLATES TAB ==========
    with tab2:
        if filtered_html:
            # Display in grid format
            html_list = list(filtered_html.items())
            for i in range(0, len(html_list), 3):
                cols = st.columns(3)
                row_items = html_list[i:i + 3]
                
                for col_idx, (template_id, template_data) in enumerate(row_items):
                    with cols[col_idx]:
                        render_saved_html_template_card(template_id, template_data)
        else:
            st.info("📂 No HTML templates found. Upload one to get started!")
            if st.button("📤 Upload HTML Template", key="upload_html_from_tab", type="primary"):
                st.session_state.show_upload_interface = True
                st.rerun()
    
    # ========== WORD TEMPLATES TAB ==========
    with tab3:
        if filtered_doc:
            # Display in grid format
            doc_list = list(filtered_doc.items())
            for i in range(0, len(doc_list), 3):
                cols = st.columns(3)
                row_items = doc_list[i:i + 3]
                
                for col_idx, (template_id, template_data) in enumerate(row_items):
                    with cols[col_idx]:
                        render_saved_doc_template_card(template_id, template_data)
        else:
            st.info("📂 No Word templates found. Upload one to get started!")
            if st.button("📤 Upload Word Template", key="upload_doc_from_tab", type="primary"):
                st.session_state.show_upload_interface = True
                st.rerun()
    
    # ========== POWERPOINT TEMPLATES TAB ==========
    with tab4:
        if filtered_ppt:
            # Display in grid format
            ppt_list = list(filtered_ppt.items())
            for i in range(0, len(ppt_list), 3):
                cols = st.columns(3)
                row_items = ppt_list[i:i + 3]
                
                for col_idx, (template_id, template_data) in enumerate(row_items):
                    with cols[col_idx]:
                        render_saved_ppt_template_card(template_id, template_data)
        else:
            st.info("📂 No PowerPoint templates found. Upload one to get started!")
            if st.button("📤 Upload PowerPoint Template", key="upload_ppt_from_tab", type="primary"):
                st.session_state.show_upload_interface = True
                st.rerun()

def render_saved_html_template_card(template_id, template_data):
    """Render a saved HTML template card."""
    st.markdown(f"""
    <div class="template-card">
        <div class="template-preview">
            <div style="width: 85%; height: 90%; background: white; 
                 border-radius: 4px; padding: 1rem; box-shadow: 0 2px 8px rgba(0,0,0,0.1);">
                <div style="width: 60%; height: 10px; background: #FF6B35; 
                     border-radius: 2px; margin-bottom: 0.5rem;"></div>
                <div style="width: 40%; height: 6px; background: #E9ECEF; 
                     border-radius: 2px; margin-bottom: 1rem;"></div>
                <div style="width: 80%; height: 6px; background: #E9ECEF; 
                     border-radius: 2px; margin-bottom: 0.5rem;"></div>
            </div>
            <div class="template-badge">Custom HTML</div>
        </div>
        <div style="padding: 1rem;">
            <div style="font-size: 1rem; font-weight: 600; color: #212529; 
                 margin-bottom: 0.3rem;">📄 {template_data['name']}</div>
            <div style="font-size: 0.75rem; color: #6c757d;">
                {template_data.get('uploaded_at', 'N/A')}
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        if st.button("Use Template", key=f"use_html_{template_id}", type="primary", use_container_width=True):
            try:
                # Get resume data
                clear_template_state()
                user_resume = st.session_state.get('final_resume_data') or st.session_state.get('enhanced_resume') or {}
                
                # Ensure enhanced_resume is initialized
                if st.session_state.get('enhanced_resume') is None:
                    st.session_state['enhanced_resume'] = user_resume.copy()
                
                # Generate fresh resume data
                resume_data = generate_resume_for_template()
                
                # Get template CSS and generate HTML
                css = template_data.get('css', '')
                html_content = generate_generic_html(resume_data)
                
                # Create full HTML preview
                full_html = f"""
                <div class="ats-page">
                    {html_content}
                </div>
                """
                
                # Store everything in session state
                st.session_state.selected_template = template_data['name']
                st.session_state.template_preview_html = full_html
                st.session_state.template_preview_css = css
                st.session_state.final_resume_data = resume_data
                st.session_state['enhanced_resume'] = resume_data
                st.session_state['template_source'] = 'html_saved'  # Set template source
                
                # Calculate ATS score if job description exists
                jd_data = st.session_state.get('job_description')
                if jd_data and isinstance(jd_data, dict) and len(jd_data) > 0:
                    try:
                        st.session_state['ats_result'] = ai_ats_score(resume_data, jd_data)
                    except Exception as e:
                        st.warning(f"Could not calculate ATS score: {str(e)}")
                        st.session_state['ats_result'] = {}
                else:
                    st.session_state['ats_result'] = {}
                
                # Switch to visual editor
                st.session_state.show_template_selector = False
                st.session_state.show_visual_editor = True
                
                st.success(f"✅ Using template: {template_data['name']}")
                time.sleep(0.5)
                st.rerun()
                
            except Exception as e:
                st.error(f"❌ Error loading template: {str(e)}")
                import traceback
                st.error(f"Details: {traceback.format_exc()}")
    
    with col2:
        if st.button("🗑️", key=f"delete_html_{template_id}", use_container_width=True):
            del st.session_state.uploaded_templates[template_id]
            save_user_templates(st.session_state.logged_in_user, st.session_state.uploaded_templates)
            st.success(f"✅ Deleted '{template_data['name']}'")
            time.sleep(0.5)
            st.rerun()

def render_saved_doc_template_card(template_id, template_data):
    """Render a saved Word template card."""
    st.markdown(f"""
    <div class="template-card">
        <div class="template-preview">
            <div style="width: 85%; height: 90%; background: white; 
                 border-radius: 4px; padding: 1rem; box-shadow: 0 2px 8px rgba(0,0,0,0.1);">
                <div style="width: 70%; height: 12px; background: #2B579A; 
                     border-radius: 2px; margin-bottom: 0.5rem;"></div>
                <div style="width: 50%; height: 8px; background: #E9ECEF; 
                     border-radius: 2px; margin-bottom: 0.8rem;"></div>
                <div style="width: 90%; height: 6px; background: #E9ECEF; 
                     border-radius: 2px; margin-bottom: 0.4rem;"></div>
                <div style="width: 85%; height: 6px; background: #E9ECEF; 
                     border-radius: 2px;"></div>
            </div>
            <div class="template-badge">Word Doc</div>
        </div>
        <div style="padding: 1rem;">
            <div style="font-size: 1rem; font-weight: 600; color: #212529; 
                 margin-bottom: 0.3rem;">📝 {template_data['name']}</div>
            <div style="font-size: 0.75rem; color: #6c757d;">
                {template_data.get('uploaded_at', 'N/A')}
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        if st.button("Use Template", key=f"use_doc_{template_id}", type="primary", use_container_width=True):
            try:
                import io
                clear_template_state()
                # Get resume data
                final_data = st.session_state.get('final_resume_data') or st.session_state.get('enhanced_resume') or {}
                
                # Ensure enhanced_resume is initialized
                if st.session_state.get('enhanced_resume') is None:
                    st.session_state['enhanced_resume'] = final_data.copy()
                
                # Generate fresh resume data
                resume_data = generate_resume_for_template()
                
                # Get the ORIGINAL template bytes
                template_bytes = template_data['doc_data']
                
                if not isinstance(template_bytes, bytes):
                    st.error("❌ Invalid template data format")
                    return
                
                with st.spinner("Processing template..."):
                    # Extract text from template
                    uploadtext = extract_temp_from_docx(io.BytesIO(template_bytes))
                    
                    # Generate mapping with user resume
                    mapped_data = ask_ai_for_mapping(uploadtext, resume_data)
                    
                    # Ensure mapping is a dictionary
                    if isinstance(mapped_data, list):
                        mapped_data = {
                            item["template"]: item["new"]
                            for item in mapped_data
                            if "template" in item and "new" in item
                        }
                    
                    # Process the template with user data
                    output_doc = auto_process_docx(
                        io.BytesIO(template_bytes),
                        mapped_data
                    )
                
                # Store everything - CRITICAL
                st.session_state['generated_docx'] = output_doc.getvalue()
                st.session_state['enhanced_resume'] = resume_data
                st.session_state.final_resume_data = resume_data
                st.session_state.selected_doc_template_id = template_id
                st.session_state.selected_doc_template = template_data
                st.session_state.selected_template = template_data['name']
                st.session_state['template_source'] = 'doc_saved'
                st.session_state['doc_original_bytes'] = template_bytes
                st.session_state['doc_original_filename'] = template_data['original_filename']
                st.session_state['mapping'] = mapped_data
                st.session_state['template_text'] = uploadtext
                
                # Calculate ATS score if job description exists
                jd_data = st.session_state.get('job_description')
                if jd_data and isinstance(jd_data, dict) and len(jd_data) > 0:
                    try:
                        st.session_state['ats_result'] = ai_ats_score(resume_data, jd_data)
                    except Exception as e:
                        st.warning(f"Could not calculate ATS score: {str(e)}")
                        st.session_state['ats_result'] = {}
                else:
                    st.session_state['ats_result'] = {}
                
                # Switch to visual editor to show preview
                st.session_state.show_template_selector = False
                st.session_state.show_visual_editor = True
                
                st.success(f"✅ Using template: {template_data['name']}")
                time.sleep(0.5)
                st.rerun()
                
            except Exception as e:
                st.error(f"❌ Error loading template: {str(e)}")
                import traceback
                st.error(f"Details: {traceback.format_exc()}")
    
    with col2:
        if st.button("🗑️", key=f"delete_doc_{template_id}", use_container_width=True):

            success = delete_user_doc_template(
                st.session_state.logged_in_user,
                template_id
            )

            if success:
                # Reload from DB (single source of truth)
                st.session_state.doc_templates = load_user_doc_templates(
                    st.session_state.logged_in_user
                )

                st.success(f"✅ Deleted '{template_data['name']}'")
                time.sleep(0.5)
                st.rerun()



def render_saved_ppt_template_card(template_id, template_data):
    """Render a saved PowerPoint template card."""
    st.markdown(f"""
    <div class="template-card">
        <div class="template-preview">
            <div style="width: 85%; height: 90%; background: linear-gradient(135deg, #D04A02 0%, #F97316 100%); 
                 border-radius: 4px; padding: 1rem; box-shadow: 0 2px 8px rgba(0,0,0,0.1);">
                <div style="width: 60%; height: 14px; background: white; 
                     border-radius: 2px; margin-bottom: 0.7rem;"></div>
                <div style="width: 80%; height: 8px; background: rgba(255,255,255,0.7); 
                     border-radius: 2px; margin-bottom: 0.5rem;"></div>
                <div style="width: 75%; height: 8px; background: rgba(255,255,255,0.7); 
                     border-radius: 2px;"></div>
            </div>
            <div class="template-badge">PowerPoint</div>
        </div>
        <div style="padding: 1rem;">
            <div style="font-size: 1rem; font-weight: 600; color: #212529; 
                 margin-bottom: 0.3rem;">📊 {template_data['name']}</div>
            <div style="font-size: 0.75rem; color: #6c757d;">
                {template_data.get('uploaded_at', 'N/A')}
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        if st.button("Use Template", key=f"use_ppt_{template_id}", type="primary", use_container_width=True):
            try:
                import io
                from pptx import Presentation
                clear_template_state()
                # Get resume data
                user_resume = st.session_state.get('final_resume_data') or st.session_state.get('enhanced_resume') or {}
                
                # Ensure enhanced_resume is initialized
                if st.session_state.get('enhanced_resume') is None:
                    st.session_state['enhanced_resume'] = user_resume.copy()
                with st.spinner("Processing template..."):
                    # Generate fresh resume data
                    resume_data = generate_resume_for_template()
                    st.spinner("Processing template...")
                    working_prs = Presentation(io.BytesIO(template_data['ppt_data']))
                    prs = Presentation(io.BytesIO(template_data['ppt_data']))
                    
                    slide_texts = []
                    for slide_idx, slide in enumerate(prs.slides):
                        text_blocks = []
                        for shape_idx, shape in enumerate(slide.shapes):
                            if shape.has_text_frame and shape.text.strip():
                                text_blocks.append({
                                    "index": shape_idx,
                                    "text": shape.text.strip(),
                                    "position": {"x": shape.left, "y": shape.top}
                                })
                        
                        if text_blocks:
                            text_blocks.sort(key=lambda x: (x["position"]["y"], x["position"]["x"]))
                            slide_texts.append({
                                "slide_number": slide_idx + 1,
                                "text_blocks": text_blocks
                            })
                    
                    structured_slides = analyze_slide_structure(slide_texts)
                    generated_sections = generate_ppt_sections(resume_data, structured_slides)
                    
                    text_elements = template_data.get('text_elements', [])
                    if not text_elements:
                        # Regenerate text elements if not stored
                        text_elements = []
                        for slide_idx, slide in enumerate(prs.slides):
                            for shape_idx, shape in enumerate(slide.shapes):
                                if shape.has_text_frame and shape.text.strip():
                                    text_elements.append({
                                        'slide': slide_idx + 1,
                                        'shape': shape_idx,
                                        'original_text': shape.text.strip(),
                                        'shape_type': type(shape).__name__
                                    })
                    
                    content_mapping, heading_shapes, basic_info_shapes = match_generated_to_original(
                        text_elements, generated_sections, prs)
                    
                    edits = {}
                    for element in text_elements:
                        key = f"{element['slide']}_{element['shape']}"
                        if key not in heading_shapes:
                            edits[key] = content_mapping.get(key, element['original_text'])
                    
                    for element in text_elements:
                        key = f"{element['slide']}_{element['shape']}"
                        if key not in heading_shapes and key in edits:
                            slide_idx = element['slide'] - 1
                            shape_idx = element['shape']
                            
                            if slide_idx < len(working_prs.slides):
                                slide = working_prs.slides[slide_idx]
                                if shape_idx < len(slide.shapes):
                                    shape = slide.shapes[shape_idx]
                                    if shape.has_text_frame:
                                        clear_and_replace_text(shape, edits[key])
                    
                    output = io.BytesIO()
                    working_prs.save(output)
                    output.seek(0)
                    
                # Store everything
                st.session_state['generated_ppt'] = output.getvalue()
                st.session_state['enhanced_resume'] = resume_data
                st.session_state.final_resume_data = resume_data
                st.session_state.selected_ppt_template_id = template_id
                st.session_state.selected_ppt_template = template_data
                st.session_state['template_source'] = 'ppt_saved'
                st.session_state.selected_template = template_data['name']
                
                # Calculate ATS score if job description exists
                jd_data = st.session_state.get('job_description')
                if jd_data and isinstance(jd_data, dict) and len(jd_data) > 0:
                    try:
                        st.session_state['ats_result'] = ai_ats_score(resume_data, jd_data)
                    except Exception as e:
                        st.warning(f"Could not calculate ATS score: {str(e)}")
                        st.session_state['ats_result'] = {}
                else:
                    st.session_state['ats_result'] = {}
                
                # Switch to visual editor to show preview
                st.session_state.show_template_selector = False
                st.session_state.show_visual_editor = True
                
                st.success(f"✅ Using template: {template_data['name']}")
                time.sleep(0.5)
                st.rerun()
                
            except Exception as e:
                st.error(f"Error loading template: {str(e)}")
                import traceback
                st.error(f"Details: {traceback.format_exc()}")
    
    with col2:
        if st.button("🗑️", key=f"delete_ppt_{template_id}", use_container_width=True):

            success = delete_user_ppt_template(
                st.session_state.logged_in_user,
                template_id
            )

            if success:
                # Reload templates from DB (source of truth)
                st.session_state.ppt_templates = load_user_ppt_templates(
                    st.session_state.logged_in_user
                )

                st.success(f"✅ Deleted '{template_data['name']}'")
                time.sleep(0.5)
                st.rerun()


# def generate_basic_docx(resume_data):
#     """Generate a basic DOCX from resume data."""
#     doc = Document()
    
#     # Header - Name
#     name = doc.add_paragraph(resume_data.get('name', 'Your Name'))
#     name.alignment = WD_ALIGN_PARAGRAPH.CENTER
#     name.runs[0].font.size = Pt(16)
#     name.runs[0].font.bold = True
    
#     # Contact Info
#     contact_parts = []
#     if resume_data.get('email'):
#         contact_parts.append(resume_data['email'])
#     if resume_data.get('phone'):
#         contact_parts.append(resume_data['phone'])
#     if resume_data.get('location'):
#         contact_parts.append(resume_data['location'])
    
#     if contact_parts:
#         contact = doc.add_paragraph(' | '.join(contact_parts))
#         contact.alignment = WD_ALIGN_PARAGRAPH.CENTER
#         contact.runs[0].font.size = Pt(10)
    
#     doc.add_paragraph()  # Spacing
    
#     # Summary
#     if resume_data.get('summary'):
#         doc.add_heading('Professional Summary', level=2)
#         doc.add_paragraph(resume_data['summary'])
    
#     # Experience
#     if resume_data.get('experience'):
#         doc.add_heading('Experience', level=2)
#         for exp in resume_data['experience']:
#             p = doc.add_paragraph()
#             p.add_run(f"{exp.get('position', '')} - {exp.get('company', '')}").bold = True
#             p.add_run(f"\n{exp.get('start_date', '')} to {exp.get('end_date', '')}")
            
#             if exp.get('description'):
#                 for desc in exp['description']:
#                     doc.add_paragraph(desc, style='List Bullet')
    
#     # Education
#     if resume_data.get('education'):
#         doc.add_heading('Education', level=2)
#         for edu in resume_data['education']:
#             p = doc.add_paragraph()
#             p.add_run(f"{edu.get('degree', '')}").bold = True
#             p.add_run(f"\n{edu.get('institution', '')}")
#             p.add_run(f"\n{edu.get('start_date', '')} to {edu.get('end_date', '')}")
    
#     # Skills
#     if resume_data.get('skills'):
#         doc.add_heading('Skills', level=2)
#         skills_text = ', '.join(resume_data['skills'])
#         doc.add_paragraph(skills_text)
    
#     # Projects
#     if resume_data.get('projects'):
#         doc.add_heading('Projects', level=2)
#         for proj in resume_data['projects']:
#             p = doc.add_paragraph()
#             p.add_run(proj.get('name', '')).bold = True
#             if proj.get('description'):
#                 for desc in proj['description']:
#                     doc.add_paragraph(desc, style='List Bullet')
    
#     # Certifications
#     if resume_data.get('certifications'):
#         doc.add_heading('Certifications', level=2)
#         for cert in resume_data['certifications']:
#             doc.add_paragraph(
#                 f"{cert.get('name', '')} - {cert.get('issuer', '')} ({cert.get('completed_date', '')})",
#                 style='List Bullet'
#             )
    
#     # Save to BytesIO
#     output = BytesIO()
#     doc.save(output)
#     output.seek(0)
#     return output.getvalue()

def show_visual_editor_with_tools():
    """Show the visual editor with resume tools on the side."""
    if not st.session_state.selected_template:
        st.session_state.show_template_selector = True
        st.session_state.show_visual_editor = False
        st.rerun()
        return

    # Header
    # st.markdown('<div class="editor-header">', unsafe_allow_html=True)
    col1, col2 = st.columns([2, 5], gap='medium')
    with col1:
        if st.button("← Back to Templates", type="primary", use_container_width=True):
            st.session_state.show_template_selector = True
            st.session_state.show_visual_editor = False
            st.rerun()

    with col2:
        st.markdown(
            f"<h3 style='margin: 0; color: #212529;'>📄 Resume Visual Editor - {st.session_state.selected_template}</h3>",
            unsafe_allow_html=True
        )
    st.markdown('</div>', unsafe_allow_html=True)

    editor_col, tools_col = st.columns([7.5, 2.5], gap="medium")

    with editor_col:
        resume_data = st.session_state.final_resume_data or {}

        if st.session_state.get('enhanced_resume') is None:
            st.session_state['enhanced_resume'] = resume_data.copy()

        resume_data = st.session_state['enhanced_resume']
        
      
  
        template_source = st.session_state.get('template_source', 'html_saved')

 
        if template_source in ['doc_saved', 'ppt_saved']:
            is_edit_mode = False 
        else:
            is_edit_mode = st.checkbox("⚙️ **Enable Edit Mode**", key='edit_toggle')
            
        if is_edit_mode:
            with st.container():
                st.markdown(
                    "<h3 style='text-align:center;color:#6b7280;margin-bottom:1.5rem;'>✏️ Content Editor</h3>",
                    unsafe_allow_html=True
                )

                # ---------------- CONTENT EDITOR ----------------
                render_basic_details(resume_data, is_edit=is_edit_mode)
                st.session_state.resume_dirty = True if is_edit_mode else st.session_state.get("resume_dirty", False)

                rendered_keys = set()
                standard_keys = get_standard_keys()

                for key in RESUME_ORDER:
                    if key in resume_data and resume_data[key]:
                        rendered_keys.add(key)
                        if key == "skills":
                            render_skills_section(resume_data, is_edit=is_edit_mode)
                        else:
                            render_generic_section(key, resume_data[key], is_edit=is_edit_mode)

                for key, value in resume_data.items():
                    if key not in rendered_keys and key not in standard_keys:
                        if isinstance(value, list) and value:
                            rendered_keys.add(key)
                            render_generic_section(key, value, is_edit=is_edit_mode)

                for key, value in resume_data.items():
                    if key not in rendered_keys and key not in standard_keys and isinstance(value, str):
                        st.markdown("<div class='resume-section'>", unsafe_allow_html=True)
                        st.markdown(f"<h3 class='custom-section-header'>{key}</h3>", unsafe_allow_html=True)

                        if is_edit_mode:
                            new_val = st.text_area(
                                f"Edit {key}",
                                value=value.strip(),
                                key=f"edit_custom_{key}",
                                height=200
                            )
                            resume_data[key] = new_val.strip()
                            st.session_state['enhanced_resume'] = resume_data
                            st.session_state.resume_dirty = True

                            if st.button(f"🗑️ Delete '{key}' Section", key=f"delete_{key}", type="secondary"):
                                del resume_data[key]
                                st.session_state['enhanced_resume'] = resume_data
                                st.session_state.resume_dirty = True
                                st.rerun()
                        else:
                            st.markdown(
                                f"<div class='custom-section-content'>{value.strip()}</div>",
                                unsafe_allow_html=True
                            )
                        st.markdown("</div>", unsafe_allow_html=True)

        # ---------------- PREVIEW REGENERATION ----------------
        if st.session_state.get("resume_dirty"):
            new_html, new_css = regenerate_live_preview()
            if new_html:
                st.session_state.template_preview_html = new_html
                st.session_state.template_preview_css = new_css
            st.session_state.resume_dirty = False

        html_content = st.session_state.template_preview_html or ""
        css_content = st.session_state.template_preview_css or ""

        # Add separator
        st.markdown("---")

        # ✅ FIX: Show appropriate editor based on template type
        if template_source == 'doc_saved':
            st.markdown("<h3 style='text-align:center;color:#6b7280;margin-top:2rem;margin-bottom:1rem;'>📄 Word Document Preview</h3>", unsafe_allow_html=True)
            
            # Show inline editor if enabled
            if st.session_state.get('show_inline_doc_editor', False):
                show_inline_doc_mapping_editor()
                st.markdown("---")
            
            # Show preview
            if st.session_state.get('generated_docx'):
                try:
                    preview_html = docx_to_html_preview(io.BytesIO(st.session_state['generated_docx']))
                    st.components.v1.html(preview_html, height=800, scrolling=True)
                except Exception as e:
                    st.error(f"Preview error: {str(e)}")
            else:
                st.info("📄 Click **'Save & Auto-Improve'** to generate the Word document.")

        elif template_source == 'ppt_saved':
            st.markdown("<h3 style='text-align:center;color:#6b7280;margin-top:2rem;margin-bottom:1rem;'>📊 PowerPoint Preview</h3>", unsafe_allow_html=True)
            
            # Show inline editor if enabled
            if st.session_state.get('show_inline_ppt_editor', False):
                show_inline_ppt_content_editor()
                st.markdown("---")
            
            # Show preview
            if st.session_state.get('generated_ppt'):
                show_ppt_preview_inline(st.session_state['generated_ppt'])
            else:
                st.warning("No presentation generated yet. Click 'Save & Auto-Improve' to generate.")

        else:
            st.markdown("<h3 style='text-align:center;color:#6b7280;margin-top:2rem;margin-bottom:1rem;'>👁️ Live Preview</h3>", unsafe_allow_html=True)
            
            # Inject data into editor
            editor_html_with_data = VISUAL_EDITOR_HTML
            
            inject_script = f"""
            <script>
                window.templateContent = `{html_content}`;
                window.templateCss = `{css_content}`;
                window.templateName = `{st.session_state.selected_template}`;
            </script>
            """
            
            editor_html_with_data = editor_html_with_data.replace('<body>', f'<body>\n{inject_script}')
            
            # Display the visual editor
            components.html(editor_html_with_data, height=1400, scrolling=True)

    with tools_col:
        with st.container():
            st.title("Resume Tools")
            
            # ========== INLINE EDITOR TOGGLE BUTTONS ==========
            template_source = st.session_state.get('template_source', 'html_saved')
            
            if template_source == 'doc_saved':
                st.markdown("---")
                # Toggle button for Word editor
                current_state = st.session_state.get('show_inline_doc_editor', False)
                button_label = "✏️ Hide Document Editor" if current_state else "✏️ Edit Document Mapping"
                
                if st.button(button_label, use_container_width=True, type="primary" if not current_state else "secondary"):
                    st.session_state.show_inline_doc_editor = not current_state
                    st.rerun()
            
            elif template_source == 'ppt_saved':
                st.markdown("---")
                # Toggle button for PPT editor
                current_state = st.session_state.get('show_inline_ppt_editor', False)
                button_label = "✏️ Hide Content Editor" if current_state else "✏️ Edit Presentation Content"
                
                if st.button(button_label, use_container_width=True, type="primary" if not current_state else "secondary"):
                    st.session_state.show_inline_ppt_editor = not current_state
                    st.rerun()
            
            # ========== DOWNLOAD BUTTONS ==========
            if template_source == 'doc_saved' and st.session_state.get('generated_docx'):
                st.markdown("---")
                filename = f"Resume_{st.session_state.final_resume_data.get('name', 'User').replace(' ', '_')}.docx"
                st.download_button(
                    label="⬇️ Download Word Document",
                    data=st.session_state['generated_docx'],
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    type="primary"
                )
            
            if template_source == 'ppt_saved' and st.session_state.get('generated_ppt'):
                st.markdown("---")
                filename = f"Resume_{st.session_state.final_resume_data.get('name', 'User').replace(' ', '_')}.pptx"
                st.download_button(
                    label="⬇️ Download PowerPoint",
                    data=st.session_state['generated_ppt'],
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    type="primary"
                )
            
            st.markdown("---")
            
            # ========== REST OF TOOLS (Save & Auto-Improve, etc.) ==========
            loading_placeholder = st.empty()
            
            # Save & Auto-Improve button
            if st.button("✨ **Save & Auto-Improve**", type="primary", use_container_width=True):
                loading_placeholder.markdown("""
                    <div id="overlay-loader">
                        <div class="loader-spinner"></div>
                        <p>Performing auto-improvement...</p>
                    </div>
                    <style>
                        #overlay-loader {
                            position: fixed;
                            top: 0;
                            left: 0;
                            width: 100vw;
                            height: 100vh;
                            background: rgba(255, 255, 255, 0.95);
                            backdrop-filter: blur(6px);
                            display: flex;
                            flex-direction: column;
                            justify-content: center;
                            align-items: center;
                            z-index: 9999;
                            font-size: 1.2rem;
                            font-weight: 500;
                        }
                        .loader-spinner {
                            border: 5px solid rgba(232, 117, 50, 0.2);
                            border-top: 5px solid #e87532;
                            border-radius: 50%;
                            width: 70px;
                            height: 70px;
                            animation: spin 1s linear infinite;
                            margin-bottom: 20px;
                        }
                        @keyframes spin {
                            0% { transform: rotate(0deg); }
                            100% { transform: rotate(360deg); }
                        }
                        #overlay-loader p {
                            color: #1f2937;
                            font-size: 1.1rem;
                            letter-spacing: 0.5px;
                        }
                    </style>
                    """, unsafe_allow_html=True)
                save_custom_sections()
                st.session_state.final_resume_data = st.session_state.enhanced_resume.copy()
                save_and_improve()
                loading_placeholder.empty()

            if template_source not in ['doc_saved', 'ppt_saved']: 
                if not st.session_state.get('edit_toggle', False):
                    st.info("⚠️ Enable Edit Mode to add new items.\n\nFor saving newly added content, disable Edit Mode after making changes.")
                else:
                    st.markdown("---")
                    st.subheader("➕ Add New Section Items")
                    st.button(
                        "Add New Experience",
                        on_click=add_new_item,
                        args=('experience', {
                            "position": "New Job Title",
                            "company": "New Company",
                            "start_date": "2025-01-01",
                            "end_date": "2025-12-31",
                            "description": ["New responsibility 1."]
                        }),
                        type="secondary",
                        use_container_width=True
                    )
                    st.button(
                        "Add New Education",
                        on_click=add_new_item,
                        args=('education', {
                            "institution": "New University",
                            "degree": "New Degree",
                            "start_date": "2025-01-01",
                            "end_date": "2025-12-31"
                        }),
                        type="secondary",
                        use_container_width=True
                    )
                    st.button(
                        "Add New Certification",
                        on_click=add_new_item,
                        args=('certifications', {
                            "name": "New Certification Name",
                            "issuer": "Issuing Body",
                            "completed_date": "2025-01-01"
                        }),
                        type="secondary",
                        use_container_width=True
                    )
                    st.button(
                        "Add New Project",
                        on_click=add_new_item,
                        args=('projects', {
                            "name": "New Project Title",
                            "description": ["Project detail"]
                        }),
                        type="secondary",
                        use_container_width=True
                    )
            # Continue with HTML/PDF downloads and ATS analysis...
            # (Rest of your existing tools_col code)
# In your show_visual_editor_with_tools() function, replace the PDF download section with:

        st.markdown("---")

        html_content = st.session_state.template_preview_html or ""
        css_content = st.session_state.template_preview_css or ""
        template_source = st.session_state.get('template_source', 'html_saved')

        if template_source not in ['doc_saved', 'ppt_saved']: 
            # Single HTML for both downloads
            download_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <style>
                    @page {{ size: A4; margin: 0.5in; }}
                    * {{ -webkit-print-color-adjust: exact !important; print-color-adjust: exact !important; }}
                    {css_content}
                </style>
            </head>
            <body style="margin: 0; padding: 0; background: #fff;">
                <div style="width: 8.5in; min-height: 11in; margin: 0 auto; padding: 1in;">
                    {html_content}
                </div>
                <script>
                    window.onload = function() {{ setTimeout(function() {{ window.print(); }}, 500); }};
                </script>
            </body>
            </html>
            """
            
            col1, col2 = st.columns(2)
            
            with col1:
                st.download_button(
                    label="⬇️ Download HTML",
                    data=download_html,
                    file_name=f"resume_{st.session_state.selected_template.replace(' ', '_')}.html",
                    mime="text/html",
                    use_container_width=True,
                    type="primary",
                    key="download_html_btn"
                )
            
            with col2:
                filename = f"Resume_{resume_data.get('name', 'User').replace(' ', '_')}.docx"

                # Convert on-the-fly
                with st.spinner("Preparing DOCX..."):
                    docx_data = convert_html_to_docx_spire(html_content, css_content)

                if docx_data:
                    st.download_button(
                        label="📄 Download DOCX",
                        data=docx_data,
                        file_name=filename,
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True,
                        type="primary",
                        key="download_docx_btn"
                    )
                else:
                    st.error("❌ Failed to convert HTML to DOCX")
            st.info("📄 **To save as PDF:** Open downloaded file → Print dialog appears → Select 'Save as PDF' → Save")

        st.markdown("---")

        # try:
        #     template_source = st.session_state.get('template_source', 'html_saved')
        #     resume_data = st.session_state.get('enhanced_resume', {})
            
        #     # Check if we have a Word document to download
        #     if template_source == 'doc_saved' and st.session_state.get('generated_docx'):
        #         # Already have DOCX from saved template
        #         docx_data = st.session_state['generated_docx']
                
        #         # Create filename
        #         filename = f"Resume_{resume_data.get('name', 'User').replace(' ', '_')}.docx"
                
        #         # Download button
        #         st.download_button(
        #             label="⬇️ Download Word Document",
        #             data=docx_data,
        #             file_name=filename,
        #             mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        #             use_container_width=True,
        #             type="primary",
        #             key="download_doc_saved"
        #         )
                
               
            
        # except Exception as e:
        #     st.error(f"Error with DOCX: {str(e)}")
        #     import traceback
        #     st.error(traceback.format_exc())
        # st.markdown("---")
        
        # ats_data = st.session_state.get('ats_result', {})
        # st.write("### ATS Analysis")
        ats_data = st.session_state.get("ats_result", {})

        if isinstance(ats_data, dict) and "overall_score" in ats_data:

            # Strengths
            st.markdown("### ✅ Strengths")
            strengths = ats_data.get("strengths", [])
            if strengths:
                for s in strengths:
                    st.markdown(f"- {s}")
            else:
                st.write("No strengths identified.")

            # Weaknesses
            st.markdown("### ⚠️ Weaknesses")
            weaknesses = ats_data.get("weaknesses", [])
            if weaknesses:
                for w in weaknesses:
                    st.markdown(f"- {w}")
            else:
                st.write("No weaknesses identified.")

            # Explanation
            explanation = ats_data.get("explanation")
            if explanation:
                st.markdown("### 🧠 Explanation")
                st.write(explanation)

        if ats_data and ats_data.get("overall_score", 0) > 0:
                score = ats_data.get("overall_score", 0)  # Changed from "score" to "overall_score"
                label = get_score_label(score)
                color = get_score_color(score)
                
                st.markdown(f"### ATS Score: {score}/100")
                st.markdown(f"**{label}**", unsafe_allow_html=True)
                st.markdown(f"""
                <div style="
                    text-align:center;
                    padding:15px;
                    margin-bottom:18px;
                    background:white;
                    border-radius:16px;
                    box-shadow:0 4px 12px rgba(0,0,0,0.08);
                    border:1px solid #e8e8e8;
                ">
                    <div style="position: relative; width:140px; height:140px; margin:auto;">
                        <svg width="140" height="140">
                            <circle cx="70" cy="70" r="60" stroke="#f0f0f0" stroke-width="12" fill="none"/>
                            <circle cx="70" cy="70" r="60"
                                stroke="{color}" stroke-width="12" fill="none"
                                stroke-linecap="round"
                                stroke-dasharray="{round(score*3.6)}, 360"
                                transform="rotate(-90 70 70)"
                            />
                            <text x="70" y="78" text-anchor="middle" font-size="28" font-weight="700" fill="#333">{score}%</text>
                        </svg>
                    </div>
                    <div style="font-size:15px; margin-top:8px; font-weight:600; color:#333;">
                        {label}
                    </div>
                </div>
                """, unsafe_allow_html=True)

                # ================= Keyword Table Expander ====================
                with st.expander("🔎 View ATS Keyword Analysis"):
                    # ---------- Technical Skills ----------
                    st.markdown("## 🛠 Technical Skills")

                    st.markdown("### 🟢 Matched")
                    tech_matched = ats_data.get("technical_skills", {}).get("matched", [])
                    if tech_matched:
                        st.write(", ".join(tech_matched))
                    else:
                        st.write("None")

                    st.markdown("### 🔴 Missing")
                    tech_missing = ats_data.get("technical_skills", {}).get("missing", [])
                    if tech_missing:
                        st.write(", ".join(tech_missing))
                    else:
                        st.write("None")

                    # ---------- Soft Skills ----------
                    st.markdown("---")
                    st.markdown("## 🤝 Soft Skills")

                    st.markdown("### 🟢 Matched")
                    soft_matched = ats_data.get("soft_skills", {}).get("matched", [])
                    if soft_matched:
                        st.write(", ".join(soft_matched))
                    else:
                        st.write("None")

                    st.markdown("### 🔴 Missing")
                    soft_missing = ats_data.get("soft_skills", {}).get("missing", [])
                    if soft_missing:
                        st.write(", ".join(soft_missing))
                    else:
                        st.write("None")

# Main app flow
def main():
    # Check if user needs to be redirected to resume creation first
    if not user_resume or len(user_resume) == 0:
        st.warning("Please create a resume first!")
        if st.button("Go to Resume Creator"):
            st.switch_page("pages/main.py")
        return
    
    # Determine what to show based on session state
    if st.session_state.show_visual_editor:
        show_visual_editor_with_tools()
    else:
        show_template_selector()

if __name__ == "__main__":
    main()