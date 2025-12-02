import streamlit as st
from PIL import Image
from utils import chatbot,show_login_modal, get_user_resume, load_users, load_user_templates, load_user_doc_templates, save_user_templates, replace_content, save_user_doc_templates, load_user_ppt_templates, analyze_slide_structure, generate_ppt_sections, match_generated_to_original, clear_and_replace_text, save_user_ppt_templates
from streamlit_extras.stylable_container import stylable_container
from pages.download import SYSTEM_TEMPLATES, generate_generic_html




st.set_page_config(
    page_title="Resume Creator",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="collapsed"
)





if "logged_in_user" not in st.session_state or st.session_state.logged_in_user is None:
    logged_user = st.query_params.get("user")
    if logged_user:
        st.session_state.logged_in_user = logged_user
        
        st.query_params["user"] = logged_user
    else:
        st.session_state.logged_in_user = None
else:
    
    if st.session_state.logged_in_user:
        st.query_params["user"] = st.session_state.logged_in_user




if "show_login_modal" not in st.session_state:
    st.session_state.show_login_modal = False

if "show_template_modal" not in st.session_state:
    st.session_state.show_template_modal = False

if "scroll_to" not in st.session_state:
    st.session_state.scroll_to = None

if "template_view_mode" not in st.session_state:
    st.session_state.template_view_mode = "custom"

if 'current_template_type' not in st.session_state:
    st.session_state.current_template_type = "html"




st.markdown("""
<style>
#MainMenu, footer, header, button[kind="header"] {visibility: hidden;}
.stMainBlockContainer, div.block-container, [data-testid="stMainBlockContainer"] {
    padding-top: 0rem !important;
    margin-top: 0rem !important;
}

.nav-wrapper {
    position: fixed;
    top: 20px;
    left: 50%;
    transform: translateX(-50%);
    width: 90%;
    max-width: 1200px;
    z-index: 99999 !important;
    background-color: white !important;
    padding: 0.8rem 2rem;
    box-shadow: 0 2px 20px rgba(0,0,0,0.1);
    display: flex;
    align-items: center;
    justify-content: space-between;
    border-radius: 50px;
}

.logo {
    font-size: 24px;
    font-weight: 400;
    color: #2c3e50;
    font-family: 'Nunito Sans', sans-serif !important;
    letter-spacing: -0.5px;
}

.nav-menu {
    display: flex;
    gap: 2rem;
    align-items: center;
}

.nav-item { position: relative; }

.nav-link {
    color: #000000 !important;
    text-decoration: none !important;
    font-size: 1rem;
    font-family: 'Nunito Sans', sans-serif;
    padding: 0.5rem 1rem;
    border-radius: 8px;
    cursor: pointer;
    transition: all 0.3s ease;
}

.nav-link:visited {
    color: #000000 !important;
}

.nav-link:hover {
    background-color: #fff5f0;
    color: #ff8c42 !important;
}

.hero-title {
    font-size: 3rem;
    font-weight: 500;
    color: #0a0f14;
    margin-bottom: 1rem;
    margin-top: 1rem !important;
    font-family: 'Nunito Sans', sans-serif !important;
    margin-left: 1rem !important;
}

.hero-subtitle {
    font-size: 1.2rem;
    color: #64748b;
    margin-top: 2rem;
    font-family: 'Nunito Sans', sans-serif !important;
    margin-left: 1rem !important;
}

.about-image-container {
    height: 650px;
    overflow: hidden;
    border-radius: 16px;
    box-shadow: 0 4px 6px rgba(0,0,0,0.1);
    margin-top:-100px !important;
}

.about-image-container img {
    width: 100%;
    height: 100%;
    object-fit: cover;
    object-position: center;
}

.about-card {
    background: white;
    border-radius: 16px;
    padding: 2rem;
    box-shadow: 0 4px 6px rgba(0,0,0,0.1);
    max-width: 800px;
}

.about-header {
    color: #ff7043;
    font-size: 1.2rem;
    font-weight: 600;
    margin-bottom: 1.5rem;
    font-family: 'Nunito Sans', sans-serif;
}

.about-title {
    color: #0a0f14;
    font-size: 2rem;
    font-weight: 600;
    margin-bottom: 1rem;
    font-family: 'Nunito Sans', sans-serif;
}

.about-description {
    color: #64748b;
    font-size: 1rem;
    line-height: 1.6;
    font-family: 'Nunito Sans', sans-serif;
    margin-bottom: 1.2rem;
}

.about-description:only-child {
    margin-bottom: 0 !important;
}

.info-grid {
    display: grid;
    grid-template-columns: repeat(2, 1fr);
    gap: 1.5rem;
    background: #f8fafc;
    padding: 1.5rem;
    border-radius: 12px;
}

.info-item {
    display: flex;
    flex-direction: column;
}

.info-label {
    color: #64748b;
    font-size: 0.9rem;
    margin-bottom: 0.3rem;
    font-family: 'Nunito Sans', sans-serif;
}

.info-value {
    color: #0a0f14;
    font-size: 1.1rem;
    font-weight: 600;
    font-family: 'Nunito Sans', sans-serif;
}

.resume-container {
    max-width: 1000px;
    margin: 0 auto;
    padding: 2rem;
}

.section-header {
    text-align: center;
    color: #8b6f47;
    font-size: 2.5rem;
    font-weight: 500;
    margin-bottom: 0.5rem;
    font-family: 'Playfair Display', serif !important;
    letter-spacing: 0.5px;
}

.section-divider-wave {
    width: 200px;
    height: 20px;
    margin: 0.5rem auto 2rem;
    background-image: url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 200 20'%3E%3Cpath d='M0,10 Q25,0 50,10 T100,10 T150,10 T200,10' stroke='%23c89665' stroke-width='3' fill='none'/%3E%3C/svg%3E");
    background-repeat: no-repeat;
    background-position: center;
    background-size: contain;
}

.section-subtitle {
    text-align: center;
    color: #64748b;
    font-size: 1rem;
    margin-bottom: 3rem;
    font-family: 'Nunito Sans', sans-serif;
}

.work-experience-section {
    margin-top: 3rem;
}

.work-experience-title {
    color: #1e3a5f;
    font-size: 1.8rem;
    font-weight: 600;
    margin-bottom: 1rem;
    font-family: 'Nunito Sans', sans-serif;
}

.work-experience-subtitle {
    color: #64748b;
    font-size: 1rem;
    margin-bottom: 2rem;
    font-family: 'Nunito Sans', sans-serif;
}

.experience-item {
    display: flex;
    gap: 2rem;
    margin-bottom: 2rem;
    padding-bottom: 2rem;
    border-bottom: 1px solid #e2e8f0;
}

.experience-timeline {
    flex-shrink: 0;
}

.company-name {
    color: #1e3a5f;
    font-size: 1.1rem;
    font-weight: 600;
    margin-bottom: 0.3rem;
    font-family: 'Nunito Sans', sans-serif;
}

.experience-date {
    color: #ff7043;
    font-size: 0.9rem;
    font-family: 'Nunito Sans', sans-serif;
}

.experience-content {
    flex: 1;
}

.experience-role {
    color: #1e3a5f;
    font-size: 1.3rem;
    font-weight: 600;
    margin-bottom: 0.8rem;
    font-family: 'Nunito Sans', sans-serif;
}

.experience-description {
    color: #64748b;
    font-size: 0.95rem;
    line-height: 1.6;
    font-family: 'Nunito Sans', sans-serif;
}

.timeline-dot {
    width: 12px;
    height: 12px;
    background: #ff7043;
    border-radius: 50%;
    margin-top: 6px;
}

.no-resume-message {
    text-align: center;
    color: #64748b;
    font-size: 1.1rem;
    padding: 3rem;
    background: #f8fafc;
    border-radius: 12px;
    margin: 2rem auto;
    max-width: 600px;
}

.template-section {
    padding: 100px 0;
    min-height: 100vh;
    background: #f8fafc;
}

.template-header {
    text-align: center;
    margin-bottom: 3rem;
}

.template-title {
    font-size: 2.5rem;
    font-weight: 600;
    color: #0a0f14;
    font-family: 'Nunito Sans', sans-serif;
    margin-bottom: 1rem;
}

.template-subtitle {
    font-size: 1.1rem;
    color: #64748b;
    font-family: 'Nunito Sans', sans-serif;
}

.template-toggle-container {
    display: flex;
    justify-content: center;
    gap: 1rem;
    margin: 3rem auto;
    max-width: 500px;
}

.template-toggle-btn {
    padding: 12px 40px;
    border-radius: 50px;
    font-size: 1rem;
    font-weight: 600;
    font-family: 'Nunito Sans', sans-serif;
    cursor: pointer;
    transition: all 0.3s ease;
    border: 2px solid #e2e8f0;
    background: white;
    color: #64748b;
}

.template-toggle-btn.active {
    background: #e87532;
    color: white;
    border-color: #e87532;
}

.template-toggle-btn:hover {
    transform: translateY(-2px);
    box-shadow: 0 4px 12px rgba(255, 107, 74, 0.25);
    border-color: #e87532;
    color: #e87532;
}

.template-toggle-btn.active:hover {
    color: white;
}

.template-type-tabs {
    display: flex;
    justify-content: center;
    gap: 1rem;
    margin: 2rem auto;
    flex-wrap: wrap;
    max-width: 800px;
}

.template-type-btn {
    padding: 10px 30px;
    border-radius: 50px;
    font-size: 0.95rem;
    font-weight: 500;
    font-family: 'Nunito Sans', sans-serif;
    cursor: pointer;
    transition: all 0.3s ease;
    border: 2px solid #e2e8f0;
    background: white;
    color: #64748b;
}

.template-type-btn.active {
    background: #e87532;
    color: white;
    border-color: #e87532;
}

.template-type-btn:hover {
    border-color: #e87532;
    color: #e87532;
    background: #fff5f2;
}

.template-type-btn.active:hover {
    color: white;
    background: #e87532;
}

.template-grid {
    display: grid;
    grid-template-columns: repeat(auto-fill, minmax(280px, 1fr));
    gap: 2rem;
    max-width: 1200px;
    margin: 0 auto;
    padding: 2rem;
}

.template-card-new {
    background: white;
    border-radius: 16px;
    overflow: visible;
    box-shadow: none;
    transition: all 0.3s ease;
    border: 3px solid #e87532;
    display: flex;
    flex-direction: column;
}

.template-card-new:hover {
    transform: translateY(-8px);
    box-shadow: 0 12px 24px rgba(232, 117, 50, 0.25);
    border-color: #d66629;
}

.template-image {
    display: none;
}

.template-info {
    padding: 2rem 1.5rem;
    background: white;
    text-align: center;
    min-height: 140px;
    display: flex;
    flex-direction: column;
    justify-content: center;
    border-radius: 13px 13px 0 0;
}

.template-category {
    color: #e87532;
    font-size: 0.75rem;
    font-weight: 700;
    text-transform: uppercase;
    letter-spacing: 1px;
    margin-bottom: 0.5rem;
    font-family: 'Nunito Sans', sans-serif;
}

.template-name {
    color: #0a0f14;
    font-size: 1.15rem;
    font-weight: 700;
    margin-bottom: 0.5rem;
    font-family: 'Nunito Sans', sans-serif;
    line-height: 1.4;
}

.template-description {
    color: #64748b;
    font-size: 0.85rem;
    line-height: 1.5;
    margin-bottom: 0;
    font-family: 'Nunito Sans', sans-serif;
}

.template-meta {
    display: none;
}

.template-actions {
    display: flex;
    gap: 0.5rem;
    padding: 0 1.5rem 1.5rem;
}

.template-action-icon {
    width: 40px;
    height: 40px;
    border-radius: 50%;
    display: flex;
    align-items: center;
    justify-content: center;
    background: #f8fafc;
    color: #64748b;
    transition: all 0.3s ease;
    cursor: pointer;
}

.template-action-icon:hover {
    background: #e87532;
    color: white;
}

.template-action-icon.primary {
    background: #e87532;
    color: white;
}

.template-action-icon.primary:hover {
    background: #ff5733;
}

div[data-testid="stButton"] > button:hover,
.stButton > button:hover {
    background-color: white !important;
    color: #ff8c42 !important;
    border: 2px solid #ff8c42 !important;
}

/* Template Toggle Buttons - Default State (Not Selected) */
div[data-testid="stButton"] > button {
    border: 2px solid #e87532 !important;
}

/* Scroll to Top Button */
.scroll-to-top {
    bottom: 30px;
    right: 30px;
    width: 50px;
    height: 50px;
    background: #e87532;
    color: white;
    border-radius: 50%;
    display: flex;
    align-items: center;
    justify-content: center;
    cursor: pointer;
    box-shadow: 0 4px 12px rgba(232, 117, 50, 0.4);
    transition: all 0.3s ease;
    opacity: 0;
    visibility: hidden;
    z-index: 9998;
    font-size: 24px;
    border: none;
}

.scroll-to-top.show {
    opacity: 1;
    visibility: visible;
}

.scroll-to-top:hover {
    background: white !important;
    color: #ff8c42 !important;
    border: 2px solid #ff8c42 !important;
    transform: translateY(-5px);
    box-shadow: 0 6px 16px rgba(232, 117, 50, 0.6);
}

.scroll-to-top:active {
    transform: translateY(-2px);
}
</style>
""", unsafe_allow_html=True)





if st.query_params.get("logout") == "true":
    st.session_state.logged_in_user = None
    for key in list(st.session_state.keys()):
        del st.session_state[key]
    st.query_params.clear()
    st.rerun()


if st.query_params.get("home") == "true":
    if "home" in st.query_params:
        del st.query_params["home"]
    
    if st.session_state.logged_in_user:
        st.query_params["user"] = st.session_state.logged_in_user
    st.rerun()




if st.session_state.get('logged_in_user') and not st.query_params.get("user"):
    st.query_params["user"] = st.session_state.logged_in_user




current_user = st.session_state.get('logged_in_user', '')
is_logged_in = bool(current_user)





current_user = st.session_state.get('logged_in_user', '')
is_logged_in = bool(current_user)

# Build the home URL - only add user param if logged in, otherwise just scroll to top
if is_logged_in and current_user:
    home_url = f"?home=true&user={current_user}"
else:
    home_url = "#Home"  # Just scroll to home section, don't trigger any navigation

if is_logged_in:
    auth_button = '<div class="nav-item"><a class="nav-link" href="?logout=true" target="_self">Logout</a></div>'
else:
    auth_button = '<div class="nav-item"><a class="nav-link" data-section="Login" href="#Login">Login</a></div>'

st.markdown(f"""
<div class="nav-wrapper">
    <div class="logo">Resume Creator</div>
    <div class="nav-menu">
        <div class="nav-item">
            <a class="nav-link" href="{home_url}" target="_self">Home</a>
        </div>
        <div class="nav-item">
            <a class="nav-link" data-section="Templates" href="#Templates">Templates</a>
        </div>
        {auth_button}
    </div>
</div>
""", unsafe_allow_html=True)

st.markdown('<div style="height: 80px;"></div>', unsafe_allow_html=True)

# Handle the home button click ONLY for logged-in users
if st.query_params.get("home") == "true":
    # Check if user is logged in
    if st.session_state.get('logged_in_user'):
        # User is logged in, allow home navigation
        if "home" in st.query_params:
            del st.query_params["home"]
        
        if st.session_state.logged_in_user:
            st.query_params["user"] = st.session_state.logged_in_user
        st.rerun()
    else:
        # User is NOT logged in, remove the home param and don't unlock
        if "home" in st.query_params:
            del st.query_params["home"]
        # Don't rerun, just remove the param
        st.rerun()




st.markdown('<div id="Home"></div>', unsafe_allow_html=True)

with st.container():
    col1, col2 = st.columns([1, 1], gap="large")

    with col1:
        st.markdown("""
        <div style="padding: 2rem 0;">
            <div class="hero-title">Creating Resume with Help of AI</div>
            <div class="hero-subtitle">
                Transforming Data into elegant solutions through creative design and innovative development
            </div>
        </div>
        """, unsafe_allow_html=True)

        btn1, btn2 = st.columns([0.9, 1])

        with btn1:
            with stylable_container(
                "create-resume-btn",
                css_styles="""
                    button {
                        background-color: #e87532 !important;
                        color: #ffffff !important;
                        padding: 15px 35px !important;
                        border-radius: 50px !important;
                        font-weight: 600 !important;
                        border: 2px solid #e87532 !important;
                    }
                    button:hover {
                        background-color: #ffffff !important;
                        color: #e87532 !important;
                    }
                """
            ):
                create_resume_clicked = st.button("Create Resume")

        with btn2, stylable_container(
            "template-btn",
            css_styles="""
                button {
                    background-color: #ffffff !important;
                    color: #e87532 !important;
                    padding: 15px 35px !important;
                    border-radius: 50px !important;
                    font-weight: 600 !important;
                    border: 1px solid #e87532 !important;
                }
                button:hover {
                    background-color: #e87532 !important;
                    color: #ffffff !important;
                }
            """
        ):
            template_clicked = st.button("Show Template")

    with col2:
        try:
            img = Image.open(r"C:\\ask.ai\\image\\image.png")
            st.image(img, use_container_width=True)
        except:
            st.info("Image not found. Check the path.")

st.markdown('<div style="padding: 10px 0; min-height: 50px;">', unsafe_allow_html=True)


if create_resume_clicked:
    if st.session_state.logged_in_user is None:
        st.warning("🔒 Please login first to create a resume.")
        st.session_state.show_login_modal = True
    else:
        email = st.session_state.logged_in_user
        users = load_users()
        user_entry = users.get(email)
        
        
        st.query_params["user"] = email
        
        if isinstance(user_entry, dict):
            st.session_state.username = email.split('@')[0]
            user_resume = get_user_resume(email)
            
            if user_resume and len(user_resume) > 0:
                st.session_state.resume_source = user_resume
                st.session_state.input_method = user_resume.get("input_method", "Manual Entry")
                st.switch_page("pages/job.py")
            else:
                st.switch_page("pages/main.py")
        else:
            st.switch_page("pages/main.py")




if st.session_state.logged_in_user is None:
    st.markdown(
        "<p style='text-align:center; color:#e87532; font-weight:600; margin-top:1rem;margin-bottom:10rem;'>🔒 Please login to access all sections.</p>",
        unsafe_allow_html=True
    )
    st.markdown('<div id="Login"></div>', unsafe_allow_html=True)
    show_login_modal()
    st.stop()




email = st.session_state.logged_in_user


if not st.query_params.get("user"):
    st.query_params["user"] = email


user_resume = get_user_resume(email)
has_resume = user_resume and len(user_resume) > 0
chatbot(user_resume)

st.markdown('<div id="Templates"></div>', unsafe_allow_html=True)


st.markdown("""
<div class="template-header">
    <div class="section-header">Templates</div>
    <div class="section-divider-wave"></div>
    <div class="template-subtitle">Choose from our collection of professional templates</div>
</div>
""", unsafe_allow_html=True)

# Template Mode Toggle
col1, col2, col3 = st.columns([1, 1, 1])
with col2:
    toggle_col1, toggle_col2 = st.columns(2)
    with toggle_col1:
        custom_style = """
            button {
                background: """ + ("#e87532 !important" if st.session_state.template_view_mode == "custom" else "white !important") + """;
                color: """ + ("white !important" if st.session_state.template_view_mode == "custom" else "#64748b !important") + """;
                border: 2px solid """ + ("#e87532 !important" if st.session_state.template_view_mode == "custom" else "#e2e8f0 !important") + """;
                border-radius: 50px !important;
                padding: 2px 50px !important;
                font-weight: 600 !important;
                transition: all 0.3s ease !important;
            }
            button:hover {
                border-color: #e87532 !important;
                color: """ + ("white !important" if st.session_state.template_view_mode == "custom" else "#e87532 !important") + """;
                background: """ + ("#e87532 !important" if st.session_state.template_view_mode == "custom" else "#fff5f2 !important") + """;
            }
        """
        with stylable_container("custom_toggle", css_styles=custom_style):
            if st.button("Custom Templates", key="custom_toggle_btn", use_container_width=True):
                st.session_state.template_view_mode = "custom"
                st.rerun()
    with toggle_col2:
        system_style = """
            button {
                background: """ + ("#e87532 !important" if st.session_state.template_view_mode == "system" else "white !important") + """;
                color: """ + ("white !important" if st.session_state.template_view_mode == "system" else "#64748b !important") + """;
                border: 2px solid """ + ("#e87532 !important" if st.session_state.template_view_mode == "system" else "#e2e8f0 !important") + """;
                border-radius: 50px !important;
                padding: 2px 40px !important;
                font-weight: 600 !important;
                transition: all 0.3s ease !important;
            }
            button:hover {
                border-color: #e87532 !important;
                color: """ + ("white !important" if st.session_state.template_view_mode == "system" else "#e87532 !important") + """;
                background: """ + ("#e87532 !important" if st.session_state.template_view_mode == "system" else "#fff5f2 !important") + """;
            }
        """
        with stylable_container("system_toggle", css_styles=system_style):
            if st.button("System Templates", key="system_toggle_btn", use_container_width=True):
                st.session_state.template_view_mode = "system"
                st.rerun()

st.markdown("<br>", unsafe_allow_html=True)


if st.session_state.template_view_mode == "custom":
    if 'uploaded_templates' not in st.session_state:
        st.session_state.uploaded_templates = load_user_templates(st.session_state.logged_in_user)
    
    if 'doc_templates' not in st.session_state:
        st.session_state.doc_templates = load_user_doc_templates(st.session_state.logged_in_user)
    
    if 'ppt_templates' not in st.session_state:
        st.session_state.ppt_templates = load_user_ppt_templates(st.session_state.logged_in_user)

    
    st.markdown('<div class="template-type-tabs">', unsafe_allow_html=True)
    type_col1, type_col2, type_col3 = st.columns([1, 1, 1])
    
    with type_col1:
        html_style = """
            button {
                background: """ + ("#e87532 !important" if st.session_state.current_template_type == "html" else "white !important") + """;
                color: """ + ("white !important" if st.session_state.current_template_type == "html" else "#64748b !important") + """;
                border: 2px solid """ + ("#e87532 !important" if st.session_state.current_template_type == "html" else "#e2e8f0 !important") + """;
                border-radius: 50px !important;
                padding: 10px 30px !important;
                font-weight: 500 !important;
                transition: all 0.3s ease !important;
            }
            button:hover {
                border-color: #e87532 !important;
                color: """ + ("white !important" if st.session_state.current_template_type == "html" else "#e87532 !important") + """;
                background: """ + ("#e87532 !important" if st.session_state.current_template_type == "html" else "#fff5f2 !important") + """;
            }
        """
        with stylable_container("html_type_btn", css_styles=html_style):
            if st.button("HTML Templates", key="html_type", use_container_width=True):
                st.session_state.current_template_type = "html"
                st.rerun()
    
    with type_col2:
        word_style = """
            button {
                background: """ + ("#e87532 !important" if st.session_state.current_template_type == "word" else "white !important") + """;
                color: """ + ("white !important" if st.session_state.current_template_type == "word" else "#64748b !important") + """;
                border: 2px solid """ + ("#e87532 !important" if st.session_state.current_template_type == "word" else "#e2e8f0 !important") + """;
                border-radius: 50px !important;
                padding: 10px 30px !important;
                font-weight: 500 !important;
                transition: all 0.3s ease !important;
            }
            button:hover {
                border-color: #e87532 !important;
                color: """ + ("white !important" if st.session_state.current_template_type == "word" else "#e87532 !important") + """;
                background: """ + ("#e87532 !important" if st.session_state.current_template_type == "word" else "#fff5f2 !important") + """;
            }
        """
        with stylable_container("word_type_btn", css_styles=word_style):
            if st.button("Word Templates", key="word_type", use_container_width=True):
                st.session_state.current_template_type = "word"
                st.rerun()
    
    with type_col3:
        ppt_style = """
            button {
                background: """ + ("#e87532 !important" if st.session_state.current_template_type == "ppt" else "white !important") + """;
                color: """ + ("white !important" if st.session_state.current_template_type == "ppt" else "#64748b !important") + """;
                border: 2px solid """ + ("#e87532 !important" if st.session_state.current_template_type == "ppt" else "#e2e8f0 !important") + """;
                border-radius: 50px !important;
                padding: 10px 30px !important;
                font-weight: 500 !important;
                transition: all 0.3s ease !important;
            }
            button:hover {
                border-color: #e87532 !important;
                color: """ + ("white !important" if st.session_state.current_template_type == "ppt" else "#e87532 !important") + """;
                background: """ + ("#e87532 !important" if st.session_state.current_template_type == "ppt" else "#fff5f2 !important") + """;
            }
        """
        with stylable_container("ppt_type_btn", css_styles=ppt_style):
            if st.button("PowerPoint Templates", key="ppt_type", use_container_width=True):
                st.session_state.current_template_type = "ppt"
                st.rerun()
    
    st.markdown('</div>', unsafe_allow_html=True)

    # Display templates based on selected type
    st.markdown('<div class="template-grid">', unsafe_allow_html=True)
    
    if st.session_state.current_template_type == "html":
        if st.session_state.uploaded_templates:
            cols = st.columns(4)
            for idx, (template_id, template_data) in enumerate(st.session_state.uploaded_templates.items()):
                with cols[idx % 4]:
                    st.markdown(f"""
                    <div class="template-card-new">
                        <div class="template-info">
                            <div class="template-category">HTML</div>
                            <div class="template-name">{template_data['name']}</div>
                            <div class="template-description">{template_data['original_filename']}</div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)

                    with stylable_container(
                        f"use_html_{template_id}",
                        css_styles="""
                            button {
                                background-color: #e87532 !important;
                                color: white !important;
                                border: none !important;
                                border-radius: 12px !important;
                                padding: 12px 20px !important;
                                font-weight: 600 !important;
                                font-size: 0.95rem !important;
                                margin-top: 0.5rem !important;
                                width: 100% !important;
                            }
                            button:hover {
                                background-color: #d66629 !important;
                                transform: translateY(-2px);
                            }
                        """
                    ):
                        if st.button("Preview", key=f"use_html_{template_id}", use_container_width=True):
                            if 'temp_upload_config' in st.session_state:
                                del st.session_state.temp_upload_config
                            
                            st.session_state.selected_template = template_data['name']
                            st.session_state.selected_template_config = template_data
                            st.session_state.template_source = 'saved'
                            st.session_state.current_upload_id = template_id
                            
                            
                            st.query_params["user"] = st.session_state.logged_in_user
                            st.switch_page("pages/template_preview.py")
                    
                    with stylable_container(
                        f"delete_html_{template_id}",
                        css_styles="""
                            button {
                                background-color: #dc3545 !important;
                                color: white !important;
                                border: none !important;
                                border-radius: 8px !important;
                                padding: 8px 16px !important;
                                font-weight: 500 !important;
                                font-size: 0.85rem !important;
                                margin-top: 0.5rem !important;
                                width: 100% !important;
                            }
                            button:hover {
                                background-color: #c82333 !important;
                            }
                        """
                    ):
                        if st.button("Delete", key=f"delete_html_{template_id}", use_container_width=True):
                            if st.session_state.get('current_upload_id') == template_id:
                                st.session_state.pop('selected_template_preview', None)
                                st.session_state.pop('selected_template', None)
                                st.session_state.pop('selected_template_config', None)
                                st.session_state.pop('current_upload_id', None)
                            
                            del st.session_state.uploaded_templates[template_id]
                            save_user_templates(st.session_state.logged_in_user, st.session_state.uploaded_templates)
                            
                            
                            st.query_params["user"] = st.session_state.logged_in_user
                            st.success(f"Deleted: {template_data['name']}")
                            st.rerun()
        else:
            st.markdown("""
            <div class="no-resume-message">
                <h3 style="color: #e87532;">No HTML Templates</h3>
                <p>Upload your first HTML template to get started</p>
            </div>
            """, unsafe_allow_html=True)

    elif st.session_state.current_template_type == "word":
        if st.session_state.doc_templates:
            cols = st.columns(4)
            for idx, (template_id, template_data) in enumerate(st.session_state.doc_templates.items()):
                with cols[idx % 4]:
                    st.markdown(f"""
                    <div class="template-card-new">
                        <div class="template-info">
                            <div class="template-category">WORD</div>
                            <div class="template-name">{template_data['name']}</div>
                            <div class="template-description">{template_data['original_filename']}</div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    with stylable_container(
                        f"use_doc_{template_id}",
                        css_styles="""
                            button {
                                background-color: #e87532 !important;
                                color: white !important;
                                border: none !important;
                                border-radius: 12px !important;
                                padding: 12px 20px !important;
                                font-weight: 600 !important;
                                font-size: 0.95rem !important;
                                margin-top: 0.5rem !important;
                                width: 100% !important;
                            }
                            button:hover {
                                background-color: #d66629 !important;
                                transform: translateY(-2px);
                            }
                        """
                    ):
                        if st.button("Preview", key=f"use_doc_{template_id}", use_container_width=True):
                            try:
                                import io
                                from docx import Document
                                
                                doc_stream = io.BytesIO(template_data['doc_data'])
                                doc = Document(doc_stream)
                                structure = template_data.get('structure', [])
                                output, replaced, removed = replace_content(doc, structure, user_resume)
                                
                                st.session_state.generated_doc = output.getvalue()
                                st.session_state.selected_doc_template_id = template_id
                                st.session_state.selected_doc_template = template_data
                                st.session_state.template_source = 'doc_saved'
                                st.session_state.selected_template = template_data['name']
                                st.session_state.selected_template_config = template_data
                                
                                
                                st.query_params["user"] = st.session_state.logged_in_user
                                st.switch_page("pages/template_preview.py")
                            except Exception as e:
                                st.error(f"Error: {str(e)}")

                    with stylable_container(
                        f"delete_doc_{template_id}",
                        css_styles="""
                            button {
                                background-color: #dc3545 !important;
                                color: white !important;
                                border: none !important;
                                border-radius: 8px !important;
                                padding: 8px 16px !important;
                                font-weight: 500 !important;
                                font-size: 0.85rem !important;
                                margin-top: 0.5rem !important;
                                width: 100% !important;
                            }
                            button:hover {
                                background-color: #c82333 !important;
                            }
                        """
                    ):
                        if st.button("Delete", key=f"delete_doc_{template_id}", use_container_width=True):
                            if st.session_state.get('selected_doc_template_id') == template_id:
                                st.session_state.pop('generated_doc', None)
                                st.session_state.pop('selected_doc_template_id', None)
                                st.session_state.pop('selected_doc_template', None)
                                st.session_state.pop('doc_template_source', None)
                            
                            del st.session_state.doc_templates[template_id]
                            save_user_doc_templates(st.session_state.logged_in_user, st.session_state.doc_templates)
                            
                            
                            st.query_params["user"] = st.session_state.logged_in_user
                            st.success(f"Deleted: {template_data['name']}")
                            st.rerun()
        else:
            st.markdown("""
            <div class="no-resume-message">
                <h3 style="color: #e87532;">No Word Templates</h3>
                <p>Upload your first Word template to get started</p>
            </div>
            """, unsafe_allow_html=True)

    elif st.session_state.current_template_type == "ppt":
        if st.session_state.ppt_templates:
            cols = st.columns(4)
            for idx, (template_id, template_data) in enumerate(st.session_state.ppt_templates.items()):
                with cols[idx % 4]:
                    st.markdown(f"""
                    <div class="template-card-new">
                        <div class="template-info">
                            <div class="template-category">POWERPOINT</div>
                            <div class="template-name">{template_data['name']}</div>
                            <div class="template-description">{template_data['original_filename']}</div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    with stylable_container(
                        f"use_ppt_{template_id}",
                        css_styles="""
                            button {
                                background-color: #e87532 !important;
                                color: white !important;
                                border: none !important;
                                border-radius: 12px !important;
                                padding: 12px 20px !important;
                                font-weight: 600 !important;
                                font-size: 0.95rem !important;
                                margin-top: 0.5rem !important;
                                width: 100% !important;
                            }
                            button:hover {
                                background-color: #d66629 !important;
                                transform: translateY(-2px);
                            }
                        """
                    ):
                        if st.button("Preview", key=f"use_ppt_{template_id}", use_container_width=True):
                            try:
                                import io
                                from pptx import Presentation
                                
                                working_prs = Presentation(io.BytesIO(template_data['ppt_data']))
                                prs = Presentation(io.BytesIO(template_data['ppt_data']))
                                
                                slide_texts = []
                                for slide_idx, slide in enumerate(prs.slides):
                                    text_blocks = []
                                    for shape_idx, shape in enumerate(slide.shapes):
                                        if shape.has_text_frame and shape.text.strip():
                                            text_blocks.append({
                                                "index": shape_idx,
                                                "text": shape.text.strip(),
                                                "position": {"x": shape.left, "y": shape.top}
                                            })
                                    
                                    if text_blocks:
                                        text_blocks.sort(key=lambda x: (x["position"]["y"], x["position"]["x"]))
                                        slide_texts.append({
                                            "slide_number": slide_idx + 1,
                                            "text_blocks": text_blocks
                                        })
                                
                                structured_slides = analyze_slide_structure(slide_texts)
                                generated_sections = generate_ppt_sections(user_resume, structured_slides)
                                
                                text_elements = template_data['text_elements']
                                content_mapping, heading_shapes, basic_info_shapes = match_generated_to_original(
                                    text_elements, generated_sections, prs)
                                
                                edits = {}
                                for element in text_elements:
                                    key = f"{element['slide']}_{element['shape']}"
                                    if key not in heading_shapes:
                                        edits[key] = content_mapping.get(key, element['original_text'])
                                
                                success_count = 0
                                for element in text_elements:
                                    key = f"{element['slide']}_{element['shape']}"
                                    if key not in heading_shapes and key in edits:
                                        slide_idx = element['slide'] - 1
                                        shape_idx = element['shape']
                                        
                                        if slide_idx < len(working_prs.slides):
                                            slide = working_prs.slides[slide_idx]
                                            if shape_idx < len(slide.shapes):
                                                shape = slide.shapes[shape_idx]
                                                if shape.has_text_frame:
                                                    clear_and_replace_text(shape, edits[key])
                                                    success_count += 1
                                
                                output = io.BytesIO()
                                working_prs.save(output)
                                output.seek(0)
                                
                                st.session_state.generated_ppt = output.getvalue()
                                st.session_state.selected_ppt_template_id = template_id
                                st.session_state.selected_ppt_template = template_data
                                st.session_state.template_source = 'ppt_saved'
                                st.session_state.selected_template = template_data['name']
                                st.session_state.selected_template_config = template_data
                                
                                
                                st.query_params["user"] = st.session_state.logged_in_user
                                st.switch_page("pages/template_preview.py")
                            except Exception as e:
                                st.error(f"Error: {str(e)}")

                    with stylable_container(
                        f"delete_ppt_{template_id}",
                        css_styles="""
                            button {
                                background-color: #dc3545 !important;
                                color: white !important;
                                border: none !important;
                                border-radius: 8px !important;
                                padding: 8px 16px !important;
                                font-weight: 500 !important;
                                font-size: 0.85rem !important;
                                margin-top: 0.5rem !important;
                                width: 100% !important;
                            }
                            button:hover {
                                background-color: #c82333 !important;
                            }
                        """
                    ):
                        if st.button("Delete", key=f"delete_ppt_{template_id}", use_container_width=True):
                            if st.session_state.get('selected_ppt_template_id') == template_id:
                                st.session_state.pop('generated_ppt', None)
                                st.session_state.pop('selected_ppt_template_id', None)
                                st.session_state.pop('selected_ppt_template', None)
                                st.session_state.pop('ppt_template_source', None)
                            
                            del st.session_state.ppt_templates[template_id]
                            save_user_ppt_templates(st.session_state.logged_in_user, st.session_state.ppt_templates)
                            
                            
                            st.query_params["user"] = st.session_state.logged_in_user
                            st.success(f"Deleted: {template_data['name']}")
                            st.rerun()
        else:
            st.markdown("""
            <div class="no-resume-message">
                <h3 style="color: #e87532;">No PowerPoint Templates</h3>
                <p>Upload your first PowerPoint template to get started</p>
            </div>
            """, unsafe_allow_html=True)

    st.markdown('</div>', unsafe_allow_html=True)


elif st.session_state.template_view_mode == "system":
    st.markdown('<div class="template-grid">', unsafe_allow_html=True)
    
    system_template_names = list(SYSTEM_TEMPLATES.keys())
    cols = st.columns(4)
    for idx, template_name in enumerate(system_template_names):
        with cols[idx % 4]:
            st.markdown(f"""
            <div class="template-card-new">
                <div class="template-info">
                    <div class="template-category">SYSTEM</div>
                    <div class="template-name">{template_name}</div>
                    <div class="template-description">Professional template</div>
                </div>
            </div>
            """, unsafe_allow_html=True)
            
            with stylable_container(
                f"preview_sys_{idx}",
                css_styles="""
                    button {
                        background-color: #e87532 !important;
                        color: white !important;
                        border: none !important;
                        border-radius: 12px !important;
                        padding: 12px 20px !important;
                        font-weight: 600 !important;
                        font-size: 0.95rem !important;
                        margin-top: 0.5rem !important;
                        width: 100% !important;
                    }
                    button:hover {
                        background-color: #d66629 !important;
                        transform: translateY(-2px);
                    }
                """
            ):
                if st.button("Preview", key=f"preview_sys_{idx}", use_container_width=True):
                    st.session_state.selected_template = template_name
                    st.session_state.selected_template_config = SYSTEM_TEMPLATES[template_name]
                    st.session_state.template_source = 'system'
                    
                    
                    st.query_params["user"] = st.session_state.logged_in_user
                    st.switch_page("pages/template_preview.py")
    
    st.markdown('</div>', unsafe_allow_html=True)

st.markdown('</div>', unsafe_allow_html=True)




if template_clicked:
    if st.session_state.logged_in_user is None:
        st.warning("🔒 Please login first to view templates.")
        st.session_state.show_login_modal = True
    else:
        st.session_state.from_template_button = True
        st.switch_page("pages/main.py")




st.markdown(f"""
<!-- Scroll to Top Button -->
<button class="scroll-to-top" id="scrollToTop" onclick="scrollToTop()">
    ↑
</button>

<script>
// CRITICAL: Preserve user login on page refresh and navigation
(function() {{
    const urlParams = new URLSearchParams(window.location.search);
    const userParam = urlParams.get('user');
    
    if (userParam) {{
        // Store in sessionStorage for backup on refresh
        sessionStorage.setItem('logged_in_user', userParam);
    }} else {{
        // Check if we have a stored user (for page refresh scenarios)
        const storedUser = sessionStorage.getItem('logged_in_user');
        if (storedUser && !window.location.search.includes('logout=true')) {{
            // Restore user param to URL
            const separator = window.location.search ? '&' : '?';
            const newUrl = window.location.pathname + window.location.search + separator + 'user=' + encodeURIComponent(storedUser);
            window.history.replaceState({{}}, '', newUrl);
            // Reload to apply the param
            window.location.reload();
        }}
    }}
    
    // Clear storage on logout
    if (window.location.search.includes('logout=true')) {{
        sessionStorage.removeItem('logged_in_user');
    }}
    
    // Add user param to all internal links that don't have logout
    document.addEventListener('DOMContentLoaded', function() {{
        const currentUser = urlParams.get('user') || sessionStorage.getItem('logged_in_user');
        if (currentUser) {{
            document.querySelectorAll('a').forEach(link => {{
                if (link.href && 
                    !link.href.includes('logout=true') && 
                    !link.href.includes('user=') &&
                    link.hostname === window.location.hostname) {{
                    const separator = link.href.includes('?') ? '&' : '?';
                    link.href += separator + 'user=' + encodeURIComponent(currentUser);
                }}
            }});
        }}
    }});
}})();

// Scroll to Top Button Functionality
const scrollToTopBtn = document.getElementById('scrollToTop');
const homeSection = document.getElementById('Home');

// Show/hide button based on scroll position
window.addEventListener('scroll', function() {{
    if (!homeSection) return;
    
    const homeSectionBottom = homeSection.offsetTop + homeSection.offsetHeight;
    const scrollPosition = window.pageYOffset || document.documentElement.scrollTop;
    
    // Show button when scrolled past the home section
    if (scrollPosition > homeSectionBottom) {{
        scrollToTopBtn.classList.add('show');
    }} else {{
        scrollToTopBtn.classList.remove('show');
    }}
}});

// Scroll to top function with smooth animation
function scrollToTop() {{
    window.scrollTo({{
        top: 0,
        behavior: 'smooth'
    }});
}}

// Optional: Hide button during scroll animation
let scrollTimeout;
scrollToTopBtn.addEventListener('click', function() {{
    clearTimeout(scrollTimeout);
    scrollTimeout = setTimeout(function() {{
        scrollToTopBtn.classList.remove('show');
    }}, 500);
}});
</script>
""", unsafe_allow_html=True)