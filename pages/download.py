import streamlit as st
import base64
from datetime import datetime
from pathlib import Path
import hashlib
import re
from streamlit_extras.switch_page_button import switch_page
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import streamlit as st
import json, os
from datetime import datetime
from utils import replace_content, extract_document_structure, save_user_ppt_templates,load_user_ppt_templates,load_user_templates,save_user_templates,get_css_sophisticated_minimal,get_css_clean_contemporary,get_css_elegant_professional,get_css_modern_minimal,get_css_date_below,get_css_classic,get_css_minimalist,get_css_horizontal,get_css_bold_title,get_css_section_box,analyze_slide_structure,generate_ppt_sections,match_generated_to_original,clear_and_replace_text


# Define the preferred display order for sections
RESUME_ORDER = ["summary", "experience", "education", "skills", "projects", "certifications", "achievements", "publications", "awards"]

# ATS-friendly color palette
ATS_COLORS = {
    "Professional Blue (Default)": "#1F497D",
    "Corporate Gray": "#4D4D4D",
    "Deep Burgundy": "#800020",
    "Navy Blue": "#000080",
    "Black":"#000000"
}



if 'uploaded_templates' not in st.session_state:
    st.session_state.uploaded_templates = {}

if "selected_template_config" not in st.session_state:
    st.session_state.selected_template_config = None

def format_section_title(key):
    """Converts keys like 'certifications' to 'Certifications'."""
    title = key.replace('_', ' ').replace('Skills', ' Skills').replace('summary', 'Summary')
    return ' '.join(word.capitalize() for word in title.split())


from datetime import datetime

def format_year_only(date_str):
    """Return only the year from a date string. If already a year, return as-is."""
    if not date_str:
        return ""
    date_str = str(date_str).strip()
    if len(date_str) == 4 and date_str.isdigit(): 
        return date_str
    try:
        dt = datetime.strptime(date_str, "%Y-%m-%d")
        return dt.strftime("%Y")
    except:
        return date_str  

def generate_generic_html(data, date_placement='right'):
    """Generates clean HTML content based on resume data, showing only years for all dates."""
    if not data: 
        return ""

    # Header
    job_title_for_header = data.get('job_title', '')
    contacts = [data.get('phone'), data.get('email'), data.get('location')]
    contacts_html = " | ".join([c for c in contacts if c])

    html = f'<div class="ats-header">'
    html += f'<h1>{data.get("name", "NAME MISSING")}</h1>'
    if job_title_for_header:
        html += f'<div class="ats-job-title-header">{job_title_for_header}</div>'
    if contacts_html:
        html += f'<div class="ats-contact">{contacts_html}</div>'
    html += '</div>'

    for key in RESUME_ORDER:
        section_data = data.get(key)
        if not section_data or (isinstance(section_data, list) and not section_data):
            continue

        title = format_section_title(key)
        html += f'<div class="ats-section-title">{title}</div>'

        # Summary
        if key == 'summary' and isinstance(section_data, str) and section_data.strip():
            html += f'<p style="margin-top:0;margin-bottom:5px;">{section_data}</p>'

        # Skills
        elif key == 'skills' and isinstance(section_data, dict):
            for skill_type, skill_list in section_data.items():
                if skill_list:
                    html += f'<div class="ats-skills-group"><strong>{format_section_title(skill_type)}:</strong> {", ".join(skill_list)}</div>'

        elif isinstance(section_data, list):
            for item in section_data:
                # Normalize string items to dict
                if isinstance(item, str) and item.strip():
                    item = {"title": item}
                
                if not isinstance(item, dict):
                    continue

                # Define comprehensive field keys for dynamic handling
                title_keys = ['position', 'title', 'name', 'degree', 'certificate_name', 'course']
                subtitle_keys = ['company', 'institution', 'issuer', 'organization', 'provider_name', 'university']
                duration_keys = ['duration', 'date', 'period', 'completed_date', 'start_date', 'end_date']
                description_keys = ['description', 'achievement', 'details', 'overview']

                # Get main title (first available)
                main_title = None
                for k in title_keys:
                    if k in item and item[k]:
                        main_title = item[k]
                        break
                
                # Get subtitle (first available that's not the same as title)
                subtitle = None
                for k in subtitle_keys:
                    if k in item and item[k] and item[k] != main_title:
                        subtitle = item[k]
                        break

                # Get duration with proper date formatting
                start = format_year_only(item.get('start_date', ''))
                end = format_year_only(item.get('end_date', ''))

                if start or end:
                    duration = f"{start} - {end}" if start and end else start or end
                else:
                    # Try other duration fields
                    duration = None
                    for k in duration_keys:
                        if k in item and item[k]:
                            duration = format_year_only(item[k])
                            break
                    duration = duration or ''

                # Skip completely empty items
                if not main_title and not subtitle and not duration:
                    # Check if there's any description content
                    has_description = False
                    for desc_key in description_keys:
                        if desc_key in item and item[desc_key]:
                            has_description = True
                            break
                    if not has_description:
                        continue

                # Build item HTML
                html += '<div class="ats-item-header">'
                if main_title or subtitle:
                    html += '<div class="ats-item-title-group">'
                    if main_title:
                        html += f'<span class="ats-item-title">{main_title}'
                        if subtitle:
                            html += f' <span class="ats-item-subtitle">{subtitle}</span>'
                        html += '</span>'
                    html += '</div>'
                if duration:
                    html += f'<div class="ats-item-duration">{duration}</div>'
                html += '</div>'

                # Description / Achievements - check all possible description fields
                description_list = None
                for desc_key in description_keys:
                    if desc_key in item and item[desc_key]:
                        description_list_raw = item[desc_key]
                        if isinstance(description_list_raw, str):
                            description_list = [description_list_raw]
                        elif isinstance(description_list_raw, list):
                            description_list = description_list_raw
                        break

                if description_list:
                    bullet_html = "".join([f"<li>{line}</li>" for line in description_list if line and str(line).strip()])
                    if bullet_html:
                        html += f'<ul class="ats-bullet-list">{bullet_html}</ul>'

    return html


SYSTEM_TEMPLATES = {
    "Minimalist (ATS Best)": {
        "html_generator": lambda data: generate_generic_html(data, date_placement='right'),
        "css_generator": get_css_minimalist,
    },
    "Horizontal Line": {
        "html_generator": lambda data: generate_generic_html(data, date_placement='right'),
        "css_generator": get_css_horizontal,
    },
    "Bold Title Accent": {
        "html_generator": lambda data: generate_generic_html(data, date_placement='right'),
        "css_generator": get_css_bold_title,
    },
    "Date Below": {
        "html_generator": lambda data: generate_generic_html(data, date_placement='below'),
        "css_generator": get_css_date_below,
    },
    "Section Box Header": {
        "html_generator": lambda data: generate_generic_html(data, date_placement='right'),
        "css_generator": get_css_section_box,
    },
    "Times New Roman Classic": {
        "html_generator": lambda data: generate_generic_html(data, date_placement='right'),
        "css_generator": get_css_classic,
    },
    "Sohisticated_minimal": {
        "html_generator": lambda data: generate_generic_html(data, date_placement='right'),
        "css_generator": get_css_sophisticated_minimal,
    },
        "Clean look": {
        "html_generator": lambda data: generate_generic_html(data, date_placement='right'),
        "css_generator": get_css_clean_contemporary,
    },
            "Elegant": {
        "html_generator": lambda data: generate_generic_html(data, date_placement='right'),
        "css_generator": get_css_elegant_professional,
    },
            "Mordern Minimal": {
        "html_generator": lambda data: generate_generic_html(data, date_placement='right'),
        "css_generator": get_css_modern_minimal,
    },
    
}

# def parse_uploaded_file(uploaded_file):
#     """Parse uploaded resume file and extract template styling."""
#     file_type = uploaded_file.name.split('.')[-1].lower()
    
#     try:
#         if file_type == 'html':
#             import chardet
#             raw_data = uploaded_file.read()
#             detected = chardet.detect(raw_data)
#             encoding = detected["encoding"] or "utf-8"
#             content = raw_data.decode(encoding, errors="ignore")

#             return extract_template_from_html(content)
#         elif file_type in ['doc', 'docx']:
#             st.warning("DOC/DOCX parsing coming soon. Using default template for now.")
#             return None
#         elif file_type == 'pdf':
#             st.warning("PDF parsing coming soon. Using default template for now.")
#             return None
#         elif file_type in ['ppt', 'pptx']:
#             st.warning("PPT/PPTX parsing coming soon. Using default template for now.")
#             return None
#         else:
#             st.error(f"Unsupported file type: {file_type}")
#             return None
#     except Exception as e:
#         st.error(f"Error parsing file: {str(e)}")
#         return None

def extract_template_from_html(html_content):
    """Extract CSS and structure from uploaded HTML."""

    css_match = re.search(r'<style[^>]*>(.*?)</style>', html_content, re.DOTALL)
    css = css_match.group(1) if css_match else ""
    
    template_id = hashlib.md5(html_content.encode()).hexdigest()[:8]
    
    return {
        'id': template_id,
        'name': f'Uploaded Template {template_id}',
        'css': css,
        'html': html_content,
        'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    }


def generate_html_content(data, color, template_generator, css_generator):
    """Generates the full HTML document and returns base64-encoded string."""
    css = css_generator(color)
    html_content = template_generator(data)
    
    title = f"{data.get('name', 'Resume')}"
    full_document_html = f"<html><head><meta charset='UTF-8'>{css}<title>{title}</title></head><body><div class='ats-page'>{html_content}</div></body></html>" 
    return base64.b64encode(full_document_html.encode('utf-8')).decode()

def get_html_download_link(data, color, template_config, filename_suffix=""):
    """Generates a download link for the styled HTML file."""
    if 'html_generator' in template_config and 'css_generator' in template_config:
        b64_data = generate_html_content(
            data, 
            color, 
            template_config['html_generator'], 
            template_config['css_generator']
        )
    else:
        # For uploaded templates
        css = template_config.get('css', '')
        html_body = generate_generic_html(data)
        full_html = f"<html><head><meta charset='UTF-8'><style>{css}</style></head><body><div class='ats-page'>{html_body}</div></body></html>"
        b64_data = base64.b64encode(full_html.encode('utf-8')).decode()
    
    filename = f"Resume_{data.get('name', 'User').replace(' ', '_')}{filename_suffix}.html"
    href = f'<a href="data:text/html;base64,{b64_data}" download="{filename}" style="font-size: 0.95em; text-decoration: none; padding: 10px 15px; background-color: #00FA9A; color: white; border-radius: 5px; display: inline-block; margin-top: 10px; width: 100%; text-align: center;"><strong>⬇️ Download HTML (.html)</strong></a>'
    return href



def generate_doc_html(data):
    """Generate a simple HTML that can be saved as .doc."""
    html = f"""
    <html xmlns:o='urn:schemas-microsoft-com:office:office' xmlns:w='urn:schemas-microsoft-com:office:word' xmlns='http://www.w3.org/TR/REC-html40'>
    <head>
        <meta charset='utf-8'>
        <title>Resume</title>
        <style>
            @page {{
                size: 8.5in 11in;
                margin: 0.5in;
            }}
            body {{ 
                font-family: Calibri, Arial, sans-serif; 
                font-size: 10pt; 
                line-height: 1.1; 
                margin: 0;
                padding: 0;
            }}
            h1 {{ font-size: 16pt; margin: 0 0 3pt 0; text-align: center; }}
            h2 {{ 
                font-size: 11pt; 
                margin-top: 8pt; 
                margin-bottom: 3pt; 
                border-bottom: 1px solid black; 
                padding-bottom: 1pt;
            }}
            p {{ margin: 3pt 0; }}
            ul {{ margin: 3pt 0; padding-left: 18pt; }}
            li {{ margin: 1pt 0; }}
            .header {{ text-align: center; margin-bottom: 8pt; }}
            .contact {{ font-size: 9pt; }}
            .job-title {{ font-size: 11pt; margin: 2pt 0; color: #555; }}
            .item-header {{ margin: 2pt 0; }}
            .item-title {{ font-weight: bold; }}
            .item-subtitle {{ font-style: italic; color: #555; }}
            .skills-group {{ margin: 3pt 0; }}
        </style>
    </head>
    <body>
        <div class='header'>
            <h1>{data.get('name', 'NAME MISSING')}</h1>
            <p class='job-title'>{data.get('job_title', '')}</p>
            <p class='contact'>{data.get('phone', '')} | {data.get('email', '')} | {data.get('location', '')}</p>
        </div>
    """
    
    for key in RESUME_ORDER:
        section_data = data.get(key)
        
        if not section_data or (isinstance(section_data, list) and not section_data) or (key == 'summary' and not section_data):
            continue
            
        title = format_section_title(key)
        html += f'<h2>{title}</h2>'
        
        if key == 'summary' and isinstance(section_data, str):
            html += f'<p>{section_data}</p>'
        
        elif key == 'skills' and isinstance(section_data, dict):
            for skill_type, skill_list in section_data.items():
                if skill_list:
                    html += f'<p class="skills-group"><strong>{format_section_title(skill_type)}:</strong> '
                    html += ", ".join(skill_list)
                    html += '</p>'
        
        elif isinstance(section_data, list):
            for item in section_data:
                if isinstance(item, str):
                    html += f'<ul><li>{item}</li></ul>'
                    continue
                
                if not isinstance(item, dict):
                    continue
                
                title_keys = ['title', 'name', 'degree']
                subtitle_keys = ['company', 'institution', 'issuer', 'organization']
                duration_keys = ['duration', 'date', 'period']
                
                main_title = next((item[k] for k in title_keys if k in item and item[k]), '')
                subtitle = next((item[k] for k in subtitle_keys if k in item and item[k] != main_title and item[k]), '')
                duration = next((item[k] for k in duration_keys if k in item and item[k]), '')

                html += f'<p class="item-header"><span class="item-title">{main_title}</span>'
                if subtitle:
                    html += f' <span class="item-subtitle">({subtitle})</span>'
                if duration:
                    html += f' | {duration}'
                html += '</p>'
                
                description_list_raw = item.get('description') or item.get('achievement') or item.get('details') 

                if description_list_raw:
                    if isinstance(description_list_raw, str):
                        description_list = [description_list_raw]
                    elif isinstance(description_list_raw, list):
                        description_list = description_list_raw
                    else:
                        description_list = None
                        
                    if description_list:
                        html += '<ul>'
                        for line in description_list:
                            html += f'<li>{line}</li>'
                        html += '</ul>'

    html += '</body></html>'
    return html

def get_doc_download_link(data, color, template_config, filename_suffix=""):
    """
    Generates a DOC download link using the selected template's HTML/CSS.

    Note: The 'DOC' file generated is an HTML file with a .doc extension and 
    Microsoft Word-specific XML/CSS (like @page rules) added to the beginning, 
    allowing it to open and be formatted correctly in Word.
    """

    if 'html_generator' in template_config and 'css_generator' in template_config:
        css = template_config['css_generator'](color)
        html_body = template_config['html_generator'](data)
    else:
       
        css = template_config.get('css', '')
        html_body = generate_generic_html(data) #
    word_doc_header = f"""
<html xmlns:o='urn:schemas-microsoft-com:office:office'
xmlns:w='urn:schemas-microsoft-com:office:word'
xmlns='http://www.w3.org/TR/REC-html40'>
<head>
    <meta charset='UTF-8'>
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{data.get('name', 'Resume')}</title>
    <style>
        @page {{ 
            size: 8.5in 11in; 
            margin: 0.5in; 
        }}
        {css}
    </style>
</head>
<body>
    <div class='ats-page'>
        {html_body}
    </div>
</body>
</html>
"""

   
    try:
        b64_data = base64.b64encode(word_doc_header.encode('utf-8')).decode()
    except NameError:
        print("Error: 'base64' module not found. Please import it.")
        return ""

    filename = f"Resume_{data.get('name', 'User').replace(' ', '_')}_Template_DOC{filename_suffix}.doc"

    doc_html = f"""
    <a href="data:application/msword;base64,{b64_data}" download="{filename}"
       style="font-size: 0.95em; text-decoration: none; padding: 10px 15px; 
              background-color: #00FFFF; color: white; border-radius: 5px; 
              display: inline-block; margin-top: 10px; width: 100%; text-align: center;">
        <strong>📄 Download Template DOC (.doc)</strong>
    </a>
    """
    return doc_html


def generate_markdown_text(data):
    """Generates a plain markdown/text version of the resume."""
    text = ""
    
    text += f"{data.get('name', 'NAME MISSING').upper()}\n"
    if data.get('job_title'):
        text += f"{data.get('job_title')}\n"
    contact_parts = [data.get('phone', ''), data.get('email', ''), data.get('location', '')]
    text += " | ".join(filter(None, contact_parts)) + "\n"
    text += "=" * 50 + "\n\n"
    
    for key in RESUME_ORDER:
        section_data = data.get(key)
        
        if not section_data or (isinstance(section_data, list) and not section_data) or (key == 'summary' and not section_data):
            continue
            
        title = format_section_title(key).upper()
        text += f"{title}\n"
        text += "-" * len(title) + "\n"
        
        if key == 'summary' and isinstance(section_data, str):
            text += section_data + "\n\n"

        elif key == 'skills' and isinstance(section_data, dict):
            for skill_type, skill_list in section_data.items():
                if skill_list:
                    text += f"{format_section_title(skill_type)}: "
                    text += ", ".join(skill_list) + "\n"
            text += "\n"
        
        elif isinstance(section_data, list):
            for item in section_data:
                if not isinstance(item, dict):
                    text += f" - {item}\n"
                    continue
                
                title_keys = ['title', 'name', 'degree']
                subtitle_keys = ['company', 'institution', 'issuer']
                duration_keys = ['duration', 'date']
                
                main_title = next((item[k] for k in title_keys if k in item and item[k]), '')
                subtitle = next((item[k] for k in subtitle_keys if k in item and item[k] != main_title and item[k]), '')
                duration = next((item[k] for k in duration_keys if k in item and item[k]), '')

                line = f"{main_title}"
                if subtitle:
                    line += f" ({subtitle})"
                if duration:
                    line += f" | {duration}"
                text += line + "\n"
                
                description_list = item.get('description') or item.get('achievement') or item.get('details')
                if description_list and isinstance(description_list, list):
                    for line in description_list:
                        text += f" - {line}\n"
            text += "\n"

    return text

def get_text_download_link(data, filename_suffix=""):
    """Generates a download link for a plain text file."""
    text_content = generate_markdown_text(data)
    b64_data = base64.b64encode(text_content.encode('utf-8')).decode()
    
    filename = f"Resume_{data.get('name', 'User').replace(' ', '_')}_ATS_Plain_Text{filename_suffix}.txt"
    href = f'<a href="data:text/plain;base64,{b64_data}" download="{filename}" style="font-size: 0.95em; text-decoration: none; padding: 10px 15px; background-color: #40E0D0; color: white; border-radius: 5px; display: inline-block; margin-top: 10px; width: 100%; text-align: center;"><strong>📋 Download Plain Text (.txt)</strong></a>'
    return href



def app_download():
    st.set_page_config(layout="wide", page_title="Download Resume")
    ACCENT_COLOR = "#6ea8fe" # Softer, bright blue

    st.markdown(f"""
    <style>
        /* ============================
        🌟 Dark Blue-Green Theme
        ============================ */
        :root {{
            --primary-color: #00BFFF; /* Deep sky blue */
            --secondary-color: #000000; /* Black background */
            --accent-color: #00FF7F; /* Spring green */
            --accent-light: #66FFB2; /* Lighter green-blue */
            --accent-ice: #0a0a0a; /* Dark icy color for borders */
            --text-dark: #FFFFFF; /* Text on dark background */
            --text-light: #FFFFFF;
            --card-bg: rgba(20, 20, 20, 0.9); /* Dark card background */
            --card-border: rgba(0, 255, 127, 0.5); /* Soft green border */
            --soft-shadow: rgba(0, 255, 127, 0.15); /* Soft shadow for depth */
            --button-gradient: -webkit-linear-gradient(45deg, #00BFFF, #00FF7F);
        }}

        /* HIDE NAVIGATION */
        [data-testid="stSidebarNav"] {{
            display: none !important;
        }}

        .stApp {{
            background-color: var(--secondary-color);
            background-attachment: fixed;
            min-height: 100vh;
            color: var(--text-dark);
            font-family: "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif;
        }}

        /* PREVIEW IFRAME BACKGROUND FIX */
        iframe {{
            background: linear-gradient(135deg, #e6f9ff 0%, #f0fff4 100%) !important;
            border-radius: 12px;
            border: 2px solid var(--accent-color);
            box-shadow: 0 8px 25px var(--soft-shadow);
        }}
        /* TEMPLATE CARDS */


.template-card {{
    background: rgba(20, 20, 20, 0.85);
    border: 1px solid rgba(0, 255, 200, 0.15);
    border-radius: 12px;
    padding: 14px;
    box-shadow: 0 3px 10px rgba(0, 0, 0, 0.4);
    margin-bottom: 18px;
    text-align: center;
    min-height: 160px;
    transition: all 0.25s ease-in-out;
    display: flex;
    flex-direction: column;
    justify-content: space-between;
}}
.template-card:hover {{
    border-color: rgba(0, 255, 180, 0.6);
    box-shadow: 0 6px 18px rgba(0, 255, 200, 0.25);
    transform: translateY(-4px);
}}
.template-card h4 {{
    font-size: 1em;
    font-weight: 600;
    color: #00f0ff;
    margin-bottom: 4px;
}}
.template-card p {{
    font-size: 0.8em;
    color: #aaa;
    margin-bottom: 8px;
}}

/* Only style buttons INSIDE .template-card */
.template-card .stButton > button {{
    background: linear-gradient(90deg, #00e0ff, #00ffa3);
    color: black;
    border: none;
    padding: 6px 0;
    border-radius: 8px;
    font-weight: 600;
    font-size: 0.85em;
    width: 100%;
    transition: all 0.3s ease;
}}
.template-card .stButton > button:hover {{
    background: linear-gradient(90deg, #00ffa3, #00e0ff);
    transform: scale(1.03);
}}



        /* HEADINGS */
        h1, h2, h3, h4 {{
            color: var(--primary-color);
            font-weight: 700;
        }}

        .template-card h3 {{
            color: var(--primary-color);
            font-weight: 600;
            margin-bottom: 5px;
            font-size: 1.4em;
        }}

        /* BADGE */
        .saved-badge {{
            background: var(--accent-light);
            color: var(--primary-color);
            padding: 5px 12px;
            border-radius: 15px;
            font-weight: 700;
            box-shadow: 0 2px 10px rgba(102, 255, 178, 0.5);
        }}

        /* BUTTONS */
        .stButton>button {{
            background: var(--button-gradient);
            color: var(--text-light);
            border: none;
            border-radius: 12px;
            padding: 0.8rem 1.8rem;
            font-weight: 700;
            text-transform: uppercase;
            letter-spacing: 0.8px;
            box-shadow: 0 5px 18px rgba(0, 255, 127, 0.4);
            transition: all 0.3s ease;
            width: 100%;
            max-width: 280px;
            margin-top: auto;
            margin-left: auto;
            margin-right: auto;
            display: block;
        }}

        .stButton>button:hover {{
            background: -webkit-linear-gradient(45deg, #00FF7F, #00BFFF);
            box-shadow: 0 8px 25px rgba(0, 255, 127, 0.8);
            transform: translateY(-4px) scale(1.02);
            color: var(--text-dark);
        }}

        /* INPUT FIELDS */
        .stTextInput > div > div > input,
        .stTextArea > div > div > textarea {{
            background-color: #1a1a1a; /* Dark input background */
            color: var(--text-light);
            border: 1px solid var(--accent-ice);
            border-radius: 10px;
            padding: 14px;
            box-shadow: inset 0 1px 4px rgba(0, 0, 0, 0.2);
        }}

        .stTextInput > div > div > input:focus,
        .stTextArea > div > div > textarea:focus {{
            border-color: var(--accent-color);
            box-shadow: 0 0 0 3px rgba(0, 255, 127, 0.4);
        }}
                
                /* ============================
        Dark Sidebar Styling
        ============================ */
        /* Sidebar background */
        [data-testid="stSidebar"] {{
            background-color: #000000 !important; /* Black background */
            color: #FFFFFF; /* White text by default */
            padding: 1rem;
        }}

        /* Sidebar headings, subheaders */
        [data-testid="stSidebar"] h2,
        [data-testid="stSidebar"] h3,
        [data-testid="stSidebar"] h4,
        [data-testid="stSidebar"] .css-1d391kg p {{
            color: #FFFFFF !important;
        }}

        /* Sidebar selectbox and color picker */
        [data-testid="stSidebar"] .stSelectbox, 
        [data-testid="stSidebar"] .stColorPicker {{
            background-color: #1a1a1a !important; /* Dark input background */
            color: #FFFFFF !important; /* White text */
        }}

        /* Sidebar buttons */
        [data-testid="stSidebar"] button {{
            background: linear-gradient(45deg, #00BFFF, #00FF7F);
            color: #FFFFFF;
            border-radius: 10px;
            padding: 10px 15px;
            margin-bottom: 10px;
            font-weight: 600;
            border: none;
            width: 100%;
            transition: all 0.3s ease;
        }}

        [data-testid="stSidebar"] button:hover {{
            background: linear-gradient(45deg, #00FF7F, #00BFFF);
            color: #000000;
            transform: scale(1.02);
        }}

        /* Sidebar horizontal lines */
        [data-testid="stSidebar"] hr {{
            border-top: 1px solid rgba(255, 255, 255, 0.2);
        }}

        /* Sidebar captions */
        [data-testid="stSidebar"] .stCaption {{
            color: #FFFFFF !important;
        }}
                    /* Make selectbox label text white in sidebar */
        [data-testid="stSidebar"] label, 
        [data-testid="stSidebar"] .stSelectbox label {{
            color: #FFFFFF !important;
        }}

        /* Make selectbox options text dark (optional) if using dark dropdown) */
        [data-testid="stSidebar"] .stSelectbox div[role="listbox"] div {{
            color: #000000;
        }}
        


    </style>
""", unsafe_allow_html=True)




    final_data = st.session_state.get('final_resume_data')
    # st.json(final_data)

    if final_data is None:
        st.error("❌ Resume data not found. Please return to the editor to finalize your resume.")
        if st.button("⬅️ Go Back to Editor"):
            switch_page("main")
        return

    if isinstance(final_data, str):
        try:
            final_data = json.loads(final_data)
        except json.JSONDecodeError:
            st.error("❌ Error: Could not parse resume data.")
            if st.button("⬅️ Go Back to Editor"):
                switch_page("main")
            return
            
    if not isinstance(final_data, dict):
        st.error("❌ Resume data is in an unusable format.")
        if st.button("⬅️ Go Back to Editor"):
            switch_page("main")
        return

    st.markdown("""
    <div style="text-align: center;">
        <h1>📄 Resume Template Manager</h1>
        <p style="color: white; font-size: 16px;">
            <b>System Templates • Saved Templates • Upload Custom Templates</b>
        </p>
    </div>
    """, unsafe_allow_html=True)
#     
    tab1, tab3 = st.tabs(["🎨 System Templates", "📤 Upload New Template"])

    # --- TAB 1: SYSTEM TEMPLATES ---
    with tab1:
        st.markdown("### Available System Templates")
        st.caption("Choose from professionally designed ATS-friendly templates")
        
        # Template selection
        cols = st.columns(3)
        for idx, (template_name, template_config) in enumerate(SYSTEM_TEMPLATES.items()):
            with cols[idx % 3]:
                with st.container():
                    st.markdown(f"""
                    <div class="template-card">
                        <h4>{template_name}</h4>
                        <p style="font-size: 0.9em; color: #666;">Click to preview and customize</p>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    if st.button(f"Select {template_name}", key=f"sys_template_{idx}", use_container_width=True):
                        st.session_state.selected_template = template_name
                        st.session_state.selected_template_config = template_config
                        st.session_state.template_source = 'system'
                        st.rerun()
        
        # Show preview and color selection if a template is selected
        if st.session_state.get('selected_template') and st.session_state.get('template_source') == 'system':
            st.markdown("---")
            st.markdown(f"### Preview: {st.session_state.selected_template}")
            
            # Color selection in main body
            col1, col2 = st.columns([3, 1])
            with col1:
                color_name = st.selectbox(
                    'Choose Accent Color:',
                    list(ATS_COLORS.keys()),
                    key='sys_color_select'
                )
                primary_color = ATS_COLORS[color_name]
            
            with col2:
                custom_color = st.color_picker(
                    'Custom Color:',
                    primary_color,
                    key='sys_custom_color'
                )
                if custom_color != primary_color:
                    primary_color = custom_color
            
            # Store selected color in session state
            st.session_state.selected_color = primary_color
            
            # Generate preview with selected color
            template_config = st.session_state.selected_template_config
            css = template_config['css_generator'](primary_color)
            html_content = template_config['html_generator'](final_data)
            
            full_html = f"""
            {css}
            <div class="ats-page">
                {html_content}
            </div>
            """
            
            st.components.v1.html(full_html, height=1000, scrolling=True)
            

    
    with tab3:
        if 'uploaded_templates' not in st.session_state:
            st.session_state.uploaded_templates = load_user_templates(st.session_state.logged_in_user)

        if st.session_state.uploaded_templates:
            st.markdown("### 🗂️ Your Saved Templates")
            cols = st.columns(3)
            for idx, (template_id, template_data) in enumerate(st.session_state.uploaded_templates.items()):
                with cols[idx % 3]:
                    st.markdown(f"""
                    <div class="template-card" style="border:1px solid #ccc; padding:10px; border-radius:10px; background:#fafafa;">
                        <h4>{template_data['name']}</h4>
                        <p style="font-size:0.85em; color:#555;">File: {template_data['original_filename']}</p>
                        <p style="font-size:0.8em; color:#888;">Uploaded: {template_data['uploaded_at']}</p>
                    </div>
                    """, unsafe_allow_html=True)

                    col1, col2 = st.columns(2)
                    with col1:
                        if st.button(f"Use", key=f"use_{template_id}", use_container_width=True):
                            # Clear temp upload config when using a saved template
                            if 'temp_upload_config' in st.session_state:
                                del st.session_state.temp_upload_config
                            
                            st.session_state.selected_template_preview = f"""
                                <style>{template_data['css']}</style>
                                <div class="ats-page">{generate_generic_html(final_data)}</div>
                            """
                            st.session_state.selected_template = template_data['name']
                            st.session_state.selected_template_config = template_data
                            st.session_state.template_source = 'saved'
                            st.session_state.current_upload_id = template_id
                            st.rerun()

                    with col2:
                        if st.button(f"Delete", key=f"delete_{template_id}", use_container_width=True):
                            # Clear selection if deleting currently selected template
                            if st.session_state.get('current_upload_id') == template_id:
                                st.session_state.pop('selected_template_preview', None)
                                st.session_state.pop('selected_template', None)
                                st.session_state.pop('selected_template_config', None)
                                st.session_state.pop('current_upload_id', None)
                            
                            del st.session_state.uploaded_templates[template_id]
                            save_user_templates(st.session_state.logged_in_user, st.session_state.uploaded_templates)
                            st.success(f"✅ Deleted '{template_data['name']}'")
                            st.rerun()

            st.markdown("---")
        else:
            st.info("📂 No saved templates yet. Upload a template below to get started!")
            st.markdown("---")

        # 2️⃣ Upload Section
        st.markdown("### 📤 Upload New Template")
        uploaded_file = st.file_uploader(
            "Upload an HTML file",
            type=['html','pptx','docx', 'doc'],
            key="template_upload"
        )

        if uploaded_file is not None:
            st.success(f"✅ File uploaded: {uploaded_file.name}")

            with st.spinner("Parsing template..."):
                file_type = uploaded_file.name.split('.')[-1].lower()

                if file_type == 'html':
                    import chardet
                    raw_data = uploaded_file.read()
                    detected = chardet.detect(raw_data)
                    encoding = detected["encoding"] or "utf-8"
                    content = raw_data.decode(encoding, errors="ignore")

                    parsed_template = extract_template_from_html(content)

                    # Store parsed template for preview and download (before saving)
                    st.session_state.temp_upload_config = {
                        'name': f"Uploaded_{uploaded_file.name.split('.')[0]}",
                        'css': parsed_template.get('css', ''),
                        'html': parsed_template.get('html', ''),
                        'original_filename': uploaded_file.name
                    }

                    # Template name input and save button
                    col1, col2 = st.columns([2, 1])
                    with col1:
                        template_name = st.text_input(
                            "Template Name:",
                            value=f"Uploaded_{uploaded_file.name.split('.')[0]}",
                            key="upload_template_name"
                        )

                    with col2:
                        if st.button("💾 Save Template", use_container_width=True):
                            template_id = f"upload_{datetime.now().strftime('%Y%m%d%H%M%S')}"
                            st.session_state.uploaded_templates[template_id] = {
                                'name': template_name,
                                'css': parsed_template.get('css', ''),
                                'html': parsed_template.get('html', ''),
                                'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                                'original_filename': uploaded_file.name
                            }

                            # Save to JSON
                            save_user_templates(st.session_state.logged_in_user, st.session_state.uploaded_templates)
                            st.success(f"✅ Template '{template_name}' saved!")

                            # Set the saved template as selected
                            st.session_state.selected_template_config = st.session_state.uploaded_templates[template_id]
                            st.session_state.selected_template = template_name
                            st.session_state.template_source = 'saved'
                            st.session_state.current_upload_id = template_id
                            
                            # Update preview to show saved template
                            st.session_state.selected_template_preview = f"""
                                <style>{parsed_template.get('css', '')}</style>
                                <div class="ats-page">{generate_generic_html(final_data)}</div>
                            """
                            
                            # Clear temp upload config since it's now saved
                            if 'temp_upload_config' in st.session_state:
                                del st.session_state.temp_upload_config
                            
                            st.rerun()

                    # Show preview for uploaded file (before saving) - Always show when file is uploaded
                    preview_html = f"""
                        <style>{parsed_template.get('css', '')}</style>
                        <div class="ats-page">{generate_generic_html(final_data)}</div>
                    """
                    st.markdown("### 🔍 Template Preview (Not Saved Yet)")
                    st.components.v1.html(preview_html, height=1000, scrolling=True)
                    
                    # Enable downloads for unsaved upload
                    st.session_state.selected_template_config = st.session_state.temp_upload_config
                    st.session_state.template_source = 'temp_upload'

                elif file_type in ['ppt', 'pptx']:
                    # Store uploaded file in session state for preview
                    st.session_state.ppt_uploaded_file = uploaded_file.getvalue()
                    st.session_state.ppt_original_filename = uploaded_file.name
                    
                    # Process the presentation
                    prs = Presentation(io.BytesIO(st.session_state.ppt_uploaded_file))
                    slide_texts = []
                    for slide_idx, slide in enumerate(prs.slides):
                        text_blocks = []
                        for shape_idx, shape in enumerate(slide.shapes):
                            if shape.has_text_frame and shape.text.strip():
                                text_blocks.append({
                                    "index": shape_idx,
                                    "text": shape.text.strip(),
                                    "position": {"x": shape.left, "y": shape.top}
                                })
                        
                        if text_blocks:
                            text_blocks.sort(key=lambda x: (x["position"]["y"], x["position"]["x"]))
                            slide_texts.append({
                                "slide_number": slide_idx + 1,
                                "text_blocks": text_blocks
                            })
                    
                    if slide_texts:
                        structured_slides = analyze_slide_structure(slide_texts)
                        generated_sections = generate_ppt_sections(final_data, structured_slides)

                        if generated_sections:
                            text_elements = []
                            for slide_idx, slide in enumerate(prs.slides):
                                for shape_idx, shape in enumerate(slide.shapes):
                                    if shape.has_text_frame and shape.text.strip():
                                        text_elements.append({
                                            'slide': slide_idx + 1,
                                            'shape': shape_idx,
                                            'original_text': shape.text.strip(),
                                            'shape_type': type(shape).__name__
                                        })
                            
                            content_mapping, heading_shapes, basic_info_shapes = match_generated_to_original(
                                text_elements, generated_sections, prs)
                            
                            # Store in session state
                            st.session_state.ppt_content_mapping = content_mapping
                            st.session_state.ppt_heading_shapes = heading_shapes
                            st.session_state.ppt_basic_info_shapes = basic_info_shapes
                            st.session_state.ppt_text_elements = text_elements
                            
                            # Auto-generate preview immediately
                            with st.spinner("🔄 Generating preview..."):
                                working_prs = Presentation(io.BytesIO(st.session_state.ppt_uploaded_file))
                                edits = {}
                                
                                for element in text_elements:
                                    key = f"{element['slide']}_{element['shape']}"
                                    if key not in heading_shapes:
                                        edits[key] = content_mapping.get(key, element['original_text'])
                                
                                success_count = 0
                                for element in text_elements:
                                    key = f"{element['slide']}_{element['shape']}"
                                    if key not in heading_shapes and key in edits:
                                        slide_idx = element['slide'] - 1
                                        shape_idx = element['shape']
                                        
                                        if slide_idx < len(working_prs.slides):
                                            slide = working_prs.slides[slide_idx]
                                            if shape_idx < len(slide.shapes):
                                                shape = slide.shapes[shape_idx]
                                                if shape.has_text_frame:
                                                    clear_and_replace_text(shape, edits[key])
                                                    success_count += 1
                                
                                # Save preview to session state
                                output = io.BytesIO()
                                working_prs.save(output)
                                output.seek(0)
                                st.session_state.generated_ppt = output.getvalue()
                            
                            # Show preview section
                            st.markdown("### 🔍 PowerPoint Preview")
                            
                            # Display slide thumbnails and content
                            preview_prs = Presentation(io.BytesIO(st.session_state.generated_ppt))
                            
                            for slide_idx, slide in enumerate(preview_prs.slides):
                                with st.expander(f"📊 Slide {slide_idx + 1}", expanded=slide_idx == 0):
                                    # Extract slide content for display
                                    slide_content = []
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text") and shape.text.strip():
                                            slide_content.append(shape.text.strip())
                                    
                                    if slide_content:
                                        for content in slide_content:
                                            st.text_area(
                                                f"Content - Slide {slide_idx + 1}",
                                                value=content,
                                                height=100,
                                                key=f"preview_slide_{slide_idx}_{hash(content)}",
                                                disabled=True
                                            )
                                    else:
                                        st.info("No text content in this slide")
                            
                            # Save PPT Template button
                            st.markdown("---")
                            ppt_name = st.text_input(
                                "PPT Template Name:",
                                value=f"PPT_{uploaded_file.name.split('.')[0]}",
                                key="ppt_template_name"
                            )
                            
                            if st.button("💾 Save PPT Template", use_container_width=True, type="primary"):
                                # Initialize ppt_templates if not exists
                                if 'ppt_templates' not in st.session_state:
                                    st.session_state.ppt_templates = load_user_ppt_templates(st.session_state.logged_in_user)
                                
                                ppt_id = f"ppt_{datetime.now().strftime('%Y%m%d%H%M%S')}"
                                st.session_state.ppt_templates[ppt_id] = {
                                    'name': ppt_name,
                                    'ppt_data': st.session_state.ppt_uploaded_file,
                                    'edits': edits,
                                    'content_mapping': st.session_state.get('ppt_content_mapping', {}),
                                    'heading_shapes': st.session_state.get('ppt_heading_shapes', set()),
                                    'basic_info_shapes': st.session_state.get('ppt_basic_info_shapes', set()),
                                    'text_elements': text_elements,
                                    'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                                    'original_filename': uploaded_file.name
                                }
                                
                                # Save to JSON
                                save_user_ppt_templates(st.session_state.logged_in_user, st.session_state.ppt_templates)
                                st.success(f"✅ PPT Template '{ppt_name}' saved!")
                                st.rerun()
                            
                            # Download button
                            st.markdown("---")
                            st.download_button(
                                label="📥 Download Enhanced PowerPoint",
                                data=st.session_state.generated_ppt,
                                file_name="ai_enhanced_presentation.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True
                            )

                elif file_type in ['docx', 'doc']:  
                    st.session_state.json_data = json.dumps(final_data, indent=2)  
                    
                    # Process document and generate preview immediately
                    uploaded_file.seek(0)
                    doc, structure = extract_document_structure(uploaded_file)
                    output, replaced, removed = replace_content(doc, structure, final_data)
                    
                    # Store in session state
                    st.session_state.generated_doc = output.getvalue()
                    st.session_state.doc_structure = structure
                    st.session_state.doc_replaced = replaced
                    st.session_state.doc_removed = removed
                    
                    # Show preview section
                    st.markdown("### 🔍 Document Preview")
                    
                    # Display document sections
                    for i, section in enumerate(structure):
                        with st.expander(f"📝 {section.get('section', 'Content')} Section", expanded=i == 0):
                            col1, col2 = st.columns([1, 1])
                            
                            with col1:
                                st.caption("Original Content:")
                                st.text_area(
                                    f"Original {i}",
                                    value=section.get('text', ''),
                                    height=150,
                                    key=f"doc_original_{i}",
                                    disabled=True
                                )
                            
                            with col2:
                                st.caption("AI-Generated Content:")
                                # Find replaced content for this section
                                replaced_text = "No replacement found"
                                for replace_item in replaced:
                                    if (replace_item.get('section') == section.get('section') and 
                                        replace_item.get('original_text') == section.get('text')):
                                        replaced_text = replace_item.get('replaced_with', replaced_text)
                                        break
                                
                                st.text_area(
                                    f"Generated {i}",
                                    value=replaced_text,
                                    height=150,
                                    key=f"doc_generated_{i}",
                                    disabled=True
                                )
                    
                    # Show statistics
                    st.metric("Sections Processed", len(replaced))
                    st.metric("Content Blocks Removed", len(removed))
                    
                    # Download button
                    filename = f"{final_data.get('name', 'Resume').replace(' ', '_')}_Final.docx"
                    st.download_button(
                        label="📥 Download Final Document",
                        data=st.session_state.generated_doc,
                        file_name=filename,
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )

        st.markdown("---")
        
        # 3️⃣ Preview Section - Show saved template preview ONLY when no upload is active
        if uploaded_file is None and st.session_state.get("selected_template_preview") and st.session_state.get("template_source") == 'saved':
            st.markdown(f"### 🔍 Template Preview — **{st.session_state.selected_template}**")
            st.components.v1.html(st.session_state.selected_template_preview, height=1000, scrolling=True)

        st.markdown("---")
        if st.button("⬅️ Go Back to Editor", use_container_width=True):
            switch_page("create")

    # --- Sidebar ---
    with st.sidebar:
        st.subheader("⚙️ Download Options")
        
        # Download buttons
        if st.session_state.get("selected_template_config"):
            current_template = st.session_state.selected_template_config
            
            # Use color from session state if available, otherwise use default
            download_color = st.session_state.get('selected_color', ATS_COLORS["Professional Blue (Default)"])
            
            # Safely get CSS and HTML with defaults
            template_css = current_template.get('css', '')
            template_html = current_template.get('html', '')

            st.markdown(get_html_download_link(
                final_data, 
                download_color, 
                st.session_state.selected_template_config
            ), unsafe_allow_html=True)
            
            st.markdown(get_doc_download_link(
                final_data, 
                download_color, 
                st.session_state.selected_template_config
            ), unsafe_allow_html=True)
            
            st.markdown(get_text_download_link(final_data), unsafe_allow_html=True)
        else:
            st.info("👆 Please select a template first to enable download links.")

        st.markdown("---")
        st.caption("### 💡 Download Tips:")
        st.caption("**HTML:** Best for web viewing")
        st.caption("**PDF:** Open HTML → Print → Save as PDF")
        st.caption("**DOC:** Edit in Microsoft Word")
        st.caption("**TXT:** Maximum ATS compatibility")

if __name__ == '__main__':
    app_download()