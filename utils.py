from openpyxl import load_workbook
import streamlit as st
import pdfplumber
from docx import Document
from typing import List, Dict,Any
import json
import requests
import os
import io
from reportlab.lib.pagesizes import A4
from pathlib import Path
import base64
import requests
from difflib import SequenceMatcher
import re
import json
import requests
from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt
from docx import Document
from docx.shared import RGBColor
from collections import Counter
import hashlib
from copy import deepcopy
import time
import base64
from io import BytesIO
from datetime import datetime
import requests
import json
from lxml import etree
import mammoth
from configuration.db import get_connection


# Recommended models:
FAST_MODEL = "meta/llama3-8b-instruct"
HEAVY_MODEL = "meta/llama3-70b-instruct"



API_URL = "https://integrate.api.nvidia.com/v1/chat/completions"
API_KEY = "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"

url = API_URL
    
headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }
# TO EXTRACT DATA FROM THE UPLOADED FILE#

TEMPLATES_DIR = "uploaded_templates.json"
users_file = Path(__file__).parent / "users.json"
user_data_file = Path(__file__).parent/"user_resume_data.json"


@st.cache_data
def get_base64_of_file(file_path):
    with open(file_path, "rb") as f:
        data = f.read()
    return base64.b64encode(data).decode()

def set_image_as_bg(image_file):
    # Get file extension to set correct MIME type
    ext = os.path.splitext(image_file)[1].lower()
    if ext == ".png":
        mime_type = "image/png"
    elif ext == ".jpg" or ext == ".jpeg":
        mime_type = "image/jpeg"
    else:
        st.error("Unsupported file type!")
        return

    bin_str = get_base64_of_file(image_file)
    page_bg_img = f"""
    <style>
    .stApp {{
        background-image: url("data:{mime_type};base64,{bin_str}");
        background-size: cover;
        background-position: center;
        background-repeat: no-repeat;
    }}
    </style>
    """
    st.markdown(page_bg_img, unsafe_allow_html=True)



import time
from typing import Dict, Any
import requests
import json

def extract_details_from_text(extracted_text: str, max_retries: int = 3) -> Dict[str, Any]:
    """
    Extract structured details from resume text with retry logic and timeout handling.
    """
    
    # Truncate if text is too long (keep first 15000 chars to avoid timeouts)
    if len(extracted_text) > 15000:
        print(f"⚠️ Resume text truncated from {len(extracted_text)} to 15000 characters")
        extracted_text = extracted_text[:15000]
    
    SCHEMA_DEFINITION = """{
  "name": "string or null",
  "email": "string or null",
  "phone": "string or null",
  "location": "string or null",
  "url": "string or null",
  "summary": "string or null",
  "skills": ["array of strings"],
  "experience": [
    {
      "company": "string or null",
      "position": "string or null",
      "location": "string or null",
      "start_date": "string or null (MM/YYYY format)",
      "end_date": "string or null (MM/YYYY or Present)",
      "exp_skills": ["array of strings"],
      "description": ["array of strings - bullet points"]
    }
  ],
  "education": [
    {
      "course": "string or null (degree/program name)",
      "university": "string or null (institution name)",
      "start_date": "string or null (MM/YYYY format)",
      "end_date": "string or null (MM/YYYY or Present)"
    }
  ],
  "certificate": [
    {
      "certificate_name": "string or null",
      "provider_name": "string or null",
      "completed_date": "string or null (MM/YYYY format)"
    }
  ],
  "project": [
    {
      "projectname": "string or null",
      "tools": "string or null (comma-separated)",
      "decription": "string or null (key achievements)"
    }
  ],
  "custom_sections": {
    "section_title": "section content as string"
  }
}"""

    prompt_parts = [
        "Extract structured information from the resume text below.",
        "Return ONLY valid JSON. Do NOT include explanations, markdown, or code blocks.",
        "",
        "CRITICAL RULES:",
        "1. Do NOT invent or infer any information that is not explicitly present.",
        "2. If a field is missing, return null or an empty array [].",
        "3. Preserve exact wording of skills, titles, companies, and achievements.",
        "4. For dates, use MM/YYYY format (e.g., '01/2020').",
        "5. For current positions, use 'Present' as end_date.",
        "6. Split multi-line descriptions into arrays of strings.",
        "",
        "SCHEMA TO FOLLOW:",
        SCHEMA_DEFINITION,
        "",
        "Resume text:",
        extracted_text
    ]
    
    prompt = "\n".join(prompt_parts)

    # Retry logic with exponential backoff
    for attempt in range(max_retries):
        try:
            # Increase timeout with each retry
            timeout = 30 + (attempt * 20)  # 30s, 50s, 70s
            
            print(f"🔄 Attempt {attempt + 1}/{max_retries} (timeout: {timeout}s)...")
            
            payload = {
                "model": FAST_MODEL,
                "messages": [
                    {
                        "role": "system", 
                        "content": "You are an expert resume parser. Extract information exactly as it appears. Return only valid JSON with no additional text or markdown formatting."
                    },
                    {
                        "role": "user", 
                        "content": prompt
                    }
                ],
                "temperature": 0.1,
                "max_tokens": 3000,  # Reduced for faster response
                "top_p": 0.7
            }

            # Make API request
            response = requests.post(
                url, 
                headers=headers, 
                json=payload, 
                timeout=timeout
            )
            response.raise_for_status()
            result = response.json()
            
            if 'choices' not in result or not result['choices']:
                raise ValueError("API returned empty response")
            
            response_text = result['choices'][0]['message']['content'].strip()

            # Clean up response - remove markdown code blocks if present
            response_text = response_text.replace('```json', '').replace('```', '').strip()

            # Extract JSON from response
            json_start = response_text.find('{')
            json_end = response_text.rfind('}') + 1
            
            if json_start == -1 or json_end <= json_start:
                raise ValueError("No valid JSON object found in API response")

            json_str = response_text[json_start:json_end]
            data = json.loads(json_str)

            # Post-processing
            if "experience" in data and isinstance(data["experience"], list):
                for exp in data["experience"]:
                    desc = exp.get("description", [])
                    if isinstance(desc, str):
                        exp["description"] = [line.strip() for line in desc.split("\n") if line.strip()]
                    elif not isinstance(desc, list):
                        exp["description"] = []
                    
                    if "exp_skills" not in exp or not isinstance(exp["exp_skills"], list):
                        exp["exp_skills"] = []

            if "project" in data and isinstance(data["project"], list):
                for proj in data["project"]:
                    desc = proj.get("decription", "")
                    if isinstance(desc, list):
                        proj["decription"] = " ".join(desc)
                    elif not isinstance(desc, str):
                        proj["decription"] = ""

            list_fields = ["skills", "experience", "education", "certificate", "project"]
            for field in list_fields:
                if field not in data:
                    data[field] = []
                elif not isinstance(data[field], list):
                    data[field] = [data[field]] if data[field] else []

            string_fields = ["name", "email", "phone", "location", "url", "summary"]
            for field in string_fields:
                if field not in data:
                    data[field] = ""
                elif data[field] is None:
                    data[field] = ""

            if "custom_sections" not in data:
                data["custom_sections"] = {}
            elif not isinstance(data["custom_sections"], dict):
                data["custom_sections"] = {}

            if not data.get("name") or not data.get("email"):
                raise ValueError("Resume must contain at least a name and email address")

            print(f"✅ Successfully parsed resume on attempt {attempt + 1}")
            return data

        except requests.exceptions.Timeout:
            print(f"⏱️ Timeout on attempt {attempt + 1}/{max_retries}")
            if attempt < max_retries - 1:
                wait_time = 2 ** attempt  # Exponential backoff: 1s, 2s, 4s
                print(f"⏳ Waiting {wait_time}s before retry...")
                time.sleep(wait_time)
                continue
            else:
                # Last attempt failed - return fallback
                print("❌ All retries exhausted. Using fallback parser...")
                return _fallback_parse(extracted_text)
        
        except requests.exceptions.RequestException as e:
            print(f"🌐 API error on attempt {attempt + 1}: {str(e)}")
            if attempt < max_retries - 1:
                wait_time = 2 ** attempt
                print(f"⏳ Waiting {wait_time}s before retry...")
                time.sleep(wait_time)
                continue
            else:
                return _fallback_parse(extracted_text)
        
        except json.JSONDecodeError as e:
            print(f"📄 JSON parse error: {str(e)}")
            if attempt < max_retries - 1:
                continue
            else:
                return _fallback_parse(extracted_text)
        
        except Exception as e:
            print(f"❌ Unexpected error: {str(e)}")
            if attempt < max_retries - 1:
                continue
            else:
                return _fallback_parse(extracted_text)


def _fallback_parse(text: str) -> Dict[str, Any]:
    """
    Simple regex-based fallback parser when API fails.
    Extracts basic information using pattern matching.
    """
    import re
    
    print("🔧 Using fallback regex parser...")
    
    data = {
        "name": "",
        "email": "",
        "phone": "",
        "location": "",
        "url": "",
        "summary": "",
        "skills": [],
        "experience": [],
        "education": [],
        "certificate": [],
        "project": [],
        "custom_sections": {}
    }
    
    # Extract email
    email_match = re.search(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', text)
    if email_match:
        data["email"] = email_match.group()
    
    # Extract phone
    phone_match = re.search(r'[\+\(]?[0-9][0-9 .\-\(\)]{8,}[0-9]', text)
    if phone_match:
        data["phone"] = phone_match.group()
    
    # Extract name (first line typically)
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    if lines:
        # Name is usually first non-empty line
        data["name"] = lines[0]
    
    # Extract skills section
    skills_section = re.search(r'(?:SKILLS?|TECHNICAL SKILLS?|CORE SKILLS?)[:\s]+(.*?)(?=\n[A-Z]{3,}|$)', 
                                text, re.IGNORECASE | re.DOTALL)
    if skills_section:
        skills_text = skills_section.group(1)
        # Split by common delimiters
        skills = re.split(r'[,•|]', skills_text)
        data["skills"] = [s.strip() for s in skills if s.strip() and len(s.strip()) > 2][:20]
    
    # Add generic experience entry
    exp_section = re.search(r'(?:EXPERIENCE|WORK EXPERIENCE)[:\s]+(.*?)(?=\n[A-Z]{3,}|EDUCATION|$)', 
                            text, re.IGNORECASE | re.DOTALL)
    if exp_section:
        data["experience"].append({
            "company": "See resume for details",
            "position": "Multiple positions",
            "location": "",
            "start_date": "",
            "end_date": "Present",
            "exp_skills": data["skills"][:5] if data["skills"] else [],
            "description": ["Experience details extracted from resume"]
        })
    
    # Add generic education entry
    edu_section = re.search(r'(?:EDUCATION)[:\s]+(.*?)(?=\n[A-Z]{3,}|$)', 
                            text, re.IGNORECASE | re.DOTALL)
    if edu_section:
        data["education"].append({
            "course": "See resume for details",
            "university": "Educational background available",
            "start_date": "",
            "end_date": ""
        })
    
    # Ensure we have at least name and email
    if not data["name"]:
        data["name"] = "Resume Holder"
    if not data["email"]:
        data["email"] = "email@example.com"
    
    print(f"✅ Fallback parser extracted: {len(data['skills'])} skills, "
          f"{len(data['experience'])} experience entries")
    
    return data

def call_llm(payload):
    response = requests.post(API_URL, json=payload, headers=headers)
    response.raise_for_status()
    data = response.json()
    llm_text = data['choices'][0]['message']['content']
    return llm_text



@st.cache_data(show_spinner=False)
def load_skills_from_json():
    # """Load skills from the cached JSON file"""
    # get_all_skills_from_llm()
    try:
        if os.path.exists(SKILL_CACHE_FILE):
            with open(SKILL_CACHE_FILE, "r") as f:
                return json.load(f)
        else:
            # If file doesn't exist, generate it
            return get_all_skills_from_llm()
    except Exception as e:
        st.error(f"Error loading skills: {e}")
        # Fallback to basic skills if there's an error
        return ["Python", "JavaScript", "SQL", "Cloud Computing", "Project Management", "Data Analysis"]


SKILL_CACHE_FILE = "../ask.ai/json/skill.json"

@st.cache_data(show_spinner=False)
def get_all_skills_from_llm():
    if os.path.exists(SKILL_CACHE_FILE):
        with open(SKILL_CACHE_FILE, "r") as f:
            return json.load(f)

    all_skills = set()
    
    # Define multiple focused prompts to get comprehensive coverage
    prompt_categories = [
        {
            "category": "Manual Labor & Service",
            "prompt": """List 200+ specific skills for: cleaning staff, janitors, dishwashers, laborers, warehouse workers, 
            packers, stock clerks, maintenance workers, landscapers, groundskeepers, movers, delivery personnel, 
            sanitation workers, custodians, housekeepers, laundry workers, kitchen helpers."""
        },
        {
            "category": "Agriculture & Farming",
            "prompt": """List 200+ specific skills for: farmers, farm workers, crop managers, livestock handlers, 
            agricultural technicians, irrigation specialists, harvesters, dairy workers, poultry workers, 
            farm equipment operators, greenhouse workers, vineyard workers, aquaculture workers, agronomists."""
        },
        {
            "category": "Transportation & Logistics",
            "prompt": """List 200+ specific skills for: truck drivers, taxi drivers, bus drivers, delivery drivers, 
            forklift operators, crane operators, pilots, ship captains, train conductors, dispatchers, 
            logistics coordinators, supply chain managers, freight handlers, route planners, fleet managers."""
        },
        {
            "category": "Construction & Trades",
            "prompt": """List 200+ specific skills for: carpenters, electricians, plumbers, welders, masons, 
            HVAC technicians, roofers, painters, drywall installers, tile setters, glaziers, insulators, 
            heavy equipment operators, pipefitters, ironworkers, sheet metal workers, construction managers."""
        },
        {
            "category": "Food Service & Hospitality",
            "prompt": """List 200+ specific skills for: chefs, cooks, bakers, bartenders, servers, baristas, 
            food prep workers, line cooks, pastry chefs, sous chefs, catering managers, restaurant managers, 
            hotel managers, front desk agents, concierges, housekeeping managers, event coordinators."""
        },
        {
            "category": "Retail & Sales",
            "prompt": """List 200+ specific skills for: retail associates, cashiers, sales representatives, 
            merchandisers, store managers, buyers, inventory specialists, loss prevention, visual merchandisers, 
            account executives, business development, inside sales, outside sales, territory managers."""
        },
        {
            "category": "Healthcare & Medical",
            "prompt": """List 200+ specific skills for: nurses, doctors, surgeons, medical assistants, 
            paramedics, EMTs, pharmacists, pharmacy technicians, radiologists, lab technicians, 
            phlebotomists, dental assistants, physical therapists, occupational therapists, medical coders."""
        },
        {
            "category": "Education & Training",
            "prompt": """List 200+ specific skills for: teachers, professors, tutors, instructional designers, 
            training specialists, curriculum developers, education administrators, school counselors, 
            librarians, teaching assistants, special education teachers, ESL teachers, corporate trainers."""
        },
        {
            "category": "Manufacturing & Production",
            "prompt": """List 200+ specific skills for: machine operators, production workers, assemblers, 
            quality inspectors, manufacturing engineers, production supervisors, CNC operators, fabricators, 
            industrial maintenance, production planners, process engineers, lean specialists, Six Sigma experts."""
        },
        {
            "category": "IT & Software Development",
            "prompt": """List 300+ specific technical skills including: programming languages (Python, Java, C++, C#, 
            JavaScript, TypeScript, Ruby, Go, Rust, PHP, Swift, Kotlin, R, MATLAB, Scala, etc.), frameworks 
            (React, Angular, Vue, Django, Flask, Spring, .NET, Node.js, Express, Laravel, etc.), databases, 
            cloud platforms, DevOps tools, version control, testing frameworks, APIs, microservices."""
        },
        {
            "category": "Cybersecurity & Network",
            "prompt": """List 200+ specific skills for: security analysts, penetration testers, ethical hackers, 
            security engineers, network administrators, firewall specialists, SIEM analysts, SOC analysts, 
            malware analysts, cryptographers, compliance officers, risk assessors, incident responders."""
        },
        {
            "category": "Data Science & Analytics",
            "prompt": """List 200+ specific skills including: data analysis tools, statistical methods, 
            machine learning algorithms, deep learning, NLP, computer vision, big data technologies (Hadoop, Spark), 
            visualization tools (Tableau, Power BI, Looker), SQL variants, Python libraries (pandas, numpy, scikit-learn), 
            R packages, data modeling, ETL processes, data warehousing."""
        },
        {
            "category": "Business & Finance",
            "prompt": """List 200+ specific skills for: accountants, bookkeepers, financial analysts, auditors, 
            tax professionals, investment bankers, wealth managers, controllers, CFOs, budget analysts, 
            credit analysts, financial planners, actuaries, payroll specialists, accounts payable/receivable."""
        },
        {
            "category": "Marketing & Advertising",
            "prompt": """List 200+ specific skills including: SEO, SEM, PPC, social media platforms (Facebook Ads, 
            Google Ads, LinkedIn Ads, TikTok Ads), email marketing, content marketing, copywriting, brand management, 
            market research, analytics tools (Google Analytics, Adobe Analytics), CRM systems, marketing automation, 
            growth hacking, influencer marketing, affiliate marketing, video marketing."""
        },
        {
            "category": "Creative & Design",
            "prompt": """List 200+ specific skills including: Adobe Creative Suite (Photoshop, Illustrator, InDesign, 
            After Effects, Premiere Pro), Figma, Sketch, UI/UX design, graphic design, web design, motion graphics, 
            3D modeling (Blender, Maya, 3ds Max), video editing, photography, illustration, typography, 
            color theory, wireframing, prototyping."""
        },
        {
            "category": "Engineering",
            "prompt": """List 200+ specific skills for: mechanical engineers, electrical engineers, civil engineers, 
            chemical engineers, aerospace engineers, industrial engineers, environmental engineers, software engineers, 
            systems engineers, including CAD tools (AutoCAD, SolidWorks, CATIA, Revit), simulation software, 
            technical drawing, project engineering, quality engineering."""
        },
        {
            "category": "Legal & Compliance",
            "prompt": """List 150+ specific skills for: lawyers, paralegals, legal assistants, compliance officers, 
            contract managers, legal researchers, court reporters, mediators, arbitrators, patent attorneys, 
            corporate counsel, litigation specialists, regulatory experts, legal technology."""
        },
        {
            "category": "HR & Recruitment",
            "prompt": """List 150+ specific skills for: HR managers, recruiters, talent acquisition, HR generalists, 
            compensation specialists, benefits administrators, HR analytics, HRIS systems (Workday, SAP SuccessFactors, 
            Oracle HCM), employee relations, performance management, training coordinators, organizational development."""
        },
        {
            "category": "Executive & Leadership",
            "prompt": """List 200+ specific skills for: CEOs, COOs, CFOs, CIOs, CTOs, VPs, directors, senior managers 
            including: strategic planning, business strategy, P&L management, M&A, stakeholder management, 
            corporate governance, change management, executive communication, board relations, investor relations, 
            crisis management, succession planning, organizational design."""
        },
        {
            "category": "Soft Skills & General",
            "prompt": """List 200+ specific soft skills and general competencies including: communication types 
            (verbal, written, presentation, public speaking, negotiation), leadership styles, emotional intelligence, 
            problem-solving approaches, critical thinking, creativity, adaptability, time management, organization, 
            collaboration, conflict resolution, decision-making, customer service, active listening."""
        }
    ]
    
    print("🔄 Generating comprehensive skill database across all sectors...")
    
    for idx, category_data in enumerate(prompt_categories, 1):
        print(f"📊 Processing category {idx}/{len(prompt_categories)}: {category_data['category']}")
        
        full_prompt = f"""{category_data['prompt']}

Return ONLY a JSON array of individual skill names. Be very specific and comprehensive. 
Include technical skills, tools, software, methodologies, certifications, and specific abilities.
No categories, no explanations, no numbering - just skill names."""

        payload = {
            "model": "meta/llama-3.1-70b-instruct",
            "messages": [
                {"role": "system", "content": "You are a skill taxonomy expert. Return a comprehensive JSON array of specific, actionable skill names."},
                {"role": "user", "content": full_prompt}
            ],
            "temperature": 0.3,
            "max_tokens": 8000
        }

        try:
            llm_response = call_llm(payload)
            
            try:
                skills_list = json.loads(llm_response)
            except json.JSONDecodeError:
                # Robust extraction
                cleaned_text = re.sub(r'[0-9]+\.?\s*', '', llm_response)
                cleaned_text = re.sub(r'```json|```', '', cleaned_text)
                skills_list = [
                    skill.strip().strip('"').strip("'").strip(',').strip('[]') 
                    for skill in re.split(r'[\n,•*-]+', cleaned_text) 
                    if skill.strip()
                ]
            
            # Clean and add skills
            excluded_keywords = ["category", "skills:", "example", "e.g.", "etc.", "including", "such as"]
            
            for skill in skills_list:
                skill = skill.strip().strip('"').strip("'").strip(',').strip()
                
                # Skip if contains excluded keywords
                if any(keyword in skill.lower() for keyword in excluded_keywords):
                    continue
                
                # Skip very short or very long entries
                if len(skill) < 2 or len(skill) > 70:
                    continue
                
                # Skip entries with parentheses at the end
                if re.search(r'\([^)]*$', skill):
                    continue
                
                # Skip entries that look like descriptions
                if skill.count(' ') > 6:
                    continue
                
                if skill:
                    all_skills.add(skill.title())
            
            print(f"   ✅ Added skills. Total so far: {len(all_skills)}")
            
        except Exception as e:
            print(f"   ⚠️ Error in category {category_data['category']}: {e}")
            continue
    
    # Add comprehensive baseline skills to ensure coverage
    baseline_skills = [
        # Manual & Service
        "Cleaning", "Janitorial Work", "Floor Cleaning", "Window Cleaning", "Carpet Cleaning",
        "Dishwashing", "Laundry", "Ironing", "Mopping", "Vacuuming", "Sweeping", "Dusting",
        "Waste Management", "Recycling", "Sanitation", "Disinfection", "Pest Control",
        
        # Agriculture
        "Farming", "Crop Management", "Livestock Care", "Tractor Operation", "Harvesting",
        "Planting", "Irrigation", "Soil Testing", "Fertilization", "Crop Rotation",
        "Animal Husbandry", "Dairy Farming", "Poultry Farming", "Beekeeping", "Composting",
        
        # Transportation
        "Driving", "Truck Driving", "Bus Driving", "Taxi Driving", "Forklift Operation",
        "CDL License", "Vehicle Maintenance", "Route Planning", "GPS Navigation", "Loading",
        "Unloading", "Delivery", "Dispatch", "Logistics", "Fleet Management",
        
        # Construction
        "Carpentry", "Plumbing", "Electrical Work", "Welding", "Masonry", "HVAC",
        "Roofing", "Painting", "Drywall", "Tiling", "Flooring", "Concrete Work",
        "Blueprint Reading", "Framing", "Demolition", "Scaffolding",
        
        # Food Service
        "Cooking", "Baking", "Food Preparation", "Food Safety", "Menu Planning",
        "Knife Skills", "Grilling", "Frying", "Sautéing", "Plating", "Bartending",
        "Coffee Making", "Food Handling", "Kitchen Management", "Recipe Development",
        
        # Retail
        "Customer Service", "Cash Handling", "POS Systems", "Inventory Management",
        "Merchandising", "Sales", "Upselling", "Stock Replenishment", "Loss Prevention",
        "Visual Merchandising", "Product Knowledge", "Returns Processing",
        
        # Healthcare
        "Patient Care", "First Aid", "CPR", "BLS", "ACLS", "Vital Signs",
        "Medical Terminology", "Phlebotomy", "IV Insertion", "Wound Care",
        "Medication Administration", "EMR Systems", "HIPAA Compliance",
        
        # IT & Programming
        "Python", "Java", "JavaScript", "C++", "C#", "PHP", "Ruby", "Go", "Rust",
        "TypeScript", "Swift", "Kotlin", "R", "MATLAB", "SQL", "HTML", "CSS",
        "React", "Angular", "Vue.js", "Node.js", "Django", "Flask", "Spring Boot",
        "Git", "Docker", "Kubernetes", "AWS", "Azure", "GCP", "Linux", "Windows Server",
        
        # Data & Analytics
        "Data Analysis", "Excel", "Power BI", "Tableau", "SQL", "Python", "R",
        "Statistics", "Data Visualization", "ETL", "Data Modeling", "Machine Learning",
        "Deep Learning", "TensorFlow", "PyTorch", "Pandas", "NumPy", "Scikit-learn",
        
        # Cybersecurity
        "Network Security", "Penetration Testing", "Ethical Hacking", "Firewall Management",
        "SIEM", "IDS/IPS", "Vulnerability Assessment", "Security Auditing", "Cryptography",
        "Incident Response", "Malware Analysis", "Security Compliance",
        
        # Business & Finance
        "Accounting", "Bookkeeping", "Financial Analysis", "Budgeting", "Forecasting",
        "Tax Preparation", "Auditing", "QuickBooks", "SAP", "Oracle Financials",
        "Financial Modeling", "Excel VBA", "Financial Reporting", "AP/AR",
        
        # Marketing
        "SEO", "SEM", "PPC", "Google Ads", "Facebook Ads", "Social Media Marketing",
        "Content Marketing", "Email Marketing", "Copywriting", "Google Analytics",
        "Marketing Automation", "CRM", "Salesforce", "HubSpot", "Mailchimp",
        
        # Design
        "Graphic Design", "UI/UX Design", "Adobe Photoshop", "Adobe Illustrator",
        "Figma", "Sketch", "InDesign", "After Effects", "Premiere Pro", "Video Editing",
        "Photography", "Photo Editing", "Web Design", "Mobile Design", "Wireframing",
        
        # Soft Skills
        "Communication", "Leadership", "Teamwork", "Problem Solving", "Critical Thinking",
        "Time Management", "Organization", "Adaptability", "Creativity", "Collaboration",
        "Negotiation", "Conflict Resolution", "Decision Making", "Public Speaking",
        "Presentation Skills", "Active Listening", "Emotional Intelligence", "Work Ethic",
        "Attention to Detail", "Multitasking", "Stress Management", "Customer Focus",
        
        # Project Management
        "Project Management", "Agile", "Scrum", "Kanban", "Waterfall", "JIRA", "Asana",
        "MS Project", "Risk Management", "Stakeholder Management", "Resource Planning",
        "PMP", "PRINCE2", "Change Management",
        
        # Engineering
        "AutoCAD", "SolidWorks", "CATIA", "Revit", "MATLAB", "Simulink", "ANSYS",
        "CAD", "CAM", "FEA", "CFD", "Technical Drawing", "GD&T", "Quality Control",
        
        # Executive
        "Strategic Planning", "Business Strategy", "P&L Management", "M&A",
        "Corporate Governance", "Stakeholder Management", "Board Relations",
        "Executive Leadership", "Change Management", "Business Development"
    ]
    
    # Add all baseline skills
    all_skills.update([skill.title() for skill in baseline_skills])
    
    # Convert to sorted list
    unique_skills = sorted(list(all_skills))
    
    # Save to cache
    os.makedirs(os.path.dirname(SKILL_CACHE_FILE), exist_ok=True)
    with open(SKILL_CACHE_FILE, "w") as f:
        json.dump(unique_skills, f, indent=2)

    print(f"\n✅ Generated {len(unique_skills)} unique skills covering all job sectors and levels!")
    return unique_skills


# ALL THE JOB ROLE LIST#


CACHE_FILE = "../ask.ai/json/job_roles.json"

def get_all_roles_from_llm():
    if os.path.exists(CACHE_FILE):
        with open(CACHE_FILE, "r") as f:
            return json.load(f)
    prompt = (
    "You are a global expert in careers. Provide an exhaustive list of job roles "
    "across all industries worldwide. Include detailed job roles from major fields like "
    "software development, software testing, hardware engineering, IT infrastructure, networking, "
    "cybersecurity, manufacturing, healthcare, finance, marketing, and more. "
    "List entry-level, mid-level, senior, and specialized roles such as 'Software Tester', "
    "'Automation Tester', 'Hardware Design Engineer', 'Firmware Developer', 'Network Administrator', "
    "'Security Analyst', and all variations in these fields. "
    "Return only a JSON array of job role names without extra text."
)


    payload = {
        "model": "meta/llama-3.1-70b-instruct",
        "messages": [
            {"role": "system", "content": "You are a global expert in careers worldwide."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.1,
        "max_tokens": 3500
    }

    llm_response = call_llm(payload)

    try:
        roles_list = json.loads(llm_response)
        roles = sorted(set(role.strip().title() for role in roles_list if role.strip()))
    except json.JSONDecodeError:
        cleaned = llm_response.replace("[", "").replace("]", "").replace('"', "").replace("\n", "")
        roles = [role.strip().title() for role in cleaned.split(",") if role.strip()]
        roles = sorted(set(roles))

    with open(CACHE_FILE, "w") as f:
        json.dump(roles, f, indent=2)

    return roles






# exract excel column contents
def extract_skills_from_column(ws, example_col_index):
    all_skills = []
    seen = set()
    for row in ws.iter_rows(min_row=2):
        skill = row[example_col_index].value
        if skill and skill not in seen:
            seen.add(skill)
            all_skills.append(skill)
    return all_skills

# extract the excel file
@st.cache_data(show_spinner=False)
def load_excel_data(filename):
    wb = load_workbook(filename=filename)
    ws = wb.active
    header = [cell.value for cell in next(ws.iter_rows(min_row=1, max_row=1))]
    return ws, header

# read the pdf file
def extract_text_from_pdf(pdf_file):
    """Extract text from PDF file"""
    try:
        with pdfplumber.open(pdf_file) as pdf:
            text = ""
            for page in pdf.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        st.error(f"Error reading PDF: {e}")
        return ""
    
import docx2txt

def extract_text_from_docx(docx_file):
    try:
        text = docx2txt.process(docx_file)
        return text.strip()
    except Exception as e:
        st.error(f"Error reading DOCX: {e}")
        return ""
""




def extract_details_from_jd(job_description: str, max_retries: int = 3) -> Dict[str, Any]:
    """
    Extract structured details from job description with retry logic.
    """
    
    # Truncate if too long
    if len(job_description) > 10000:
        print(f"⚠️ JD truncated from {len(job_description)} to 10000 characters")
        job_description = job_description[:10000]
    
    SCHEMA_DEFINITION = """{
  "job_title": "string",
  "company": "string or null",
  "location": "string or null",
  "required_skills": ["array of strings - technical and soft skills"],
  "required_experience": "string or null (e.g., '3-5 years')",
  "education": "string or null (e.g., 'Bachelor's in CS')",
  "responsibilities": ["array of strings - main duties"],
  "qualifications": ["array of strings - requirements"],
  "nice_to_have": ["array of strings - optional skills"],
  "benefits": "string or null"
}"""

    prompt = f"""Extract structured information from the job description below.
Return ONLY valid JSON. Do NOT include explanations, markdown, or code blocks.

RULES:
1. Extract only information explicitly mentioned.
2. Use null for missing fields.
3. Return empty arrays [] for missing list fields.
4. Extract ALL skills mentioned (technical, tools, soft skills).

SCHEMA TO FOLLOW:
{SCHEMA_DEFINITION}

Job Description:
{job_description}"""

    for attempt in range(max_retries):
        try:
            timeout = 30 + (attempt * 15)
            print(f"🔄 Parsing JD - Attempt {attempt + 1}/{max_retries} (timeout: {timeout}s)...")
            
            payload = {
                "model": HEAVY_MODEL,
                "messages": [
                    {
                        "role": "system",
                        "content": "You are an expert job description parser. Extract information exactly as it appears. Return only valid JSON."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                "temperature": 0.1,
                "max_tokens": 2000,
                "top_p": 0.7
            }

            response = requests.post(url, headers=headers, json=payload, timeout=timeout)
            response.raise_for_status()
            result = response.json()
            
            if 'choices' not in result or not result['choices']:
                raise ValueError("API returned empty response")
            
            response_text = result['choices'][0]['message']['content'].strip()
            response_text = response_text.replace('```json', '').replace('```', '').strip()

            json_start = response_text.find('{')
            json_end = response_text.rfind('}') + 1
            
            if json_start == -1 or json_end <= json_start:
                raise ValueError("No valid JSON found")

            json_str = response_text[json_start:json_end]
            data = json.loads(json_str)

            # Ensure all fields exist
            list_fields = ["required_skills", "responsibilities", "qualifications", "nice_to_have"]
            for field in list_fields:
                if field not in data:
                    data[field] = []
                elif not isinstance(data[field], list):
                    data[field] = [data[field]] if data[field] else []

            string_fields = ["job_title", "company", "location", "required_experience", "education", "benefits"]
            for field in string_fields:
                if field not in data:
                    data[field] = ""
                elif data[field] is None:
                    data[field] = ""

            print(f"✅ Successfully parsed JD on attempt {attempt + 1}")
            return data

        except requests.exceptions.Timeout:
            print(f"⏱️ JD timeout on attempt {attempt + 1}/{max_retries}")
            if attempt < max_retries - 1:
                time.sleep(2 ** attempt)
                continue
            else:
                return _fallback_parse_jd(job_description)
        
        except Exception as e:
            print(f"❌ JD parse error on attempt {attempt + 1}: {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(2 ** attempt)
                continue
            else:
                return _fallback_parse_jd(job_description)


def extract_details_from_jd(job_description: str, max_retries: int = 3) -> Dict[str, Any]:
    """
    Extract structured details from job description with retry logic.
    """
    
    # Truncate if too long
    if len(job_description) > 10000:
        print(f"⚠️ JD truncated from {len(job_description)} to 10000 characters")
        job_description = job_description[:10000]
    
    SCHEMA_DEFINITION = """{
  "job_title": "string",
  "company": "string or null",
  "location": "string or null",
  "required_skills": ["array of strings - technical and soft skills"],
  "required_experience": "string or null (e.g., '3-5 years')",
  "education": "string or null (e.g., 'Bachelor's in CS')",
  "responsibilities": ["array of strings - main duties"],
  "qualifications": ["array of strings - requirements"],
  "nice_to_have": ["array of strings - optional skills"],
  "benefits": "string or null"
}"""

    # ✅ FIXED: More explicit prompt with examples
    prompt = f"""You are a JSON parser. Extract job description information and return ONLY a valid JSON object.

CRITICAL RULES:
- Return ONLY the JSON object, nothing else
- No markdown, no code blocks, no explanations
- No ```json``` tags
- Start with {{ and end with }}

SCHEMA:
{SCHEMA_DEFINITION}

EXAMPLE OUTPUT:
{{
  "job_title": "Software Engineer",
  "company": "Tech Corp",
  "location": "Remote",
  "required_skills": ["Python", "Django", "REST APIs"],
  "required_experience": "3-5 years",
  "education": "Bachelor's in Computer Science",
  "responsibilities": ["Build APIs", "Write tests"],
  "qualifications": ["Strong Python skills", "Team player"],
  "nice_to_have": ["AWS experience"],
  "benefits": "Health insurance, 401k"
}}

Now extract from this job description:
{job_description}

Return only the JSON object:"""

    for attempt in range(max_retries):
        try:
            timeout = 30 + (attempt * 15)
            print(f"🔄 Parsing JD - Attempt {attempt + 1}/{max_retries} (timeout: {timeout}s)...")
            
            # ✅ FIXED: Check if URL and headers are properly configured
            if not url or not headers:
                print("❌ API configuration missing (url or headers)")
                return _fallback_parse_jd(job_description)
            
            payload = {
                "model": HEAVY_MODEL,  # ⚠️ Make sure this is correct (e.g., "meta/llama-3.1-70b-instruct")
                "messages": [
                    {
                        "role": "system",
                        "content": "You are a JSON extraction tool. You ONLY output valid JSON objects. Never use markdown or code blocks."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                "temperature": 0.1,
                "max_tokens": 2000,
                "top_p": 0.7
            }

            response = requests.post(url, headers=headers, json=payload, timeout=timeout)
            
            # ✅ FIXED: Better error handling
            if response.status_code == 404:
                print(f"❌ API endpoint not found (404). Check your URL and model name.")
                print(f"   URL: {url}")
                print(f"   Model: {HEAVY_MODEL}")
                return _fallback_parse_jd(job_description)
            
            response.raise_for_status()
            result = response.json()
            
            if 'choices' not in result or not result['choices']:
                raise ValueError("API returned empty response")
            
            response_text = result['choices'][0]['message']['content'].strip()
            
            # ✅ FIXED: Better JSON extraction
            # Remove common wrapper patterns
            response_text = response_text.replace('```json', '').replace('```', '').strip()
            response_text = re.sub(r'^[^{]*', '', response_text)  # Remove everything before first {
            response_text = re.sub(r'[^}]*$', '', response_text)  # Remove everything after last }
            
            # Find JSON boundaries
            json_start = response_text.find('{')
            json_end = response_text.rfind('}') + 1
            
            if json_start == -1 or json_end <= json_start:
                print(f"⚠️ No JSON found in response. First 200 chars: {response_text[:200]}")
                raise ValueError("No valid JSON found")

            json_str = response_text[json_start:json_end]
            
            # ✅ FIXED: Try to fix common JSON issues
            try:
                data = json.loads(json_str)
            except json.JSONDecodeError as e:
                print(f"⚠️ JSON decode error: {str(e)}")
                print(f"   Problematic JSON: {json_str[:500]}")
                raise

            # Ensure all fields exist
            list_fields = ["required_skills", "responsibilities", "qualifications", "nice_to_have"]
            for field in list_fields:
                if field not in data:
                    data[field] = []
                elif not isinstance(data[field], list):
                    data[field] = [data[field]] if data[field] else []

            string_fields = ["job_title", "company", "location", "required_experience", "education", "benefits"]
            for field in string_fields:
                if field not in data:
                    data[field] = ""
                elif data[field] is None:
                    data[field] = ""

            print(f"✅ Successfully parsed JD on attempt {attempt + 1}")
            print(f"   Extracted {len(data.get('required_skills', []))} skills")
            return data

        except requests.exceptions.Timeout:
            print(f"⏱️ JD timeout on attempt {attempt + 1}/{max_retries}")
            if attempt < max_retries - 1:
                time.sleep(2 ** attempt)
                continue
            else:
                return _fallback_parse_jd(job_description)
        
        except requests.exceptions.HTTPError as e:
            print(f"❌ HTTP error on attempt {attempt + 1}: {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(2 ** attempt)
                continue
            else:
                return _fallback_parse_jd(job_description)
        
        except Exception as e:
            print(f"❌ JD parse error on attempt {attempt + 1}: {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(2 ** attempt)
                continue
            else:
                return _fallback_parse_jd(job_description)


def _fallback_parse_jd(text: str) -> Dict[str, Any]:
    """
    Enhanced fallback parser for job descriptions.
    """
    import re
    
    print("🔧 Using fallback JD parser...")
    
    data = {
        "job_title": "",
        "company": "",
        "location": "",
        "required_skills": [],
        "required_experience": "",
        "education": "",
        "responsibilities": [],
        "qualifications": [],
        "nice_to_have": [],
        "benefits": ""
    }
    
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    
    # ✅ FIXED: Better job title extraction
    if lines:
        # Try to find "Job Title:" or similar
        for i, line in enumerate(lines[:5]):
            if any(keyword in line.lower() for keyword in ['job title', 'position', 'role']):
                if ':' in line:
                    data["job_title"] = line.split(':', 1)[1].strip()
                elif i + 1 < len(lines):
                    data["job_title"] = lines[i + 1]
                break
        if not data["job_title"]:
            data["job_title"] = lines[0]
    
    # ✅ FIXED: Enhanced skill extraction
    # Expanded skill keywords
    tech_skills = [
        # Programming
        "Python", "Java", "JavaScript", "TypeScript", "C++", "C#", "Ruby", "Go", "Rust", "PHP", "Swift", "Kotlin",
        # Web
        "React", "Angular", "Vue", "Node.js", "Express", "Django", "Flask", "FastAPI", "Spring", "ASP.NET",
        # Data
        "SQL", "PostgreSQL", "MySQL", "MongoDB", "Redis", "Cassandra", "Elasticsearch",
        "Pandas", "NumPy", "Scikit-learn", "TensorFlow", "PyTorch", "Keras",
        # Cloud & DevOps
        "AWS", "Azure", "GCP", "Docker", "Kubernetes", "Jenkins", "GitLab", "CircleCI", "Terraform",
        # Tools
        "Git", "Jira", "Confluence", "Tableau", "Power BI", "Excel", "Figma", "Postman",
        # Concepts
        "Machine Learning", "Deep Learning", "Data Analysis", "ETL", "CI/CD", "Microservices",
        "REST API", "GraphQL", "Agile", "Scrum", "TDD", "DevOps"
    ]
    
    soft_skills = [
        "Communication", "Leadership", "Problem Solving", "Team Player", "Analytical",
        "Critical Thinking", "Creativity", "Adaptability", "Time Management"
    ]
    
    all_skills = tech_skills + soft_skills
    found_skills = set()
    
    text_lower = text.lower()
    text_words = set(re.findall(r'\b\w+\b', text_lower))
    
    # ✅ FIXED: Better matching
    for skill in all_skills:
        skill_lower = skill.lower()
        # Check exact phrase match or word match
        if skill_lower in text_lower or any(word in text_words for word in skill_lower.split()):
            found_skills.add(skill)
    
    # ✅ FIXED: Extract custom skills from text
    # Look for bullet points or lines that might contain skills
    skill_sections = re.findall(
        r'(?:skills?|technologies?|requirements?|qualifications?)[\s:]*\n((?:[-•*]\s*.+\n?)+)',
        text,
        re.IGNORECASE | re.MULTILINE
    )
    
    for section in skill_sections:
        # Extract items from bullets
        items = re.findall(r'[-•*]\s*(.+)', section)
        for item in items:
            # Clean and add if looks like a skill
            item = item.strip().rstrip('.,;')
            if 3 <= len(item) <= 50 and not any(word in item.lower() for word in ['year', 'degree', 'bachelor']):
                found_skills.add(item)
    
    data["required_skills"] = list(found_skills)[:30]  # Limit to 30 skills
    
    # ✅ FIXED: Better experience extraction
    exp_patterns = [
        r'(\d+[\+\-]?\s*(?:to\s+)?\d*\s*years?)',
        r'(\d+\+?\s*years?)',
        r'(minimum\s+\d+\s+years?)',
        r'(at least\s+\d+\s+years?)'
    ]
    for pattern in exp_patterns:
        exp_match = re.search(pattern, text, re.IGNORECASE)
        if exp_match:
            data["required_experience"] = exp_match.group(1)
            break
    
        # ✅ FIXED: Better education extraction
    edu_patterns = [
        r"(bachelor's?\s+(?:degree\s+)?(?:in\s+)?[\w\s]+)",
        r"(master's?\s+(?:degree\s+)?(?:in\s+)?[\w\s]+)",
        r"(phd\s+(?:in\s+)?[\w\s]+)",
        r"(degree\s+in\s+[\w\s]+)"
    ]

    for pattern in edu_patterns:
        edu_match = re.search(pattern, text, re.IGNORECASE)
        if edu_match:
            data["education"] = edu_match.group(1).strip()
            break
    
    # ✅ FIXED: Extract responsibilities
    resp_section = re.search(
        r"(?:responsibilities|duties|what you'll do)[\s:]*\n((?:[-•*]\s*.+\n?)+)",
        text,
        re.IGNORECASE | re.MULTILINE
    )

    if resp_section:
        data["responsibilities"] = [
            r.strip().rstrip('.,;') 
            for r in re.findall(r'[-•*]\s*(.+)', resp_section.group(1))
        ][:10]
    
    # ✅ FIXED: Extract qualifications
    qual_section = re.search(
        r'(?:qualifications|requirements|must have)[\s:]*\n((?:[-•*]\s*.+\n?)+)',
        text,
        re.IGNORECASE | re.MULTILINE
    )
    if qual_section:
        data["qualifications"] = [
            q.strip().rstrip('.,;') 
            for q in re.findall(r'[-•*]\s*(.+)', qual_section.group(1))
        ][:10]
    
    print(f"✅ Fallback JD parser extracted:")
    print(f"   - {len(data['required_skills'])} skills")
    print(f"   - {len(data['responsibilities'])} responsibilities")
    print(f"   - {len(data['qualifications'])} qualifications")
    
    return data




def lower_set(arr: List[str]) -> set:
    return set([s.lower().strip() for s in arr])

def highlight_missing_items(source: List[str], required: List[str]) -> List[str]:
    source_set = lower_set(source)
    required_set = lower_set(required)
    missing = required_set - source_set
    return list(missing)

def enhance_resume_with_jd_skills(resume_skills: List[str], jd_skills: List[str]) -> List[str]:
    updated = resume_skills.copy()
    resume_lower = lower_set(updated)
    for skill in jd_skills:
        if skill.lower().strip() not in resume_lower:
            updated.append(skill)
    return updated



def rewrite_resume_for_job(resume_data: dict, jd_data: dict) -> dict:
    """
    Enhance the existing candidate resume using the JD,
    but DO NOT add new fake roles, skills, or projects that
    are not present in user's resume. Only highlight and optimize.
    """
    # st.write(resume_data)
    # st.write(jd_data)
    # ============================================
    # 🔧 SAFETY CHECK
    # ============================================
    if resume_data is None:
        raise ValueError("resume_data cannot be None")
    
    if not isinstance(resume_data, dict):
        raise ValueError(f"resume_data must be a dictionary, got {type(resume_data)}")
    
    if jd_data is None:
        jd_data = {}
    
    # ============================================
    # 🔧 FIX: Handle nested structure
    # If resume_data contains a nested "resume_data" key, flatten it
    # ============================================
    # ============================================
# 🔧 FIX: Handle nested structure
    # ============================================
    if "resume_data" in resume_data:
        nested_data = resume_data["resume_data"]
        
        # If it's a string, parse it as JSON
        if isinstance(nested_data, str):
            print("⚠️ Detected JSON string, parsing...")
            actual_resume = json.loads(nested_data)
        # If it's already a dict, use it
        elif isinstance(nested_data, dict):
            print("⚠️ Detected nested dict, flattening...")
            actual_resume = nested_data
        else:
            actual_resume = resume_data
    else:
        actual_resume = resume_data
    
    # ============================================
    # START WITH COMPLETELY NEW DICT
    # ============================================
    rewritten_resume = {}

    # ============================================
    # PRESERVE BASIC INFO - NEVER CHANGE THESE
    # ============================================
    rewritten_resume["name"] = actual_resume.get("name", "")
    rewritten_resume["email"] = actual_resume.get("email", "")
    rewritten_resume["phone"] = actual_resume.get("phone", "")
    rewritten_resume["location"] = actual_resume.get("location", "")
    rewritten_resume["url"] = actual_resume.get("url", "")

    # ============================================
    # 1. SEPARATE EDUCATION AND CERTIFICATIONS
    # ============================================
    education_all = actual_resume.get("education", [])
    education_list = []
    certification_list = []
    
    for item in education_all:
        if any(k in item for k in ["duration", "gpa", "degree", "course", "university", "institution"]):
            education_list.append(item)
        elif any(k in item for k in ["issuer", "name", "provider_name", "certificate_name"]):
            certification_list.append(item)
        else:
            education_list.append(item)

    # Add certifications from separate field
    cert_data = actual_resume.get("certifications", []) or actual_resume.get("certificate", [])
    if isinstance(cert_data, list):
        certification_list.extend(cert_data)

    rewritten_resume["education"] = education_list
    rewritten_resume["certifications"] = certification_list

    # ============================================
    # 2. NORMALIZE CANDIDATE SKILLS
    # ============================================
    candidate_skills = actual_resume.get("skills", [])
    if isinstance(candidate_skills, list) and candidate_skills:
        if candidate_skills and isinstance(candidate_skills[0], dict):
            flat_skills = []
            for group in candidate_skills:
                if isinstance(group, dict):
                    for vals in group.values():
                        if isinstance(vals, list):
                            flat_skills.extend([str(v) for v in vals if v])
            candidate_skills = flat_skills
        else:
            candidate_skills = [str(s) for s in candidate_skills if s]
    else:
        candidate_skills = []

    # ============================================
    # 3. GET EXPERIENCE AND PROJECTS
    # ============================================
    experience = actual_resume.get("experience", [])
    projects = actual_resume.get("projects", []) or actual_resume.get("project", [])

    all_exp_desc = []
    for exp in experience:
        desc = exp.get("description", [])
        if isinstance(desc, str):
            desc = [line.strip() for line in desc.split("\n") if line.strip()]
        all_exp_desc.extend(desc if isinstance(desc, list) else [])

    # ============================================
    # 4. EXTRACT JD INFO
    # ============================================
    job_title = jd_data.get("job_title", "")
    responsibilities = jd_data.get("responsibilities", [])
    required_skills = jd_data.get("required_skills", [])

    responsibilities = [str(r) for r in responsibilities if r] if isinstance(responsibilities, list) else []
    required_skills = [str(s) for s in required_skills if s] if isinstance(required_skills, list) else []

    # ============================================
    # 5. REWRITE SUMMARY
    # ============================================
    original_summary = actual_resume.get("summary", "")
    
    if original_summary or candidate_skills or all_exp_desc:
        summary_prompt = f"""
Rewrite this resume summary for a {job_title} position. 

CRITICAL RULES:
1. Write in THIRD PERSON - NO "I", "my", "me" anywhere
2. Use PAST TENSE for experience (e.g., "worked with", "built", "tested")
3. Write 2-3 sentences MAXIMUM
4. Only mention what's actually in their resume - no fake skills
5. DO NOT include any preamble - return ONLY the summary text
6. DO NOT use [Name] placeholder - just describe the person's experience

THEIR RESUME INFO:
Name: {actual_resume.get('name','')}
Current Summary: {original_summary}
Work Experience: {all_exp_desc[:10]}
Skills: {', '.join(candidate_skills[:15])}

JD Context: {', '.join(required_skills[:10])}

Return ONLY the summary text, nothing else.
"""

        raw_summary = call_llm_api(summary_prompt, 200)
        
        # Clean up AI preambles
        clean_summary = raw_summary.strip()
        
        preambles = [
            "Here is the rewritten summary:",
            "Here's the rewritten summary:",
            "The rewritten summary:",
            "Rewritten summary:",
            "Summary:",
        ]
        
        for preamble in preambles:
            if clean_summary.lower().startswith(preamble.lower()):
                clean_summary = clean_summary[len(preamble):].strip()
        
        clean_summary = clean_summary.strip('"\'')
        
        rewritten_resume["summary"] = clean_summary
    else:
        rewritten_resume["summary"] = ""
    
    rewritten_resume["job_title"] = job_title if job_title else actual_resume.get("job_title", "")

    # ============================================
    # 6. REWRITE EXPERIENCE (ONLY IF EXISTS)
    # ============================================
    if experience and len(experience) > 0:
        exp_prompt = f"""
Rewrite work experience bullet points for a {job_title} application.

ABSOLUTE REQUIREMENTS:
1. THIRD PERSON ONLY - NO "I", "my", "me", "we"
2. Start bullets with PAST TENSE VERBS: "Transported", "Maintained", "Ensured", "Provided"
3. Keep company names, positions, dates, locations EXACTLY as original
4. Only improve description bullets
5. Keep any numbers/percentages from original
6. Return ONLY JSON array, no markdown, no explanation

ORIGINAL EXPERIENCE:
{json.dumps(experience, indent=2)}

JD CONTEXT:
{', '.join(required_skills[:10])}

Return ONLY this JSON (no ```json``` wrapper):
[
  {{
    "position": "exact from original",
    "company": "exact from original",
    "start_date": "exact from original",
    "end_date": "exact from original",
    "location": "exact from original",
    "description": ["improved bullet 1", "improved bullet 2", "improved bullet 3"]
  }}
]
"""
        rewritten_exp_text = call_llm_api(exp_prompt, 600)
        try:
            rewritten_exp_text = rewritten_exp_text.strip().replace("```json", "").replace("```", "").strip()
            json_start = rewritten_exp_text.find('[')
            json_end = rewritten_exp_text.rfind(']') + 1
            
            if json_start >= 0 and json_end > json_start:
                rewritten_resume["experience"] = json.loads(rewritten_exp_text[json_start:json_end])
            else:
                rewritten_resume["experience"] = experience
        except Exception as e:
            print(f"Error parsing experience: {e}")
            rewritten_resume["experience"] = experience
    else:
        rewritten_resume["experience"] = experience if experience else []

    # ============================================
    # 7. REWRITE PROJECTS (ONLY IF EXISTS)
    # ============================================
    if projects and len(projects) > 0:
        proj_prompt = f"""
Rewrite project descriptions for a {job_title} role.

RULES:
1. THIRD PERSON - No "I", "my", "me", "we"
2. Start with PAST TENSE VERBS
3. Keep project names exactly as original
4. Only improve description bullets
5. Return ONLY JSON array, no markdown

ORIGINAL PROJECTS:
{json.dumps(projects, indent=2)}

Return ONLY this JSON:
[
  {{
    "name": "exact from original",
    "description": ["improved bullet 1", "improved bullet 2"]
  }}
]
"""
        rewritten_proj_text = call_llm_api(proj_prompt, 400)
        try:
            rewritten_proj_text = rewritten_proj_text.strip().replace("```json", "").replace("```", "").strip()
            json_start = rewritten_proj_text.find('[')
            json_end = rewritten_proj_text.rfind(']') + 1
            
            if json_start >= 0 and json_end > json_start:
                rewritten_resume["projects"] = json.loads(rewritten_proj_text[json_start:json_end])
            else:
                rewritten_resume["projects"] = projects
        except Exception as e:
            print(f"Error parsing projects: {e}")
            rewritten_resume["projects"] = projects
    else:
        rewritten_resume["projects"] = projects if projects else []

    # ============================================
    # 8. CATEGORIZE SKILLS (NO NEW SKILLS)
    # ============================================
    if candidate_skills and len(candidate_skills) > 0:
        skills_prompt = f"""
Organize these exact skills into categories. DO NOT add any new skills.

THEIR SKILLS (use ONLY these):
{', '.join(candidate_skills)}

Return ONLY this JSON (no markdown):
{{
  "technicalSkills": ["skill1"],
  "tools": ["tool1"],
  "softSkills": ["soft1"]
}}

Only include categories with items from above list.
"""
        skills_text = call_llm_api(skills_prompt, 300)
        try:
            skills_text = skills_text.strip().replace("```json", "").replace("```", "").strip()
            json_start = skills_text.find('{')
            json_end = skills_text.rfind('}') + 1
            
            if json_start >= 0 and json_end > json_start:
                parsed_skills = json.loads(skills_text[json_start:json_end])
                rewritten_resume["skills"] = {k: v for k, v in parsed_skills.items() if v}
            else:
                rewritten_resume["skills"] = candidate_skills
        except Exception as e:
            print(f"Error parsing skills: {e}")
            rewritten_resume["skills"] = candidate_skills
    else:
        rewritten_resume["skills"] = candidate_skills if candidate_skills else []

    # ============================================
    # 9. PRESERVE OTHER FIELDS
    # ============================================
    if "custom_sections" in actual_resume:
        rewritten_resume["custom_sections"] = actual_resume["custom_sections"]
    
    if "achievements" in actual_resume:
        rewritten_resume["achievements"] = actual_resume["achievements"]
    
    if "languages" in actual_resume:
        rewritten_resume["languages"] = actual_resume["languages"]

    # ============================================
    # FINAL: Ensure no nested structure
    # ============================================
    if "resume_data" in rewritten_resume:
        del rewritten_resume["resume_data"]

    return rewritten_resume



def generate_basic_description(position: str, company: str, candidate_skills: list, call_llm_api) -> str:
    key_skills = candidate_skills[:5]
    prompt = f"""
    You are a resume filler AI. Generate a single, short paragraph (2-3 sentences maximum) 
    describing the general responsibilities for a '{position}' role at '{company}'. 
    
    The description MUST incorporate these key skills: {', '.join(key_skills)}. 
    
    Start with an action verb and provide a summary of work done. Do not include bullet points.
    Return ONLY the paragraph.
    """
    return call_llm_api(prompt, 100)

# --- Helper function for LLM API ---
def call_llm_api(prompt: str, max_tokens: int = 200) -> str:
  

    payload = {
        "model": FAST_MODEL,
        "messages": [
            {"role": "system", "content": "You are an expert AI assistant for resume optimization and job matching."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.1,
        "max_tokens": max_tokens
    }

    # Try a few times if the server gives 500 or times out
    for attempt in range(3):
        try:
            response = requests.post(url, headers=headers, json=payload, timeout=60)
            response.raise_for_status()
            result = response.json()
            return result["choices"][0]["message"]["content"]

        except requests.exceptions.Timeout:
            print(f"⏳ Timeout on attempt {attempt + 1}, retrying...")
            time.sleep(2 ** attempt)

        except requests.exceptions.HTTPError as e:
            if response.status_code >= 500:
                print(f"⚠️ Server error 500 on attempt {attempt + 1}, retrying...")
                time.sleep(2 ** attempt)
            else:
                # Show more details for debugging
                print("❌ Client error:", response.text)
                break

        except Exception as e:
            print("Unexpected error:", e)
            break

    return "⚠️ Sorry, the AI service is currently unavailable. Please try again later."





def analyze_and_improve_resume(resume_data: Dict[str, Any], job_description: str = "") -> Dict[str, Any]:
    """
    Automatically improve resume content with better grammar, stronger action verbs,
    and professional formatting while preserving all information, using an external API.
    """
    

    resume_text = json.dumps(resume_data, indent=2)
    
    # The prompt now includes the job_description for context
    prompt = f"""
    Improve the following resume while maintaining all factual information and structure.
    
    Make these improvements:
    1. Fix all grammar, spelling, and punctuation errors
    2. Replace weak verbs with strong action verbs (Led, Spearheaded, Optimized, etc.)
    3. Ensure consistent tense (past for previous jobs, present for current)
    4. Improve clarity and impact of descriptions
    5. Ensure professional tone throughout
    6. Make descriptions more concise and impactful
    7. Ensure consistent formatting (dates, capitalization, etc.)
    
    **Job Description for Context (if available):**
    {job_description or "N/A"}
    
    IMPORTANT:
    - Keep the exact same JSON structure
    - Do NOT remove or add any sections
    - Do NOT change facts, dates, or company names
    - Only improve the language and presentation
    
    Resume Data:
    {resume_text}

    Return the improved resume in the EXACT SAME JSON structure, with only the text content improved.
    Return only valid JSON.
    """

    payload = {
        "model": FAST_MODEL,
        "messages": [
            {"role": "system", "content": "You are an expert resume writer who improves content while preserving structure and facts. Return improved resume in exact same JSON format."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.3,
        "max_tokens": 5000
    }

    try:
        # Perform API call
        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()
        result = response.json()
        
        # NOTE: Assumes the Llama model response structure wraps the content
        response_text = result['choices'][0]['message']['content']

        # Extract JSON (robustly handle text wrapping)
        json_start = response_text.find('{')
        json_end = response_text.rfind('}') + 1
        
        if json_start == -1 or json_end <= json_start:
            st.error("Could not find valid JSON in improved resume response")
            return resume_data

        improved_data = json.loads(response_text[json_start:json_end])
        return improved_data

    except requests.exceptions.RequestException as e:
        # Crucial: Check API config and network connectivity
        st.error(f"External API request failed: {str(e)}. Check your `url` and `headers` config.")
        return resume_data
    except json.JSONDecodeError as e:
        st.error(f"Failed to parse JSON response: {str(e)}. Using original data.")
        return resume_data
    except KeyError as e:
        st.error(f"Unexpected response structure from external API. Missing key: {str(e)}. Using original data.")
        return resume_data
    except Exception as e:
        st.error(f"Error improving resume: {str(e)}. Using original data.")
        return resume_data



def check_specific_section(section_name: str, section_data: Any) -> Dict[str, Any]:
    """
    Check a specific section for grammar, clarity, and professional quality.
    Useful for real-time checking as user edits.
    """
    
    section_text = json.dumps({section_name: section_data}, indent=2)
    
    prompt = f"""
    Review this resume section for quality and provide quick feedback.
    
    Check for:
    - Grammar and spelling errors
    - Clarity and professionalism
    - Strong action verbs
    - Consistency
    
    Section:
    {section_text}

    Return feedback in JSON format:
    {{
        "has_issues": true/false,
        "issues": ["list of specific issues found"],
        "suggestions": ["list of improvement suggestions"],
        "quality_score": <number 0-10>
    }}

    Return only valid JSON.
    """

    payload = {
        "model": "meta/llama-3.1-70b-instruct",
        "messages": [
            {"role": "system", "content": "You are a resume quality checker. Provide concise, actionable feedback."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.2,
        "max_tokens": 1000
    }

    try:
        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()
        result = response.json()
        response_text = result['choices'][0]['message']['content']

        json_start = response_text.find('{')
        json_end = response_text.rfind('}') + 1
        if json_start == -1 or json_end <= json_start:
            return {"has_issues": False, "issues": [], "suggestions": [], "quality_score": 5}

        feedback = json.loads(response_text[json_start:json_end])
        return feedback

    except Exception as e:
        return {"has_issues": False, "issues": [], "suggestions": [], "quality_score": 5}



def rewrite_resume_for_job_manual(resume_data: dict, jd_data: dict) -> dict:
    """
    Enhance the existing candidate resume using JD guidance,
    without adding any fake skills, tools, achievements or content
    that does not exist in the original resume.
    """
    rewritten_resume = resume_data.copy()
    rewritten_resume.pop("input_method", None)

    # -------------------- Education and Certifications (NO CHANGE) --------------------
    education_all = resume_data.get("education", [])
    education_list, certification_list = [], []

    for item in education_all:
        if any(k in item for k in ["duration", "gpa", "degree", "course", "university"]):
            education_list.append(item)
        elif any(k in item for k in ["issuer", "name", "provider_name", "certificate_name"]):
            certification_list.append(item)
        else:
            education_list.append(item)

    cert_data = resume_data.get("certificate", [])
    if isinstance(cert_data, list):
        certification_list.extend(cert_data)

    rewritten_resume["education"] = education_list
    rewritten_resume["certifications"] = certification_list

    # -------------------- Flatten Candidate Skills (NO CHANGE) --------------------
    candidate_skills = resume_data.get("skills", [])
    if isinstance(candidate_skills, list) and candidate_skills:
        if isinstance(candidate_skills[0], dict):
            flat_skills = []
            for group in candidate_skills:
                if isinstance(group, dict):
                    for vals in group.values():
                        if isinstance(vals, list):
                            flat_skills.extend([str(v) for v in vals if v])
            candidate_skills = flat_skills
        else:
            candidate_skills = [str(s) for s in candidate_skills if s]
    else:
        candidate_skills = []

    # -------------------- Job Info (NO CHANGE) --------------------
    job_title = jd_data.get("job_title", "")
    responsibilities = jd_data.get("responsibilities", [])
    required_skills = jd_data.get("required_skills", [])
    rewritten_resume["job_title"] = job_title

    responsibilities = [str(r) for r in responsibilities if r] if isinstance(responsibilities, list) else []
    required_skills = [str(s) for s in required_skills if s] if isinstance(required_skills, list) else []

    # -------------------- Merge Experience (NO CHANGE) --------------------
    experience_all = []
    if "experience" in resume_data and isinstance(resume_data["experience"], list):
        experience_all.extend(resume_data["experience"])
    if "professional_experience" in resume_data and isinstance(resume_data["professional_experience"], list):
        experience_all.extend(resume_data["professional_experience"])

    rewritten_resume.pop("experience", None)
    rewritten_resume.pop("professional_experience", None)
    rewritten_resume.pop("input_method", None)

    seen = set()
    merged_experience = []
    for exp in experience_all:
        key = (exp.get("company", ""), exp.get("position", ""))
        if key not in seen:
            merged_experience.append(exp)
            seen.add(key)

    rewritten_resume["experience"] = merged_experience
    rewritten_resume["total_experience_count"] = len(merged_experience)

    # -------------------- Generate Experience (UPDATED LOGIC) --------------------
    rewritten_experience = []
    all_exp_desc = []

    for exp in merged_experience:
        position = exp.get("position", "Professional")
        company = exp.get("company", "A Company")
        start_date = exp.get("start_date", "")
        end_date = exp.get("end_date", "")
        exp_skills = exp.get("exp_skills", [])

        exp_prompt = f"""
        You are a professional resume writer. Rewrite ONLY the existing content for this experience
        to align with the {job_title} role.

        STRICT RULES:
        - DO NOT add new tools, tasks, achievements, metrics or claims.
        - Use only the information already present in the role.
        - You may rephrase content and highlight relevancy using existing skills only.
        - Highlight only if a skill matches BOTH resume and JD.
        - Make description strong but 100% factual.

        USER ROLE: {position} at {company} ({start_date} to {end_date})
        USER EXISTING SKILLS: {', '.join(exp_skills)}
        JD SKILLS TO HIGHLIGHT ONLY IF THEY ALREADY EXIST ABOVE: {', '.join(required_skills)}

        Return ONLY JSON array with one object containing:
        {{"description": ["sentence1", "sentence2", ...]}}
        """

        rewritten_exp_text = call_llm_api(exp_prompt, 500)

        try:
            json_start = rewritten_exp_text.find('[')
            json_end = rewritten_exp_text.rfind(']') + 1
            rewritten_exp_list = json.loads(rewritten_exp_text[json_start:json_end])
            rewritten_exp_obj = rewritten_exp_list[0] if rewritten_exp_list else {}
            desc = rewritten_exp_obj.get("description", [])
            if isinstance(desc, str):
                desc = [line.strip() for line in desc.split("\n") if line.strip()]

            new_exp = {
                "company": company,
                "position": position,
                "start_date": start_date,
                "end_date": end_date,
                "description": desc if desc else [f"Worked as a {position} at {company}."]
            }
            rewritten_experience.append(new_exp)
            all_exp_desc.extend(desc)

        except Exception:
            fallback_desc = [f"Worked as a {position} at {company}, using {', '.join(exp_skills)}."]
            new_exp = {
                "company": company,
                "position": position,
                "start_date": start_date,
                "end_date": end_date,
                "description": fallback_desc
            }
            rewritten_experience.append(new_exp)
            all_exp_desc.extend(fallback_desc)

    rewritten_resume["experience"] = rewritten_experience
    rewritten_resume["total_experience_count"] = len(rewritten_experience)

    # -------------------- Professional Summary (UPDATED LOGIC) --------------------
    summary_prompt = f"""
    Rewrite the professional summary for {resume_data.get('name','')} applying for {job_title}.

    STRICT RULES:
    - Do NOT invent new skills, tools, achievements or claims.
    - Only use what exists in resume & experience.
    - Highlight only skills that match both resume and JD.
    - Keep 100% factual but strongly aligned to the job.
    - Write in first person implied (no "I", "me", "my").

    USER SUMMARY: {resume_data.get('summary','')}
    EXPERIENCE POINTS: {all_exp_desc[:10]}
    JD SKILLS (highlight only if they exist in resume): {', '.join(required_skills)}

    Return ONLY the rewritten summary text. 
    Do NOT add titles, headers, labels, notes, explanations, quotes, or extra lines.
    """

    rewritten_resume["summary"] = call_llm_api(summary_prompt, 200)

    # -------------------- Projects (UPDATED LOGIC) --------------------
    if "project" in resume_data:
        projects = resume_data.get("project", [])
        rewritten_projects = []

        for proj in projects:
            project_name = proj.get("projectname", proj.get("name", "Project"))
            tools_used = proj.get("tools", [])
            original_desc = proj.get("description", proj.get("decription", ""))

            proj_prompt = f"""
            Rewrite this project description using ONLY the existing information provided.

            ALLOWED:
            - You may elaborate, expand, and improve clarity.
            - You may make the sentences more descriptive and professional.
            - You may highlight relevant responsibilities already present.
            - You may emphasize important parts using stronger wording.

            STRICT RULES:
            - Do NOT introduce or assume new technologies, tools, achievements, metrics, or responsibilities.
            - If a JD skill matches the project and ALREADY exists in the description or tools, highlight it naturally.
            - The final output must remain factual based only on given content.

            Project Name: {project_name}
            Original Description: {original_desc}
            Tools Already Used: {', '.join(tools_used)}
            JD Skills to highlight ONLY if already present: {', '.join(required_skills)}

            Return ONLY a JSON array in this exact structure:
            [{{ 
                "name": "{project_name}", 
                "description": ["sentence1", "sentence2", ...], 
                "overview": "" 
            }}]

            Do NOT include notes, labels, headers, comments, or explanations.
            """


            rewritten_proj_text = call_llm_api(proj_prompt, 300)

            try:
                json_start = rewritten_proj_text.find('[')
                json_end = rewritten_proj_text.rfind(']') + 1
                parsed_proj = json.loads(rewritten_proj_text[json_start:json_end])

                proj_obj = parsed_proj[0] if parsed_proj else {}
                desc = proj_obj.get("description", [])
                if isinstance(desc, str):
                    proj_obj["description"] = [line.strip() for line in desc.split("\n") if line.strip()]

                proj_obj["name"] = project_name
                proj_obj["overview"] = proj_obj.get("overview", "")
                rewritten_projects.append(proj_obj)

            except:
                fallback_desc = (
                    [line.strip() for line in original_desc.split("\n") if line.strip()]
                    if isinstance(original_desc, str)
                    else [f"Worked on {project_name} using {', '.join(tools_used)}."]
                )
                rewritten_projects.append({
                    "name": project_name,
                    "description": fallback_desc,
                    "overview": ""
                })

        rewritten_resume["projects"] = rewritten_projects

    # -------------------- Skill Categorization (UPDATED LOGIC) --------------------
    skills_prompt = f"""
    Categorize ONLY the candidate's EXISTING skills to align with {job_title}.

    STRICT RULES:
    - DO NOT add or assume any new skills.
    - If a skill matches the JD AND exists in user's list, move it to the top.
    - Categorize only, do not invent.

    USER SKILLS: {', '.join(candidate_skills[:30])}
    JD SKILLS to highlight ONLY if also in user skills: {', '.join(required_skills)}

    Return ONLY JSON with:
    technicalSkills, tools, cloudSkills, softSkills, languages
    """
    skills_text = call_llm_api(skills_prompt, 300)

    try:
        json_start = skills_text.find('{')
        json_end = skills_text.rfind('}') + 1
        parsed_skills = json.loads(skills_text[json_start:json_end])
        categorized_skills = {k: parsed_skills.get(k, []) for k in 
                              ["technicalSkills","tools","cloudSkills","softSkills","languages"]}
        rewritten_resume["skills"] = {k: v for k,v in categorized_skills.items() if v}
    except:
        rewritten_resume["skills"] = candidate_skills

    # -------------------- Handle Custom Sections (NO CHANGE) --------------------
    standard_keys = {
        "name", "email", "phone", "location", "url", "summary", "job_title",
        "education", "experience", "skills", "projects", "certifications", 
        "achievements", "total_experience_count", "professional_experience",
        "certificate", "project", "input_method"
    }

    custom_sections = resume_data.get("custom_sections", {})
    if isinstance(custom_sections, dict):
        for title, description in custom_sections.items():
            clean_title = str(title).strip()
            clean_description = description.strip() if isinstance(description, str) else (
                "\n".join(map(str, description)) if isinstance(description, list) else str(description)
            )
            if clean_title and clean_title not in standard_keys:
                rewritten_resume[clean_title] = clean_description

    for key, value in resume_data.items():
        if key not in standard_keys and isinstance(value, str):
            rewritten_resume[key] = value.strip()

    if "project" in rewritten_resume and "projects" in rewritten_resume:
        del rewritten_resume["project"]
    if "certificate" in rewritten_resume and "certifications" in rewritten_resume:
        del rewritten_resume["certificate"]
    if "custom_sections" in rewritten_resume:
        del rewritten_resume["custom_sections"]

    rewritten_resume["projects"] = rewritten_resume.get("projects", [])
    rewritten_resume["certifications"] = rewritten_resume.get("certifications", [])

    return rewritten_resume


def is_valid_email(email):
    pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
    return re.match(pattern, email) is not None

def is_valid_phone(phone):
    pattern = r'^\+?\d{7,15}$'
    return re.match(pattern, phone) is not None





def save_user_resume(email, resume_data, input_method=None):

    def convert_dates(obj):
        if isinstance(obj, dict):
            return {k: convert_dates(v) for k, v in obj.items()}
        elif isinstance(obj, list):
            return [convert_dates(item) for item in obj]
        elif hasattr(obj, "isoformat"):
            return obj.isoformat()
        return obj

    resume_data = convert_dates(resume_data)

    try:
        conn = get_connection()
        cursor = conn.cursor()

        cursor.execute("""
            MERGE user_resumes AS target
            USING (
                SELECT 
                    ? AS email,
                    ? AS resume_data,
                    ? AS input_method
            ) AS source
            ON target.email = source.email
            WHEN MATCHED THEN
                UPDATE SET
                    resume_data = source.resume_data,
                    input_method = source.input_method,
                    updated_at = SYSDATETIME()
            WHEN NOT MATCHED THEN
                INSERT (email, resume_data, input_method, created_at, updated_at)
                VALUES (
                    source.email,
                    source.resume_data,
                    source.input_method,
                    SYSDATETIME(),
                    SYSDATETIME()
                );
        """,
        email,
        json.dumps(resume_data),
        input_method
        )

        conn.commit()
        return True

    except Exception as e:
        st.error(f"Error saving resume data: {e}")
        return False

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass




def get_user_template_path(email):
    conn = get_connection()
    cursor = conn.cursor()

    cursor.execute("""
        SELECT template_key, template_name, html, css, uploaded_at, original_filename
        FROM user_templates
        WHERE email = ?
        ORDER BY uploaded_at DESC
    """, email)

    rows = cursor.fetchall()

    cursor.close()
    conn.close()

    templates = {}
    for row in rows:
        templates[row[0]] = {
            "name": row[1],
            "html": row[2],
            "css": row[3],
            "uploaded_at": row[4],
            "original_filename": row[5]
        }

    return templates



# def get_user_template_path(user_email):
#     """Return the JSON path for the given user's templates."""
#     if not user_email:
#         raise ValueError("❌ user_email is None — make sure the user is logged in before saving templates.")
    
#     os.makedirs(TEMPLATES_DIR, exist_ok=True)
#     safe_email = user_email.replace("@", "_at_").replace(".", "_dot_")
#     return os.path.join(TEMPLATES_DIR, f"{safe_email}.json")








def load_user_templates(user_email):
    """Load templates for the logged-in user from database."""
    templates = {}

    try:
        conn = get_connection()
        cursor = conn.cursor()

        cursor.execute("""
            SELECT template_key,
                   template_name,
                   html,
                   css,
                   uploaded_at,
                   original_filename
            FROM user_templates
            WHERE email = ?
            ORDER BY uploaded_at DESC
        """, user_email)

        for row in cursor.fetchall():
            templates[row[0]] = {
                "name": row[1],
                "html": row[2],
                "css": row[3],
                "uploaded_at": row[4],
                "original_filename": row[5]
            }

    except Exception as e:
        st.error(f"Error loading templates: {e}")

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass

    return templates




def save_user_templates(user_email, templates):
    """Save or update templates for the logged-in user in database."""
    try:
        conn = get_connection()
        cursor = conn.cursor()

        for template_key, data in templates.items():
            cursor.execute("""
                MERGE user_templates AS target
                USING (
                    SELECT
                        ? AS email,
                        ? AS template_key,
                        ? AS template_name,
                        ? AS html,
                        ? AS css,
                        ? AS uploaded_at,
                        ? AS original_filename
                ) AS source
                ON target.email = source.email
                   AND target.template_key = source.template_key
                WHEN MATCHED THEN
                    UPDATE SET
                        template_name = source.template_name,
                        html = source.html,
                        css = source.css,
                        uploaded_at = source.uploaded_at,
                        original_filename = source.original_filename,
                        updated_at = SYSDATETIME()
                WHEN NOT MATCHED THEN
                    INSERT (
                        email,
                        template_key,
                        template_name,
                        html,
                        css,
                        uploaded_at,
                        original_filename,
                        created_at,
                        updated_at
                    )
                    VALUES (
                        source.email,
                        source.template_key,
                        source.template_name,
                        source.html,
                        source.css,
                        source.uploaded_at,
                        source.original_filename,
                        SYSDATETIME(),
                        SYSDATETIME()
                    );
            """,
            user_email,
            template_key,
            data.get("name"),
            data.get("html"),
            data.get("css"),
            data.get("uploaded_at"),
            data.get("original_filename")
            )

        conn.commit()
        return True

    except Exception as e:
        st.error(f"Error saving templates: {e}")
        return False

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass


# import base64

# def get_user_ppt_template_path(user_email):
#     """Get the file path for user's PPT templates."""
#     return os.path.join("user_data", f"{user_email}_ppt_templates.json")

# def load_user_ppt_templates(user_email):
#     """Load PowerPoint templates from JSON file for the logged-in user."""
#     path = get_user_ppt_template_path(user_email)
#     if os.path.exists(path):
#         with open(path, "r", encoding="utf-8") as f:
#             templates = json.load(f)
#             # Convert base64 back to binary for ppt_data and convert lists back to sets
#             for key in templates:
#                 if 'ppt_data' in templates[key]:
#                     templates[key]['ppt_data'] = base64.b64decode(templates[key]['ppt_data'])
#                 # Convert lists back to sets
#                 if 'heading_shapes' in templates[key]:
#                     templates[key]['heading_shapes'] = set(templates[key]['heading_shapes'])
#                 if 'basic_info_shapes' in templates[key]:
#                     templates[key]['basic_info_shapes'] = set(templates[key]['basic_info_shapes'])
#             return templates
#     return {}

# def save_user_ppt_templates(user_email, templates):
#     """Save PowerPoint templates to JSON file for the logged-in user."""
#     path = get_user_ppt_template_path(user_email)
#     os.makedirs(os.path.dirname(path), exist_ok=True)
    
#     # Create a copy to avoid modifying the original
#     templates_copy = {}
#     for key, val in templates.items():
#         templates_copy[key] = val.copy()
#         # Convert binary ppt_data to base64 for JSON storage
#         if 'ppt_data' in val and isinstance(val['ppt_data'], bytes):
#             templates_copy[key]['ppt_data'] = base64.b64encode(val['ppt_data']).decode('utf-8')
#         # Convert sets to lists for JSON serialization
#         if 'heading_shapes' in val and isinstance(val['heading_shapes'], set):
#             templates_copy[key]['heading_shapes'] = list(val['heading_shapes'])
#         if 'basic_info_shapes' in val and isinstance(val['basic_info_shapes'], set):
#             templates_copy[key]['basic_info_shapes'] = list(val['basic_info_shapes'])
    
#     with open(path, "w", encoding="utf-8") as f:
#         json.dump(templates_copy, f, indent=4, ensure_ascii=False)




def extract_json(text):
    """
    Extracts JSON from LLM response, handling markdown fences and extra text.
    """
    if not text:
        raise ValueError("Empty response – cannot extract JSON")

    # Remove markdown code fences
    text = re.sub(r'```json\s*|\s*```', '', text, flags=re.IGNORECASE)

    # Try multiple extraction strategies
    strategies = [
        # Strategy 1: Find outermost balanced braces
        lambda t: re.search(r'\{(?:[^{}]|(?:\{[^{}]*\}))*\}', t, re.DOTALL),
        
        # Strategy 2: Non-greedy match
        lambda t: re.search(r'\{.*?\}(?=\s*$)', t, re.DOTALL),
        
        # Strategy 3: Greedy match
        lambda t: re.search(r'\{.*\}', t, re.DOTALL)
    ]

    for strategy in strategies:
        match = strategy(text)
        if match:
            json_text = match.group(0)
            try:
                return json.loads(json_text)
            except json.JSONDecodeError:
                continue

    # Last resort: try the entire text
    try:
        return json.loads(text.strip())
    except json.JSONDecodeError as e:
        print("=" * 60)
        print("FAILED TO EXTRACT JSON")
        print("=" * 60)
        print(f"Response length: {len(text)}")
        print(f"First 500 chars: {text[:500]}")
        print("=" * 60)
        raise ValueError(f"No valid JSON found in response: {str(e)}")


def call_llm_for_json(prompt: str, max_tokens: int = 2000) -> dict:
    """
    Wrapper around call_llm_api that handles JSON extraction.
    Returns a dict instead of a string.
    """
    try:
        # Call the existing function
        result = call_llm_api(prompt, max_tokens=max_tokens)
        
        print("\n" + "=" * 80)
        print("🔍 Raw LLM Response:")
        print(result[:500])
        print("=" * 80)
        
        # Check for error message
        if "⚠️" in result or "unavailable" in result.lower():
            raise ValueError(f"API service error: {result}")
        
        # Extract JSON from the response
        return extract_json(result)
        
    except Exception as e:
        print(f"\n❌ Error in call_llm_for_json: {str(e)}")
        raise

# Add this to your utils.py file

import json

def ai_ats_score(resume_text, jd_text, max_retries=3):
    """
    Evaluates resume-job fit using AI and returns structured scores with detailed skill breakdown.
    Includes validation and retry logic for consistency.
    """
    
    def safe_truncate(text, max_length=3000):
        """Safely convert to string and truncate"""
        if text is None:
            return ""
        if isinstance(text, dict):
            text = json.dumps(text)
        return str(text)[:max_length]
    
    resume_truncated = safe_truncate(resume_text)
    jd_truncated = safe_truncate(jd_text)
    
    if not resume_truncated.strip():
        raise ValueError("Resume text is empty")
    if not jd_truncated.strip():
        raise ValueError("Job description is empty")
    
    prompt = f"""
You are an ATS scoring engine designed to evaluate resume–job fit using semantic understanding.

Your job is to:
1. Compare the resume with the job description comprehensively
2. Score the match on a scale of 0–100
3. Categorize skills into technical and soft skills
4. Identify matched and missing skills in both categories
5. Evaluate semantic similarity, not just keyword overlap
6. Consider experience level and relevance

Return ONLY valid JSON in the following format:

{{
  "overall_score": number (0-100),
  "skill_match_score": number (0-100),
  "experience_match_score": number (0-100),
  "semantic_alignment_score": number (0-100),
  "technical_skills": {{
    "matched": [string] (list of technical skills found in resume that match JD),
    "missing": [string] (list of technical skills from JD not found in resume),
    "match_percentage": number (0-100)
  }},
  "soft_skills": {{
    "matched": [string] (list of soft skills demonstrated in resume),
    "missing": [string] (list of soft skills from JD not demonstrated),
    "match_percentage": number (0-100)
  }},
  "experience_analysis": {{
    "years_required": string (e.g., "5-7 years"),
    "years_present": string (e.g., "6 years"),
    "level_match": string (e.g., "Good Match", "Senior Level", "Entry Level"),
    "relevant_experience": [string] (list of relevant work experiences)
  }},
  "strengths": [string] (3-5 key strengths of the resume),
  "weaknesses": [string] (3-5 areas that need improvement),
  "recommendations": [string] (5-7 specific actionable recommendations),
  "explanation": string (2-3 paragraph detailed explanation of the analysis)
}}

Guidelines:
- Technical skills include: programming languages, frameworks, tools, technologies, platforms, databases, methodologies
- Soft skills include: leadership, communication, teamwork, problem-solving, time management, adaptability
- Be thorough and specific in identifying skills
- Provide actionable, specific recommendations
- Consider industry context and role requirements

CRITICAL REQUIREMENTS FOR SCORE CONSISTENCY:
- The overall_score MUST be logically consistent with matched skills
- If matched technical_skills is empty or has fewer than 3 items, overall_score should be below 40
- If overall_score is above 60, there MUST be at least 5+ matched technical skills
- If overall_score is above 80, there MUST be at least 10+ matched technical skills
- The match_percentage in technical_skills and soft_skills must align with the number of matched vs missing skills
- Ensure the response is COMPLETE with all fields properly filled
- Double-check all scores for logical consistency before returning

Resume:
{resume_truncated}

Job Description:
{jd_truncated}

IMPORTANT: Return ONLY the JSON object. No markdown formatting, no code blocks, no explanation text.
"""

    for attempt in range(max_retries):
        try:
            print(f"\n🔄 ATS Analysis Attempt {attempt + 1}/{max_retries}")
            
            # Increased max_tokens to prevent truncation
            result = call_llm_for_json(prompt, max_tokens=4500)
            
            # Validate and set defaults for all required fields
            required_fields = {
                "overall_score": 0,
                "skill_match_score": 0,
                "experience_match_score": 0,
                "semantic_alignment_score": 0,
                "technical_skills": {
                    "matched": [],
                    "missing": [],
                    "match_percentage": 0
                },
                "soft_skills": {
                    "matched": [],
                    "missing": [],
                    "match_percentage": 0
                },
                "experience_analysis": {
                    "years_required": "Not specified",
                    "years_present": "Not specified",
                    "level_match": "Unknown",
                    "relevant_experience": []
                },
                "strengths": [],
                "weaknesses": [],
                "recommendations": [],
                "explanation": "Analysis completed"
            }
            
            # Ensure all fields exist with proper defaults
            for field, default in required_fields.items():
                if field not in result or result[field] is None:
                    result[field] = default
                elif isinstance(default, dict):
                    # For nested dicts, ensure sub-fields exist
                    for sub_field, sub_default in default.items():
                        if sub_field not in result[field] or result[field][sub_field] is None:
                            result[field][sub_field] = sub_default
            
            # Check if response is complete
            is_complete = (
                len(result.get('explanation', '')) > 50 and
                len(result.get('recommendations', [])) >= 3 and
                len(result.get('strengths', [])) >= 2
            )
            
            if not is_complete:
                print(f"⚠️ Incomplete response detected on attempt {attempt + 1}")
                if attempt < max_retries - 1:
                    continue
            
            # Validate score consistency
            tech_matched = len(result['technical_skills']['matched'])
            soft_matched = len(result['soft_skills']['matched'])
            overall_score = result['overall_score']
            
            # Check for logical inconsistencies
            inconsistent = False
            
            if overall_score > 60 and tech_matched < 5:
                print(f"⚠️ Score Inconsistency: Overall={overall_score}% but only {tech_matched} technical skills matched")
                inconsistent = True
            
            if overall_score > 80 and tech_matched < 10:
                print(f"⚠️ Score Inconsistency: Overall={overall_score}% but only {tech_matched} technical skills matched")
                inconsistent = True
            
            if overall_score < 40 and tech_matched > 10:
                print(f"⚠️ Score Inconsistency: Overall={overall_score}% but {tech_matched} technical skills matched")
                inconsistent = True
            
            if tech_matched == 0 and soft_matched == 0 and overall_score > 30:
                print(f"⚠️ Major Inconsistency: Overall={overall_score}% but NO skills matched")
                inconsistent = True
            
            # If inconsistent and not last attempt, retry
            if inconsistent and attempt < max_retries - 1:
                print(f"🔄 Retrying due to score inconsistency...")
                continue
            
            # Calculate a more accurate overall score based on components
            if inconsistent:
                print("⚙️ Recalculating overall score based on skill matches...")
                
                # Calculate match percentages if not set
                tech_total = tech_matched + len(result['technical_skills']['missing'])
                soft_total = soft_matched + len(result['soft_skills']['missing'])
                
                tech_match_pct = (tech_matched / tech_total * 100) if tech_total > 0 else 0
                soft_match_pct = (soft_matched / soft_total * 100) if soft_total > 0 else 0
                
                result['technical_skills']['match_percentage'] = round(tech_match_pct)
                result['soft_skills']['match_percentage'] = round(soft_match_pct)
                
                # Weighted calculation: 50% technical, 20% soft, 30% experience
                calculated_overall = (
                    tech_match_pct * 0.50 +
                    soft_match_pct * 0.20 +
                    result['experience_match_score'] * 0.30
                )
                
                result['overall_score'] = round(calculated_overall)
                result['skill_match_score'] = round((tech_match_pct * 0.7) + (soft_match_pct * 0.3))
                
                print(f"✅ Adjusted Overall Score: {result['overall_score']}%")
            
            # Debug output
            print("\n" + "=" * 80)
            print("✅ ATS Analysis Results:")
            print(f"Overall Score: {result['overall_score']}%")
            print(f"Technical Skills - Matched: {tech_matched}, Missing: {len(result['technical_skills']['missing'])}")
            print(f"Soft Skills - Matched: {soft_matched}, Missing: {len(result['soft_skills']['missing'])}")
            print(f"Recommendations: {len(result['recommendations'])}")
            print(f"Explanation Length: {len(result['explanation'])} characters")
            print("=" * 80)
            
            return result
            
        except Exception as e:
            print(f"\n❌ Error on attempt {attempt + 1}: {str(e)}")
            if attempt < max_retries - 1:
                print("🔄 Retrying...")
                continue
            else:
                import traceback
                traceback.print_exc()
                
                # Return error structure with all required fields
                return {
                    "overall_score": 0,
                    "skill_match_score": 0,
                    "experience_match_score": 0,
                    "semantic_alignment_score": 0,
                    "technical_skills": {
                        "matched": [],
                        "missing": [],
                        "match_percentage": 0
                    },
                    "soft_skills": {
                        "matched": [],
                        "missing": [],
                        "match_percentage": 0
                    },
                    "experience_analysis": {
                        "years_required": "Error",
                        "years_present": "Error",
                        "level_match": "Error",
                        "relevant_experience": []
                    },
                    "strengths": ["Unable to complete analysis"],
                    "weaknesses": ["API error occurred"],
                    "recommendations": ["Please try again or contact support"],
                    "explanation": f"Error occurred during analysis after {max_retries} attempts: {str(e)}"
                }


def get_score_color(score):
    """Return color based on score."""
    if score >= 80:
        return "#00FF7F"  # Green
    elif score >= 60:
        return "#FFD700"  # Gold
    elif score >= 40:
        return "#FFA500"  # Orange
    else:
        return "#FF6347"  # Red


def get_score_label(score):
    """Return label based on score."""
    if score >= 80:
        return "Excellent Match"
    elif score >= 60:
        return "Good Match"
    elif score >= 40:
        return "Fair Match"
    else:
        return "Needs Improvement"


def validate_ats_result(result):
    """
    Additional validation function to check ATS result consistency.
    Returns (is_valid, issues_list)
    """
    issues = []
    
    overall = result.get('overall_score', 0)
    tech_matched = len(result.get('technical_skills', {}).get('matched', []))
    soft_matched = len(result.get('soft_skills', {}).get('matched', []))
    
    # Rule 1: High score requires matched skills
    if overall > 60 and tech_matched < 5:
        issues.append(f"Overall score {overall}% is too high for only {tech_matched} matched technical skills")
    
    # Rule 2: Low score with many matched skills
    if overall < 40 and tech_matched > 8:
        issues.append(f"Overall score {overall}% is too low for {tech_matched} matched technical skills")
    
    # Rule 3: No matched skills with moderate/high score
    if tech_matched == 0 and soft_matched == 0 and overall > 30:
        issues.append(f"Overall score {overall}% but no skills matched at all")
    
    # Rule 4: Check if explanation is meaningful
    if len(result.get('explanation', '')) < 50:
        issues.append("Explanation is too short or missing")
    
    # Rule 5: Check recommendations
    if len(result.get('recommendations', [])) < 3:
        issues.append("Insufficient recommendations provided")
    
    return len(issues) == 0, issues

def analyze_slide_structure(slide_texts):
    """Use LLM to find headings, subheadings, and related contents - INCLUDING BASIC INFO"""
    print("analysing slide structure with basic info...")
    system_prompt = """
You are a presentation content analyzer.

You will receive a list of text boxes per slide, each with coordinates (x, y) and text content.

Your goal:
- Identify **basic information** (name, title, contact info, summary) - these should be marked as editable
- Identify **main headings** (e.g., "Selected Experiences", "Core Competencies", "Profile")
- Identify **subheadings** under a main heading
- Group all descriptive text under the closest subheading

CRITICAL: Basic information like name, job title, contact info, and summary SHOULD be marked as editable content, NOT locked headings.

Return **strictly valid JSON**:
[
  {
    "slide_number": 1,
    "sections": [
      {
        "heading": "Profile",
        "heading_shape_idx": 0,
        "is_basic_info": true,
        "subsections": [
          {
            "subheading": "",
            "subheading_shape_idx": null,
            "content_shape_idx": 1,
            "content_text": "Sudheer Varma\\nOracle Cloud HCM Techno Functional Consultant\\n8+ years of experience",
            "content_type": "basic_info"
          }
        ]
      }
    ]
  }
]
"""
    user_prompt = f"Analyze these slide texts and extract ALL content including basic information:\n{json.dumps(slide_texts, indent=2)}"

    payload = {
        "model": "meta/llama-3.1-70b-instruct",
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "temperature": 0.3,
        "max_tokens": 2500
    }

    try:
        response = requests.post(API_URL, headers=headers, json=payload)
        result_text = response.json()["choices"][0]["message"]["content"]
        
        try:
            structured = json.loads(result_text)
        except json.JSONDecodeError:
            match = re.search(r'\[.*\]', result_text, re.DOTALL)
            structured = json.loads(match.group(0)) if match else []
        
        return structured
    except:
        return []
    
def analyze_content_length(text):
    """Analyze content length and return category"""
    word_count = len(text.split())
    char_count = len(text)
    
    if word_count < 30 or char_count < 150:
        return "very_short", word_count, char_count
    elif word_count < 60 or char_count < 400:
        return "short", word_count, char_count
    elif word_count < 100 or char_count < 700:
        return "medium", word_count, char_count
    else:
        return "long", word_count, char_count



def generate_ppt_sections(resume_data, structured_slides):
    """Generate AI content for ALL sections including basic info"""
    print("generating ppt sections with basic info...")
    sections = []
    for slide in structured_slides:
        for section in slide.get("sections", []):
            # Check if this is basic info section
            is_basic_info = section.get("is_basic_info", False)
            
            entry = {
                "slide_number": slide.get("slide_number"),
                "heading": section.get("heading", ""),
                "heading_shape_idx": section.get("heading_shape_idx"),
                "is_basic_info": is_basic_info,
                "subsections": []
            }

            for subsec in section.get("subsections", []):
                original_text = subsec.get("content_text", "")
                content_type = subsec.get("content_type", "normal")
                
                # For basic info, we need special handling
                if is_basic_info or content_type == "basic_info":
                    # Basic info should be short and concise
                    target_word_count = min(len(original_text.split()), 50)
                    target_char_count = min(len(original_text), 300)
                    length_category = "short"
                else:
                    length_category, word_count, char_count = analyze_content_length(original_text)
                    target_word_count = word_count
                    target_char_count = char_count
                
                entry["subsections"].append({
                    "subheading": subsec.get("subheading", ""),
                    "subheading_shape_idx": subsec.get("subheading_shape_idx"),
                    "content_shape_idx": subsec.get("content_shape_idx"),
                    "original_content": original_text,
                    "target_word_count": target_word_count,
                    "target_char_count": target_char_count,
                    "length_category": length_category,
                    "content_type": content_type,
                    "is_basic_info": is_basic_info
                })

            sections.append(entry)

    system_prompt = """
You are a professional AI assistant that generates polished PowerPoint content from resume data.

CRITICAL REQUIREMENTS:
1. Match the target length specified for each section PRECISELY
2. For BASIC INFORMATION sections (name, title, contact, summary), replace with actual resume data
3. For EXPERIENCE sections, generate professional content that matches the target length

SPECIAL HANDLING FOR BASIC INFORMATION:
- Replace name with: {name}
- Replace title/position with actual experience titles
- Update contact info: {email}, {phone}, {location}
- For profile/summary: use the resume summary

TASK:
- For basic info sections: DIRECTLY REPLACE with actual resume data
- For other sections: Generate professional content that MATCHES target word count
- Always write in a formal, presentation-friendly tone
- Rewrite resume data to sound concise, polished, and visually engaging
- DO NOT use bullet points - write in paragraph form
- IMPORTANT: Respect the target_word_count for each section

OUTPUT FORMAT (strict JSON):
[
  {
    "slide_number": 1,
    "heading": "Profile",
    "is_basic_info": true,
    "subsections": [
      {
        "subheading": "",
        "content_shape_idx": 1,
        "target_word_count": 25,
        "generated_content": "Anaha Joy\\nMagento 2 Developer\\n2+ years of experience\\n+91 8304078233 | anahajoy2022@gmail.com"
      }
    ]
  }
]
"""

    user_prompt = f"""
Resume Data:
{json.dumps(resume_data, indent=2)}

PPT Structure with Target Lengths:
{json.dumps(sections, indent=2)}

Generate content for ALL sections:
- For BASIC INFO: Directly replace with actual resume data
- For EXPERIENCE: Generate professional content matching target length
"""

    payload = {
        "model": "meta/llama-3.1-70b-instruct",
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "temperature": 0.45,
        "max_tokens": 4000
    }

    try:
        response = requests.post(API_URL, headers=headers, json=payload)
        result_text = response.json()["choices"][0]["message"]["content"]
        
        try:
            structured_output = json.loads(result_text)
        except json.JSONDecodeError:
            match = re.search(r'\[.*\]', result_text, re.DOTALL)
            structured_output = json.loads(match.group(0)) if match else []
        
        return structured_output
    except:
        return []
    
def trim_content_to_length(content, target_word_count, target_char_count):
    """Trim content to match target length while maintaining coherence"""
    current_words = len(content.split())
    current_chars = len(content)
    
    # If content is within acceptable range (±20%), return as is
    word_diff = abs(current_words - target_word_count) / max(target_word_count, 1)
    if word_diff < 0.2:
        return content
    
    # If too long, trim sentences
    if current_words > target_word_count * 1.2:
        sentences = content.split('. ')
        trimmed = []
        word_count = 0
        
        for sentence in sentences:
            sentence_words = len(sentence.split())
            if word_count + sentence_words <= target_word_count * 1.1:
                trimmed.append(sentence)
                word_count += sentence_words
            else:
                break
        
        result = '. '.join(trimmed)
        if result and not result.endswith('.'):
            result += '.'
        return result
    
    return content
def similarity_score(a, b):
    """Calculate similarity between two strings"""
    return SequenceMatcher(None, a.lower(), b.lower()).ratio()

def match_generated_to_original(original_elements, generated_sections, prs):
    """
    Smart matching system that maps generated content to original PPT positions
    NOW ALLOWS BASIC INFO TO BE EDITABLE
    """
    print("started content mapping")
    content_mapping = {}
    heading_shapes = set()
    basic_info_shapes = set()
    
    for gen_section in generated_sections:
        slide_num = gen_section.get("slide_number", 1)
        is_basic_info = gen_section.get("is_basic_info", False)
        
        # Mark heading as non-editable UNLESS it's basic info
        heading_shape_idx = gen_section.get("heading_shape_idx")
        if heading_shape_idx is not None and not is_basic_info:
            heading_shapes.add(f"{slide_num}_{heading_shape_idx}")
        
        for subsec in gen_section.get("subsections", []):
            # Mark subheading as non-editable
            subheading_shape_idx = subsec.get("subheading_shape_idx")
            if subheading_shape_idx is not None:
                heading_shapes.add(f"{slide_num}_{subheading_shape_idx}")
            
            # Map generated content to content shape
            content_shape_idx = subsec.get("content_shape_idx")
            generated_content = subsec.get("generated_content", "")
            target_word_count = subsec.get("target_word_count", 0)
            target_char_count = subsec.get("target_char_count", 0)
            
            # For basic info, always replace
            if is_basic_info or subsec.get("is_basic_info", False):
                basic_info_shapes.add(f"{slide_num}_{content_shape_idx}")
            
            # Trim content to match target length
            if generated_content and target_word_count:
                generated_content = trim_content_to_length(
                    generated_content, 
                    target_word_count, 
                    target_char_count
                )
            
            if content_shape_idx is not None and generated_content:
                key = f"{slide_num}_{content_shape_idx}"
                content_mapping[key] = generated_content
            elif generated_content:
                # Fallback: find best matching content by similarity
                subheading = subsec.get("subheading", "")
                original_content = subsec.get("original_content", "")
                
                # Find shapes on this slide
                best_match_key = None
                best_score = 0
                
                for elem in original_elements:
                    if elem['slide'] == slide_num:
                        elem_key = f"{elem['slide']}_{elem['shape']}"
                        
                        # Skip if already mapped or is heading
                        if elem_key in content_mapping or elem_key in heading_shapes:
                            continue
                        
                        # Check if this shape contains similar content
                        score = similarity_score(elem['original_text'], original_content)
                        
                        if score > best_score and score > 0.3:
                            best_score = score
                            best_match_key = elem_key
                
                if best_match_key:
                    content_mapping[best_match_key] = generated_content
    
    return content_mapping, heading_shapes, basic_info_shapes

def clear_and_replace_text(shape, new_text):
    """Replace text while preserving formatting"""
    print("replacing the text")
    if not shape.has_text_frame:
        return
    
    # Store original formatting
    original_font = None
    if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
        first_run = shape.text_frame.paragraphs[0].runs[0]
        original_font = {
            'name': first_run.font.name,
            'size': first_run.font.size,
            'bold': first_run.font.bold,
            'italic': first_run.font.italic,
        }
        try:
            if first_run.font.color.type == 1:
                original_font['color'] = first_run.font.color.rgb
        except:
            pass
    
    # Clear all text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""
    
    # Replace with new text
    if shape.text_frame.paragraphs:
        first_paragraph = shape.text_frame.paragraphs[0]
        if first_paragraph.runs:
            first_run = first_paragraph.runs[0]
        else:
            first_run = first_paragraph.add_run()
        
        first_run.text = new_text
        
        # Restore formatting
        if original_font:
            if original_font.get('name'):
                first_run.font.name = original_font['name']
            if original_font.get('size'):
                first_run.font.size = original_font['size']
            if original_font.get('bold') is not None:
                first_run.font.bold = original_font['bold']
            if original_font.get('italic') is not None:
                first_run.font.italic = original_font['italic']
            if original_font.get('color'):
                try:
                    first_run.font.color.rgb = original_font['color']
                except:
                    pass
        
        # Remove extra paragraphs
        while len(shape.text_frame.paragraphs) > 1:
            p = shape.text_frame.paragraphs[-1]
            shape.text_frame._element.remove(p._element)

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import io
import json
import re

# ============= HELPER FUNCTIONS =============

def is_heading(para):
    """Check if paragraph is a heading"""
    if para.style.name.startswith('Heading') or para.style.name == 'Title':
        return True
    
    if para.runs and para.runs[0].bold and len(para.text.strip().split()) <= 10:
        return True
    
    return False

def detect_section(text):
    """Detect which section a paragraph belongs to"""
    text_lower = text.lower().strip()
    
    section_keywords = {
        'experience': ['experience', 'employment', 'work history', 'professional experience'],
        'education': ['education', 'academic', 'qualification'],
        'skills': ['skills', 'technical skills', 'competencies', 'expertise', 'tools', 'languages'],
        'projects': ['projects', 'key projects'],
        'certifications': ['certifications', 'certificates', 'licenses'],
        'achievements': ['achievements', 'accomplishments', 'awards'],
        'summary': ['summary', 'profile', 'about', 'objective', 'professional summary'],
    }
    
    for section, keywords in section_keywords.items():
        if any(keyword == text_lower or keyword in text_lower for keyword in keywords):
            return section
    
    return None

def extract_document_structure(uploaded_file):
    """Extract document structure with sections"""
    doc = Document(uploaded_file)
    structure = []
    current_section = None
    section_content_start = None
    
    for idx, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        
        # Check if it's a header
        is_header = is_heading(para)
        detected_section = detect_section(text) if is_header else None
        
        # If we found a new section
        if detected_section:
            # Save previous section if exists
            if current_section and section_content_start is not None:
                structure.append({
                    'section': current_section,
                    'header_idx': section_content_start - 1,
                    'content_start': section_content_start,
                    'content_end': idx - 1
                })
            
            current_section = detected_section
            section_content_start = idx + 1
        
        # Handle header content (name, title, contact)
        elif idx <= 3 and not current_section:
            if idx == 0:
                structure.append({'section': 'name', 'para_idx': idx})
            elif idx == 1:
                structure.append({'section': 'job_title', 'para_idx': idx})
            elif idx == 2 or idx == 3:
                structure.append({'section': 'contact', 'para_idx': idx})
    
    # Save last section
    if current_section and section_content_start is not None:
        structure.append({
            'section': current_section,
            'header_idx': section_content_start - 1,
            'content_start': section_content_start,
            'content_end': len(doc.paragraphs) - 1
        })
    
    return doc, structure

def clean_list_representation(data):
    """Clean Python list/dict string representations"""
    if isinstance(data, str):
        # Try to parse as JSON
        try:
            return json.loads(data.replace("'", '"'))
        except:
            pass
        
        # Remove Python dict/list syntax
        data = re.sub(r"^\{|\}$", "", data)  # Remove outer braces
        data = re.sub(r"^\[|\]$", "", data)  # Remove outer brackets
        data = re.sub(r"['\"]", "", data)     # Remove quotes
        
    return data

def format_skills(skills_data):
    """Format skills data properly"""
    if isinstance(skills_data, dict):
        lines = []
        for key, value in skills_data.items():
            key_title = key.replace('_', ' ').title()
            
            if isinstance(value, list):
                # Join list items with commas
                value_str = ", ".join(str(v) for v in value)
                lines.append(f"{key_title}: {value_str}")
            elif isinstance(value, str):
                # Clean string representation
                cleaned = clean_list_representation(value)
                if isinstance(cleaned, list):
                    value_str = ", ".join(str(v) for v in cleaned)
                    lines.append(f"{key_title}: {value_str}")
                else:
                    lines.append(f"{key_title}: {cleaned}")
            else:
                lines.append(f"{key_title}: {value}")
        
        return "\n\n".join(lines)
    
    elif isinstance(skills_data, list):
        return ", ".join(str(s) for s in skills_data)
    
    elif isinstance(skills_data, str):
        # Clean any Python syntax
        cleaned = clean_list_representation(skills_data)
        if isinstance(cleaned, (list, dict)):
            return format_skills(cleaned)
        return cleaned
    
    return str(skills_data)

def format_experience(exp_list):
    """Format experience data with bold subheadings"""
    result = []
    for exp in exp_list:
        job_title = exp.get('job_title', exp.get('position', ''))
        company = exp.get('company', '')
        location = exp.get('location', '')
        start_date = exp.get('start_date', '')
        end_date = exp.get('end_date', '')
        
        # Format date range
        duration = f"{start_date} - {end_date}" if start_date or end_date else ""
        
        # Get responsibilities
        responsibilities = exp.get('responsibilities', exp.get('description', []))
        if isinstance(responsibilities, str):
            # Clean string representation
            cleaned = clean_list_representation(responsibilities)
            if isinstance(cleaned, list):
                responsibilities = cleaned
            else:
                responsibilities = [cleaned]
        
        result.append({
            'job_title': job_title,
            'company': company,
            'duration': duration,
            'location': location,
            'bullets': [str(r).strip() for r in responsibilities if r]
        })
    
    return result

def format_education(edu_list):
    """Format education data"""
    result = []
    for edu in edu_list:
        degree = edu.get('degree', '')
        institution = edu.get('institution', '')
        location = edu.get('location', '')
        year = edu.get('year', '')
        start_date = edu.get('start_date', '')
        end_date = edu.get('end_date', '')
        
        # Use year or date range
        duration = year if year else (f"{start_date} - {end_date}" if start_date or end_date else "")
        
        result.append({
            'degree': degree,
            'institution': institution,
            'year': duration,
            'location': location
        })
    
    return result

def format_projects(proj_list):
    """Format projects data with bold titles"""
    result = []
    for proj in proj_list:
        title = proj.get('title', proj.get('name', ''))
        description = proj.get('description', [])
        
        if isinstance(description, str):
            cleaned = clean_list_representation(description)
            if isinstance(cleaned, list):
                description = cleaned
            else:
                description = [cleaned]
        
        result.append({
            'title': title,
            'bullets': [str(d).strip() for d in description if d]
        })
    
    return result

def format_certifications(cert_list):
    """Format certifications"""
    result = []
    for cert in cert_list:
        if isinstance(cert, str):
            result.append(cert)
        elif isinstance(cert, dict):
            name = cert.get('name', cert.get('title', ''))
            issuer = cert.get('issuer', cert.get('institution', ''))
            if name:
                line = name
                if issuer:
                    line += f" {issuer}"
                result.append(line)
    
    return result

def replace_paragraph_preserving_format(para, new_text, keep_bold=False):
    """Replace paragraph text while preserving all formatting"""
    if not para.runs:
        para.add_run(new_text)
        return
    
    # Store first run's formatting
    first_run = para.runs[0]
    font_name = first_run.font.name
    font_size = first_run.font.size
    bold = first_run.font.bold if keep_bold else False
    italic = first_run.font.italic
    
    try:
        color = first_run.font.color.rgb
    except:
        color = None
    
    # Clear all runs
    for run in para.runs:
        run.text = ""
    
    # Remove extra runs
    while len(para.runs) > 1:
        para._element.remove(para.runs[-1]._element)
    
    # Set new text and restore formatting
    para.runs[0].text = new_text
    if font_name:
        para.runs[0].font.name = font_name
    if font_size:
        para.runs[0].font.size = font_size
    if color:
        try:
            para.runs[0].font.color.rgb = color
        except:
            pass
    para.runs[0].font.bold = bold
    para.runs[0].font.italic = italic

def add_formatted_run(para, text, bold=False, italic=False, ref_run=None):
    """Add a run with specific formatting"""
    run = para.add_run(text)
    
    if ref_run:
        if ref_run.font.name:
            run.font.name = ref_run.font.name
        if ref_run.font.size:
            run.font.size = ref_run.font.size
        try:
            if ref_run.font.color.rgb:
                run.font.color.rgb = ref_run.font.color.rgb
        except:
            pass
    
    run.font.bold = bold
    run.font.italic = italic
    
    return run

def replace_experience_content(doc, start_idx, end_idx, exp_data, ref_para):
    """Replace experience content with bold job titles and italic companies"""
    # Get reference formatting
    ref_run = ref_para.runs[0] if ref_para and ref_para.runs else None
    
    # Build text with proper formatting
    text_parts = []
    
    for exp in exp_data:
        # Format: **Job Title** *Company*
        job_company = f"{exp['job_title']} {exp['company']}".strip()
        
        text_parts.append(job_company)
        
        # Add bullets
        for bullet in exp['bullets']:
            text_parts.append(f"• {bullet}")
        
        text_parts.append("")  # Empty line between experiences
    
    # Join and replace
    full_text = "\n".join(text_parts).strip()
    
    if start_idx < len(doc.paragraphs):
        para = doc.paragraphs[start_idx]
        
        # Clear existing content
        for run in para.runs:
            run.text = ""
        while len(para.runs) > 1:
            para._element.remove(para.runs[-1]._element)
        
        # Add content with mixed formatting
        current_text = ""
        for exp in exp_data:
            # Add job title (bold) and company (italic)
            if exp['job_title']:
                job_run = para.add_run(exp['job_title'] + " ")
                job_run.font.bold = True
                if ref_run:
                    if ref_run.font.name:
                        job_run.font.name = ref_run.font.name
                    if ref_run.font.size:
                        job_run.font.size = ref_run.font.size
            
            if exp['company']:
                company_run = para.add_run(exp['company'])
                company_run.font.italic = True
                if ref_run:
                    if ref_run.font.name:
                        company_run.font.name = ref_run.font.name
                    if ref_run.font.size:
                        company_run.font.size = ref_run.font.size
            
            # Add line break
            para.add_run("\n")
            
            # Add bullets
            for bullet in exp['bullets']:
                bullet_run = para.add_run(f"• {bullet}\n")
                if ref_run:
                    if ref_run.font.name:
                        bullet_run.font.name = ref_run.font.name
                    if ref_run.font.size:
                        bullet_run.font.size = ref_run.font.size
            
            # Add spacing between experiences
            para.add_run("\n")
    
    return list(range(start_idx + 1, end_idx + 1))

def replace_projects_content(doc, start_idx, end_idx, proj_data, ref_para):
    """Replace projects content with bold titles"""
    ref_run = ref_para.runs[0] if ref_para and ref_para.runs else None
    
    if start_idx < len(doc.paragraphs):
        para = doc.paragraphs[start_idx]
        
        # Clear existing content
        for run in para.runs:
            run.text = ""
        while len(para.runs) > 1:
            para._element.remove(para.runs[-1]._element)
        
        # Add content with bold titles
        for proj in proj_data:
            # Add project title (bold)
            if proj['title']:
                title_run = para.add_run(proj['title'])
                title_run.font.bold = True
                if ref_run:
                    if ref_run.font.name:
                        title_run.font.name = ref_run.font.name
                    if ref_run.font.size:
                        title_run.font.size = ref_run.font.size
                
                para.add_run("\n")
            
            # Add bullets
            for bullet in proj['bullets']:
                bullet_run = para.add_run(f"• {bullet}\n")
                if ref_run:
                    if ref_run.font.name:
                        bullet_run.font.name = ref_run.font.name
                    if ref_run.font.size:
                        bullet_run.font.size = ref_run.font.size
            
            # Add spacing
            para.add_run("\n")
    
    return list(range(start_idx + 1, end_idx + 1))

def replace_section_content(doc, start_idx, end_idx, formatted_data, ref_para, section_type):
    """Replace section content with new data"""
    # Handle special formatting for experience and projects
    if section_type == 'experience' and isinstance(formatted_data, list):
        return replace_experience_content(doc, start_idx, end_idx, formatted_data, ref_para)
    
    if section_type == 'projects' and isinstance(formatted_data, list):
        return replace_projects_content(doc, start_idx, end_idx, formatted_data, ref_para)
    
    # Regular content replacement
    ref_run = ref_para.runs[0] if ref_para and ref_para.runs else None
    
    # Build text content
    text_lines = []
    
    if isinstance(formatted_data, list):
        if formatted_data and isinstance(formatted_data[0], dict):
            # Education or other structured data
            for item in formatted_data:
                parts = []
                if item.get('degree'):
                    parts.append(item['degree'])
                if item.get('institution'):
                    parts.append(item['institution'])
                
                header = " ".join(parts)
                if header:
                    text_lines.append(header)
                
                if item.get('year'):
                    text_lines.append(item['year'])
                
                text_lines.append("")
        else:
            # Simple list (certifications)
            for item in formatted_data:
                text_lines.append(f"• {item}")
    
    elif isinstance(formatted_data, str):
        text_lines.append(formatted_data)
    
    # Join text
    full_text = "\n".join(text_lines).strip()
    
    # Replace content
    if start_idx < len(doc.paragraphs):
        para = doc.paragraphs[start_idx]
        replace_paragraph_preserving_format(para, full_text, keep_bold=False)
    
    return list(range(start_idx + 1, end_idx + 1))

def replace_content(doc, structure, final_data):
    """Main function to replace content"""
    paragraphs_to_remove = set()
    replaced_count = 0
    sections_to_remove = set()
    
    # First pass: identify sections not in final_data
    for item in structure:
        section = item['section']
        if section not in ['name', 'job_title', 'contact']:
            if section not in final_data or not final_data[section]:
                # Mark entire section for removal
                if 'header_idx' in item:
                    sections_to_remove.add(item['header_idx'])
                if 'content_start' in item:
                    for idx in range(item['content_start'], item['content_end'] + 1):
                        sections_to_remove.add(idx)
    
    # Second pass: replace content
    for item in structure:
        section = item['section']
        
        # Skip sections marked for removal
        if 'header_idx' in item and item['header_idx'] in sections_to_remove:
            continue
        
        # Handle header sections
        if section == 'name':
            if 'para_idx' in item and 'name' in final_data:
                para_idx = item['para_idx']
                if para_idx < len(doc.paragraphs):
                    replace_paragraph_preserving_format(doc.paragraphs[para_idx], final_data['name'], keep_bold=True)
                    replaced_count += 1
        
        elif section == 'job_title':
            if 'para_idx' in item and 'job_title' in final_data:
                para_idx = item['para_idx']
                if para_idx < len(doc.paragraphs):
                    replace_paragraph_preserving_format(doc.paragraphs[para_idx], final_data['job_title'], keep_bold=False)
                    replaced_count += 1
        
        elif section == 'contact':
            if 'para_idx' in item:
                contact_parts = []
                if final_data.get('phone'):
                    contact_parts.append(final_data['phone'])
                if final_data.get('email'):
                    contact_parts.append(final_data['email'])
                if final_data.get('location'):
                    contact_parts.append(final_data['location'])
                
                if contact_parts:
                    para_idx = item['para_idx']
                    if para_idx < len(doc.paragraphs):
                        replace_paragraph_preserving_format(doc.paragraphs[para_idx], " | ".join(contact_parts), keep_bold=False)
                        replaced_count += 1
        
        # Handle content sections
        else:
            if 'content_start' in item and section in final_data and final_data[section]:
                content_start = item['content_start']
                content_end = item['content_end']
                ref_para = doc.paragraphs[content_start] if content_start < len(doc.paragraphs) else None
                
                # Format data based on section
                formatted_data = None
                
                if section == 'experience':
                    formatted_data = format_experience(final_data['experience'])
                elif section == 'education':
                    formatted_data = format_education(final_data['education'])
                elif section == 'skills':
                    formatted_data = format_skills(final_data['skills'])
                elif section == 'projects':
                    formatted_data = format_projects(final_data['projects'])
                elif section == 'certifications':
                    formatted_data = format_certifications(final_data['certifications'])
                elif section == 'summary':
                    formatted_data = final_data['summary']
                
                if formatted_data:
                    to_remove = replace_section_content(doc, content_start, content_end, formatted_data, ref_para, section)
                    paragraphs_to_remove.update(to_remove)
                    replaced_count += 1
    
    # Combine all paragraphs to remove
    paragraphs_to_remove.update(sections_to_remove)
    
    # Remove marked paragraphs
    for idx in sorted(paragraphs_to_remove, reverse=True):
        if idx < len(doc.paragraphs):
            p = doc.paragraphs[idx]._element
            p.getparent().remove(p)
    
    # Save document
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    
    return output, replaced_count, len(paragraphs_to_remove)



def load_user_doc_templates(username):
    """Load user's saved document templates from database"""
    templates = {}

    try:
        conn = get_connection()
        cursor = conn.cursor()

        cursor.execute("""
            SELECT template_key,
                   template_name,
                   doc_data,
                   doc_text,
                   uploaded_at,
                   original_filename
            FROM user_doc_templates
            WHERE email = ?
            ORDER BY uploaded_at DESC
        """, username)

        rows = cursor.fetchall()

        for row in rows:
            template_key = row[0]
            template_name = row[1]
            doc_data = row[2]
            doc_text = row[3]
            uploaded_at = row[4]
            original_filename = row[5]

            # Decide which content to return
            if doc_data is not None:
                doc_content = bytes(doc_data)   # VARBINARY → bytes
            else:
                doc_content = doc_text           # NVARCHAR text

            templates[template_key] = {
                "name": template_name,
                "doc_data": doc_content,
                "uploaded_at": uploaded_at,
                "original_filename": original_filename
            }

    except Exception as e:
        print(f"Error loading doc templates: {e}")

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass

    return templates


import psycopg2


def save_user_doc_templates(username, templates):
    """Save user's document templates to database"""
    try:
        conn = get_connection()
        cursor = conn.cursor()

        for template_id, template_data in templates.items():
            doc_data = template_data.get("doc_data")

            # Decide storage (SQL Server)
            if isinstance(doc_data, (bytes, bytearray)):
                doc_binary = doc_data          # VARBINARY(MAX)
                doc_text = None
            else:
                doc_binary = None
                doc_text = doc_data            # NVARCHAR(MAX)

            cursor.execute("""
                MERGE user_doc_templates AS target
                USING (
                    SELECT
                        ? AS email,
                        ? AS template_key,
                        ? AS template_name,
                        ? AS doc_data,
                        ? AS doc_text,
                        ? AS uploaded_at,
                        ? AS original_filename
                ) AS source
                ON target.email = source.email
                   AND target.template_key = source.template_key
                WHEN MATCHED THEN
                    UPDATE SET
                        template_name = source.template_name,
                        doc_data = source.doc_data,
                        doc_text = source.doc_text,
                        uploaded_at = source.uploaded_at,
                        original_filename = source.original_filename,
                        updated_at = SYSDATETIME()
                WHEN NOT MATCHED THEN
                    INSERT (
                        email,
                        template_key,
                        template_name,
                        doc_data,
                        doc_text,
                        uploaded_at,
                        original_filename,
                        created_at,
                        updated_at
                    )
                    VALUES (
                        source.email,
                        source.template_key,
                        source.template_name,
                        source.doc_data,
                        source.doc_text,
                        source.uploaded_at,
                        source.original_filename,
                        SYSDATETIME(),
                        SYSDATETIME()
                    );
            """,
            username,
            template_id,
            template_data.get("name"),
            doc_binary,
            doc_text,
            template_data.get("uploaded_at"),
            template_data.get("original_filename")
            )

        conn.commit()
        return True

    except Exception as e:
        print(f"❌ Error saving doc templates: {e}")
        return False

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass




def delete_user_doc_template(username, template_id):
    try:
        conn = get_connection()
        cursor = conn.cursor()

        cursor.execute("""
            DELETE FROM user_doc_templates
            WHERE email = ? AND template_key = ?
        """, username, template_id)

        conn.commit()
        return True

    except Exception as e:
        print(f"Error deleting doc template: {e}")
        return False

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass


# ============= PPT TEMPLATE STORAGE FUNCTIONS =============



def load_user_ppt_templates(username):
    """Load user's saved PPT templates from database"""
    templates = {}

    try:
        conn = get_connection()
        cursor = conn.cursor()

        cursor.execute("""
            SELECT template_key,
                   template_name,
                   ppt_data,
                   heading_shapes,
                   basic_info_shapes,
                   uploaded_at,
                   original_filename
            FROM user_ppt_templates
            WHERE email = ?
            ORDER BY uploaded_at DESC
        """, username)

        rows = cursor.fetchall()

        for row in rows:
            template_key = row[0]
            template_name = row[1]
            ppt_data = row[2]
            heading_shapes_json = row[3]
            basic_info_shapes_json = row[4]
            uploaded_at = row[5]
            original_filename = row[6]

            templates[template_key] = {
                "name": template_name,
                "ppt_data": bytes(ppt_data),  # VARBINARY → bytes
                "heading_shapes": set(json.loads(heading_shapes_json)) if heading_shapes_json else set(),
                "basic_info_shapes": set(json.loads(basic_info_shapes_json)) if basic_info_shapes_json else set(),
                "uploaded_at": uploaded_at,
                "original_filename": original_filename
            }

    except Exception as e:
        print(f"Error loading PPT templates: {e}")

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass

    return templates


import psycopg2



def save_user_ppt_templates(username, templates):
    """Save user's PPT templates to database"""
    try:
        conn = get_connection()
        cursor = conn.cursor()

        for template_id, template_data in templates.items():
            ppt_bytes = template_data.get("ppt_data")

            # Convert sets/lists → JSON strings for SQL Server
            heading_shapes_json = json.dumps(
                list(template_data.get("heading_shapes", []))
            )
            basic_info_shapes_json = json.dumps(
                list(template_data.get("basic_info_shapes", []))
            )

            cursor.execute("""
                MERGE user_ppt_templates AS target
                USING (
                    SELECT
                        ? AS email,
                        ? AS template_key,
                        ? AS template_name,
                        ? AS ppt_data,
                        ? AS heading_shapes,
                        ? AS basic_info_shapes,
                        ? AS uploaded_at,
                        ? AS original_filename
                ) AS source
                ON target.email = source.email
                   AND target.template_key = source.template_key
                WHEN MATCHED THEN
                    UPDATE SET
                        template_name = source.template_name,
                        ppt_data = source.ppt_data,
                        heading_shapes = source.heading_shapes,
                        basic_info_shapes = source.basic_info_shapes,
                        uploaded_at = source.uploaded_at,
                        original_filename = source.original_filename,
                        updated_at = SYSDATETIME()
                WHEN NOT MATCHED THEN
                    INSERT (
                        email,
                        template_key,
                        template_name,
                        ppt_data,
                        heading_shapes,
                        basic_info_shapes,
                        uploaded_at,
                        original_filename,
                        created_at,
                        updated_at
                    )
                    VALUES (
                        source.email,
                        source.template_key,
                        source.template_name,
                        source.ppt_data,
                        source.heading_shapes,
                        source.basic_info_shapes,
                        source.uploaded_at,
                        source.original_filename,
                        SYSDATETIME(),
                        SYSDATETIME()
                    );
            """,
            username,
            template_id,
            template_data.get("name"),
            ppt_bytes,                      # VARBINARY(MAX)
            heading_shapes_json,            # JSON string
            basic_info_shapes_json,         # JSON string
            template_data.get("uploaded_at"),
            template_data.get("original_filename")
            )

        conn.commit()
        return True

    except Exception as e:
        print(f"Error saving PPT templates: {e}")
        return False

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass






def delete_user_ppt_template(username, template_id):
    try:
        conn = get_connection()
        cursor = conn.cursor()

        cursor.execute("""
            DELETE FROM user_ppt_templates
            WHERE email = ? AND template_key = ?
        """, username, template_id)

        conn.commit()
        return True

    except Exception as e:
        st.error(f"Error deleting PPT template: {e}")
        return False

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass



def get_resume_hash(resume_data):
    """Generate a hash of resume data to detect changes"""
    resume_str = json.dumps(resume_data, sort_keys=True, default=str)
    return hashlib.md5(resume_str.encode()).hexdigest()

def should_regenerate_resume():
    """Check if we need to regenerate the enhanced resume"""
    current_user = st.session_state.get('logged_in_user')
    resume_data = st.session_state.get('resume_source')
    jd_data = st.session_state.get('job_description')
    
    if 'enhanced_resume' not in st.session_state:
        return True
    
   
    if st.session_state.get('last_resume_user') != current_user:
        return True
    
    
    current_resume_hash = get_resume_hash(resume_data) if resume_data else None
    if st.session_state.get('last_resume_hash') != current_resume_hash:
        return True

    current_jd_hash = get_resume_hash(jd_data) if jd_data else None
    if st.session_state.get('last_jd_hash') != current_jd_hash:
        return True
    
    return False

def generate_enhanced_resume(resume_data=None, jd_data=None):
    """Generate enhanced resume and store metadata"""
    # resume_data = st.session_state.get('resume_source')
    # jd_data = st.session_state.get('job_description')
    current_user = st.session_state.get('logged_in_user') or st.query_params.get('user', '')
    # st.write(resume_data)
    # st.write(jd_data)
    # ============================================
    # 🔧 FIX: Ensure resume_data is loaded
    # ============================================
    # If resume_data is None, try multiple sources
    if resume_data is None:
        # Try final_resume_data first
        resume_data = st.session_state.get('final_resume_data')
        
        # Try enhanced_resume next
        if resume_data is None:
            resume_data = st.session_state.get('enhanced_resume')
        
        # Finally, try to fetch from database
        if resume_data is None and current_user:
            try:
                users = load_users()
                user_entry = users.get(current_user)
                
                if isinstance(user_entry, dict):
                    user_resume = get_user_resume(current_user)
                    
                    if user_resume and len(user_resume) > 0:
                        resume_data = user_resume
                        st.session_state['resume_source'] = user_resume
                        st.session_state['final_resume_data'] = user_resume  # Also set this
                        st.session_state['input_method'] = user_resume.get("input_method", "Manual Entry")
                        st.session_state['username'] = current_user.split('@')[0]
                    else:
                        st.error(f"No resume data found for user: {current_user}. Please create a resume first.")
                        if st.button("Go to Home"):
                            st.query_params["home"] = "true"
                            st.query_params["user"] = current_user
                            st.rerun()
                        st.stop()
                else:
                    st.error(f"User not found: {current_user}")
                    st.stop()
            except Exception as e:
                st.error(f"Error fetching resume data: {str(e)}")
                import traceback
                st.error(traceback.format_exc())
                st.stop()
        elif resume_data is None:
            st.error("No resume data found and no user logged in. Please go back to the main page.")
            if st.button("Go to Login"):
                st.switch_page("pages/main.py")
            st.stop()
    
    # ============================================
    # CRITICAL: Verify resume_data is not None before proceeding
    # ============================================
    if resume_data is None:
        st.error("❌ Failed to load resume data from any source!")
        st.write("Debug info:")
        st.write(f"- resume_source: {st.session_state.get('resume_source')}")
        st.write(f"- final_resume_data: {st.session_state.get('final_resume_data')}")
        st.write(f"- enhanced_resume: {st.session_state.get('enhanced_resume')}")
        st.write(f"- current_user: {current_user}")
        st.stop()
    
    # Verify it's a dictionary
    if not isinstance(resume_data, dict):
        st.error(f"❌ Resume data is not a dictionary! Type: {type(resume_data)}")
        st.write("Data:", resume_data)
        st.stop()
    
    # If jd_data is None, use empty dict
    if jd_data is None:
        st.warning("⚠️ No job description found. Using default optimization.")
        jd_data = {}
    
    # Safely get input_method with fallback
    input_method = st.session_state.get(
        "input_method", 
        resume_data.get("input_method", "Manual Entry")
    )
   
    # ============================================
    # Call the appropriate rewrite function
    # ============================================
    try:
        if input_method == "Manual Entry":
            enhanced_resume = rewrite_resume_for_job_manual(resume_data, jd_data)
        else:
            enhanced_resume = rewrite_resume_for_job(resume_data, jd_data)
            # st.write("inside upload entry")
            # st.write(enhanced_resume)
        
        # ============================================
        # 🔧 FIX: Flatten nested structure if present
        # ============================================
        if enhanced_resume and 'resume_data' in enhanced_resume:
            original_data = enhanced_resume.pop('resume_data')

            # If resume_data is JSON string → parse it
            if isinstance(original_data, str):
                try:
                    original_data = json.loads(original_data)
                except Exception:
                    original_data = {}

            # Only iterate if it's a dict
            if isinstance(original_data, dict):
                for key, value in original_data.items():
                    if key not in enhanced_resume or not enhanced_resume[key]:
                        enhanced_resume[key] = value

        
        # Ensure critical fields exist
        if not enhanced_resume.get('name'):
            enhanced_resume['name'] = resume_data.get('name', 'Your Name')
        if not enhanced_resume.get('email'):
            enhanced_resume['email'] = resume_data.get('email', '')
        if not enhanced_resume.get('phone'):
            enhanced_resume['phone'] = resume_data.get('phone', '')
        if not enhanced_resume.get('location'):
            enhanced_resume['location'] = resume_data.get('location', '')
        
        # Store in session state
        st.session_state['enhanced_resume'] = enhanced_resume
        st.session_state['last_resume_user'] = current_user
        st.session_state['last_resume_hash'] = get_resume_hash(resume_data)
        st.session_state['last_jd_hash'] = get_resume_hash(jd_data)
        
        # Store in query params
        try:
            st.query_params["enhanced_resume"] = json.dumps(enhanced_resume)
            st.query_params["last_resume_hash"] = st.session_state['last_resume_hash'] or ""
            st.query_params["last_jd_hash"] = st.session_state['last_jd_hash'] or ""
            st.query_params["last_resume_user"] = current_user
        except:
            pass  # Query params can fail if data is too large
        
        chatbot(enhanced_resume)
        
        return enhanced_resume
        
    except Exception as e:
        st.error(f"❌ Error generating enhanced resume: {str(e)}")
        import traceback
        st.error(traceback.format_exc())
        st.stop()


def format_section_title(key):
    """Converts keys like 'certifications' to 'Certifications'."""
    title = key.replace('_', ' ')
    return ' '.join(word.capitalize() for word in title.split())


def render_basic_details(data, is_edit):
    """Top header section (Name, title, contact info)."""
    if is_edit:
        st.markdown('<h2>Basic Details</h2>', unsafe_allow_html=True)
        data['name'] = st.text_input("Name", data.get('name', ''), key="edit_name")
        data['job_title'] = st.text_input("Job Title", data.get('job_title', ''), key="edit_job_title")

        col1, col2, col3,col4 = st.columns(4)
        with col1:
            data['phone'] = st.text_input("Phone", data.get('phone', ''), key="edit_phone")
        with col2:
            data['email'] = st.text_input("Email", data.get('email', ''), key="edit_email")
        with col3:
            data['location'] = st.text_input("Location", data.get('location', ''), key="edit_location")
        with col4:
            data['url'] = st.text_input("url", data.get('url', ''), key="edit_url")
 

        st.markdown('<div class="resume-section">', unsafe_allow_html=True)
        st.markdown('<h2>Summary</h2>', unsafe_allow_html=True)
        data['summary'] = st.text_area("Summary", data.get('summary', ''), height=150, key="edit_summary")
        st.markdown('</div>', unsafe_allow_html=True)
    else:
        st.markdown(f"<h1>{data.get('name', 'Name Not Found')}</h1>", unsafe_allow_html=True)
        if data.get('job_title'):
            st.markdown(f"<h3>{data['job_title']}</h3>", unsafe_allow_html=True)

        contact_html = f"""
        <div class="contact-info">
            {data.get('phone', '')} | {data.get('email', '')} | {data.get('location', '')}
        </div>
        """
        st.markdown(contact_html, unsafe_allow_html=True)

        if data.get('summary'):
            st.markdown('<div class="resume-section">', unsafe_allow_html=True)
            st.markdown('<h2>Summary</h2>', unsafe_allow_html=True)
            st.markdown(f"<p style='color:#000000;'>{data['summary']}</p>", unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)


def format_date(date_string):
    """
    Converts date from YYYY-MM-DD to MMM YYYY format.
    Example: "2025-10-14" -> "Oct 2025"
    """
    if not date_string or date_string.strip() == '':
        return ''
    
    try:
        from datetime import datetime
        date_obj = datetime.strptime(date_string, '%Y-%m-%d')
        return date_obj.strftime('%b %Y')
    except:
        return date_string  # Return original if parsing fails


def render_list_item(item, index, key_prefix, section_title, is_edit=True):
    """Generic list item renderer for both edit and view modes with dynamic field handling."""
    
    # Normalize item to dict
    if isinstance(item, str):
        item = {"title": item}
    
    # Define field priority order for display
    title_keys = ['position', 'title', 'name', 'degree', 'institution', 'company', 'certificate_name']
    subtitle_keys = ['company', 'institution', 'issuer', 'organization', 'provider_name']
    detail_keys_to_skip = ['position', 'title', 'name', 'degree', 'company', 'institution', 
                           'description', 'overview', 'issuer', 'start_date', 'end_date', 
                           'duration', 'certificate_name', 'organization', 'provider_name',
                           'details', 'achievement']

    if not is_edit:
        html_content = "<div>"
        
        # Get main title (first available from title_keys)
        main_title = None
        for key in title_keys:
            if key in item and item[key]:
                main_title = item[key]
                break
        
        if main_title:
            html_content += f'<div class="item-title">{main_title}</div>'
            
            # Get subtitle (first available from subtitle_keys, but not same as title)
            subtitle = None
            for key in subtitle_keys:
                if key in item and item[key] and item[key] != main_title:
                    subtitle = item[key]
                    break
            
            if subtitle:
                html_content += f'<div class="item-subtitle">{subtitle}</div>'

        # Format dates to MMM YYYY
        duration = item.get('duration')
        if not duration:
            start_date = format_date(item.get('start_date', ''))
            end_date = format_date(item.get('end_date', ''))
            if start_date or end_date:
                duration = f"{start_date} - {end_date}"
        
        if duration and duration.strip() != '-':
            html_content += f'<div class="item-details"><em>{duration}</em></div>'

        # Handle description/overview/details/achievement fields
        description_fields = ['description', 'overview', 'details', 'achievement']
        main_description_list = None
        for field in description_fields:
            if field in item:
                main_description_list = item[field]
                break
        
        if main_description_list:
            if isinstance(main_description_list, str):
                main_description_list = [main_description_list]
            if isinstance(main_description_list, list) and main_description_list:
                bullet_html = "".join([f"<li>{line}</li>" for line in main_description_list if line])
                html_content += f'<ul class="bullet-list">{bullet_html}</ul>'
        
        # Display any remaining fields not in skip list
        for k, v in item.items():
            if isinstance(v, str) and v.strip() and k not in detail_keys_to_skip + ['duration']:
                formatted_k = format_section_title(k)
                html_content += f'<div class="item-details">**{formatted_k}:** {v}</div>'

        html_content += "</div>"
        return html_content
    
    else:
        # EDIT MODE
        edited_item = item.copy()
        
        # Get all fields in item
        edit_fields = list(item.keys())
        
        # Priority fields for ordering
        priority_fields = ['position', 'title', 'name', 'company', 'institution', 'degree', 
                          'certificate_name', 'issuer', 'organization', 'provider_name',
                          'duration', 'start_date', 'end_date', 'description', 'overview', 
                          'details', 'achievement']
        
        # Order fields: priority first, then rest
        ordered_fields = []
        for field in priority_fields:
            if field in edit_fields:
                ordered_fields.append(field)
                edit_fields.remove(field)
        ordered_fields.extend(edit_fields)
        
        # Render input fields for each field in order
        for k in ordered_fields:
            v = item[k]
            if isinstance(v, str):
                edited_item[k] = st.text_input(
                    format_section_title(k), 
                    v, 
                    key=f"{key_prefix}_{k}_{index}"
                )
            elif isinstance(v, list):
                text = "\n".join(v)
                edited_text = st.text_area(
                    format_section_title(k), 
                    text, 
                    height=150, 
                    key=f"{key_prefix}_area_{k}_{index}"
                )
                edited_item[k] = [line.strip() for line in edited_text.split('\n') if line.strip()]
        
        return edited_item
    
def render_generic_section(section_key, data_list, is_edit):
    """Renders dynamic list sections with consistent formatting."""
    section_title = format_section_title(section_key)
    if not data_list: 
        return
    
    # Convert data_list to a list if it's not already
    if isinstance(data_list, str):
        data_list = [data_list]
        # Update the session state with the converted list
        if 'enhanced_resume' in st.session_state:
            st.session_state['enhanced_resume'][section_key] = data_list
    
    st.markdown('<div class="resume-section">', unsafe_allow_html=True)
    st.markdown(f'<h2>{section_title}</h2>', unsafe_allow_html=True)
    
    for i, item in enumerate(data_list):
        # Normalize all items to dictionaries
        if not isinstance(item, dict):
            if isinstance(item, str):
                item = {"title": item}
            else:
                item = {"title": str(item)}
            data_list[i] = item  # Now safe because data_list is guaranteed to be a list
        
        with st.container(border=False):
            # Build expander title from available fields
            expander_title_parts = [
                item.get('position'),
                item.get('title'),
                item.get('name'),
                item.get('certificate_name'),
                item.get('company'),
                item.get('institution'),
                item.get('issuer')
            ]
            expander_title = next((t for t in expander_title_parts if t), f"{section_title[:-1]} Item {i+1}")
            
            if is_edit:
                with st.expander(f"📝 Edit: **{expander_title}**", expanded=False):
                    temp_item = deepcopy(item) 
                    edited_item = render_list_item(temp_item, i, f"{section_key}_edit_{i}", section_title, is_edit=True)
                    
                    if edited_item:
                        st.session_state['enhanced_resume'][section_key][i] = edited_item
                    
                    if st.button(f"❌ Remove this {section_title[:-1]}", key=f"{section_key}_remove_{i}"):
                        st.session_state['enhanced_resume'][section_key].pop(i)
                        st.rerun()
                
                # Show view mode preview below edit section
                st.markdown(render_list_item(item, i, f"{section_key}_view_{i}", section_title, is_edit=False), unsafe_allow_html=True)
            else:
                # View mode only
                st.markdown(render_list_item(item, i, f"{section_key}_view_{i}", section_title, is_edit=False), unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)

def render_skills_section(data, is_edit):
    """Handles the nested dictionary structure of the 'skills' section."""
    skills_data = data.get('skills', {})
    if not skills_data: return

    st.markdown('<div class="resume-section">', unsafe_allow_html=True)
    st.markdown('<h2>Skills</h2>', unsafe_allow_html=True)

    # Handle both dict and list formats
    if isinstance(skills_data, dict):
        # Dictionary format: {'technical': ['Python', 'Java'], 'soft': ['Leadership']}
        if is_edit:
            with st.expander("📝 Edit Skills (Separate by Line)", expanded=False):
                for skill_type, skill_list in skills_data.items():
                    st.subheader(format_section_title(skill_type))
                    
                    # Ensure skill_list is actually a list
                    if isinstance(skill_list, list):
                        skill_text = "\n".join(skill_list)
                    else:
                        skill_text = str(skill_list)
                    
                    edited_text = st.text_area(f"Edit {skill_type}", skill_text, height=100, key=f"skills_edit_{skill_type}")
                    
                    st.session_state['enhanced_resume']['skills'][skill_type] = [line.strip() for line in edited_text.split('\n') if line.strip()]
        
        for skill_type, skill_list in skills_data.items():
            if skill_list:
                st.markdown(f"**{format_section_title(skill_type)}:**", unsafe_allow_html=True)
                
                # Ensure skill_list is actually a list
                if isinstance(skill_list, list):
                    skills_html = "".join([f'<li class="skill-item">{s}</li>' for s in skill_list])
                else:
                    skills_html = f'<li class="skill-item">{skill_list}</li>'
                
                st.markdown(f'<ul class="skill-list">{skills_html}</ul>', unsafe_allow_html=True)
    
    elif isinstance(skills_data, list):
        # List format: ['Python', 'Java', 'Leadership']
        if is_edit:
            with st.expander("📝 Edit Skills (Separate by Line)", expanded=False):
                skill_text = "\n".join(skills_data)
                edited_text = st.text_area("Edit Skills", skill_text, height=150, key="skills_edit_all")
                st.session_state['enhanced_resume']['skills'] = [line.strip() for line in edited_text.split('\n') if line.strip()]
        
        # Display as a single list
        if skills_data:
            skills_html = "".join([f'<li class="skill-item">{s}</li>' for s in skills_data])
            st.markdown(f'<ul class="skill-list">{skills_html}</ul>', unsafe_allow_html=True)

    st.markdown('</div>', unsafe_allow_html=True)

def add_new_item(section_key, default_item):
    """Generic function to add a new item to any list section."""
    if section_key not in st.session_state['enhanced_resume']:
        st.session_state['enhanced_resume'][section_key] = []
            
    st.session_state['enhanced_resume'][section_key].append(default_item)
    st.rerun()

# MODIFIED: Save and Improve with hash update
def save_and_improve():
    """Calls auto_improve_resume and updates session state."""
    data = deepcopy(st.session_state['enhanced_resume'])
    user_skills_before = deepcopy(data.get('skills', {}))
    job_description = st.session_state.get('job_description', '') 

    improved_data = analyze_and_improve_resume(data, job_description)
    
    # Skills merging logic - Handle both list and dict formats
    llm_skills_after = improved_data.get('skills', {})
    
    # Check if skills are in list format or dict format
    user_is_list = isinstance(user_skills_before, list)
    llm_is_list = isinstance(llm_skills_after, list)
    
    if user_is_list and llm_is_list:
        # Both are lists - simple merge
        user_set = set(user_skills_before)
        llm_set = set(llm_skills_after)
        
        # Add new skills from LLM that user doesn't have
        for skill in llm_set:
            if skill not in user_set:
                user_set.add(skill)
        
        improved_data['skills'] = sorted(list(user_set))
        
    elif user_is_list and not llm_is_list:
        # User has list, LLM returned dict - convert user list to dict format
        merged_skills = {}
        
        # First, add all LLM categorized skills
        for category, llm_list in llm_skills_after.items():
            merged_skills[category] = list(set(llm_list))
        
        # Add uncategorized user skills to a "General" category
        if user_skills_before:
            if "General" not in merged_skills:
                merged_skills["General"] = []
            
            for skill in user_skills_before:
                # Check if skill exists in any category
                skill_exists = False
                for cat_skills in merged_skills.values():
                    if skill in cat_skills:
                        skill_exists = True
                        break
                
                if not skill_exists:
                    merged_skills["General"].append(skill)
            
            # Remove empty General category
            if not merged_skills.get("General"):
                merged_skills.pop("General", None)
        
        improved_data['skills'] = merged_skills
        
    elif not user_is_list and llm_is_list:
        # User has dict, LLM returned list - add list items to existing categories
        merged_skills = deepcopy(user_skills_before)
        
        # Try to categorize new skills or add to "General"
        for skill in llm_skills_after:
            skill_added = False
            
            # Check if skill already exists in any category
            for category, skills_list in merged_skills.items():
                if skill in skills_list:
                    skill_added = True
                    break
            
            # If not found, add to General category
            if not skill_added:
                if "General" not in merged_skills:
                    merged_skills["General"] = []
                merged_skills["General"].append(skill)
        
        improved_data['skills'] = merged_skills
        
    else:
        # Both are dicts - original logic
        merged_skills = {}
        all_categories = set(user_skills_before.keys()) | set(llm_skills_after.keys())

        for category in all_categories:
            user_list = user_skills_before.get(category, [])
            llm_list = llm_skills_after.get(category, [])
            
            user_set = set(user_list)
            llm_set = set(llm_list)
            final_skills_set = user_set.copy()
            
            # Add new skills from LLM
            for skill in llm_set:
                if skill not in final_skills_set:
                    final_skills_set.add(skill)
                    
            merged_skills[category] = sorted(list(final_skills_set))

        improved_data['skills'] = merged_skills
    
    # Update session state
    st.session_state['enhanced_resume'] = improved_data
    
    if 'ats_score_data' in st.session_state:
        del st.session_state['ats_score_data']
    
    # Update hash so it doesn't regenerate
    st.session_state['last_resume_hash'] = get_resume_hash(st.session_state.get('resume_source'))
    
    st.success("✅ Resume content saved and improved! Check the updated details below.")
    st.rerun()



def image_to_base64_local(image):
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    img_str = base64.b64encode(buffer.getvalue()).decode()
    return img_str



def load_users():
    users = {}

    try:
        conn = get_connection()
        cursor = conn.cursor()

        cursor.execute(
            "SELECT email, password_hash, name FROM users"
        )

        for row in cursor.fetchall():
            email, password_hash, name = row
            users[email] = {
                "password": password_hash,
                "name": name
            }

    except Exception as e:
        st.error(f"Failed to load users: {e}")

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass

    return users


    


def save_users(users):
    try:
        conn = get_connection()
        cursor = conn.cursor()

        for email, data in users.items():
            cursor.execute("""
                MERGE users AS target
                USING (SELECT ? AS email, ? AS password_hash, ? AS name) AS source
                ON target.email = source.email
                WHEN MATCHED THEN
                    UPDATE SET
                        password_hash = source.password_hash,
                        name = source.name
                WHEN NOT MATCHED THEN
                    INSERT (email, password_hash, name)
                    VALUES (source.email, source.password_hash, source.name);
            """, email, data["password"], data["name"])

        conn.commit()

    except Exception as e:
        st.error(f"Failed to save users: {e}")

    finally:
        try:
            cursor.close()
            conn.close()
        except:
            pass




def load_user_resume_data(email):
    conn = get_connection()
    cur = conn.cursor()

    cur.execute("""
        SELECT resume_data, input_method
        FROM user_resumes
        WHERE email = ?
    """, email)

    row = cur.fetchone()

    cur.close()
    conn.close()

    if row:
        return {
            "resume_data": row[0],
            "input_method": row[1]
        }

    return None



def get_user_resume(email):
    """Get resume data for a specific user"""

    resume_data = load_user_resume_data(email)

    if resume_data and isinstance(resume_data, dict):
        return resume_data

    return None


def is_valid_email(email):
    pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
    return re.match(pattern, email) is not None




def show_login_modal():
    import streamlit as st
    import time 
    from PIL import Image
    from streamlit_extras.stylable_container import stylable_container
    from utils import load_users, save_users, is_valid_email, get_user_resume

    # Clear cache
    st.cache_data.clear()
    st.cache_resource.clear()
    
    # Session states
    if 'mode' not in st.session_state:
        st.session_state.mode = 'login'

    if 'page_transitioning' not in st.session_state:
        st.session_state.page_transitioning = False

    # -------------------------------------------------------------------------------------------------
    # CSS STYLES - Modern Orange & White Theme
    # -------------------------------------------------------------------------------------------------
    st.markdown("""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Archivo:wght@400;500;600;700;800;900&display=swap');
        
        /* Overall page background */
        [data-testid="stAppViewContainer"] {
            background: linear-gradient(135deg, #FAFAFA 0%, #FFFFFF 100%) !important;
            font-family: 'Inter', sans-serif !important;
        }

        /* Hide default Streamlit elements */
        #MainMenu {visibility: hidden;}
        footer {visibility: hidden;}
        header {visibility: hidden;}

        /* Main container styling */
        .block-container {
            padding-top: 5rem !important;
            max-width: 1200px !important;
        }

        /* Main heading */
        .main-heading {
            font-family: 'Archivo', sans-serif;
            font-size: 2.5rem;
            font-weight: 900;
            color: #1A1A1A;
            margin-bottom: 2rem;
            text-align: center;
            letter-spacing: -1px;
        }

        /* Subtext */
        .form-subtext {
            font-size: 0.95rem;
            color: #666666;
            text-align: center;
            margin-bottom: 2rem;
            letter-spacing: 0.3px;
            font-weight: 400;
        }

        /* Input fields styling */
        .stTextInput > div > div > input {
            background-color: #F5F5F5 !important;
            color: #1A1A1A !important;
            border: 2px solid #E5E5E5 !important;
            border-radius: 12px !important;
            padding: 16px 18px !important;
            font-size: 15px !important;
            transition: all 0.3s ease !important;
            font-family: 'Inter', sans-serif !important;
            font-weight: 500 !important;
            max-width: 500px !important;  /* ADDED: Maximum width constraint */
            margin: 0 auto !important;     /* ADDED: Center the input */
        }

        .stTextInput > div > div > input:focus {
            background-color: #FFFFFF !important;
            border-color: #FF6B35 !important;
            box-shadow: 0 0 0 4px rgba(255, 107, 53, 0.1) !important;
            outline: none !important;
        }

        .stTextInput > div > div > input::placeholder {
            color: #999999 !important;
            font-weight: 400;
        }
                 .stTextInput {
            max-width: 500px !important;  /* ADDED: Constrain text input container */
            margin: 0 auto !important;     /* ADDED: Center the container */
        }

        /* Label styling */
        .stTextInput label {
            font-size: 14px !important;
            font-weight: 600 !important;
            color: #1A1A1A !important;
            margin-bottom: 8px !important;
        }

        /* Social icons styling */
        .social-icons {
            display: flex;
            justify-content: center;
            gap: 16px;
            margin: 2rem 0;
        }

        .social-icon {
            border: 2px solid #E5E5E5;
            border-radius: 12px;
            display: inline-flex;
            justify-content: center;
            align-items: center;
            width: 48px;
            height: 48px;
            color: #666666;
            text-decoration: none;
            transition: all 0.3s ease;
            background: #FFFFFF;
        }

        .social-icon:hover {
            border-color: #FF6B35;
            color: #FF6B35;
            transform: translateY(-3px);
            box-shadow: 0 4px 12px rgba(255, 107, 53, 0.2);
        }

        /* Toggle text styling */
        .stCheckbox {
            display: flex !important;
            justify-content: center !important;
            align-items: center !important;
            max-width: 500px !important;     /* ADDED: Match form width */
            margin-left: 130px !important;     /* ADDED: Push to center/right */
            margin-right: auto !important;    /* ADDED: Keep centered */
            padding-left: 40px !important;    /* ADDED: Shift content right */
        }
        
        .stCheckbox > label {
            display: flex !important;
            width: 100% !important;
            margin: 0 !important;
            padding: 0 !important;
            font-size: 14px !important;
            color: #666666 !important;
            font-weight: 400 !important;
             justify-content: flex-start !important;  
        }
        
        .stCheckbox input[type="checkbox"] {
            display: none !important;
        }
        
        .stCheckbox > label > div:first-child {
            display: none !important;
        }
        
        .stCheckbox strong {
            color: #FF6B35 !important;
            font-weight: 600 !important;
            cursor: pointer !important;
            margin-left: 4px !important;
            transition: all 0.3s ease !important;
        }
        
        .stCheckbox strong:hover {
            color: #E85A28 !important;
            text-decoration: underline !important;
        }

        /* Divider */
        .divider-text {
            text-align: center;
            color: #999999;
            margin: 2rem 0;
            font-size: 0.9rem;
            position: relative;
        }

        .divider-text::before,
        .divider-text::after {
            content: '';
            position: absolute;
            top: 50%;
            width: 40%;
            height: 1px;
            background: #E5E5E5;
        }

        .divider-text::before {
            left: 0;
        }

        .divider-text::after {
            right: 0;
        }

        /* Welcome panel (right side toggle) */
        .welcome-panel {
            background: linear-gradient(135deg, #FF6B35 0%, #FFA500 100%);
            color: #FFFFFF;
            padding: 4rem 3rem;
            border-radius: 24px;
            text-align: center;
            min-height: 500px;
            display: flex;
            flex-direction: column;
            justify-content: center;
            box-shadow: 0 20px 60px rgba(255, 107, 53, 0.25);
            position: relative;
            overflow: hidden;
        }

        .welcome-panel::before {
            content: '';
            position: absolute;
            top: -50%;
            right: -20%;
            width: 400px;
            height: 400px;
            background: radial-gradient(circle, rgba(255, 255, 255, 0.15) 0%, transparent 70%);
            border-radius: 50%;
        }

        .welcome-heading {
            font-family: 'Archivo', sans-serif;
            font-size: 2.8rem;
            font-weight: 900;
            margin-bottom: 1rem;
            letter-spacing: -1px;
            position: relative;
            z-index: 1;
        }

        .welcome-text {
            font-size: 1.1rem;
            line-height: 1.7;
            margin-bottom: 2rem;
            opacity: 0.95;
            position: relative;
            z-index: 1;
        }

        /* Button styling - FIXED */
        [data-testid="stButton"] button {
            background: linear-gradient(135deg, #FF6B35 0%, #E85A28 100%) !important;
            color: #FFFFFF !important;
            font-size: 15px !important;
            padding: 16px 40px !important;
            border: none !important;
            border-radius: 12px !important;
            font-weight: 700 !important;
            letter-spacing: 0.5px !important;
            text-transform: uppercase !important;
            cursor: pointer !important;
            width: 100% !important;
            transition: all 0.3s ease !important;
            font-family: 'Inter', sans-serif !important;
            box-shadow: 0 6px 20px rgba(255, 107, 53, 0.3) !important;
            margin-top: 1rem !important;
            margin-left: 150px !important;
        }
        
        [data-testid="stButton"] button:hover {
            background: linear-gradient(135deg, #E85A28 0%, #D84315 100%) !important;
            transform: translateY(-3px) !important;
            box-shadow: 0 10px 30px rgba(255, 107, 53, 0.4) !important;
        }

        [data-testid="stButton"] button:active {
            transform: translateY(-1px) !important;
        }

        /* Ensure buttons are visible */
        .stButton {
            display: block !important;
            visibility: visible !important;
            max-width: 400px !important;  /* ADDED: Constrain button container */
           
        }

        /* Column alignment */
        [data-testid="column"] {
            display: flex;
            flex-direction: column;
            justify-content: center;
        }

        /* Responsive adjustments */
        @media (max-width: 768px) {
            .welcome-panel {
                display: none;
            }
            
            .main-heading {
                font-size: 2rem;
            }

            .block-container {
                padding: 2rem 1rem !important;
            }
        }
    </style>
    """, unsafe_allow_html=True)

    # -------------------------------------------------------------------------------------------------
    # PAGE LAYOUT
    # -------------------------------------------------------------------------------------------------
    
    # Create container columns
    col_form, col_welcome = st.columns([1, 1], gap="large")

    # FORM SECTION (Left side)
    with col_form:
        # Mode-specific heading
        if st.session_state.mode == 'login':
            st.markdown('<h1 class="main-heading">Sign In</h1>', unsafe_allow_html=True)
        else:
            st.markdown('<h1 class="main-heading">Create Account</h1>', unsafe_allow_html=True)
        
        # Social icons
        # st.markdown('''
        # <div class="social-icons">
        #     <a href="#" class="social-icon" title="Google">
        #         <svg width="20" height="20" viewBox="0 0 24 24" fill="currentColor">
        #             <path d="M22.56 12.25c0-.78-.07-1.53-.2-2.25H12v4.26h5.92c-.26 1.37-1.04 2.53-2.21 3.31v2.77h3.57c2.08-1.92 3.28-4.74 3.28-8.09z"/>
        #             <path d="M12 23c2.97 0 5.46-.98 7.28-2.66l-3.57-2.77c-.98.66-2.23 1.06-3.71 1.06-2.86 0-5.29-1.93-6.16-4.53H2.18v2.84C3.99 20.53 7.7 23 12 23z"/>
        #             <path d="M5.84 14.09c-.22-.66-.35-1.36-.35-2.09s.13-1.43.35-2.09V7.07H2.18C1.43 8.55 1 10.22 1 12s.43 3.45 1.18 4.93l2.85-2.22.81-.62z"/>
        #             <path d="M12 5.38c1.62 0 3.06.56 4.21 1.64l3.15-3.15C17.45 2.09 14.97 1 12 1 7.7 1 3.99 3.47 2.18 7.07l3.66 2.84c.87-2.6 3.3-4.53 6.16-4.53z"/>
        #         </svg>
        #     </a>
        #     <a href="#" class="social-icon" title="Facebook">
        #         <svg width="20" height="20" viewBox="0 0 24 24" fill="currentColor">
        #             <path d="M24 12.073c0-6.627-5.373-12-12-12s-12 5.373-12 12c0 5.99 4.388 10.954 10.125 11.854v-8.385H7.078v-3.47h3.047V9.43c0-3.007 1.792-4.669 4.533-4.669 1.312 0 2.686.235 2.686.235v2.953H15.83c-1.491 0-1.956.925-1.956 1.874v2.25h3.328l-.532 3.47h-2.796v8.385C19.612 23.027 24 18.062 24 12.073z"/>
        #         </svg>
        #     </a>
        #     <a href="#" class="social-icon" title="LinkedIn">
        #         <svg width="20" height="20" viewBox="0 0 24 24" fill="currentColor">
        #             <path d="M20.447 20.452h-3.554v-5.569c0-1.328-.027-3.037-1.852-3.037-1.853 0-2.136 1.445-2.136 2.939v5.667H9.351V9h3.414v1.561h.046c.477-.9 1.637-1.85 3.37-1.85 3.601 0 4.267 2.37 4.267 5.455v6.286zM5.337 7.433c-1.144 0-2.063-.926-2.063-2.065 0-1.138.92-2.063 2.063-2.063 1.14 0 2.064.925 2.064 2.063 0 1.139-.925 2.065-2.064 2.065zm1.782 13.019H3.555V9h3.564v11.452zM22.225 0H1.771C.792 0 0 .774 0 1.729v20.542C0 23.227.792 24 1.771 24h20.451C23.2 24 24 23.227 24 22.271V1.729C24 .774 23.2 0 22.222 0h.003z"/>
        #         </svg>
        #     </a>
        #     <a href="#" class="social-icon" title="GitHub">
        #         <svg width="20" height="20" viewBox="0 0 24 24" fill="currentColor">
        #             <path d="M12 0c-6.626 0-12 5.373-12 12 0 5.302 3.438 9.8 8.207 11.387.599.111.793-.261.793-.577v-2.234c-3.338.726-4.033-1.416-4.033-1.416-.546-1.387-1.333-1.756-1.333-1.756-1.089-.745.083-.729.083-.729 1.205.084 1.839 1.237 1.839 1.237 1.07 1.834 2.807 1.304 3.492.997.107-.775.418-1.305.762-1.604-2.665-.305-5.467-1.334-5.467-5.931 0-1.311.469-2.381 1.236-3.221-.124-.303-.535-1.524.117-3.176 0 0 1.008-.322 3.301 1.23.957-.266 1.983-.399 3.003-.404 1.02.005 2.047.138 3.006.404 2.291-1.552 3.297-1.23 3.297-1.23.653 1.653.242 2.874.118 3.176.77.84 1.235 1.911 1.235 3.221 0 4.609-2.807 5.624-5.479 5.921.43.372.823 1.102.823 2.222v3.293c0 .319.192.694.801.576 4.765-1.589 8.199-6.086 8.199-11.386 0-6.627-5.373-12-12-12z"/>
        #         </svg>
        #     </a>
        # </div>
        # ''', unsafe_allow_html=True)
        
        # Subtext
        # if st.session_state.mode == 'login':
        #     st.markdown('<p class="form-subtext">use your email and password</p>', unsafe_allow_html=True)
        # else:
        #     st.markdown('<p class="form-subtext">use your email for registration</p>', unsafe_allow_html=True)
        
        # Name field for registration
        name = ""
        if st.session_state.mode == 'register':
            name = st.text_input("Name", placeholder="Enter your full name", key="full_name")
        
        # Email field
        email = st.text_input("Email", placeholder="Enter your email", key="lemail")
        
        # Password field
        password = st.text_input("Password", placeholder="Enter your password", type="password", key="password")
        
        # Main action button (Login/Sign Up)
        button_text = "Sign Up" if st.session_state.mode == 'register' else "Sign In"
        
        if st.button(button_text, key="main_action_btn", use_container_width=True):
            if email and password:
                if not is_valid_email(email):
                    st.error("Please enter a valid email address")
                else:
                    users = load_users()

                    if st.session_state.mode == 'login':
                        # LOGIN LOGIC
                        user_entry = users.get(email)
                        stored_pw = user_entry.get("password") if isinstance(user_entry, dict) else user_entry
                        stored_name = user_entry.get("name") if isinstance(user_entry, dict) else None

                        if user_entry is None or stored_pw != password:
                            st.error("Invalid email or password")
                        else:
                            st.session_state.logged_in_user = email
                            st.query_params["user"] = email
                            st.session_state.username = stored_name or email.split('@')[0]

                            user_resume = get_user_resume(email)
                            if user_resume and len(user_resume) > 0:
                                st.session_state.resume_source = user_resume
                                st.session_state.input_method = user_resume.get("input_method", "Manual Entry")
                                st.success(f"Welcome back, {st.session_state.username}!")
                                time.sleep(0.8)
                                st.switch_page("app.py")
                            else:
                                st.success(f"Welcome, {st.session_state.username}!")
                                time.sleep(0.8)
                                st.switch_page("app.py")
                    else:
                        # REGISTRATION LOGIC
                        if email in users:
                            st.error("Email already registered. Please login.")
                            st.session_state.mode = 'login'
                            st.rerun()
                        elif len(password) < 6:
                            st.error("Password must be at least 6 characters long")
                        elif not name or len(name.strip()) == 0:
                            st.error("Please enter your full name")
                        else:
                            users[email] = {"password": password, "name": name.strip()}
                            save_users(users)
                            st.session_state.logged_in_user = email
                            st.query_params["user"] = email
                            st.session_state.username = name.strip()
                            st.session_state.mode = 'login'
                            st.success("Account created successfully!")
                            time.sleep(0.8)
                            st.switch_page("app.py")
            else:
                st.warning("Please enter both email and password")
        
        # Toggle between login/register
        st.markdown('<br>', unsafe_allow_html=True)
        if st.session_state.mode == 'login':
            toggle_signup = st.checkbox(
                "Don't have an account? **Sign up**",
                key="toggle_to_signup",
                value=False
            )
            
            if toggle_signup:
                st.session_state.mode = 'register'
                st.rerun()
                
        else:
            toggle_login = st.checkbox(
                "Already have an account? **Back to Login**",
                key="toggle_to_login", 
                value=False
            )
            
            if toggle_login:
                st.session_state.mode = 'login'
                st.rerun()

    # WELCOME PANEL (Right side)
    with col_welcome:
        if st.session_state.mode == 'login':
            st.markdown('''
            <div class="welcome-panel">
                <h1 class="welcome-heading">Hello, Friend!</h1>
                <p class="welcome-text">Register with your personal details to use all of our amazing features and start your journey with us.</p>
            </div>
            ''', unsafe_allow_html=True)
        else:
            st.markdown('''
            <div class="welcome-panel">
                <h1 class="welcome-heading">Welcome Back!</h1>
                <p class="welcome-text">Enter your personal details to continue your journey with us and access all features.</p>
            </div>
            ''', unsafe_allow_html=True)
            
            # # Toggle button
            # with stylable_container(
            #     key="login_toggle_btn",
            #     css_styles="""
            #     button {
            #         background-color: transparent !important;
            #         color: #fff !important;
            #         border: 1px solid #fff !important;
            #         border-radius: 8px !important;
            #         padding: 10px 45px !important;
            #         font-size: 12px !important;
            #         font-weight: 600 !important;
            #         letter-spacing: 0.5px !important;
            #         text-transform: uppercase !important;
            #         cursor: pointer !important;
            #         width: 60% !important;
            #         transition: all 0.3s ease !important;
            #         font-family: 'Montserrat', sans-serif !important;
            #         margin: 0 auto !important;
            #         display: block !important;
            #     }
            #     button:hover {
            #         background-color: #fff !important;
            #         color: #f12711 !important;
            #         transform: translateY(-2px) !important;
            #     }
            #     """
            # ):
            #     # if st.button("Sign In", key="toggle_login"):
            #     #     st.session_state.mode = 'login'
            #     #     st.rerun()



# =====================================
# 🤖 CHATBOT BEHAVIOR & AI RESPONSE LOGIC
# =====================================

# Enhanced ask_llama function
RESUME_ASSISTANT_PROMPT = """
You are ResumeBot, a professional resume and career assistant. Your role is to provide concise, actionable advice for resume improvement and career guidance.

CORE RESPONSIBILITIES:
- Analyze resume content and provide specific improvement suggestions
- Guide users on relevant content inclusion based on their industry/role
- Offer formatting and structuring advice
- Suggest better ways to phrase accomplishments
- Help tailor resumes for specific job applications
- Provide career development guidance
- Assist with template selection and customization
- Give interview preparation tips

RESPONSE GUIDELINES:
- Be professional, polite, and encouraging
- Keep responses concise but helpful (3-5 sentences or bullet points)
- Use bullet points for actionable advice when appropriate
- Focus on quantifiable achievements and specific examples
- Ask clarifying questions when needed to provide better guidance
- Never request personal sensitive information
- Tailor advice to the user's mentioned industry/experience level

FORMAT PREFERENCE:
- Use bullet points (•) for lists when helpful
- Use emojis sparingly to make responses engaging
- Keep paragraphs short and scannable
- End with a relevant question to continue conversation when appropriate

Example response style:
• **Quantify achievements**: Instead of "managed team" try "Led 5-person team to achieve 15% productivity increase"
• **Add relevant keywords**: Include "project management" and "stakeholder communication" from job description
• **Improve formatting**: Use consistent bullet points and clear section headers

Always be helpful and specific in your advice!"""



def ask_llama(message, resume_data=None):
    """
    Stream tokens live + fallback model + optional resume context
    """
    print(f"🔍 DEBUG: ask_llama called with message: {message[:50]}...")  # Debug log
    
    context = f"\nHere is the user's resume:\n{resume_data}\n" if resume_data else ""
    prompt = f"{RESUME_ASSISTANT_PROMPT}{context}\nUser: {message}\nAI:"

    payload = {
        "model": HEAVY_MODEL,  # try high-quality model first
        "messages": [
            {"role": "system", "content": RESUME_ASSISTANT_PROMPT},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.3,     # faster + less hallucinations
        "max_tokens": 200,      # limit output for speed
        "stream": True
    }

    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }

    print(f"🔍 DEBUG: Making API request to {API_URL}")  # Debug log
    print(f"🔍 DEBUG: Using model: {HEAVY_MODEL}")  # Debug log

    try:
        response = requests.post(API_URL, headers=headers, json=payload, stream=True, timeout=35)
        print(f"🔍 DEBUG: Response status code: {response.status_code}")  # Debug log
        
        if response.status_code != 200:
            print(f"❌ DEBUG: API Error - {response.text}")  # Debug log
            raise Exception(f"API returned {response.status_code}")
            
    except Exception as e:
        print(f"⚠️ DEBUG: Primary model failed ({e}), trying fallback...")  # Debug log
        # Fallback to fast 8B model automatically 👍
        payload["model"] = FAST_MODEL
        print(f"🔍 DEBUG: Using fallback model: {FAST_MODEL}")  # Debug log
        response = requests.post(API_URL, headers=headers, json=payload, stream=True, timeout=35)
        print(f"🔍 DEBUG: Fallback response status: {response.status_code}")  # Debug log

    # Streaming the reply token by token
    token_count = 0
    for line in response.iter_lines():
        if line:
            try:
                decoded = line.decode("utf-8").replace("data:", "").strip()
                
                # Skip [DONE] marker
                if decoded == "[DONE]":
                    print(f"✅ DEBUG: Stream complete. Total tokens: {token_count}")  # Debug log
                    continue
                    
                data = json.loads(decoded)
                token = data["choices"][0]["delta"].get("content", "")
                if token:
                    token_count += 1
                    if token_count == 1:
                        print(f"✅ DEBUG: First token received!")  # Debug log
                    yield token
            except json.JSONDecodeError as e:
                print(f"⚠️ DEBUG: JSON decode error: {e} | Line: {line}")  # Debug log
                pass
            except Exception as e:
                print(f"⚠️ DEBUG: Stream error: {e}")  # Debug log
                pass
    
    if token_count == 0:
        print(f"❌ DEBUG: No tokens received from API!")  # Debug log

    
def chatbot(user_resume):
    import streamlit as st

    # Page config
  
    # Initialize session state for messages
    if 'messages' not in st.session_state:
        st.session_state.messages = [
            {"role": "bot", "content": "Hello! How can I help you today?", "time": "10:30 AM"},
            {"role": "user", "content": "I need help with my resume", "time": "10:31 AM"},
            {"role": "bot", "content": "I'd be happy to help! Could you provide your basic information", "time": "10:31 AM"},
            {"role": "user", "content": "Sure,", "time": "10:32 AM"},
            {"role": "bot", "content": "Thank you! Let me check that for you...", "time": "10:32 AM"}
        ]

    # Custom CSS for the popover and chat styling
    st.html("""
    <style>
        /* Force popover to bottom right with higher specificity */
        [data-testid="stPopover"],
        div[data-testid="stPopover"],
        section[data-testid="stPopover"] {
            position: fixed !important;
            bottom: 20px !important;
            right: 20px !important;
            left: unset !important;
            top: unset !important;
            z-index: 999999 !important;
            margin: 0 !important;
        }
        
        /* Target the parent container */
        [data-testid="stPopover"] > div {
            position: fixed !important;
            bottom: 20px !important;
            right: 20px !important;
            left: unset !important;
        }
        
        /* Style the popover button */
        [data-testid="stPopover"] button,
        [data-testid="stPopover"] > button {
            width: 60px !important;
            height: 60px !important;
            border-radius: 50% !important;
            background: linear-gradient(135deg, #ff6b00 0%, #ff8c42 100%) !important;
            border: none !important;
            box-shadow: 0 5px 20px rgba(255, 107, 0, 0.4) !important;
            font-size: 28px !important;
            padding: 0 !important;
            transition: all 0.3s ease !important;
        }
        
        [data-testid="stPopover"] button:hover {
            transform: scale(1.1) !important;
            box-shadow: 0 7px 25px rgba(255, 107, 0, 0.5) !important;
        }
        
        /* Style the popover content */
        [data-testid="stPopover"] > div > div,
        [data-testid="stPopoverBody"] {
            background: linear-gradient(135deg, #fff5f0 0%, #ffe8db 100%) !important;
            border: 2px solid #ff8c42 !important;
            border-radius: 20px !important;
            box-shadow: 0 10px 40px rgba(255, 107, 0, 0.25) !important;
            padding: 0 !important;
            width: 380px !important;
            max-width: 380px !important;
            position: fixed !important;
            bottom: 90px !important;
            right: 20px !important;
        }
        
        /* Hide Streamlit branding */
        #MainMenu {visibility: hidden;}
        footer {visibility: hidden;}
        header {visibility: hidden;}
        
        /* Chat header styling */
        .chat-header {
            background: linear-gradient(135deg, #ff6b00 0%, #ff8c42 100%);
            padding: 18px 20px;
            color: white;
            font-weight: 600;
            font-size: 16px;
            display: flex;
            align-items: center;
            gap: 10px;
            border-radius: 18px 18px 0 0;
            margin: -1rem -1rem 1rem -1rem;
        }
        
        .chat-icon {
            width: 30px;
            height: 30px;
            background: white;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 18px;
        }
        
        /* Message styling */
        .chat-message {
            padding: 12px 16px;
            border-radius: 18px;
            margin-bottom: 12px;
            font-size: 14px;
            line-height: 1.4;
            max-width: 75%;
        }
        
        .chat-message.bot {
            background: white;
            color: #333;
            border: 1px solid #ffd4b8;
            box-shadow: 0 2px 5px rgba(255, 107, 0, 0.1);
            margin-right: auto;
        }
        
        .chat-message.user {
            background: linear-gradient(135deg, #ff6b00 0%, #ff8c42 100%);
            color: white;
            box-shadow: 0 2px 8px rgba(255, 107, 0, 0.3);
            margin-left: auto;
        }
        
        .timestamp {
            font-size: 11px;
            margin-top: 4px;
            opacity: 0.7;
        }
        
        .chat-message.bot .timestamp {
            color: #999;
        }
        
        .chat-message.user .timestamp {
            color: rgba(255, 255, 255, 0.8);
        }
        
        /* Chat messages container */
        .chat-messages-container {
            max-height: 350px;
            overflow-y: auto;
            padding: 10px;
            margin-bottom: 15px;
        }
        
        /* Scrollbar styling */
        .chat-messages-container::-webkit-scrollbar {
            width: 6px;
        }
        
        .chat-messages-container::-webkit-scrollbar-track {
            background: #fff5f0;
        }
        
        .chat-messages-container::-webkit-scrollbar-thumb {
            background: #ff8c42;
            border-radius: 10px;
        }
        
        .chat-messages-container::-webkit-scrollbar-thumb:hover {
            background: #ff6b00;
        }
        
        /* Chat input styling */
        [data-testid="stPopover"] .stChatInput {
            margin-top: 10px;
        }
        
        [data-testid="stPopover"] .stChatInput input {
            border: 2px solid #ff8c42 !important;
            border-radius: 25px !important;
            background: #fff9f5 !important;
        }
        
        [data-testid="stPopover"] .stChatInput input:focus {
            border-color: #ff6b00 !important;
            box-shadow: 0 0 0 3px rgba(255, 107, 0, 0.1) !important;
        }
    </style>

    <script>
        // Force position with JavaScript as backup
        function positionPopover() {
            const popover = document.querySelector('[data-testid="stPopover"]');
            if (popover) {
                popover.style.position = 'fixed';
                popover.style.bottom = '20px';
                popover.style.right = '20px';
                popover.style.left = 'auto';
                popover.style.top = 'auto';
            }
        }
        
        // Run on load and with observer
        window.addEventListener('load', positionPopover);
        setTimeout(positionPopover, 100);
        setTimeout(positionPopover, 500);
        setTimeout(positionPopover, 1000);
        
        // Watch for DOM changes
        const observer = new MutationObserver(positionPopover);
        observer.observe(document.body, { childList: true, subtree: true });
    </script>
    """)


    # Chat popover in bottom right
    with st.popover("💬"):
        # Chat header
        st.markdown("""
        <div class="chat-header">
            <div class="chat-icon">💬</div>
            <span>Chat Assistant</span>
        </div>
        """, unsafe_allow_html=True)
        
        # Messages container
        st.markdown('<div class="chat-messages-container">', unsafe_allow_html=True)
        
        # Display messages
        for msg in st.session_state.messages:
            if msg["role"] == "bot":
                st.markdown(f"""
                <div class="chat-message bot">
                    {msg["content"]}
                    <div class="timestamp">{msg["time"]}</div>
                </div>
                """, unsafe_allow_html=True)
            else:
                st.markdown(f"""
                <div class="chat-message user">
                    {msg["content"]}
                    <div class="timestamp">{msg["time"]}</div>
                </div>
                """, unsafe_allow_html=True)
        
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Chat input
        # Chat input
        prompt = st.chat_input("Type your message here...")
        if prompt:
            from datetime import datetime
            

            # Add user message
            current_time = datetime.now().strftime("%I:%M %p")
            st.session_state.messages.append({
                "role": "user",
                "content": prompt,
                "time": current_time
            })

            # Call ResumeBot (LLaMA)
            with st.chat_message("assistant"):
                placeholder = st.empty()
                reply = ""
                for token in ask_llama(prompt, user_resume):
                    reply += token
                    placeholder.markdown(reply)

            # Add bot response to UI
            st.session_state.messages.append({
                "role": "bot",
                "content": reply,
                "time": current_time
            })

            st.rerun()



def format_section_title(key):
    """Converts keys like 'certifications' to 'Certifications'."""
    title = key.replace('_', ' ').replace('Skills', ' Skills').replace('summary', 'Summary')
    return ' '.join(word.capitalize() for word in title.split())

def get_standard_keys():
    """Return set of standard resume keys that should not be treated as custom sections."""
    return {
        "name", "email", "phone", "location", "url", "summary", "job_title",
        "education", "experience", "skills", "projects", "certifications", 
        "achievements", "total_experience_count"
    }

def format_year_only(date_str):
    """Return only the year from a date string. If already a year, return as-is."""
    if not date_str:
        return ""
    date_str = str(date_str).strip()
    if len(date_str) == 4 and date_str.isdigit(): 
        return date_str
    try:
        dt = datetime.strptime(date_str, "%Y-%m-%d")
        return dt.strftime("%Y")
    except:
        return date_str  


RESUME_ORDER = ["summary", "experience", "education", "skills", "projects", "certifications", "achievements", "publications", "awards"]

# ATS-friendly color palette



def generate_generic_html(data, date_placement='right'):
    # Normalise all keys to lowercase
    normalized_data = {k.lower(): v for k, v in data.items()}

    html = ""

    # ===============================
    # HELPER FUNCTION: Extract Year
    # ===============================
    def extract_year(date_str):
        """Extract year from date string. Handles formats like '2024-01-01', '2024', 'Jan 2024', etc."""
        if not date_str:
            return ""
        
        date_str = str(date_str).strip()
        
        # If already just a year (4 digits)
        if date_str.isdigit() and len(date_str) == 4:
            return date_str
        
        # If it's "Present" or similar, return as-is
        if date_str.lower() in ['present', 'current', 'ongoing']:
            return date_str.title()
        
        # Try to extract 4-digit year from string
        import re
        year_match = re.search(r'\b(19|20)\d{2}\b', date_str)
        if year_match:
            return year_match.group(0)
        
        # If format is like "2024-01-01", split by dash and take first part
        if '-' in date_str:
            return date_str.split('-')[0]
        
        # If format is like "01/01/2024", split and take last part
        if '/' in date_str:
            parts = date_str.split('/')
            for part in parts:
                if len(part) == 4 and part.isdigit():
                    return part
        
        # Default: return as-is
        return date_str

    # ===============================
    # 1. HEADER
    # ===============================
    html += '<div class="ats-header">'

    html += f'<h1>{normalized_data.get("name", "")}</h1>'

    contacts = [
        normalized_data.get("email"),
        normalized_data.get("phone"),
        normalized_data.get("url"),
        normalized_data.get("location")
    ]

    contacts = [c for c in contacts if c]
    if contacts:
        html += f'<div class="ats-contact">{" | ".join(contacts)}</div>'

    html += "</div>"

    # ===============================
    # 2. SECTION DEFINITIONS
    # ===============================

    SECTION_TITLES = {
        "summary": "Summary",
        "experience": "Experience",
        "skills": "Skills",
        "education": "Education",
        "projects": "Projects",
        "project": "Project",
        "certifications": "Certifications",
        "achievements": "Achievements",
        "awards": "Awards",
        "interest": "Interests",
        "interests": "Interests",
        "languages": "Languages",
        "publications": "Publications"
    }

    # Helper to detect empty fields
    def is_empty(value):
        if value is None:
            return True
        if isinstance(value, str) and value.strip() == "":
            return True
        if isinstance(value, list) and all(is_empty(v) for v in value):
            return True
        if isinstance(value, dict) and all(is_empty(v) for v in value.values()):
            return True
        return False

    # ===============================
    # 3. AUTO PROCESS SECTIONS
    # ===============================
    for key, title in SECTION_TITLES.items():

        if key not in normalized_data:
            continue

        value = normalized_data[key]

        if is_empty(value):
            continue

        html += f'<div class="ats-section-title">{title}</div>'

        # -----------------------------------------
        # CASE A: Simple paragraph
        # -----------------------------------------
        if isinstance(value, str):
            html += f"<p>{value}</p>"
            continue

        # -----------------------------------------
        # CASE B: Dictionary (Skills, etc.)
        # -----------------------------------------
        if isinstance(value, dict):
            for subkey, subvalue in value.items():
                sub_title = subkey.replace("_", " ").title()
                html += f'<div class="ats-subtitle"><strong>{sub_title}:</strong> '

                if isinstance(subvalue, list):
                    html += ", ".join([str(v) for v in subvalue])
                else:
                    html += str(subvalue)

                html += "</div>"

            continue

        # -----------------------------------------
        # CASE C: List
        # -----------------------------------------
        if isinstance(value, list):

            # -------- CASE C1: List of plain strings --------
            if all(isinstance(x, str) for x in value):
                html += "<ul>"
                for item in value:
                    html += f"<li>{item}</li>"
                html += "</ul>"
                continue

            # -------- CASE C2: List of objects --------
            for item in value:
                if not isinstance(item, dict):
                    continue

                # ✅ FIXED: Auto-detect title, subtitle, duration with ALL field support
                title_val = (
                    item.get("position") or 
                    item.get("role") or 
                    item.get("projectname") or
                    item.get("name") or 
                    item.get("degree") or  # Education
                    item.get("course") or 
                    item.get("certificate_name") or
                    item.get("title")
                )

                # ✅ FIXED: Include issuer for certifications
                subtitle_val = (
                    item.get("company") or 
                    item.get("institution") or  # Education
                    item.get("university") or 
                    item.get("issuer") or  # ✅ ADD THIS for certifications
                    item.get("tools") or
                    item.get("provider_name")
                )

                # ✅ IMPROVED: Handle dates with year extraction
                start_date = item.get("start_date", "")
                end_date = item.get("end_date", "")
                completed_date = item.get("completed_date", "")
                duration = item.get("duration", "")
                
                # Determine which dates to use based on section
                if key in ["education", "certifications"]:
                    # For education and certifications, extract years only
                    if start_date or end_date:
                        start_year = extract_year(start_date)
                        end_year = extract_year(end_date)
                        
                        if start_year and end_year:
                            duration_val = f"{start_year} - {end_year}"
                        elif start_year:
                            duration_val = start_year
                        elif end_year:
                            duration_val = end_year
                        else:
                            duration_val = ""
                    elif completed_date:
                        duration_val = extract_year(completed_date)
                    elif duration:
                        duration_val = duration
                    else:
                        duration_val = ""
                else:
                    # For other sections (experience, projects), keep full dates
                    if start_date or end_date:
                        if start_date and end_date:
                            duration_val = f"{start_date} - {end_date}"
                        else:
                            duration_val = start_date or end_date
                    elif completed_date:
                        duration_val = completed_date
                    elif duration:
                        duration_val = duration
                    else:
                        duration_val = ""

                # Item header
                html += '<div class="ats-item-header">'

                if title_val:
                    html += '<div class="ats-item-title-group">'
                    html += f'<span class="ats-item-title">{title_val}'

                    if subtitle_val:
                        html += f' <span class="ats-item-subtitle">{subtitle_val}</span>'

                    html += '</span></div>'

                if duration_val:
                    html += f'<div class="ats-item-duration">{duration_val}</div>'

                html += '</div>'

                # Description or details
                desc = item.get("description") or item.get("details") or item.get("overview")

                if isinstance(desc, list):
                    html += "<ul>"
                    for d in desc:
                        html += f"<li>{d}</li>"
                    html += "</ul>"

                elif isinstance(desc, str):
                    html += f"<p>{desc}</p>"

    return html


def extract_template_from_html(html_content):
    """Extract CSS and structure from uploaded HTML."""

    css_match = re.search(r'<style[^>]*>(.*?)</style>', html_content, re.DOTALL)
    css = css_match.group(1) if css_match else ""
    
    template_id = hashlib.md5(html_content.encode()).hexdigest()[:8]
    
    return {
        'id': template_id,
        'name': f'Uploaded Template {template_id}',
        'css': css,
        'html': html_content,
        'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    }

def generate_html_content(data, color, template_generator, css_generator):
    """Generates the full HTML document and returns base64-encoded string."""
    css = css_generator(color)
    html_content = template_generator(data)
    
    title = f"{data.get('name', 'Resume')}"
    full_document_html = f"<html><head><meta charset='UTF-8'>{css}<title>{title}</title></head><body><div class='ats-page'>{html_content}</div></body></html>" 
    return base64.b64encode(full_document_html.encode('utf-8')).decode()

def get_html_download_link(data, color, template_config, filename_suffix=""):
    """Generates a download link for the styled HTML file."""
    if 'html_generator' in template_config and 'css_generator' in template_config:
        b64_data = generate_html_content(
            data, 
            color, 
            template_config['html_generator'], 
            template_config['css_generator']
        )
    else:
        # For uploaded templates
        css = template_config.get('css', '')
        html_body = generate_generic_html(data)
        full_html = f"<html><head><meta charset='UTF-8'><style>{css}</style></head><body><div class='ats-page'>{html_body}</div></body></html>"
        b64_data = base64.b64encode(full_html.encode('utf-8')).decode()
    
    filename = f"Resume_{data.get('name', 'User').replace(' ', '_')}{filename_suffix}.html"
    href = f'<a href="data:text/html;base64,{b64_data}" download="{filename}" style="font-size: 0.95em; text-decoration: none; padding: 10px 15px; background-color: #0077B6; color: white; border-radius: 5px; display: inline-block; margin-top: 10px; width: 100%; text-align: center;"><strong> HTML (.html)</strong></a>'
    return href



# UPDATED generate_doc_html function
def generate_doc_html(data):
    """Generate a simple HTML that can be saved as .doc."""
    html = f"""
    <html xmlns:o='urn:schemas-microsoft-com:office:office' xmlns:w='urn:schemas-microsoft-com:office:word' xmlns='http://www.w3.org/TR/REC-html40'>
    <head>
        <meta charset='utf-8'>
        <title>Resume</title>
        <style>
            @page {{
                size: 8.5in 11in;
                margin: 0.5in;
            }}
            body {{ 
                font-family: Calibri, Arial, sans-serif; 
                font-size: 10pt; 
                line-height: 1.1; 
                margin: 0;
                padding: 0;
            }}
            h1 {{ font-size: 16pt; margin: 0 0 3pt 0; text-align: center; }}
            h2 {{ 
                font-size: 11pt; 
                margin-top: 8pt; 
                margin-bottom: 3pt; 
                border-bottom: 1px solid black; 
                padding-bottom: 1pt;
            }}
            p {{ margin: 3pt 0; }}
            ul {{ margin: 3pt 0; padding-left: 18pt; }}
            li {{ margin: 1pt 0; }}
            .header {{ text-align: center; margin-bottom: 8pt; }}
            .contact {{ font-size: 9pt; }}
            .job-title {{ font-size: 11pt; margin: 2pt 0; color: #555; }}
            .item-header {{ margin: 2pt 0; }}
            .item-title {{ font-weight: bold; }}
            .item-subtitle {{ font-style: italic; color: #555; }}
            .skills-group {{ margin: 3pt 0; }}
            .custom-section {{ margin: 8pt 0; line-height: 1.4; }}
        </style>
    </head>
    <body>
        <div class='header'>
            <h1>{data.get('name', 'NAME MISSING')}</h1>
            <p class='job-title'>{data.get('job_title', '')}</p>
            <p class='contact'>{data.get('phone', '')} | {data.get('email', '')} | {data.get('location', '')}</p>
        </div>
    """
    
    # Standard sections
    for key in RESUME_ORDER:
        section_data = data.get(key)
        
        if not section_data or (isinstance(section_data, list) and not section_data) or (key == 'summary' and not section_data):
            continue
            
        title = format_section_title(key)
        html += f'<h2>{title}</h2>'
        
        if key == 'summary' and isinstance(section_data, str):
            html += f'<p>{section_data}</p>'
        
        elif key == 'skills' and isinstance(section_data, dict):
            for skill_type, skill_list in section_data.items():
                if skill_list:
                    html += f'<p class="skills-group"><strong>{format_section_title(skill_type)}:</strong> '
                    html += ", ".join(skill_list)
                    html += '</p>'
        
        elif isinstance(section_data, list):
            for item in section_data:
                if isinstance(item, str):
                    html += f'<ul><li>{item}</li></ul>'
                    continue
                
                if not isinstance(item, dict):
                    continue
                
                title_keys = ['title', 'name', 'degree', 'position', 'course']
                subtitle_keys = ['company', 'institution', 'issuer', 'organization', 'university']
                duration_keys = ['duration', 'date', 'period']
                
                main_title = next((item[k] for k in title_keys if k in item and item[k]), '')
                subtitle = next((item[k] for k in subtitle_keys if k in item and item[k] != main_title and item[k]), '')
                duration = next((item[k] for k in duration_keys if k in item and item[k]), '')

                html += f'<p class="item-header"><span class="item-title">{main_title}</span>'
                if subtitle:
                    html += f' <span class="item-subtitle">({subtitle})</span>'
                if duration:
                    html += f' | {duration}'
                html += '</p>'
                
                description_list_raw = item.get('description') or item.get('achievement') or item.get('details') 

                if description_list_raw:
                    if isinstance(description_list_raw, str):
                        description_list = [description_list_raw]
                    elif isinstance(description_list_raw, list):
                        description_list = description_list_raw
                    else:
                        description_list = None
                        
                    if description_list:
                        html += '<ul>'
                        for line in description_list:
                            html += f'<li>{line}</li>'
                        html += '</ul>'

    # ========== ADD CUSTOM SECTIONS HERE ==========
    standard_keys = get_standard_keys()
    
    for key, value in data.items():
        if key not in standard_keys and isinstance(value, str) and value.strip():
            title = format_section_title(key)
            html += f'<h2>{title}</h2>'
            # Replace newlines with <br> for proper rendering
            formatted_value = value.strip().replace('\n', '<br>')
            html += f'<p class="custom-section">{formatted_value}</p>'

    html += '</body></html>'
    return html


def get_doc_download_link(data, color, template_config, filename_suffix=""):
    """
    Generates a DOC download link using the selected template's HTML/CSS.

    Note: The 'DOC' file generated is an HTML file with a .doc extension and 
    Microsoft Word-specific XML/CSS (like @page rules) added to the beginning, 
    allowing it to open and be formatted correctly in Word.
    """

    if 'html_generator' in template_config and 'css_generator' in template_config:
        css = template_config['css_generator'](color)
        html_body = template_config['html_generator'](data)
    else:
       
        css = template_config.get('css', '')
        html_body = generate_generic_html(data) #
    word_doc_header = f"""
<html xmlns:o='urn:schemas-microsoft-com:office:office'
xmlns:w='urn:schemas-microsoft-com:office:word'
xmlns='http://www.w3.org/TR/REC-html40'>
<head>
    <meta charset='UTF-8'>
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{data.get('name', 'Resume')}</title>
    <style>
        @page {{ 
            size: 8.5in 11in; 
            margin: 0.5in; 
        }}
        {css}
    </style>
</head>
<body>
    <div class='ats-page'>
        {html_body}
    </div>
</body>
</html>
"""

   
    try:
        b64_data = base64.b64encode(word_doc_header.encode('utf-8')).decode()
    except NameError:
        print("Error: 'base64' module not found. Please import it.")
        return ""

    filename = f"Resume_{data.get('name', 'User').replace(' ', '_')}_Template_DOC{filename_suffix}.doc"

    doc_html = f"""
    <a href="data:application/msword;base64,{b64_data}" download="{filename}"
       style="font-size: 0.95em; text-decoration: none; padding: 10px 15px; 
              background-color: #00B4D8; color: white; border-radius: 5px; 
              display: inline-block; margin-top: 10px; width: 100%; text-align: center;">
        <strong> Word Document(.doc)</strong>
    </a>
    """
    return doc_html





# UPDATED generate_markdown_text function
def generate_markdown_text(data):
    """Generates a plain markdown/text version of the resume."""
    text = ""
    
    text += f"{data.get('name', 'NAME MISSING').upper()}\n"
    if data.get('job_title'):
        text += f"{data.get('job_title')}\n"
    contact_parts = [data.get('phone', ''), data.get('email', ''), data.get('location', '')]
    text += " | ".join(filter(None, contact_parts)) + "\n"
    text += "=" * 50 + "\n\n"
    
    # Standard sections
    for key in RESUME_ORDER:
        section_data = data.get(key)
        
        if not section_data or (isinstance(section_data, list) and not section_data) or (key == 'summary' and not section_data):
            continue
            
        title = format_section_title(key).upper()
        text += f"{title}\n"
        text += "-" * len(title) + "\n"
        
        if key == 'summary' and isinstance(section_data, str):
            text += section_data + "\n\n"

        elif key == 'skills' and isinstance(section_data, dict):
            for skill_type, skill_list in section_data.items():
                if skill_list:
                    text += f"{format_section_title(skill_type)}: "
                    text += ", ".join(skill_list) + "\n"
            text += "\n"
        
        elif isinstance(section_data, list):
            for item in section_data:
                if not isinstance(item, dict):
                    text += f" - {item}\n"
                    continue
                
                title_keys = ['title', 'name', 'degree', 'position', 'course']
                subtitle_keys = ['company', 'institution', 'issuer', 'university']
                duration_keys = ['duration', 'date']
                
                main_title = next((item[k] for k in title_keys if k in item and item[k]), '')
                subtitle = next((item[k] for k in subtitle_keys if k in item and item[k] != main_title and item[k]), '')
                duration = next((item[k] for k in duration_keys if k in item and item[k]), '')

                line = f"{main_title}"
                if subtitle:
                    line += f" ({subtitle})"
                if duration:
                    line += f" | {duration}"
                text += line + "\n"
                
                description_list = item.get('description') or item.get('achievement') or item.get('details')
                if description_list and isinstance(description_list, list):
                    for line in description_list:
                        text += f" - {line}\n"
            text += "\n"

    # ========== ADD CUSTOM SECTIONS HERE ==========
    standard_keys = get_standard_keys()
    
    for key, value in data.items():
        if key not in standard_keys and isinstance(value, str) and value.strip():
            title = format_section_title(key).upper()
            text += f"{title}\n"
            text += "-" * len(title) + "\n"
            text += value.strip() + "\n\n"

    return text

def get_text_download_link(data, filename_suffix=""):
    """Generates a download link for a plain text file."""
    text_content = generate_markdown_text(data)
    b64_data = base64.b64encode(text_content.encode('utf-8')).decode()
    
    filename = f"Resume_{data.get('name', 'User').replace(' ', '_')}_ATS_Plain_Text{filename_suffix}.txt"
    href = f'<a href="data:text/plain;base64,{b64_data}" download="{filename}" style="font-size: 0.95em; text-decoration: none; padding: 10px 15px; background-color: #40E0D0; color: white; border-radius: 5px; display: inline-block; margin-top: 10px; width: 100%; text-align: center;"><strong> Plain Text (.txt)</strong></a>'
    return href




# utils.py
import zipfile
import re
import json
import requests
from io import BytesIO
# import aspose.words as aw


# ---------------------------
# 1. Extract all visible text
# ---------------------------
def extract_all_text(docx_bytes):
    """
    Extract visible text from word/document.xml.
    This is only for AI mapping — not for replacement.
    """
    with zipfile.ZipFile(BytesIO(docx_bytes), 'r') as zf:
        xml = zf.read('word/document.xml').decode('utf-8')

    pattern = r'<w:t[^>]*>(.*?)</w:t>'
    matches = re.findall(pattern, xml, flags=re.S)

    result = []
    for text in matches:
        clean = text.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">").strip()
        if clean:
            result.append(clean)

    return result


def identify_headings(extracted_texts):
    """
    Identify which paragraphs are section headings vs content.
    Returns tuple: (headings_list, content_list)
    """
    HEADING_KEYWORDS = {
        "EXPERIENCE", "EDUCATION", "SKILLS", "PROJECTS", "CERTIFICATIONS",
        "SUMMARY", "OBJECTIVE", "PROFILE", "ACHIEVEMENTS", "AWARDS",
        "WORK EXPERIENCE", "PROFESSIONAL EXPERIENCE", "TECHNICAL SKILLS",
        "SOFT SKILLS", "PROFESSIONAL SUMMARY", "CAREER OBJECTIVE",
        "CONTACT", "ABOUT", "LANGUAGES", "INTERESTS", "REFERENCES"
    }
    
    headings = []
    content = []
    
    for para in extracted_texts:
        para_upper = para.strip().upper()
        
        # Check if it's a heading
        is_heading = (
            # Exact match with common headings
            para_upper in HEADING_KEYWORDS or
            # Short and all uppercase
            (len(para) < 50 and para.strip().isupper()) or
            # Contains heading keyword and is relatively short
            (len(para) < 50 and any(kw in para_upper for kw in HEADING_KEYWORDS))
        )
        
        if is_heading:
            headings.append(para)
        else:
            content.append(para)
    
    return headings, content


def ask_ai_for_mapping(extracted_texts, resume_data):
    """
    Creates direct paragraph-to-paragraph mapping from template to new resume data.
    Section headings are preserved unchanged, only content is mapped.
    """
    
    # Identify headings vs content
    headings, content_paragraphs = identify_headings(extracted_texts)
    
    prompt = f"""
You are creating a direct mapping from old resume content paragraphs to new resume data.

IMPORTANT: Section headings have been filtered out. You are ONLY mapping content paragraphs, NOT headings.

### TEMPLATE CONTENT PARAGRAPHS (exact keys to use):
{json.dumps(content_paragraphs, indent=2)}

### NEW RESUME DATA:
{json.dumps(resume_data, indent=2)}

### MAPPING INSTRUCTIONS:

1. **Name paragraph**: 
   - Map to: resume_data["name"]
   - Example: "Andree Rocher": "Rahul Menon"

2. **Contact info paragraph**:
   - Map to: phone + " | " + email
   - Example: "Philadelphia, PA | 705.555.0121 | andree@example.com": "+91-9812345678 | rahul@example.com"

3. **Job title paragraph**:
   - Map to: resume_data["job_title"]
   - Example: "Data Scientist": "Robotics Engineer"

4. **Summary/Objective content paragraph**:
   - Map to: resume_data["summary"]
   - This is the paragraph text UNDER the "Summary" heading
   - Example: "To obtain a challenging position...": "Results-driven Robotics Engineer..."

5. **Experience company/position paragraph**:
   - Format: resume_data["experience"][0]["company"] + " | " + resume_data["experience"][0]["position"]
   - Example: "FlueroGen | Data Scientist 20XX – 20XX": "Quartrdesign.com | Quality Assurance Tester"
   
6. **Experience bullet points**:
   - If template has MULTIPLE separate bullet paragraphs, DISTRIBUTE the description bullets across them
   - Each template bullet should get ONE description bullet (one-to-one mapping)
   - DO NOT combine all descriptions into one paragraph
   - Example with 3 template bullets:
     * Template bullet 1: "Increased retention by 20%" → "Spearheaded the creation of detailed Functional Requirement Documents (FRDs) from customer requirements, ensuring 95% accuracy in implementation and setting the standard for future projects."
     * Template bullet 2: "Developed new features" → "Collaborated with cross-functional teams to guarantee successful completion of requirements, resulting in a 25% reduction in project timelines and improved overall efficiency."
     * Template bullet 3: "Improved performance" → "Developed and executed comprehensive test cases using HTML, CSS, JavaScript, and React, achieving 99% test coverage and identifying critical defects that significantly improved product quality."
   - If template has only ONE bullet paragraph, put ALL descriptions there joined with \\n: "• Desc1\\n• Desc2\\n• Desc3"
   - IMPORTANT: Match the number of template bullets - don't put everything in the first bullet

7. **Skills content paragraph**:
   - If there are MULTIPLE skill paragraphs in template, DISTRIBUTE skills across them
   - Each paragraph should get DIFFERENT skills (no repetition)
   - Example with 3 skill bullets:
     * First bullet: "Data Modeling\\nData Warehousing\\nData Validation\\nData Cleansing"
     * Second bullet: "Azure Data Factory (ADF)\\nSSIS\\nDBT\\nPower BI\\nTableau"
     * Third bullet: "Communication\\nInterpersonal Skills\\nOrganizational Skills\\nSQL"
   - If template has only ONE skill paragraph, put ALL skills there
   - NEVER repeat the same skill in multiple paragraphs

8. **Project name**:
   - Map to: resume_data["projects"][0]["name"]
   
9. **Project descriptions**:
   - Map to: ALL project descriptions joined with \\n
   - Format: "• Desc 1\\n• Desc 2"

10. **Certifications**:
    - Format: "Name - Issuer" for each, joined with \\n
    - Example: "Database Management System - IIT Madras\\nGoogle Analytics - Google"

11. **Education**:
    - If resume_data["education"] is empty: map to ""
    - Otherwise: format education details

12. **No matching data**:
    - Map to: ""

### DATA STRUCTURE REFERENCE:
```
experience[0] = {{
  "company": "Quartrdesign.com",
  "position": "Quality Assurance Tester",
  "duration": "",
  "description": ["bullet 1", "bullet 2", ...]
}}

skills = {{
  "technicalSkills": ["Data Modeling", "Data Warehousing", ...],
  "tools": ["Azure Data Factory (ADF)", "SSIS", "DBT", ...],
  "cloudSkills": ["Azure", "Azure Storage", ...],
  "softSkills": ["Communication", "Interpersonal Skills", ...],
  "languages": ["SQL"]
}}

To map skills: Combine ALL arrays from skills dict into one list with \\n
```

### CRITICAL RULES:
- Keys MUST exactly match template paragraphs (character-for-character)
- Use \\n to join multi-line content (NOT actual newlines)
- Add bullet points (•) before each description item
- Combine all related data together

### OUTPUT (JSON only, no markdown):
{{
  "Andree Rocher": "Rahul Menon",
  "Data Scientist": "Robotics Engineer",
  "Philadelphia, PA | 705.555.0121 | andree@example.com": "+91-9812345678 | rahul@example.com",
  "To obtain a challenging position...": "Results-driven Robotics Engineer with expertise...",
  "FlueroGen | Data Scientist 20XX": "Quartrdesign.com | Quality Assurance Tester",
  "Increased retention by 20%": "Spearheaded the creation of detailed Functional Requirement Documents (FRDs) from customer requirements, ensuring 95% accuracy in implementation and setting the standard for future projects.",
  "Developed new features": "Collaborated with cross-functional teams to guarantee successful completion of requirements, resulting in a 25% reduction in project timelines and improved overall efficiency.",
  "Improved performance": "Developed and executed comprehensive test cases using HTML, CSS, JavaScript, and React, achieving 99% test coverage and identifying critical defects that significantly improved product quality.",
  "Data Modeling": "Data Modeling\\nData Warehousing\\nData Validation\\nData Cleansing\\nData Optimization",
  "Problem solving": "Unit Testing\\nDebugging\\nProblem Solving\\nBusiness Intelligence\\nData Management",
  "Communication": "Communication\\nInterpersonal Skills\\nOrganizational Skills\\nSQL"
}}

CRITICAL: 
- For EXPERIENCE bullets: Each template bullet gets ONE description (distribute across bullets, don't combine)
- For SKILLS bullets: Distribute different skills to each bullet (no repetition)

NOW MAP ALL CONTENT PARAGRAPHS:
"""
    
    payload = {
        "model":"meta/llama-3.1-70b-instruct",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.0,
        "max_tokens": 2048,
        "top_p": 0.9
    }
    
    try:
        response = requests.post(API_URL, headers=headers, json=payload)
        response.raise_for_status()
        raw = response.json()["choices"][0]["message"]["content"]
        
        # Clean markdown if present
        raw_clean = raw.strip()
        if raw_clean.startswith("```"):
            raw_clean = re.sub(r'^```(?:json)?\n', '', raw_clean)
            raw_clean = re.sub(r'\n```$', '', raw_clean)
        
        # Extract JSON
        match = re.search(r'\{.*\}', raw_clean, re.DOTALL)
        if not match:
            raise ValueError(f"LLM returned non-JSON response: {raw[:200]}")
        
        content_mapping = json.loads(match.group(0))
        
        # Post-process: Ensure skills are properly mapped without repetition
        all_skills_text = format_all_skills(resume_data.get('skills', {}))
        all_skills_list = [skill.strip() for skill in all_skills_text.split('\n') if skill.strip()]
        
        # Track which skills have been used
        used_skills = set()
        skills_paragraphs = []
        
        # Identify all skills-related paragraphs
        for para in content_paragraphs:
            para_lower = para.lower()
            if any(kw in para_lower for kw in ['skill', 'management', 'technical', 'tools', 'language', 'abilities']):
                skills_paragraphs.append(para)
        
        # Build complete mapping: headings unchanged + content mapped
        complete_mapping = {}
        
        # Add all headings (map to themselves to keep unchanged)
        for heading in headings:
            complete_mapping[heading] = heading
        
        # Add content with AI mappings
        for para in content_paragraphs:
            mapped_value = content_mapping.get(para, "")
            
            # Check if this is a skills paragraph
            if para in skills_paragraphs:
                # If AI didn't map it or mapped poorly
                if mapped_value == "" or len(mapped_value) < 20:
                    # Distribute remaining skills
                    remaining_skills = [s for s in all_skills_list if s not in used_skills]
                    
                    if remaining_skills:
                        # Take a reasonable chunk (e.g., 5-10 skills per paragraph)
                        chunk_size = max(5, len(remaining_skills) // len(skills_paragraphs))
                        skills_chunk = remaining_skills[:chunk_size]
                        
                        mapped_value = '\n'.join(skills_chunk)
                        used_skills.update(skills_chunk)
                    else:
                        # All skills used, leave empty to remove
                        mapped_value = ""
                else:
                    # AI provided mapping, track which skills were used
                    for skill in all_skills_list:
                        if skill in mapped_value:
                            used_skills.add(skill)
            
            complete_mapping[para] = mapped_value
        
        return complete_mapping
        
    except requests.exceptions.RequestException as e:
        raise Exception(f"API request failed: {str(e)}")
    except json.JSONDecodeError as e:
        raise Exception(f"Failed to parse JSON response: {str(e)}\nRaw: {raw[:500]}")


def format_experience_bullets(descriptions):
    """Format experience descriptions with bullet points"""
    return "\n".join([f"• {desc}" for desc in descriptions])


def format_all_skills(skills_dict):
    """Combine all skill categories into one list (supports dict or list safely)"""

    all_skills = []

    # ✅ CASE 1: Proper dictionary format
    if isinstance(skills_dict, dict):
        for category, skills in skills_dict.items():
            if isinstance(skills, list):
                all_skills.extend(skills)

    # ✅ CASE 2: Direct list format (DOCX / user input case)
    elif isinstance(skills_dict, list):
        all_skills.extend(skills_dict)

    # ✅ Final clean output
    return "\n".join(str(skill) for skill in all_skills if skill)



def format_certifications(certs_list):
    """Format certifications as 'Name - Issuer'"""
    return "\n".join([f"{cert['name']} - {cert['issuer']}" for cert in certs_list])


def format_projects(projects_list):
    """Format project descriptions"""
    formatted = []
    for project in projects_list:
        formatted.append(project['name'])
        if 'description' in project:
            for desc in project['description']:
                formatted.append(f"• {desc}")
    return "\n".join(formatted)


def extract_temp_from_docx(uploaded_file):
    """Extract all text paragraphs from DOCX file"""
    file_bytes = uploaded_file.read()
    with zipfile.ZipFile(BytesIO(file_bytes)) as docx:
        xml = docx.read("word/document.xml")

    tree = etree.fromstring(xml)
    paragraphs = []
    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"

    for p in tree.iter(f"{ns}p"):
        text = "".join((t.text or "") for t in p.iter(f"{ns}t")).strip()
        if text:
            paragraphs.append(text)

    return paragraphs




def get_text_runs(tree):
    """Get all text run nodes from XML tree"""
    return [node for node in tree.iter() if node.tag.endswith("}t")]


def safe_replace(xml_content, old, new):
    """
    Safely replace text in DOCX XML without corrupting formatting.
    Uses run-by-run replacement to preserve document structure.
    """
    try:
        tree = etree.fromstring(xml_content)
    except:
        return xml_content, False

    runs = get_text_runs(tree)
    if not runs:
        return xml_content, False

    run_texts = [r.text or "" for r in runs]
    full = "".join(run_texts)
    old_lower = old.lower()

    # Find match position (case-insensitive)
    match_start = full.lower().find(old_lower)
    if match_start == -1:
        return xml_content, False

    match_end = match_start + len(old)

    # Replace text across runs
    pos = 0
    for i, r in enumerate(runs):
        t = run_texts[i]
        start_in_run = None
        end_in_run = None

        for j, ch in enumerate(t):
            if pos == match_start:
                start_in_run = j
            if pos == match_end - 1:
                end_in_run = j
            pos += 1

        if start_in_run is not None and end_in_run is not None:
            # Match contained in this run
            runs[i].text = t[:start_in_run] + new + t[end_in_run + 1:]
        elif start_in_run is not None:
            # Match starts here
            runs[i].text = t[:start_in_run] + new
        elif end_in_run is not None:
            # Match ends here
            runs[i].text = t[end_in_run + 1:]
        elif match_start < pos - len(t) and match_end > pos - len(t):
            # Run is inside match
            runs[i].text = ""

    return etree.tostring(tree, encoding="utf-8", xml_declaration=True), True


def remove_empty_bullets(xml_content):
    """
    Remove paragraphs that contain only bullet points with no actual text content.
    Preserves divider lines, horizontal rules, and other formatting elements.
    """
    try:
        tree = etree.fromstring(xml_content)
    except:
        return xml_content

    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    paragraphs_to_remove = []

    for p in tree.iter(f"{ns}p"):
        # Check if paragraph has numbering/bullet formatting
        has_bullet_formatting = False
        pPr = p.find(f"{ns}pPr")
        if pPr is not None:
            numPr = pPr.find(f"{ns}numPr")
            if numPr is not None:
                has_bullet_formatting = True
        
        # Get text content
        text = "".join((t.text or "") for t in p.iter(f"{ns}t")).strip()
        
        # Only remove if:
        # 1. Has bullet formatting AND
        # 2. Is empty OR contains only bullet character
        if has_bullet_formatting and (not text or text in ['•', '-', '*', '▪', '◦', '■', '●', '○', '–', '—']):
            paragraphs_to_remove.append(p)
        # Also remove standalone bullet characters without formatting (edge case)
        elif text in ['•', '-', '*', '▪', '◦', '■', '●', '○'] and len(text) == 1:
            paragraphs_to_remove.append(p)

    # Remove empty bullet paragraphs
    for p in paragraphs_to_remove:
        parent = p.getparent()
        if parent is not None:
            parent.remove(p)

    return etree.tostring(tree, encoding="utf-8", xml_declaration=True)


def auto_process_docx(file, mapping):
    """
    Process DOCX file by applying paragraph mapping and removing empty bullets.
    
    Args:
        file: Uploaded DOCX file
        mapping: Dictionary mapping old paragraphs to new content
    
    Returns:
        BytesIO object containing the updated DOCX file
    """
    out = BytesIO()
    
    with zipfile.ZipFile(file, "r") as src:
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
            for item in src.infolist():
                data = src.read(item.filename)

                if item.filename == "word/document.xml":
                    xml = data
                    
                    # Apply all mappings
                    for old_text, new_text in mapping.items():
                        if new_text.strip() == "":
                            # Remove text completely
                            xml, _ = safe_replace(xml, old_text, "")
                        else:
                            # Replace with new content
                            xml, _ = safe_replace(xml, old_text, new_text)
                    
                    # Clean up empty bullet points
                    xml = remove_empty_bullets(xml)
                    
                    dst.writestr(item.filename, xml)
                else:
                    # Copy other files unchanged
                    dst.writestr(item.filename, data)

    out.seek(0)
    return out


def extract_docx_xml(uploaded_file):
    """Extract raw XML (word/document.xml) from DOCX"""
    uploaded_file.seek(0)
    file_bytes = uploaded_file.read()

    with zipfile.ZipFile(BytesIO(file_bytes)) as docx:
        xml = docx.read("word/document.xml")

    return xml  # ✅ This is BYTES


def docx_to_html_preview(docx_bytes):
    """
    Convert DOCX BytesIO to HTML for preview with enhanced style preservation
    
    Args:
        docx_bytes: BytesIO object containing DOCX file
    
    Returns:
        HTML string for preview
    """
    try:
        docx_bytes.seek(0)
        
        # Use mammoth with style mapping to preserve more formatting
        style_map = """
        p[style-name='Heading 1'] => h1:fresh
        p[style-name='Heading 2'] => h2:fresh
        p[style-name='Heading 3'] => h3:fresh
        p[style-name='Title'] => h1.title:fresh
        p[style-name='Subtitle'] => p.subtitle:fresh
        """
        
        result = mammoth.convert_to_html(
            docx_bytes,
            style_map=style_map,
            include_default_style_map=True,
            include_embedded_style_map=True
        )
        html = result.value
        
        # Add base styling that adapts to the document
        styled_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <style>
                * {{
                    margin: 0;
                    padding: 0;
                    box-sizing: border-box;
                }}
                
                body {{
                    margin: 0;
                    padding: 0;
                    background: #f5f5f5;
                }}
                
                .preview-container {{
                    font-family: 'Calibri', 'Segoe UI', 'Arial', sans-serif;
                    line-height: 1.5;
                    padding: 60px 70px;
                    background: #ffffff;
                    max-width: 8.5in;
                    margin: 0 auto;
                    min-height: 11in;
                    box-shadow: 0 0 10px rgba(0,0,0,0.1);
                }}
                
                /* Reset default margins */
                .preview-container > * {{
                    margin-block-start: 0;
                    margin-block-end: 0;
                }}
                
                /* Headers */
                .preview-container h1 {{
                    font-size: 14pt;
                    font-weight: bold;
                    margin-top: 18px;
                    margin-bottom: 10px;
                    padding-bottom: 5px;
                }}
                
                .preview-container h2 {{
                    font-size: 12pt;
                    font-weight: bold;
                    margin-top: 12px;
                    margin-bottom: 8px;
                }}
                
                .preview-container h3 {{
                    font-size: 11pt;
                    font-weight: bold;
                    margin-top: 10px;
                    margin-bottom: 6px;
                }}
                
                /* Paragraphs */
                .preview-container p {{
                    font-size: 11pt;
                    line-height: 1.5;
                    margin-top: 6px;
                    margin-bottom: 6px;
                }}
                
                /* Lists */
                .preview-container ul,
                .preview-container ol {{
                    margin-left: 25px;
                    margin-top: 8px;
                    margin-bottom: 8px;
                }}
                
                .preview-container li {{
                    font-size: 11pt;
                    line-height: 1.5;
                    margin-top: 4px;
                    margin-bottom: 4px;
                }}
                
                /* Tables */
                .preview-container table {{
                    border-collapse: collapse;
                    width: 100%;
                    margin: 10px 0;
                }}
                
                .preview-container td,
                .preview-container th {{
                    padding: 8px;
                    border: 1px solid #ddd;
                }}
                
                /* Strong and emphasis */
                .preview-container strong,
                .preview-container b {{
                    font-weight: bold;
                }}
                
                .preview-container em,
                .preview-container i {{
                    font-style: italic;
                }}
                
                /* Links */
                .preview-container a {{
                    color: #0563C1;
                    text-decoration: underline;
                }}
                
                /* Preserve any inline styles from mammoth */
                .preview-container [style] {{
                    /* Inline styles will be preserved */
                }}
                
                /* Scrollbar */
                .preview-container::-webkit-scrollbar {{
                    width: 10px;
                }}
                
                .preview-container::-webkit-scrollbar-track {{
                    background: #f5f5f5;
                }}
                
                .preview-container::-webkit-scrollbar-thumb {{
                    background: #cccccc;
                    border-radius: 5px;
                }}
            </style>
        </head>
        <body>
            <div class="preview-container">
                {html}
            </div>
        </body>
        </html>
        """
        
        return styled_html
    except Exception as e:
        return f"""
        <!DOCTYPE html>
        <html>
        <body>
            <div style="padding: 20px; color: red; font-family: Arial;">
                <h3>Preview Generation Failed</h3>
                <p>{str(e)}</p>
                <p>The document will still be available for download.</p>
            </div>
        </body>
        </html>
        """
    



def generate_two_column_html(data):
    if not data:
        return ""

    # Normalize keys (Interest → interest, CERTIFICATIONS → certifications)
    data = {k.lower(): v for k, v in data.items()}

    # Helper to detect empty values
    def is_empty(v):
        if v is None:
            return True
        if isinstance(v, str) and not v.strip():
            return True
        if isinstance(v, list) and all(is_empty(x) for x in v):
            return True
        if isinstance(v, dict) and all(is_empty(x) for x in v.values()):
            return True
        return False

    # Section titles mapping
    SECTION_TITLES = {
        "summary": "Summary",
        "objective": "Summary",
        "skills": "Skills",
        "experience": "Experience",
        "education": "Education",
        "projects": "Projects",
        "project": "Projects",
        "certifications": "Certifications",
        "achievements": "Achievements",
        "awards": "Awards",
        "interest": "Interests",
        "interests": "Interests"
    }

    # LEFT COLUMN SECTIONS
    LEFT_SECTIONS = ["summary", "objective", "skills"]

    RIGHT_SECTIONS = [
        key for key in SECTION_TITLES.keys()
        if key not in LEFT_SECTIONS
    ]


    # Auto title/subtitle/duration detection
    def extract_title_block(item):
        title = (
            item.get("position") or item.get("role") or item.get("title") or
            item.get("projectname") or item.get("name") or
            item.get("course") or item.get("certificate_name")
        )

        subtitle = (
            item.get("company") or item.get("organization") or
            item.get("university") or item.get("school") or
            item.get("tools") or item.get("provider_name")
        )

        duration = (
            item.get("duration") or item.get("start_date") or
            item.get("completed_date") or item.get("end_date") or
            item.get("date") or item.get("year")
        )

        return title, subtitle, duration

    # ======================================================
    # HEADER
    # ======================================================

    html = '<div class="ats-two-col-container"><div class="ats-two-col-header">'

    html += f'<h1>{data.get("name", "")}</h1>'

    if data.get("job_title"):
        html += f'<div class="ats-two-col-job-title">{data["job_title"]}</div>'

    contacts = [data.get("phone"), data.get("email"), data.get("url"), data.get("location")]
    contacts = [c for c in contacts if c]

    if contacts:
        html += f'<div class="ats-two-col-contact-header">{" | ".join(contacts)}</div>'

    html += "</div>"  # End header

    # ======================================================
    # BODY (Two Columns)
    # ======================================================

    html += '<div class="ats-two-col-body">'

    # ======================================================
    # LEFT COLUMN
    # ======================================================

    html += '<div class="ats-left-column">'

    for key in LEFT_SECTIONS:
        if key not in data or is_empty(data[key]):
            continue

        section_title = SECTION_TITLES[key]
        value = data[key]

        html += f'<div class="ats-two-col-section-title">{section_title}</div>'

        # -------- SUMMARY / OBJECTIVE --------
        if isinstance(value, str):
            html += f'<div class="ats-objective-text">{value}</div>'
            continue

        # -------- SKILLS --------
        if key == "skills":
            html += '<ul class="ats-skills-list">'
            skills = value

            if isinstance(skills, dict):
                for _, items in skills.items():
                    if isinstance(items, list):
                        for s in items:
                            html += f"<li>{s}</li>"
                    elif isinstance(items, str):
                        for s in items.split(","):
                            html += f"<li>{s.strip()}</li>"

            elif isinstance(skills, list):
                for s in skills:
                    html += f"<li>{s}</li>"

            elif isinstance(skills, str):
                for s in skills.split(","):
                    html += f"<li>{s.strip()}</li>"

            html += "</ul>"

    html += "</div>"  # END LEFT COLUMN

    # ======================================================
    # RIGHT COLUMN
    # ======================================================

    html += '<div class="ats-right-column">'

    for key in RIGHT_SECTIONS:
        if key not in data or is_empty(data[key]):
            continue

        section_title = SECTION_TITLES[key]
        value = data[key]

        html += f'<div class="ats-two-col-section-title">{section_title}</div>'

        # -------- STRING SECTION --------
        if isinstance(value, str):
            html += f"<p>{value}</p>"
            continue

        # -------- DICTIONARY SECTION (Skills-like structure) --------
        if isinstance(value, dict):
            for subkey, subval in value.items():
                html += f"<strong>{subkey.replace('_', ' ').title()}:</strong> "
                if isinstance(subval, list):
                    html += ", ".join(subval)
                else:
                    html += str(subval)
                html += "<br>"
            continue

        # -------- LIST SECTION --------
        if isinstance(value, list):

            # plain list of strings
            if all(isinstance(x, str) for x in value):
                html += "<ul>"
                for item in value:
                    html += f"<li>{item}</li>"
                html += "</ul>"
                continue

            # list of objects (Experience, Education, Projects, Certs)
            for item in value:
                if not isinstance(item, dict):
                    continue

                title, subtitle, duration = extract_title_block(item)

                # Item header
                html += '<div class="ats-two-col-item-header">'
                html += '<div class="ats-two-col-item-title-row">'

                if title:
                    html += f'<span class="ats-two-col-item-title">{title}'
                    if subtitle:
                        html += f' <span class="ats-two-col-item-subtitle">/ {subtitle}</span>'
                    html += '</span>'

                if duration:
                    html += f'<span class="ats-two-col-item-duration">{duration}</span>'

                html += '</div></div>'

                # Description
                desc = (
                    item.get("description") or item.get("details") or
                    item.get("overview") or item.get("achievement")
                )

                if isinstance(desc, list):
                    html += '<ul class="ats-two-col-bullet-list">'
                    for d in desc:
                        html += f"<li>{d}</li>"
                    html += "</ul>"
                elif isinstance(desc, str):
                    html += f"<p>{desc}</p>"

    html += "</div>"  # END RIGHT COLUMN

    html += "</div></div>"  # END BODY + CONTAINER

    return html


def format_year_only(date_str):
    """Extract only the year from a date string."""
    if not date_str:
        return ""
    date_str = str(date_str).strip()
    if not date_str or date_str.lower() in ['present', 'current', 'now']:
        return 'Present' if date_str.lower() in ['present', 'current', 'now'] else date_str
    
    # Try to extract year (4 digits)
    import re
    year_match = re.search(r'\b(19|20)\d{2}\b', date_str)
    if year_match:
        return year_match.group(0)
    return date_str


def get_standard_keys():
    """Returns set of standard resume keys."""
    return {
        'name', 'email', 'phone', 'location', 'linkedin', 'website', 'github',
        'job_title', 'summary', 'objective', 'skills', 'experience', 'education',
        'certifications', 'projects', 'awards', 'achievements', 'publications', 
        'languages', 'volunteer', 'interests', 'references', 'input_method'
    }


def format_section_title(key):
    """Format section key into readable title."""
    return key.replace('_', ' ').title()





def save_custom_sections():
    """Save custom section edits to session state."""
    data = st.session_state['enhanced_resume']
    standard_keys = get_standard_keys()
    
    # Update custom sections from session state (for edit mode)
    for key in list(data.keys()):
        if key not in standard_keys and isinstance(data.get(key), str):
            edit_key = f"edit_custom_{key}"
            if edit_key in st.session_state:
                data[key] = st.session_state[edit_key].strip()
    
    # Ensure all custom sections are preserved in the data
    st.session_state['enhanced_resume'] = data


def generate_and_switch():
    """Performs final analysis and switches to download page."""
    # Save any pending custom section edits
    save_custom_sections()
    
    data = st.session_state['enhanced_resume']
    
    # Extract custom sections before analysis
    standard_keys = get_standard_keys()
    custom_sections = {k: v for k, v in data.items() 
                      if k not in standard_keys and isinstance(v, str)}
    
    # Perform analysis
    finalized_data = analyze_and_improve_resume(data)
    
    # Re-add custom sections to finalized data
    for key, value in custom_sections.items():
        if key not in finalized_data:
            finalized_data[key] = value

    default_template = "Minimalist (ATS Best)"

    st.session_state.selected_template = default_template
    st.session_state.template_source = 'saved'
    st.session_state['final_resume_data'] = finalized_data

    # Set default template config
    from templates.templateconfig import SYSTEM_TEMPLATES  # make sure this import exists
    st.session_state.selected_template_config = SYSTEM_TEMPLATES.get(default_template)

    st.switch_page("pages/template_preview.py")




    # Add this import at the top of your file
from spire.doc import *
from spire.doc.common import *
from docx import Document as DocxDocument
from docx.shared import Inches
import tempfile
import os
import io

def convert_html_to_docx_spire(html_content, css_content=""):
    """Convert HTML to DOCX using Spire.Doc with watermark removal and margin adjustment."""
    temp_html_path = None
    temp_docx_path = None
    
    try:
        # Create temporary HTML file
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as temp_html:
            # Write complete HTML with CSS - reduced margins
            complete_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <style>
                    @page {{ size: A4; margin: 0.3in; }}
                    * {{ -webkit-print-color-adjust: exact !important; print-color-adjust: exact !important; }}
                    body {{ margin: 0; padding: 0; }}
                    {css_content}
                </style>
            </head>
            <body style="margin: 0; padding: 0; background: #fff;">
                <div style="width: 100%; margin: 0; padding: 0.3in;">
                    {html_content}
                </div>
            </body>
            </html>
            """
            temp_html.write(complete_html)
            temp_html_path = temp_html.name
        
        # Create temporary DOCX file path
        temp_docx_path = tempfile.mktemp(suffix='.docx')
        
        # Convert HTML to DOCX using Spire.Doc
        document = Document()
        document.LoadFromFile(temp_html_path, FileFormat.Html, XHTMLValidationType.none)
        
        # Remove watermark
        try:
            document.Watermark.Type = WatermarkType.NoWatermark
        except:
            pass
        
        # Adjust margins using index-based access instead of iteration
        try:
            section_count = document.Sections.Count
            for i in range(section_count):
                section = document.Sections.get_Item(i)
                section.PageSetup.Margins.Left = 36.0   # 0.5 inch
                section.PageSetup.Margins.Right = 36.0  # 0.5 inch
                section.PageSetup.Margins.Top = 36.0    # 0.5 inch
                section.PageSetup.Margins.Bottom = 36.0 # 0.5 inch
        except Exception as margin_error:
            print(f"Could not adjust margins with Spire: {margin_error}")
        
        document.SaveToFile(temp_docx_path, FileFormat.Docx2016)
        document.Close()
        
        # ========== POST-PROCESS TO REMOVE WATERMARK TEXT ==========
        # Load the generated DOCX with python-docx to remove any remaining watermark text
        doc = DocxDocument(temp_docx_path)
        
        # Adjust margins using python-docx (more reliable)
        for section in doc.sections:
            section.left_margin = Inches(0.5)
            section.right_margin = Inches(0.5)
            section.top_margin = Inches(0.5)
            section.bottom_margin = Inches(0.5)
        
        # Remove paragraphs containing the watermark text
        watermark_texts = [
            "Evaluation Warning: The document was created with Spire.Doc for Python",
            "Evaluation Warning",
            "Spire.Doc for Python",
            "created with Spire",
            "This message will not appear"
        ]
        
        # Check and remove from paragraphs
        paragraphs_to_remove = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                for watermark in watermark_texts:
                    if watermark.lower() in paragraph.text.lower():
                        paragraphs_to_remove.append(paragraph)
                        break
        
        for paragraph in paragraphs_to_remove:
            p_element = paragraph._element
            p_element.getparent().remove(p_element)
        
        # Check and remove from headers/footers
        for section in doc.sections:
            # Check header
            try:
                for paragraph in section.header.paragraphs[:]:  # Create a copy of the list
                    for watermark in watermark_texts:
                        if watermark.lower() in paragraph.text.lower():
                            p_element = paragraph._element
                            p_element.getparent().remove(p_element)
                            break
            except:
                pass
            
            # Check footer
            try:
                for paragraph in section.footer.paragraphs[:]:  # Create a copy of the list
                    for watermark in watermark_texts:
                        if watermark.lower() in paragraph.text.lower():
                            p_element = paragraph._element
                            p_element.getparent().remove(p_element)
                            break
            except:
                pass
        
        # Save the cleaned document to BytesIO
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        docx_data = output.read()
        
        return docx_data
        
    except Exception as e:
        st.error(f"Error converting HTML to DOCX: {str(e)}")
        import traceback
        st.error(traceback.format_exc())
        return None
        
    finally:
        # Clean up temporary files
        if temp_html_path and os.path.exists(temp_html_path):
            try:
                os.unlink(temp_html_path)
            except Exception as e:
                print(f"Warning: Could not delete temporary HTML file: {e}")
        
        if temp_docx_path and os.path.exists(temp_docx_path):
            try:
                os.unlink(temp_docx_path)
            except Exception as e:
                print(f"Warning: Could not delete temporary DOCX file: {e}")