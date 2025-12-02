import streamlit as st
import base64
from datetime import datetime
import json
from templates.templateconfig import SYSTEM_TEMPLATES,ATS_COLORS,load_css_template
from utils import(get_user_resume,generate_markdown_text,  chatbot,
    generate_generic_html,generate_markdown_text,save_user_doc_templates,
    load_user_templates,load_user_doc_templates,save_user_templates,replace_content
    ,load_user_ppt_templates,analyze_slide_structure,generate_ppt_sections,
    match_generated_to_original,clear_and_replace_text,save_user_ppt_templates)

# ----------------------------------
# PAGE CONFIG
# ----------------------------------
st.set_page_config(
    page_title="Template Preview",
    page_icon="👁️",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# ----------------------------------
# RESTORE USER FROM QUERY PARAMS
# ----------------------------------


# ----------------------------------
# RESTORE USER FROM QUERY PARAMS
# ----------------------------------
if "logged_in_user" not in st.session_state or st.session_state.logged_in_user is None:
    logged_user = st.query_params.get("user")
    if logged_user:
        st.session_state.logged_in_user = logged_user
    else:
        st.warning("Please login first!")
        st.switch_page("app.py")

if st.session_state.logged_in_user:
    st.query_params["user"] = st.session_state.logged_in_user

current_user = st.session_state.get("logged_in_user", "")

# ----------------------------------
# LOAD USER UPLOADED TEMPLATES (ALWAYS)
# ----------------------------------
if "uploaded_templates" not in st.session_state:
    st.session_state.uploaded_templates = load_user_templates(current_user)

# ----------------------------------
# RESTORE TEMPLATE SELECTION FROM QUERY PARAMS
# ----------------------------------
DEFAULT_TEMPLATE = "Minimalist (ATS Best)"

if "selected_template" not in st.session_state or st.session_state.selected_template is None:
    template_name = st.query_params.get("template")
    template_source = st.query_params.get("source")

    if template_name:  # coming from URL
        st.session_state.selected_template = template_name
        st.session_state.template_source = template_source or "system"
    else:
        # First load → default
        st.session_state.selected_template = DEFAULT_TEMPLATE
        st.session_state.template_source = "system"

# ----------------------------------
# LOAD TEMPLATE CONFIG BASED ON SOURCE
# ----------------------------------
template_name = st.session_state.selected_template
template_source = st.session_state.template_source

if template_source == "system":
    st.session_state.selected_template_config = SYSTEM_TEMPLATES.get(template_name)

elif template_source == "user":
    uploaded = st.session_state.get("uploaded_templates", {})
    if uploaded and template_name in uploaded:
        st.session_state.selected_template_config = uploaded[template_name]

# SAFETY FALLBACK (Never allow missing config)
if not st.session_state.get("selected_template_config"):
    st.session_state.selected_template = DEFAULT_TEMPLATE
    st.session_state.template_source = "system"
    st.session_state.selected_template_config = SYSTEM_TEMPLATES[DEFAULT_TEMPLATE]

# SAVE TO URL (PERSIST ON REFRESH)
st.query_params["template"] = st.session_state.selected_template
st.query_params["source"] = st.session_state.template_source

# ----------------------------------
# RESTORE FINAL RESUME IF EXISTS
# ----------------------------------
if "final_resume_data" not in st.session_state or st.session_state.final_resume_data is None:
    resume_param = st.query_params.get("final_resume")
    if resume_param:
        try:
            st.session_state.final_resume_data = json.loads(resume_param)
        except:
            st.session_state.final_resume_data = None

if st.session_state.get("final_resume_data"):
    try:
        st.query_params["final_resume"] = json.dumps(st.session_state.final_resume_data)
    except:
        pass

# ----------------------------------
# GET USER RESUME
# ----------------------------------
email = st.session_state.logged_in_user
user_resume = st.session_state.get("final_resume_data") or get_user_resume(email)

if not user_resume:
    st.warning("No resume found. Please create one.")
    if st.button("Create Resume"):
        st.query_params["user"] = email
        st.switch_page("pages/main.py")
    st.stop()

chatbot(user_resume)

# ----------------------------------
# GET TEMPLATE SELECTION VALUES
# ----------------------------------
selected_template = st.session_state.selected_template
selected_config = st.session_state.selected_template_config
template_source = st.session_state.template_source

# ----------------------------------
# CUSTOM CSS
# ----------------------------------
# Replace the CUSTOM CSS section (starting from line 106) with this improved version:

st.markdown("""
<style>
/* CRITICAL: Prevent ALL scrolling on the page */
html, body {
    overflow: hidden !important;
    height: 100vh !important;
    width: 100vw !important;
    position: fixed !important;
    margin: 0 !important;
    padding: 0 !important;
    overscroll-behavior: none !important;
}

/* Prevent touch scrolling on mobile */
html {
    touch-action: none !important;
    -ms-touch-action: none !important;
}

/* Hide all scrollbars globally */
* {
    scrollbar-width: none !important;
    -ms-overflow-style: none !important;
}

*::-webkit-scrollbar {
    display: none !important;
}

/* Prevent scroll on all Streamlit containers */
#root, [data-testid="stAppViewContainer"], .main, .block-container {
    overflow: hidden !important;
    position: fixed !important;
    height: 100vh !important;
    width: 100% !important;
    overscroll-behavior: none !important;
}

#MainMenu, footer, header, button[kind="header"] {visibility: hidden;}

/* Streamlit app container - lock it */
.stApp {
    overflow: hidden !important;
    height: 100vh !important;
    width: 100vw !important;
    position: fixed !important;
    top: 0 !important;
    left: 0 !important;
    overscroll-behavior: none !important;
}

section[data-testid="stAppViewContainer"] {
    overflow: hidden !important;
    height: 100vh !important;
    position: fixed !important;
    width: 100% !important;
}

section[data-testid="stAppViewContainer"] > div {
    overflow: hidden !important;
    height: 100vh !important;
}

.main {
    overflow: hidden !important;
    height: 100vh !important;
    position: fixed !important;
    width: 100% !important;
}

.main .block-container {
    overflow: hidden !important;
    height: 100vh !important;
    padding-top: 100px !important;
    padding-bottom: 40px !important;
    max-height: 100vh !important;
    position: relative !important;
}

:root {
    --primary-orange: #e87532;
    --primary-dark: #d66629;
    --light-bg: #fff9f5;
    --border-color: #ffe5d9;
    --panel-bg: #ffffff;
    --section-bg: #fffbf7;
}

.nav-wrapper {
    position: fixed;
    top: 20px;
    left: 50%;
    transform: translateX(-50%);
    width: 90%;
    max-width: 1200px;
    z-index: 99999 !important;
    background-color: white !important;
    padding: 0.8rem 2rem;
    box-shadow: 0 2px 20px rgba(0,0,0,0.1);
    display: flex;
    align-items: center;
    justify-content: space-between;
    border-radius: 50px;
}

.logo {
    font-size: 24px;
    font-weight: 400;
    color: #2c3e50;
    font-family: 'Nunito Sans', sans-serif !important;
    letter-spacing: -0.5px;
}

.nav-menu {
    display: flex;
    gap: 2rem;
    align-items: center;
}

.nav-item { position: relative; }

.nav-link {
    color: #000000 !important;
    text-decoration: none !important;
    font-size: 1rem;
    font-family: 'Nunito Sans', sans-serif;
    padding: 0.5rem 1rem;
    border-radius: 8px;
    cursor: pointer;
    transition: all 0.3s ease;
}

.nav-link:visited {
    color: #000000 !important;
}

.nav-link:hover {
    background-color: #fff5f0;
    color: #ff8c42 !important;  /* Added !important to override the default color */
}

/* Main Container - Increased bottom padding for white space */
.main-content {
    margin-top: 100px;
    display: grid;
    grid-template-columns: 280px 250px 1fr;
    gap: 1.5rem;
    padding: 0 2rem 3rem 2rem;
    height: calc(100vh - 180px) !important;
    max-height: calc(100vh - 180px) !important;
    overflow: hidden !important;
}

/* Horizontal block wrapper */
div[data-testid="stHorizontalBlock"] {
    overflow: hidden !important;
    height: calc(100vh - 180px) !important;
    max-height: calc(100vh - 180px) !important;
}

/* Vertical blocks - prevent overflow */
div[data-testid="stVerticalBlock"] {
    overflow: hidden !important;
}

/* Left column (col1) - ONLY this scrolls */
div[data-testid="column"]:first-child {
    height: calc(100vh - 200px) !important;
    max-height: calc(100vh - 200px) !important;
    overflow-y: auto !important;
    overflow-x: hidden !important;
    padding-right: 10px;
    scrollbar-width: thin !important;
    -ms-overflow-style: auto !important;
    overscroll-behavior: contain !important;
}

div[data-testid="column"]:first-child::-webkit-scrollbar {
    display: block !important;
    width: 8px !important;
    height: 8px !important;
}

/* Right column (col3 - preview) - ONLY this scrolls with MORE bottom space */
div[data-testid="column"]:last-child {
    height: calc(100vh - 200px) !important;
    max-height: calc(100vh - 200px) !important;
    overflow-y: auto !important;
    overflow-x: hidden !important;
    padding-right: 10px;
    padding-bottom: 40px !important;
    scrollbar-width: thin !important;
    -ms-overflow-style: auto !important;
    overscroll-behavior: contain !important;
}

div[data-testid="column"]:last-child::-webkit-scrollbar {
    display: block !important;
    width: 8px !important;
    height: 8px !important;
}

/* Scrollbar styling for columns */
div[data-testid="column"]::-webkit-scrollbar-track {
    background: #f1f5f9;
    border-radius: 10px;
}

div[data-testid="column"]::-webkit-scrollbar-thumb {
    background: var(--primary-orange);
    border-radius: 10px;
}

div[data-testid="column"]::-webkit-scrollbar-thumb:hover {
    background: var(--primary-dark);
}

/* Panel Styling */
.control-panel {
    background: var(--section-bg);
    border-radius: 12px;
    padding: 1.5rem;
    box-shadow: 0 2px 8px rgba(232, 117, 50, 0.1);
    height: fit-content;
    border: 1px solid var(--border-color);
}

.panel-header {
    font-size: 0.85rem;
    font-weight: 600;
    color: #666;
    text-transform: uppercase;
    letter-spacing: 0.5px;
    margin-bottom: 1rem;
    padding-bottom: 0.5rem;
    border-bottom: 2px solid var(--primary-orange);
}

/* Template List Container - WITH SCROLLING */
.template-list-container {
    background: white;
    border-radius: 8px;
    padding: 1rem;
    max-height: 350px !important;
    overflow-y: auto !important;
    overflow-x: hidden !important;
    scrollbar-width: thin !important;
    -ms-overflow-style: auto !important;
    overscroll-behavior: contain !important;
}

.template-list-container::-webkit-scrollbar {
    display: block !important;
    width: 6px !important;
    height: 6px !important;
}

.template-list-container::-webkit-scrollbar-track {
    background: #f1f5f9;
    border-radius: 10px;
}

.template-list-container::-webkit-scrollbar-thumb {
    background: var(--primary-orange);
    border-radius: 10px;
}

.template-list-container::-webkit-scrollbar-thumb:hover {
    background: var(--primary-dark);
}

/* Template List Header */
.template-list-header {
    display: flex;
    justify-content: space-between;
    align-items: center;
    padding: 0.5rem 0.8rem;
    margin-bottom: 0.5rem;
    font-size: 0.75rem;
    color: #666;
    border-bottom: 1px solid #e5e7eb;
}

/* Template Item */
.template-item {
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 0.8rem;
    margin-bottom: 0.5rem;
    background: white;
    border: 1px solid #e5e7eb;
    border-radius: 6px;
    transition: all 0.2s ease;
}

.template-item:hover {
    background: #f9fafb;
    border-color: var(--primary-orange);
    box-shadow: 0 2px 4px rgba(232, 117, 50, 0.1);
}

.template-item-left {
    display: flex;
    align-items: center;
    gap: 0.8rem;
    flex: 1;
}

.template-icon {
    width: 32px;
    height: 32px;
    display: flex;
    align-items: center;
    justify-content: center;
    background: #f3f4f6;
    border-radius: 4px;
    color: #6b7280;
    font-size: 1rem;
}

.template-info {
    flex: 1;
}

.template-name-text {
    font-size: 0.9rem;
    font-weight: 500;
    color: #1f2937;
    margin-bottom: 0.2rem;
}

.template-meta {
    font-size: 0.75rem;
    color: #9ca3af;
}

.template-actions {
    display: flex;
    gap: 0.5rem;
}

/* Custom Button Styles for Template Actions */
.template-use-btn {
    background: var(--primary-orange) !important;
    color: white !important;
    border: none !important;
    padding: 0.4rem 1rem !important;
    border-radius: 6px !important;
    font-size: 0.8rem !important;
    font-weight: 500 !important;
    cursor: pointer !important;
    transition: all 0.2s ease !important;
    min-width: 60px !important;
}

.template-use-btn:hover {
    background: var(--primary-dark) !important;
    transform: translateY(-1px);
}

.template-delete-btn {
    background: transparent !important;
    color: #6b7280 !important;
    border: 1px solid #d1d5db !important;
    padding: 0.4rem 0.6rem !important;
    border-radius: 6px !important;
    font-size: 0.8rem !important;
    cursor: pointer !important;
    transition: all 0.2s ease !important;
}

.template-delete-btn:hover {
    background: #fee2e2 !important;
    color: #dc2626 !important;
    border-color: #dc2626 !important;
}

/* Color Circles */
.color-palette {
    display: flex;
    flex-wrap: wrap;
    gap: 12px;
    margin-top: 15px;
}

.color-circle {
    width: 45px;
    height: 45px;
    border-radius: 50%;
    cursor: pointer;
    transition: all 0.3s ease;
    border: 3px solid transparent;
    box-shadow: 0 2px 8px rgba(0,0,0,0.15);
}

.color-circle:hover {
    transform: scale(1.15);
    box-shadow: 0 4px 12px rgba(0,0,0,0.25);
}

.color-circle.selected {
    border-color: #333;
    box-shadow: 0 0 0 3px rgba(232, 117, 50, 0.3);
    transform: scale(1.1);
}

/* Template Selection Box */
.template-box {
    background: white;
    border: 2px solid var(--primary-orange);
    border-radius: 8px;
    padding: 1rem;
    margin-bottom: 1rem;
    cursor: pointer;
    transition: all 0.3s ease;
}

.template-box:hover {
    transform: translateY(-2px);
    box-shadow: 0 4px 12px rgba(232, 117, 50, 0.2);
}

.template-name {
    font-weight: 600;
    color: var(--primary-orange);
    font-size: 0.95rem;
}

/* Preview Panel - with bottom spacing */
.preview-panel {
    background: var(--section-bg);
    border-radius: 12px;
    padding: 1.5rem;
    padding-bottom: 2.5rem;
    box-shadow: 0 2px 8px rgba(232, 117, 50, 0.1);
    height: 100%;
    max-height: calc(100vh - 200px);
    overflow: visible;
    border: 1px solid var(--border-color);
    margin-bottom: 20px;
}

.preview-header {
    font-size: 1.2rem;
    font-weight: 600;
    color: var(--primary-orange);
    margin-bottom: 1rem;
    padding-bottom: 0.8rem;
    border-bottom: 2px solid var(--primary-orange);
}

/* Streamlit components iframe container - adjust height with bottom space */
iframe {
    max-height: calc(100vh - 280px) !important;
    margin-bottom: 20px !important;
}

/* Streamlit Overrides */
.stButton > button {
    background: var(--primary-orange) !important;
    color: white !important;
    border: none !important;
    padding: 10px 20px !important;
    border-radius: 8px !important;
    font-weight: 600 !important;
    transition: all 0.3s ease !important;
    width: 100% !important;
}

.stButton > button:hover {
    background: var(--primary-dark) !important;
    transform: translateY(-2px);
    box-shadow: 0 4px 12px rgba(232, 117, 50, 0.3) !important;
}

.stDownloadButton > button {
    background: var(--primary-orange) !important;
    color: white !important;
    border: none !important;
    padding: 10px 20px !important;
    border-radius: 8px !important;
    font-weight: 600 !important;
    width: 100% !important;
}

.stDownloadButton > button:hover {
    background: var(--primary-dark) !important;
    transform: translateY(-2px);
}

.stSelectbox > div > div {
    border: 2px solid var(--border-color) !important;
    border-radius: 8px !important;
    background: white !important;
}

/* Hide default data editor toolbar */
.stDataFrameToolbar {
    display: none !important;
}

/* Tabs styling */
.stTabs [data-baseweb="tab-list"] {
    gap: 8px;
    background-color: transparent;
    border-bottom: 2px solid var(--border-color);
}

.stTabs [data-baseweb="tab"] {
    background-color: white;
    border: 1px solid var(--border-color);
    border-radius: 8px 8px 0 0;
    padding: 8px 16px;
    font-size: 0.85rem;
    font-weight: 500;
}

.stTabs [aria-selected="true"] {
    background-color: var(--primary-orange) !important;
    color: white !important;
    border-color: var(--primary-orange) !important;
}

/* Tab panels with scrolling and bottom space */
.stTabs [data-baseweb="tab-panel"] {
    background-color: var(--section-bg);
    padding: 1.5rem;
    padding-bottom: 2rem;
    border-radius: 0 0 8px 8px;
    border: 1px solid var(--border-color);
    border-top: none;
    max-height: calc(100vh - 320px) !important;
    overflow-y: auto !important;
    overflow-x: hidden !important;
    scrollbar-width: thin !important;
    -ms-overflow-style: auto !important;
    overscroll-behavior: contain !important;
}

.stTabs [data-baseweb="tab-panel"]::-webkit-scrollbar {
    display: block !important;
    width: 6px !important;
    height: 6px !important;
}

.stTabs [data-baseweb="tab-panel"]::-webkit-scrollbar-track {
    background: #f1f5f9;
    border-radius: 10px;
}

.stTabs [data-baseweb="tab-panel"]::-webkit-scrollbar-thumb {
    background: var(--primary-orange);
    border-radius: 10px;
}

.stTabs [data-baseweb="tab-panel"]::-webkit-scrollbar-thumb:hover {
    background: var(--primary-dark);
}

/* File uploader */
.stFileUploader {
    border: 2px dashed var(--border-color) !important;
    border-radius: 8px !important;
    padding: 1rem !important;
    background: white !important;
}

.stFileUploader:hover {
    border-color: var(--primary-orange) !important;
    background: var(--light-bg) !important;
}

.stFileUploader label {
    color: #666 !important;
    font-size: 0.85rem !important;
}

/* Remove extra spacing bars */
div[data-testid="stVerticalBlock"] > div:first-child {
    gap: 0.5rem;
}

/* Column container adjustments */
div[data-testid="column"] {
    background: transparent;
    padding: 0;
}

/* Hide horizontal rules */
hr {
    display: none;
}
</style>

<script>
// JavaScript to prevent scrolling on the main window
document.addEventListener('DOMContentLoaded', function() {
    // Prevent wheel scrolling on window
    window.addEventListener('wheel', function(e) {
        // Only prevent if not scrolling inside a column or template list
        if (!e.target.closest('[data-testid="column"]') && 
            !e.target.closest('.template-list-container') &&
            !e.target.closest('[data-baseweb="tab-panel"]')) {
            e.preventDefault();
            e.stopPropagation();
        }
    }, { passive: false });
    
    // Prevent arrow key scrolling
    window.addEventListener('keydown', function(e) {
        if(['ArrowUp', 'ArrowDown', 'PageUp', 'PageDown', 'Home', 'End', ' '].includes(e.key)) {
            if (!e.target.closest('[data-testid="column"]') && 
                !e.target.closest('.template-list-container') &&
                !e.target.closest('[data-baseweb="tab-panel"]')) {
                e.preventDefault();
            }
        }
    });
    
    // Prevent touch scrolling
    document.body.addEventListener('touchmove', function(e) {
        if (!e.target.closest('[data-testid="column"]') && 
            !e.target.closest('.template-list-container') &&
            !e.target.closest('[data-baseweb="tab-panel"]')) {
            e.preventDefault();
        }
    }, { passive: false });
});
</script>
""", unsafe_allow_html=True)

# ----------------------------------
# TOP NAVIGATION
# ----------------------------------
# Build navigation menu conditionally
nav_items = f"""
    <div class="nav-item">
        <a class="nav-link" href="?home=true&user={current_user}" target="_self">Home</a>
    </div>
    <div class="nav-item">
        <a class="nav-link" href="?create=true&user={current_user}" target="_self">Create New Resume</a>
    </div>
"""

# Add Edit Content link only if final_resume_data exists
if st.session_state.get("final_resume_data"):
    nav_items += f"""
    <div class="nav-item">
        <a class="nav-link" href="?edit=true&user={current_user}" target="_self">Edit Content</a>
    </div>
"""

nav_items += f"""
    <div class="nav-item">
        <a class="nav-link" href="?addjd=true&user={current_user}" target="_self">Add New JD</a>
    </div>
    <div class="nav-item">
        <a class="nav-link" href="?logout=true" target="_self">Logout</a>
    </div>
"""

st.markdown(f"""
    <div class="nav-wrapper">
        <div class="logo">Resume Creator</div>
        <div class="nav-menu">
            {nav_items}
       
    </div>
""", unsafe_allow_html=True)

    # Handle navigation - PRESERVE USER IN SESSION STATE
if st.query_params.get("addjd") == "true":
        st.query_params.clear()
        if current_user:
            st.query_params["user"] = current_user
        st.switch_page("pages/job.py")

if st.query_params.get("create") == "true":
        st.query_params.clear()
        if current_user:
            st.query_params["user"] = current_user
        st.switch_page("pages/main.py")

if st.query_params.get("home") == "true":
        st.query_params.clear()
        if current_user:
            st.query_params["user"] = current_user
        st.switch_page("app.py")

if st.query_params.get("logout") == "true":
        # Clear session state AND user info on logout
        for key in list(st.session_state.keys()):
            del st.session_state[key]
        st.query_params.clear()
        st.switch_page("app.py")


if st.query_params.get("edit") == "true":
    st.query_params.clear()
    if current_user:
        st.query_params["user"] = current_user
    # Pass final_resume_data through query params
    if st.session_state.get("final_resume_data"):
        try:
            st.query_params["final_resume"] = json.dumps(st.session_state.final_resume_data)
        except:
            pass
    st.switch_page("pages/create.py") 
# ----------------------------------
# MAIN LAYOUT - THREE COLUMNS
# ----------------------------------
col1, col3 = st.columns([1.8, 2.5])

with col1:
    tab1, tab2, tab3, tab4 = st.tabs(["System Templates", "Custom Templates", "Download","Add New Template"])
    with tab1:
        # if template_source == 'system':
            # st.markdown('<div class="control-panel">', unsafe_allow_html=True)
            st.markdown('<div class="panel-header">🎨 Colors</div>', unsafe_allow_html=True)
            
            # Color selection dropdown (like in download.py)
            color_name = st.selectbox(
                'Choose Accent Color:',
                list(ATS_COLORS.keys()),
                key='preview_color_select'
            )
            primary_color = ATS_COLORS[color_name]
            
            # Store in session state
            st.session_state.preview_selected_color = primary_color
            
            # Display the selected color
            st.markdown(f"""
            <div style="width: 50%; height: 30px; background: {primary_color}; 
                border-radius: 8px; margin-top: 10px;margin-left:50px;
                border: 3px solid #333;
                box-shadow: 0 4px 12px rgba(0,0,0,0.2);"></div>
            """, unsafe_allow_html=True)
            
            st.markdown('</div>', unsafe_allow_html=True)
            st.markdown('<div class="panel-header">📋 Template Selection</div>', unsafe_allow_html=True)
        
            # Current template display
            st.markdown(f"""
            <div class="template-box">
                <div class="template-name">{selected_template}</div>
                <div style="font-size: 0.8rem; color: #666; margin-top: 0.3rem;">
                    {template_source.replace('_', ' ').title()}
                </div>
            </div>
            """, unsafe_allow_html=True)
            
            system_template_names = list(SYSTEM_TEMPLATES.keys())
           # Find the current template's index
            try:
                current_index = system_template_names.index(st.session_state.selected_template)
            except (ValueError, AttributeError):
                current_index = 0

            selected_system_template = st.selectbox(
                "Select a System Template:",
                system_template_names,
                index=current_index,  # Add this line
                key="system_template_dropdown",
                help="Choose a template from our professionally designed collection"
)

            if selected_system_template:
                # Check if the selection has actually changed
                if st.session_state.selected_template != selected_system_template:
                    template_config = SYSTEM_TEMPLATES[selected_system_template]
                    st.session_state.selected_template = selected_system_template
                    st.session_state.selected_template_config = template_config
                    st.session_state.template_source = 'system'
                    
                    # Update query params
                    st.query_params["template"] = selected_system_template
                    st.query_params["source"] = 'system'
                    
                    # Force rerun to update preview
                    st.rerun()
            
            st.markdown('</div>', unsafe_allow_html=True)
            
            st.markdown("<br>", unsafe_allow_html=True)
    with tab2:
        st.markdown('<div class="panel-header">🗂️ Your Saved Templates</div>', unsafe_allow_html=True)
        
        # Load templates
        if 'uploaded_templates' not in st.session_state:
            st.session_state.uploaded_templates = load_user_templates(current_user)
        
        if 'doc_templates' not in st.session_state:
            st.session_state.doc_templates = load_user_doc_templates(current_user)
        
        if 'ppt_templates' not in st.session_state:
            st.session_state.ppt_templates = load_user_ppt_templates(current_user)
        
        # Tabs for different template types
        template_tab1, template_tab2, template_tab3 = st.tabs(["📄 HTML", "📝 Word", "📊 PPT"])
        
        # ========== HTML TEMPLATES TAB ==========
        with template_tab1:
            if st.session_state.uploaded_templates:
                # Search and filter header
                st.markdown("""
                <div class="template-list-header">
                    <div style="flex: 2;">Template Name</div>
                    <div style="flex: 1; text-align: center;">Type</div>
                    <div style="flex: 1; text-align: right;">Actions</div>
                </div>
                """, unsafe_allow_html=True)
                
                st.markdown('<div class="template-list-container">', unsafe_allow_html=True)
                
                for template_id, template_data in st.session_state.uploaded_templates.items():
                    st.markdown(f"""
                    <div class="template-item">
                        <div class="template-item-left">
                            <div class="template-icon">📄</div>
                            <div class="template-info">
                                <div class="template-name-text">{template_data['name']}</div>
                                <div class="template-meta">{template_data['original_filename']}</div>
                            </div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    col1, col2 = st.columns([2, 1])
                    with col1:
                        if st.button("Use", key=f"use_html_{template_id}", type="primary", use_container_width=True):
                            if 'temp_upload_config' in st.session_state:
                                del st.session_state.temp_upload_config
                            
                            st.session_state.selected_template_preview = f"""
                                <style>{template_data['css']}</style>
                                <div class="ats-page">{generate_generic_html(user_resume)}</div>
                            """
                            st.session_state.selected_template = template_data['name']
                            st.session_state.selected_template_config = template_data
                            st.session_state.template_source = 'saved'
                            st.session_state.current_upload_id = template_id
                            st.rerun()
                    
                    with col2:
                        if st.button("🗑️", key=f"delete_html_{template_id}", use_container_width=True):
                            if st.session_state.get('current_upload_id') == template_id:
                                st.session_state.pop('selected_template_preview', None)
                                st.session_state.pop('selected_template', None)
                                st.session_state.pop('selected_template_config', None)
                                st.session_state.pop('current_upload_id', None)
                            
                            del st.session_state.uploaded_templates[template_id]
                            save_user_templates(st.session_state.logged_in_user, st.session_state.uploaded_templates)
                            st.success(f"✅ Deleted '{template_data['name']}'")
                            st.rerun()
                
                st.markdown('</div>', unsafe_allow_html=True)
            else:
                st.info("📂 No saved HTML templates yet.")
        
        # ========== WORD TEMPLATES TAB ==========
        with template_tab2:
            if st.session_state.doc_templates:
                st.markdown("""
                <div class="template-list-header">
                    <div style="flex: 2;">Template Name</div>
                    <div style="flex: 1; text-align: center;">Type</div>
                    <div style="flex: 1; text-align: right;">Actions</div>
                </div>
                """, unsafe_allow_html=True)
                
                st.markdown('<div class="template-list-container">', unsafe_allow_html=True)
                
                for template_id, template_data in st.session_state.doc_templates.items():
                    st.markdown(f"""
                    <div class="template-item">
                        <div class="template-item-left">
                            <div class="template-icon">📝</div>
                            <div class="template-info">
                                <div class="template-name-text">{template_data['name']}</div>
                                <div class="template-meta">{template_data['original_filename']}</div>
                            </div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    col1, col2 = st.columns([2, 1])
                    with col1:
                        if st.button("Use", key=f"use_doc_{template_id}", type="primary", use_container_width=True):
                            try:
                                import io
                                from docx import Document
                                
                                doc_stream = io.BytesIO(template_data['doc_data'])
                                doc = Document(doc_stream)
                                structure = template_data.get('structure', [])
                                output, replaced, removed = replace_content(doc, structure, user_resume)
                                
                                st.session_state.generated_doc = output.getvalue()
                                st.session_state.selected_doc_template_id = template_id
                                st.session_state.selected_doc_template = template_data
                                st.session_state.doc_template_source = 'saved'
                                
                                st.success(f"✅ Using template: {template_data['name']}")
                                st.rerun()
                            except Exception as e:
                                st.error(f"Error loading template: {str(e)}")
                    
                    with col2:
                        if st.button("🗑️", key=f"delete_doc_{template_id}", use_container_width=True):
                            if st.session_state.get('selected_doc_template_id') == template_id:
                                st.session_state.pop('generated_doc', None)
                                st.session_state.pop('selected_doc_template_id', None)
                                st.session_state.pop('selected_doc_template', None)
                                st.session_state.pop('doc_template_source', None)
                            
                            del st.session_state.doc_templates[template_id]
                            save_user_doc_templates(st.session_state.logged_in_user, st.session_state.doc_templates)
                            st.success(f"✅ Deleted '{template_data['name']}'")
                            st.rerun()
                
                st.markdown('</div>', unsafe_allow_html=True)
            else:
                st.info("📂 No saved Word templates yet.")
        
        # ========== POWERPOINT TEMPLATES TAB ==========
        with template_tab3:
            if st.session_state.ppt_templates:
                st.markdown("""
                <div class="template-list-header">
                    <div style="flex: 2;">Template Name</div>
                    <div style="flex: 1; text-align: center;">Type</div>
                    <div style="flex: 1; text-align: right;">Actions</div>
                </div>
                """, unsafe_allow_html=True)
                
                st.markdown('<div class="template-list-container">', unsafe_allow_html=True)
                
                for template_id, template_data in st.session_state.ppt_templates.items():
                    st.markdown(f"""
                    <div class="template-item">
                        <div class="template-item-left">
                            <div class="template-icon">📊</div>
                            <div class="template-info">
                                <div class="template-name-text">{template_data['name']}</div>
                                <div class="template-meta">{template_data['original_filename']}</div>
                            </div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    col1, col2 = st.columns([2, 1])
                    with col1:
                        if st.button("Use", key=f"use_ppt_{template_id}", type="primary", use_container_width=True):
                            try:
                                import io
                                from pptx import Presentation
                                
                                working_prs = Presentation(io.BytesIO(template_data['ppt_data']))
                                prs = Presentation(io.BytesIO(template_data['ppt_data']))
                                
                                slide_texts = []
                                for slide_idx, slide in enumerate(prs.slides):
                                    text_blocks = []
                                    for shape_idx, shape in enumerate(slide.shapes):
                                        if shape.has_text_frame and shape.text.strip():
                                            text_blocks.append({
                                                "index": shape_idx,
                                                "text": shape.text.strip(),
                                                "position": {"x": shape.left, "y": shape.top}
                                            })
                                    
                                    if text_blocks:
                                        text_blocks.sort(key=lambda x: (x["position"]["y"], x["position"]["x"]))
                                        slide_texts.append({
                                            "slide_number": slide_idx + 1,
                                            "text_blocks": text_blocks
                                        })
                                
                                structured_slides = analyze_slide_structure(slide_texts)
                                generated_sections = generate_ppt_sections(user_resume, structured_slides)
                                
                                text_elements = template_data['text_elements']
                                content_mapping, heading_shapes, basic_info_shapes = match_generated_to_original(
                                    text_elements, generated_sections, prs)
                                
                                edits = {}
                                for element in text_elements:
                                    key = f"{element['slide']}_{element['shape']}"
                                    if key not in heading_shapes:
                                        edits[key] = content_mapping.get(key, element['original_text'])
                                
                                success_count = 0
                                for element in text_elements:
                                    key = f"{element['slide']}_{element['shape']}"
                                    if key not in heading_shapes and key in edits:
                                        slide_idx = element['slide'] - 1
                                        shape_idx = element['shape']
                                        
                                        if slide_idx < len(working_prs.slides):
                                            slide = working_prs.slides[slide_idx]
                                            if shape_idx < len(slide.shapes):
                                                shape = slide.shapes[shape_idx]
                                                if shape.has_text_frame:
                                                    clear_and_replace_text(shape, edits[key])
                                                    success_count += 1
                                
                                output = io.BytesIO()
                                working_prs.save(output)
                                output.seek(0)
                                
                                st.session_state.generated_ppt = output.getvalue()
                                st.session_state.selected_ppt_template_id = template_id
                                st.session_state.selected_ppt_template = template_data
                                st.session_state.ppt_template_source = 'saved'
                                
                                st.success(f"✅ Using template: {template_data['name']}")
                                st.rerun()
                            except Exception as e:
                                st.error(f"Error loading template: {str(e)}")
                    
                    with col2:
                        if st.button("🗑️", key=f"delete_ppt_{template_id}", use_container_width=True):
                            if st.session_state.get('selected_ppt_template_id') == template_id:
                                st.session_state.pop('generated_ppt', None)
                                st.session_state.pop('selected_ppt_template_id', None)
                                st.session_state.pop('selected_ppt_template', None)
                                st.session_state.pop('ppt_template_source', None)
                            
                            del st.session_state.ppt_templates[template_id]
                            save_user_ppt_templates(st.session_state.logged_in_user, st.session_state.ppt_templates)
                            st.success(f"✅ Deleted '{template_data['name']}'")
                            st.rerun()
                
                st.markdown('</div>', unsafe_allow_html=True)
            else:
                st.info("📂 No saved PowerPoint templates yet.")
        
        st.markdown('</div>', unsafe_allow_html=True)
        






    # Add Template Section
    # st.markdown('<div class="control-panel">', unsafe_allow_html=True)
        
    with tab3:
    # Download Panel
    # st.markdown('<div class="control-panel">', unsafe_allow_html=True)
        st.markdown('<div class="panel-header">📥 Download</div>', unsafe_allow_html=True)
        
        if template_source == 'system':
            selected_color = st.session_state.get('preview_selected_color', ATS_COLORS["Professional Blue (Default)"])
            css = load_css_template(
                selected_config['css_template'],
                selected_color
            )

            html_content = selected_config['html_generator'](user_resume)
            
            # HTML Download
            full_doc = f"<html><head><meta charset='UTF-8'><style>{css}</style></head><body><div class='ats-page'>{html_content}</div></body></html>"
            html_filename = f"Resume_{user_resume.get('name', 'User').replace(' ', '_')}.html"
            st.download_button(
                label="📄 Download HTML",
                data=full_doc.encode('utf-8'),
                file_name=html_filename,
                mime="text/html",
                use_container_width=True,
                key="download_html"
            )
            
            # DOC Download
            word_doc = f"""
    <html xmlns:o='urn:schemas-microsoft-com:office:office'
    xmlns:w='urn:schemas-microsoft-com:office:word'
    xmlns='http://www.w3.org/TR/REC-html40'>
    <head>
        <meta charset='UTF-8'>
        <style>
            @page {{ size: 8.5in 11in; margin: 0.5in; }}
            {css}
        </style>
    </head>
    <body>
        <div class='ats-page'>{html_content}</div>
    </body>
    </html>
    """
            doc_filename = f"Resume_{user_resume.get('name', 'User').replace(' ', '_')}.doc"
            st.download_button(
                label="📑 Download DOC",
                data=word_doc.encode('utf-8'),
                file_name=doc_filename,
                mime="application/msword",
                use_container_width=True,
                key="download_doc"
            )
            
            # TXT Download
            txt_content = generate_markdown_text(user_resume)
            txt_filename = f"Resume_{user_resume.get('name', 'User').replace(' ', '_')}.txt"
            st.download_button(
                label="📝 Download TXT",
                data=txt_content.encode('utf-8'),
                file_name=txt_filename,
                mime="text/plain",
                use_container_width=True,
                key="download_txt"
            )
        
        elif template_source in ['saved', 'temp_upload']:
            css = selected_config.get('css', '')
            html_content = generate_generic_html(user_resume)
            full_doc = f"<html><head><meta charset='UTF-8'><style>{css}</style></head><body><div class='ats-page'>{html_content}</div></body></html>"
            html_filename = f"Resume_{user_resume.get('name', 'User').replace(' ', '_')}.html"
            
            st.download_button(
                label="📄 Download HTML",
                data=full_doc.encode('utf-8'),
                file_name=html_filename,
                mime="text/html",
                use_container_width=True,
                key="download_html_custom"
            )
        
        elif template_source == 'doc_saved' and st.session_state.get('generated_doc'):
            doc_filename = f"Resume_{user_resume.get('name', 'Resume').replace(' ', '_')}.docx"
            st.download_button(
                label="📄 Download DOCX",
                data=st.session_state.generated_doc,
                file_name=doc_filename,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True,
                key="download_docx"
            )
        
        elif template_source == 'ppt_saved' and st.session_state.get('generated_ppt'):
            ppt_filename = f"Resume_{user_resume.get('name', 'Presentation').replace(' ', '_')}.pptx"
            st.download_button(
                label="📊 Download PPTX",
                data=st.session_state.generated_ppt,
                file_name=ppt_filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key="download_pptx"
            )
        
        st.markdown('</div>', unsafe_allow_html=True)

# ----------------------------------
# COLUMN 2: COLORS
    with tab4:
        st.markdown('<div class="panel-header">➕ Add Template</div>', unsafe_allow_html=True)
        
        add_method = st.selectbox(
            "Upload Type:",
            ["Select...", "HTML", "DOCX", "PPTX"],
            key="add_method"
        )
        
        if add_method == "HTML":
            uploaded_file = st.file_uploader("Upload HTML/CSS File", type=['html', 'css'], key="html_upload")
            if uploaded_file:
                import chardet
                raw_data = uploaded_file.read()
                detected = chardet.detect(raw_data)
                encoding = detected["encoding"] or "utf-8"
                content = raw_data.decode(encoding, errors="ignore")
                
                # Extract template from HTML
                from utils import extract_template_from_html
                parsed_template = extract_template_from_html(content)
                
                # Store in temporary session state for preview
                st.session_state.temp_upload_config = {
                    'name': f"Uploaded_{uploaded_file.name.split('.')[0]}",
                    'css': parsed_template.get('css', ''),
                    'html': parsed_template.get('html', ''),
                    'original_filename': uploaded_file.name
                }
                
                # Show template name input and save button
                col1, col2 = st.columns([2, 1])
                with col1:
                    template_name = st.text_input(
                        "Template Name:",
                        value=f"Uploaded_{uploaded_file.name.split('.')[0]}",
                        key="upload_template_name"
                    )
                
                with col2:
                    if st.button("💾 Save Template", use_container_width=True):
                        if 'uploaded_templates' not in st.session_state:
                            st.session_state.uploaded_templates = load_user_templates(current_user)
                        
                        template_id = f"upload_{datetime.now().strftime('%Y%m%d%H%M%S')}"
                        st.session_state.uploaded_templates[template_id] = {
                            'name': template_name,
                            'css': parsed_template.get('css', ''),
                            'html': parsed_template.get('html', ''),
                            'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                            'original_filename': uploaded_file.name
                        }
                        
                        save_user_templates(current_user, st.session_state.uploaded_templates)
                        st.success(f"✅ Template '{template_name}' saved!")
                        
                        # Set as active template
                        st.session_state.selected_template_config = st.session_state.uploaded_templates[template_id]
                        st.session_state.selected_template = template_name
                        st.session_state.template_source = 'saved'
                        st.session_state.current_upload_id = template_id
                        
                        if 'temp_upload_config' in st.session_state:
                            del st.session_state.temp_upload_config
                        
                        st.rerun()
                
                # Apply button for immediate preview without saving
                if st.button("👁️ Apply Template Preview", use_container_width=True, type="secondary"):
                    st.session_state.selected_template_config = st.session_state.temp_upload_config
                    st.session_state.selected_template = f"Temp_{uploaded_file.name.split('.')[0]}"
                    st.session_state.template_source = 'temp_upload'
                    st.success("✅ Template applied for preview!")
                    st.rerun()
        
        elif add_method == "DOCX":
            uploaded_file = st.file_uploader("Upload DOCX File", type=['docx', 'doc'], key="docx_upload")
            if uploaded_file:
                try:
                    from utils import extract_document_structure, replace_content
                    
                    # Process document
                    uploaded_file.seek(0)
                    doc, structure = extract_document_structure(uploaded_file)
                    
                    # Store original template data
                    uploaded_file.seek(0)
                    st.session_state.temp_doc_data = uploaded_file.read()
                    st.session_state.temp_doc_filename = uploaded_file.name
                    
                    # Replace content
                    output, replaced, removed = replace_content(doc, structure, user_resume)
                    
                    # Store results
                    st.session_state.generated_doc = output.getvalue()
                    st.session_state.doc_structure = structure
                    st.session_state.doc_replaced = replaced
                    st.session_state.doc_removed = removed
                    
                    # Template name and save section
                    col1, col2 = st.columns([2, 1])
                    
                    with col1:
                        doc_template_name = st.text_input(
                            "Template Name:",
                            value=f"Doc_{uploaded_file.name.split('.')[0]}",
                            key="doc_template_name"
                        )
                    
                    with col2:
                        st.write("")
                        st.write("")
                        if st.button("💾 Save Template", use_container_width=True, type="primary"):
                            if 'doc_templates' not in st.session_state:
                                st.session_state.doc_templates = load_user_doc_templates(current_user)
                            
                            template_id = f"doc_{datetime.now().strftime('%Y%m%d%H%M%S')}"
                            st.session_state.doc_templates[template_id] = {
                                'name': doc_template_name,
                                'doc_data': st.session_state.temp_doc_data,
                                'structure': structure,
                                'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                                'original_filename': uploaded_file.name,
                                'sections_detected': [s['section'] for s in structure]
                            }
                            
                            if save_user_doc_templates(current_user, st.session_state.doc_templates):
                                st.success(f"✅ Template '{doc_template_name}' saved!")
                                st.session_state.selected_doc_template_id = template_id
                                st.session_state.selected_doc_template = st.session_state.doc_templates[template_id]
                                st.session_state.doc_template_source = 'saved'
                                st.balloons()
                                st.rerun()
                            else:
                                st.error("Failed to save template. Please try again.")
                    
                    # Apply button for immediate preview
                    if st.button("👁️ Apply Template Preview", use_container_width=True, type="secondary"):
                        st.session_state.selected_doc_template = {
                            'name': f"Temp_Doc_{uploaded_file.name.split('.')[0]}",
                            'doc_data': st.session_state.temp_doc_data,
                            'structure': structure
                        }
                        st.session_state.doc_template_source = 'temp_upload'
                        st.session_state.generated_doc = output.getvalue()
                        st.success("✅ DOCX template applied for preview!")
                        st.rerun()
                        
                except Exception as e:
                    st.error(f"❌ Error processing document: {str(e)}")
        
        elif add_method == "PPTX":
            uploaded_file = st.file_uploader("Upload PPTX File", type=['pptx'], key="pptx_upload")
            if uploaded_file:
                try:
                    import io
                    from pptx import Presentation
                    from utils import analyze_slide_structure, generate_ppt_sections, match_generated_to_original, clear_and_replace_text
                    
                    st.session_state.ppt_uploaded_file = uploaded_file.getvalue()
                    st.session_state.ppt_original_filename = uploaded_file.name
                    
                    prs = Presentation(io.BytesIO(st.session_state.ppt_uploaded_file))
                    slide_texts = []
                    for slide_idx, slide in enumerate(prs.slides):
                        text_blocks = []
                        for shape_idx, shape in enumerate(slide.shapes):
                            if shape.has_text_frame and shape.text.strip():
                                text_blocks.append({
                                    "index": shape_idx,
                                    "text": shape.text.strip(),
                                    "position": {"x": shape.left, "y": shape.top}
                                })
                        
                        if text_blocks:
                            text_blocks.sort(key=lambda x: (x["position"]["y"], x["position"]["x"]))
                            slide_texts.append({
                                "slide_number": slide_idx + 1,
                                "text_blocks": text_blocks
                            })
                    
                    if slide_texts:
                        structured_slides = analyze_slide_structure(slide_texts)
                        generated_sections = generate_ppt_sections(user_resume, structured_slides)
                        
                        text_elements = []
                        for slide_idx, slide in enumerate(prs.slides):
                            for shape_idx, shape in enumerate(slide.shapes):
                                if shape.has_text_frame and shape.text.strip():
                                    text_elements.append({
                                        'slide': slide_idx + 1,
                                        'shape': shape_idx,
                                        'original_text': shape.text.strip(),
                                        'shape_type': type(shape).__name__
                                    })
                        
                        content_mapping, heading_shapes, basic_info_shapes = match_generated_to_original(
                            text_elements, generated_sections, prs)
                        
                        working_prs = Presentation(io.BytesIO(st.session_state.ppt_uploaded_file))
                        edits = {}
                        
                        for element in text_elements:
                            key = f"{element['slide']}_{element['shape']}"
                            if key not in heading_shapes:
                                edits[key] = content_mapping.get(key, element['original_text'])
                        
                        success_count = 0
                        for element in text_elements:
                            key = f"{element['slide']}_{element['shape']}"
                            if key not in heading_shapes and key in edits:
                                slide_idx = element['slide'] - 1
                                shape_idx = element['shape']
                                
                                if slide_idx < len(working_prs.slides):
                                    slide = working_prs.slides[slide_idx]
                                    if shape_idx < len(slide.shapes):
                                        shape = slide.shapes[shape_idx]
                                        if shape.has_text_frame:
                                            clear_and_replace_text(shape, edits[key])
                                            success_count += 1
                        
                        output = io.BytesIO()
                        working_prs.save(output)
                        output.seek(0)
                        st.session_state.generated_ppt = output.getvalue()
                        
                        # Template name and save section
                        col1, col2 = st.columns([2, 1])
                        
                        with col1:
                            ppt_template_name = st.text_input(
                                "Template Name:",
                                value=f"PPT_{uploaded_file.name.split('.')[0]}",
                                key="ppt_template_name"
                            )
                        
                        with col2:
                            st.write("")
                            st.write("")
                            if st.button("💾 Save Template", use_container_width=True, type="primary"):
                                if 'ppt_templates' not in st.session_state:
                                    st.session_state.ppt_templates = load_user_ppt_templates(current_user)
                                
                                ppt_id = f"ppt_{datetime.now().strftime('%Y%m%d%H%M%S')}"
                                st.session_state.ppt_templates[ppt_id] = {
                                    'name': ppt_template_name,
                                    'ppt_data': st.session_state.ppt_uploaded_file,
                                    'edits': edits,
                                    'content_mapping': content_mapping,
                                    'heading_shapes': list(heading_shapes),
                                    'basic_info_shapes': list(basic_info_shapes),
                                    'text_elements': text_elements,
                                    'uploaded_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                                    'original_filename': uploaded_file.name
                                }
                                
                                if save_user_ppt_templates(current_user, st.session_state.ppt_templates):
                                    st.success(f"✅ PPT Template '{ppt_template_name}' saved!")
                                    st.session_state.selected_ppt_template_id = ppt_id
                                    st.session_state.selected_ppt_template = st.session_state.ppt_templates[ppt_id]
                                    st.session_state.ppt_template_source = 'saved'
                                    st.balloons()
                                    st.rerun()
                                else:
                                    st.error("Failed to save template. Please try again.")
                        
                        # Apply button for immediate preview
                        if st.button("👁️ Apply Template Preview", use_container_width=True, type="secondary"):
                            st.session_state.selected_ppt_template = {
                                'name': f"Temp_PPT_{uploaded_file.name.split('.')[0]}",
                                'ppt_data': st.session_state.ppt_uploaded_file
                            }
                            st.session_state.ppt_template_source = 'temp_upload'
                            st.session_state.generated_ppt = output.getvalue()
                            st.success("✅ PPTX template applied for preview!")
                            st.rerun()
                            
                except Exception as e:
                    st.error(f"❌ Error processing PowerPoint: {str(e)}")
        
        st.markdown('</div>', unsafe_allow_html=True)




    
    # ========== YOUR SAVED TEMPLATES ==========
    # st.markdown('<div class="control-panel">', unsafe_allow_html=True)
        

# ----------------------------------
# COLUMN 3: PREVIEW
# ----------------------------------
with col3:
    # st.markdown('<div class="preview-panel">', unsafe_allow_html=True)
    st.markdown('<div class="preview-header">Resume Preview</div>', unsafe_allow_html=True)
    
    # Generate preview based on template type
    if template_source == 'system':
        # Get the selected color from session state
        selected_color = st.session_state.get('preview_selected_color', ATS_COLORS["Professional Blue (Default)"])
        
        # Generate CSS and HTML with the selected color
        css = load_css_template(
    selected_config['css_template'],
    selected_color
)

        html_content = selected_config['html_generator'](user_resume)
        
        # Create full HTML with inline styles - using the exact format from download.py
        full_html = f"""
        {css}
        <div class="ats-page">
            {html_content}
        </div>
        """
        
        # Use a unique container key based on color to force refresh
        preview_container = st.container()
        with preview_container:
            # Add a hidden element that changes with color to force iframe reload
            st.markdown(f'<!-- color: {selected_color} -->', unsafe_allow_html=True)
            st.components.v1.html(full_html, height=1000, scrolling=True)
    
    elif template_source in ['saved', 'temp_upload']:
        css = selected_config.get('css', '')
        html_content = generate_generic_html(user_resume)
        
        full_html = f"""
        <style>{css}</style>
        <div class="ats-page">
            {html_content}
        </div>
        """
        st.components.v1.html(full_html, height=1000, scrolling=True)
    
    elif template_source == 'doc_saved' and st.session_state.get('generated_doc'):
        try:
            from docx import Document
            import io
            
            doc_stream = io.BytesIO(st.session_state.generated_doc)
            processed_doc = Document(doc_stream)
            
            html_preview = '<div style="background: white; padding: 40px; font-family: Calibri, Arial; line-height: 1.6;">'
            for para in processed_doc.paragraphs:
                if para.text.strip():
                    html_preview += f'<p style="margin: 8px 0;">{para.text.strip()}</p>'
            html_preview += '</div>'
            
            st.markdown(html_preview, unsafe_allow_html=True)
        except Exception as e:
            st.error(f"Preview error: {str(e)}")
    
    elif template_source == 'ppt_saved' and st.session_state.get('generated_ppt'):
        try:
            from pptx import Presentation
            import io
            
            preview_prs = Presentation(io.BytesIO(st.session_state.generated_ppt))
            
            for slide_idx, slide in enumerate(preview_prs.slides):
                st.markdown(f"""
                <div style="background: #f5f7fa; padding: 20px; margin: 15px 0; border-radius: 10px; border: 2px solid #ffe5d9;">
                    <div style="background: #e87532; color: white; padding: 10px; border-radius: 5px; margin-bottom: 15px; font-weight: bold;">
                        Slide {slide_idx + 1}
                    </div>
                """, unsafe_allow_html=True)
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        st.markdown(f"""
                        <div style="background: white; padding: 15px; margin: 10px 0; border-radius: 8px; border-left: 4px solid #e87532;">
                            {shape.text.strip().replace(chr(10), '<br>')}
                        </div>
                        """, unsafe_allow_html=True)
                
                st.markdown("</div>", unsafe_allow_html=True)
        except Exception as e:
            st.error(f"Preview error: {str(e)}")
    
    st.markdown('</div>', unsafe_allow_html=True)